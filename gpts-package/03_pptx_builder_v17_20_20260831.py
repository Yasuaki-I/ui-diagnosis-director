# -*- coding: utf-8 -*-
"""
紺＆クリーン スライド作成 ─ python-pptx ビルダー
==================================================

GPTs の Knowledge にアップロードして使用する PPTX 生成ライブラリ。
ChatGPT の Code Interpreter から exec() で読み込み、関数を呼び出して .pptx を組み立てます。

使い方:
    exec(open('/mnt/data/03_pptx_builder.py').read())
    prs = create_presentation()
    add_cover(prs, title='Q2 販売戦略', date='2026年7月', author='○○株式会社 営業企画部')
    add_agenda(prs, items=['市場環境', '現状課題', '施策', '計画'])
    # ... 他のレイアウト関数 ...
    prs.save('output.pptx')

必達13条項（Instructions参照）を厳格に遵守。

バージョン履歴:
    v17.1.0 (2026-08-24): v17 P2 図解パターン描画3種を統合。
        - draw_pyramid / draw_sequence / draw_framework。
        - draw_pattern を6種対応に拡張（P1 3種＋P2 3種）。
        - 要素数・軸要否は集約表（8/15）を正とする（8/24 統括指示）。
        - pyramid の台形は add_freeform を使わず幅可変の矩形段で近似（原則④）。
        - フォールバック理由を戻り値の notes に記録するよう修正（P1にも遡及）。
    v17.0.1 (2026-08-24): 実機検証（入江さん）で検出した2件を修正。
        - カード高さ上限を 230→186px（V17_CARD_H_MAX）。充填率47%→60%。
        - add_diagram_slide() を追加。ヘッダ帯とスライド内見出しの
          タイトル二重表示を構造的に防止（ヘッダは原本 `ja` を表示）。
    v17.0.0 (2026-08-23): v17 P1 図解パターン描画3種を統合。
        - draw_category / draw_breakdown / draw_comparison。
        - DIAGRAM_PATTERNS（原本）は無改変。拡張層は DIAGRAM_PATTERN_SPEC に分離。
    v15 (2026-07-12): リスト表示の「折返し行間」と「項目間余白」を分離制御。
        - _add_multi_run_box / add_paragraph_box に space_after_pt /
          space_before_pt / vertical_anchor パラメータを追加。
        - リスト表示6箇所（C-1強み/課題、C-3 slide1 Top3/強み、C-3 slide2
          コンパクト強み/課題）の line_height を 1.6-1.7 → 1.2 に、
          space_after_pt=8pt を項目間に付与。「1文の連続感」と「項目間の分離感」
          を独立制御し、折返し時の可読性を根本改善。
        - C-3 slide1 総評帯：SUMMARY_H を 62→68 に拡大、
          vertical_anchor=MIDDLE で天地余白を均等化。
    v14 (2026-07-09): 総評 1行運用確定、17字上限、SUMMARY_H=62 に最適化。
    v13 (2026-07-08): 総評本文18字、SUMMARY_H=76、ISSUE_H=190。
    v12 (2026-07-07): 強み項目 40→28字厳格化（語途中折返し防止）。
    v11 (2026-07-06): 総評⇄Top3レイアウト崩れ根治（文字数上限厳格化）。
    v10 (2026-07-05): C-1/C-2/C-3 の120字統一・word_wrap=True で文字切れ根絶。
"""
__version__ = '17.2.0'
__version_date__ = '2026-08-25'

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from copy import deepcopy

# =====================================================================
# 配色パレット（02_design_spec.md 準拠・変更禁止）
# =====================================================================
NAVY        = RGBColor(0x1C, 0x36, 0x6C)   # メインナビー
RED         = RGBColor(0xD0, 0x02, 0x1B)   # アクセント赤
LIGHT_GRAY  = RGBColor(0xF4, 0xF5, 0xF8)   # カード背景
BORDER_GRAY = RGBColor(0xD0, 0xD4, 0xDC)   # 罫線
TEXT        = RGBColor(0x40, 0x40, 0x40)   # 本文
SUB_TEXT    = RGBColor(0x60, 0x60, 0x60)   # 注釈
NAVY_LIGHT  = RGBColor(0x9D, 0xB0, 0xD6)   # ナビー帯内補助
NAVY_E6     = RGBColor(0xE6, 0xEA, 0xF3)   # ナビー帯内本文
GOLD        = RGBColor(0xFF, 0xD5, 0x4F)   # 紺帯内強調
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BORDER = RGBColor(0xCC, 0xCC, 0xCC)   # フッター境界線
PAGE_NUM    = RGBColor(0x26, 0x26, 0x26)   # ページ番号
STRIPE      = RGBColor(0xF8, 0xF9, 0xFB)   # 縞模様セル
# C-3 / UI診断ボード用 アクセントカラー
ORANGE      = RGBColor(0xF9, 0x73, 0x16)   # アクセントオレンジ（After/方向性）
ORANGE_LIGHT = RGBColor(0xFE, 0xF1, 0xE3)  # オレンジ薄帯（After背景）
RED_LIGHT   = RGBColor(0xFD, 0xEC, 0xEC)   # 赤薄帯（Before背景）
PRIO_RED    = RGBColor(0xD0, 0x02, 0x1B)   # 優先度高
PRIO_ORANGE = RGBColor(0xE8, 0x8B, 0x1F)   # 優先度中
PRIO_GRAY   = RGBColor(0x88, 0x88, 0x88)   # 優先度低

# =====================================================================
# ▼▼▼ Phase A 追加ブロック（v16.5・2026-07-28） ▼▼▼
# 目的：機能強化の下ごしらえ（辞書のみ・利用ロジック未実装）
# 既存のC-1〜C-3描画ロジックには一切影響しない（後方互換完全維持）
# 詳細設計：phase_a_design_20260727_rev2.md 参照
# =====================================================================

# ============================================================
# A-3：UI診断ディレクターの10項目評価軸の背景知識
# 出典：入江さんご自身の知見の集約（web-director.skill）
# 併用プロジェクト：マーケティングオーケストレーター
# 対応マッピング：10項目中8項目が web-director.skill 6大機能領域と「高」対応
# 詳細：phase_a_design_20260727_rev2.md A-3セクション参照
# ============================================================

# ============================================================
# A-1：デジタル庁 公式カラーパレット
# 出典：デジタル庁 ダッシュボードデザインの実践ガイドブックとデザインテンプレート
# https://www.digital.go.jp/resources/dashboard-guidebook
# ライセンス：PDL1.0（公共データ利用規約 第1.0版）
# GitHub：https://github.com/digital-go-jp/policy-dashboard-assets
# 取得日：2026-07-27（Claude-Chat先行調査）
# 用途：v3.5時点は内部保持のみ・利用ロジック未実装
#       将来（Phase B以降）で10項目スコア色分けへ応用予定
# ============================================================
DIGITAL_AGENCY_PALETTE = {
    "SolidGray": {
        "primary":  "#4D4D4D",
        "secondary":"#767676",
        "midtone":  "#999999",
        "light":    "#CCCCCC",
        "lightest": "#F2F2F2",
        "accent":   "#3460FB",  # 強調用（青）
        "warning":  "#FE3939",  # 警告用（赤）
        "bg":       "#F8F8FB",
    },
    "Blue": {
        "primary":  "#0017C1",
        "secondary":"#3460FB",
        "midtone":  "#7096F8",
        "light":    "#C5D7FB",
        "lightest": "#E8F1FE",
        "accent":   "#FE3939",
        "warning":  "#FFBBBB",
        "bg":       "#F8F8FB",
    },
    "LightBlue": {
        "primary":  "#0055AD",
        "secondary":"#008BF2",
        "midtone":  "#57B8FF",
        "light":    "#C0E4FF",
        "lightest": "#F0F9FF",
        "accent":   "#FE3939",
        "warning":  "#FFBBBB",
        "bg":       "#F8F8FB",
    },
    "Green": {
        "primary":  "#115A36",
        "secondary":"#259D63",
        "midtone":  "#51B883",
        "light":    "#9BD4B5",
        "lightest": "#E6F5EC",
        "accent":   "#666666",
        "warning":  "#CCCCCC",
        "bg":       "#F8F8FB",
    },
    "Cyan": {
        "primary":  "#006F83",
        "secondary":"#00A3BF",
        "midtone":  "#2BC8E4",
        "light":    "#99F2FF",
        "lightest": "#E9F7F9",
        "accent":   "#666666",
        "warning":  "#CCCCCC",
        "bg":       "#F8F8FB",
    },
    "Red": {
        "primary":  "#CE0000",
        "secondary":"#FE3939",
        "midtone":  "#FF7171",
        "light":    "#FFBBBB",
        "lightest": "#FDEEEE",
        "accent":   "#666666",
        "warning":  "#CCCCCC",
        "bg":       "#F8F8FB",
    },
    "Orange": {
        "primary":  "#AC3E00",
        "secondary":"#FB5B01",
        "midtone":  "#FF8D44",
        "light":    "#FFC199",
        "lightest": "#FFEEE2",
        "accent":   "#666666",
        "warning":  "#CCCCCC",
        "bg":       "#F8F8FB",
    },
}

# 全テーマ共通の閾値色（good/bad判定用）
DIGITAL_AGENCY_THRESHOLD = {
    "center":  "#E6E6E6",  # 中央値・中立表示
    # maximum / minimum は各テーマの4番目色に準ずる（テーマ依存）
}

# ============================================================
# A-2：パーツ図鑑・図解集 由来 図解パターン辞書
# 出典：うちた氏「パーツ図鑑_120種」「図解集_50種」「テンプレ大全_100枚」
# 商用利用許諾：2026-07-26 note返信にて確認済み
# ライセンス記録：UCHITA_LICENSE_RECORD_20260726.md 参照
# 用途：v3.5時点は辞書のみ・描画実装なし（Phase B以降で実装）
# ============================================================
DIAGRAM_PATTERNS = {
    'category':     {'ja': '分類',           'use': '要素を並列カテゴリで整理',      'shape': 'grid'},
    'pyramid':      {'ja': 'ピラミッド',     'use': '階層・優先順位を上下で表現',    'shape': 'triangle'},
    'comparison':   {'ja': '比較',           'use': '2〜3要素の対比',                'shape': 'side_by_side'},
    'sequence':     {'ja': '順序',           'use': 'ステップ・時系列を左→右',       'shape': 'arrow_chain'},
    'cycle':        {'ja': '循環',           'use': '反復プロセスを円環で表現',      'shape': 'circle_arrow'},
    'funnel':       {'ja': '絞り込み',       'use': '上から下へ絞り込むファネル型',  'shape': 'trapezoid'},
    'timeline':     {'ja': '時間軸',         'use': '期間別のマイルストーン',        'shape': 'horizontal_bar'},
    'breakdown':    {'ja': '分解',           'use': '全体を構成要素に分解',          'shape': 'tree'},
    'contrast':     {'ja': '対比',           'use': '対照的な2要素の並列強調',       'shape': 'split_screen'},
    'integration':  {'ja': '統合',           'use': '複数要素の統合結果',            'shape': 'merge'},
    'framework':    {'ja': 'フレームワーク', 'use': '4象限マトリクス等',             'shape': 'quadrant'},
    'network':      {'ja': 'ネットワーク',   'use': 'ノード間の関係性',              'shape': 'node_edge'},
}

# 診断結果 → 推奨図解パターンのマッピング
DIAGNOSIS_TO_PATTERN = {
    'proposal_categorization': 'category',
    'priority_ranking':        'pyramid',
    'before_after':            'comparison',
    'user_flow':               'sequence',
    'improvement_cycle':       'cycle',
    'conversion_funnel':       'funnel',
    'schedule':                'timeline',
    'score_breakdown':         'breakdown',
    'ux_contrast':             'contrast',
    'impact_cost_matrix':      'framework',
    'site_structure':          'network',
}

# =====================================================================
# ▲▲▲ Phase A 追加ブロック ここまで ▲▲▲
# =====================================================================

# キャンバス（1280×720 想定）
CANVAS_W_PX = 1280
CANVAS_H_PX = 720

# 共通フォント
FONT = 'メイリオ'

# =====================================================================
# 単位変換・ユーティリティ
# =====================================================================
def px(n):
    """1280×720 想定の px を EMU に変換 (1 px = 9525 EMU)"""
    return Emu(int(n * 9525))


def set_run(run, text, size_pt, bold=False, color=TEXT, italic=False):
    """run にテキスト・フォント設定を適用（メイリオ・14pt以上を強制）

    フォント設定は run の OOXML rPr 内の latin / ea / cs すべてに
    明示的にメイリオを書き込み、英数字も含めて完全に統一する。
    """
    if size_pt < 14:
        raise ValueError(
            f"font_size must be >= 14pt (got {size_pt}pt) ─ 条項2違反"
        )
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = False  # 条項13: 斜体強調禁止
    run.font.color.rgb = color
    _force_all_script_fonts(run, FONT)


def _force_all_script_fonts(run, font_name):
    """run の latin / eastAsia / complexScript すべてにフォントを強制設定。

    PowerPoint OOXML は run のフォントを latin（英数字）/ ea（日本語等）/
    cs（複雑表記）の3スクリプトで個別に持つ。run.font.name は latin の
    みを設定するため、日本語文字＋数字混在のテキストではフォントが
    分かれて見えることがある。このヘルパーで3スクリプトすべて明示的に
    同じフォントへ統一する。
    """
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        # 既存タグを削除してから追加（重複防止）
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
        el = rPr.makeelement(qn(tag), {'typeface': font_name})
        rPr.append(el)


# 後方互換のエイリアス
_force_east_asian_font = _force_all_script_fonts


# =====================================================================
# 文字数バリデーション（design_system.md §1 準拠・絶対遵守）
# =====================================================================
# UI診断ディレクター用の文字数規定。
# スライドはコンテナサイズ固定のため、テキスト量を物理的に制御しないと
# 必ずレイアウトが崩壊する。規定文字数を超えた場合は ValueError で停止し、
# GPTs 側で「末尾…省略ではなく、規定文字内で要約し直す」運用とする。
# 詳細は gpts-package/design_system.md §1 を参照。

LIMITS = {
    # ─── C-1 スコアカード ─────────────────────
    'service_name':       30,
    'c1_comment':         40,   # 一言所見（10項目×1コメント）
    'c1_strength':        60,   # 強み（3項目）
    'c1_issue':           60,   # 最優先課題（3項目）
    'c1_conclusion':     120,   # 結論・総評
    # ─── C-2 改善提案リスト ──────────────────
    'c2_title':           30,   # 提案タイトル
    'c2_point':          150,   # POINT本文
    'c2_priority':         6,   # 優先度バッジ
    'c2_category':        12,   # カテゴリ
    # ─── C-3 ビジュアル診断ボード ────────────
    'c3_summary':         80,   # 総評（スライド1中央）
    'c3_top_issue':       40,   # 最優先課題タイトル
    'c3_direction':      120,   # 改善方向帯（2行折り返し許容）
    'c3_section_label':   20,   # LP構造マップ各セクション
    'c3_flow_step':       25,   # 行動フローステップ
    'c3_highlight_title': 30,   # Before/Afterタイトル
    'c3_before':          40,   # Before本文
    'c3_after':           40,   # After本文
}


def validate_length(text, limit_key, field_label=None, *, allow_none=True):
    """規定文字数を超えた場合 ValueError で停止する。

    Parameters
    ----------
    text : str or None
        検査対象テキスト。None / 空文字は許容（allow_none=False で禁止）。
    limit_key : str
        LIMITS 辞書のキー。例: 'c1_comment', 'c2_point', 'c3_summary'。
    field_label : str, optional
        エラーメッセージに表示する人間向けの項目名。
        未指定なら limit_key をそのまま使う。
    allow_none : bool
        True なら None / 空文字を素通り。False なら必須扱いで ValueError。

    Returns
    -------
    str
        検査済みテキスト（None なら空文字に正規化）。

    Raises
    ------
    KeyError
        limit_key が LIMITS に存在しない場合（開発者ミス）。
    ValueError
        文字数超過、または allow_none=False で空入力の場合。
    """
    if limit_key not in LIMITS:
        raise KeyError(
            f"validate_length: 未定義の limit_key '{limit_key}'。"
            f" LIMITS 辞書に追加してください。"
        )
    max_len = LIMITS[limit_key]
    label = field_label or limit_key

    if text is None or text == '':
        if allow_none:
            return ''
        raise ValueError(
            f"\n[文字数バリデーションエラー]\n"
            f"  項目: {label}\n"
            f"  問題: 必須項目が空欄です。\n"
            f"→ GPTs側でテキストを生成してください。"
        )

    text_str = str(text)
    n = len(text_str)
    if n > max_len:
        over = n - max_len
        preview = text_str[:30] + ('...' if n > 30 else '')
        raise ValueError(
            f"\n[文字数バリデーションエラー]\n"
            f"  項目: {label}\n"
            f"  上限: {max_len}文字（design_system.md §1 準拠）\n"
            f"  実際: {n}文字（{over}文字オーバー）\n"
            f"  内容: 「{preview}」\n"
            f"→ GPTs側で規定文字数内に要約し直してください。\n"
            f"  ※末尾「…」での省略は禁止です。要点を絞り込んでください。"
        )
    return text_str


def add_text(slide, left_px, top_px, width_px, text, size_pt, *,
             bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             height_px=None, line_height=None, letter_spacing=None,
             vertical_anchor=None):
    """1行のテキストボックスを追加（内容追従、高さ自動・条項5準拠）

    vertical_anchor: MSO_ANCHOR.TOP / MIDDLE / BOTTOM（既定None=PPTXデフォルト）。
                     帯内の縦中央配置などで使用。
    """
    if height_px is None:
        height_px = max(int(size_pt * 1.6), 24)
    box = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.auto_size = MSO_AUTO_SIZE.NONE  # 高さ調整は word_wrap で
    if vertical_anchor is not None:
        tf.vertical_anchor = vertical_anchor
    para = tf.paragraphs[0]
    para.alignment = align
    if line_height:
        para.line_spacing = line_height
    run = para.add_run()
    set_run(run, text, size_pt, bold=bold, color=color)
    return box


def add_paragraph_box(slide, left_px, top_px, width_px, paragraphs, *,
                      height_px=None, default_size=16, default_color=TEXT,
                      line_height=1.6,
                      space_after_pt=None, space_before_pt=None):
    """
    複数行（段落 or 箇条書き）を1つのテキストボックスに集約（条項6準拠）

    paragraphs: list of dict or str
      dict: {'text': str, 'size': int, 'bold': bool, 'color': RGBColor,
             'align': PP_ALIGN, 'bullet': bool,
             'space_after_pt': float, 'space_before_pt': float}
      str:  既定スタイルで追加

    line_height : float
        段落**内部**の折返し行間（既定 1.6）。
    space_after_pt : float, optional
        [v15/2026-07-12] 全段落に適用する段落末尾余白（pt）。
        リスト表示で項目間の分離感を演出する用途。
    space_before_pt : float, optional
        全段落に適用する段落先頭余白（pt）。通常は未使用。
    """
    if height_px is None:
        height_px = 400  # 大きめに確保（内容で自動縮小される表示）
    box = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, p in enumerate(paragraphs):
        if isinstance(p, str):
            p = {'text': p}
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        para.line_spacing = p.get('line_height', line_height)
        # [v15] 段落間余白の分離制御
        sa = p.get('space_after_pt', space_after_pt)
        if sa is not None and sa != 0:
            para.space_after = Pt(sa)
        sb = p.get('space_before_pt', space_before_pt)
        if sb is not None and sb != 0:
            para.space_before = Pt(sb)
        run = para.add_run()
        set_run(
            run,
            p.get('text', ''),
            p.get('size', default_size),
            bold=p.get('bold', False),
            color=p.get('color', default_color),
        )
    return box


def add_shape(slide, shape_type, left_px, top_px, width_px, height_px, *,
              fill=None, line=None, line_width_pt=None, fill_alpha=None):
    """図形を追加（条項7: 図形には文字を入れない）

    Parameters
    ----------
    fill_alpha : int or None
        塗りつぶしの不透明度（0-100000、None=不透明）。
        例: 95000=95%不透明=5%透過、80000=80%不透明=20%透過。
        design_system.md §3.2 の透明度バリエーション実装に使用。
    """
    shape = slide.shapes.add_shape(
        shape_type, px(left_px), px(top_px), px(width_px), px(height_px)
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        # 透明度指定がある場合は alpha 属性を OOXML に直接書き込む
        if fill_alpha is not None:
            _set_shape_fill_alpha(shape, fill_alpha)
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    # 図形内のテキストは空に
    if shape.has_text_frame:
        shape.text_frame.text = ''
    shape.shadow.inherit = False
    return shape


def _set_shape_fill_alpha(shape, alpha):
    """図形の塗りつぶし色に透明度（alpha）を設定する。

    python-pptx は標準で fill の alpha 属性をサポートしないため、
    OOXML を直接操作して <a:srgbClr><a:alpha val=".."/></a:srgbClr>
    を組み立てる。

    Parameters
    ----------
    alpha : int
        不透明度（0=完全透過、100000=完全不透明）。
        例: 95000 = 95%不透明 = 5%透過
            80000 = 80%不透明 = 20%透過
            12000 = 12%不透明 = 88%透過（薄帯背景）
    """
    spPr = shape.fill._xPr.find(qn('a:solidFill')) if False else None
    # 上記は型ヒント目的。実体は以下：
    solidFill = shape.fill._xPr.find(qn('a:solidFill'))
    if solidFill is None:
        return  # solid 塗りでない場合はスキップ
    srgbClr = solidFill.find(qn('a:srgbClr'))
    if srgbClr is None:
        return
    # 既存の alpha タグを削除してから追加（重複防止）
    for el in srgbClr.findall(qn('a:alpha')):
        srgbClr.remove(el)
    alpha_el = srgbClr.makeelement(qn('a:alpha'), {'val': str(int(alpha))})
    srgbClr.append(alpha_el)


def _add_bg_frame(slide, left_px, top_px, width_px, height_px, *,
                  fill=None, fill_alpha=None, line=None, line_width_pt=None,
                  radius_px=8):
    """カード背景フレーム（角丸 + 任意で透明度フィル + 任意で罫線）。

    design_system.md §4.3 カード規定準拠：
    - 角丸標準 8px（radius_px で変更可）
    - フラットデザイン（影なし）
    - 罫線 1pt / BORDER_GRAY を推奨

    ビジュアル強化（フェーズ1.5）で典型的に使うパターン：

        # NAVY 5%透過カード薄塗り
        _add_bg_frame(slide, 40, 90, 1200, 200,
                      fill=NAVY, fill_alpha=5000,
                      line=BORDER_GRAY, line_width_pt=1)

        # NAVY 18%透過カード強塗り
        _add_bg_frame(slide, 40, 90, 1200, 200,
                      fill=NAVY, fill_alpha=18000)

        # 不透明白カード（標準）
        _add_bg_frame(slide, 40, 90, 1200, 200,
                      fill=WHITE, line=BORDER_GRAY, line_width_pt=1)
    """
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      left_px, top_px, width_px, height_px,
                      fill=fill, fill_alpha=fill_alpha,
                      line=line, line_width_pt=line_width_pt)
    # 角丸半径を radius_px に合わせる（PowerPoint は短辺比で 0-0.5）
    short_side = min(width_px, height_px)
    if short_side > 0:
        ratio = max(0.0, min(0.5, radius_px / short_side))
        try:
            shape.adjustments[0] = ratio
        except (IndexError, AttributeError):
            pass
    return shape


# =====================================================================
# 共通要素：ヘッダ・フッター
# =====================================================================
def _add_header(slide, title, sub_label_en=''):
    """ヘッダ帯（高さ60px）を追加"""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W_PX, 60, fill=NAVY)
    add_text(slide, 40, 18, 800, title, 22, bold=True, color=WHITE)
    if sub_label_en:
        # サブラベルは右寄せ。長い英文（'Improvement Proposals' 等）が
        # 折り返さないよう、幅を400pxに広げる（左端= CANVAS_W_PX - 40 - 400）
        add_text(slide, CANVAS_W_PX - 40 - 400, 22, 400,
                 sub_label_en, 14,
                 color=NAVY_LIGHT, align=PP_ALIGN.RIGHT,
                 height_px=24)


def _add_footer(slide, page_num, total, author='紺＆クリーン スライド作成'):
    """フッター帯（ロゴ＋ページ番号のみ・条項8準拠）"""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 660, CANVAS_W_PX, 1, fill=GRAY_BORDER)
    add_text(slide, 40, 682, 400, author, 14, bold=True, color=NAVY)
    add_text(slide, 1140, 685, 120, f'{page_num} / {total}', 14,
             color=PAGE_NUM, align=PP_ALIGN.RIGHT)


def _add_lead(slide, text, top_px=90):
    """リード文（ヘッダ直下の導入文）"""
    add_text(slide, 40, top_px, 1200, text, 18, color=TEXT,
             height_px=60, line_height=1.7)


def _add_conclusion_band(slide, top_px, body_text, label='CONCLUSION',
                          color_label=NAVY_LIGHT, color_body=WHITE,
                          height_px=92, bg=NAVY, body_size=18):
    """結論帯（本文末尾の主張帯）

    本文の下端切れを防ぐため：
      - 帯の高さを十分（既定92px）に確保
      - 本文 textbox は ラベル分(32px)を引いた残量を高さに割当
      - body_text に '\\n' が含まれる場合は複数行段落として配置
    """
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 40, top_px, 1200, height_px,
              fill=bg)
    # ラベル（上部）
    add_text(slide, 60, top_px + 10, 1160, label, 14, bold=True,
             color=color_label, letter_spacing=2, height_px=20)
    # 本文（複数行対応）
    body_top = top_px + 34
    body_h = max(int(body_size * 1.5), height_px - 38)
    # body_text に改行が含まれていれば複数段落、そうでなければ単一行
    lines = [ln for ln in body_text.split('\n') if ln]
    if not lines:
        lines = [body_text]
    paragraphs = [{'text': line, 'size': body_size, 'bold': True,
                    'color': color_body, 'line_height': 1.35}
                   for line in lines]
    add_paragraph_box(slide, 60, body_top, 1160, paragraphs,
                       height_px=body_h, line_height=1.35)


# =====================================================================
# Presentation 初期化
# =====================================================================
def create_presentation():
    """1280×720 px キャンバスの空 Presentation を生成"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _blank_slide(prs):
    """白紙レイアウトでスライドを追加"""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


# =====================================================================
# Layout 1: 表紙（add_cover）
# =====================================================================
def add_cover(prs, title, date, author, subtitle='', page_total=1):
    """
    表紙スライド。

    Args:
        title: メインタイトル
        date: 日付（"2026年7月15日" など）
        author: 著者表記（"○○株式会社 営業企画部"）
        subtitle: サブタイトル（省略可）
    """
    slide = _blank_slide(prs)

    # ナビー背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W_PX, 660, fill=NAVY)

    # プロジェクトラベル
    add_text(slide, 80, 160, 900, 'PRESENTATION', 16, bold=True,
             color=NAVY_LIGHT, letter_spacing=3)

    # メインタイトル（複数行対応）
    title_lines = title.split('\n')
    title_h = 70 * len(title_lines)
    add_paragraph_box(slide, 80, 200, 1120,
        [{'text': line, 'size': 40, 'bold': True, 'color': WHITE,
          'line_height': 1.35} for line in title_lines],
        height_px=title_h)

    # 赤ライン
    add_shape(slide, MSO_SHAPE.RECTANGLE, 80, 200 + title_h + 20, 480, 3,
              fill=RED)

    # サブタイトル
    if subtitle:
        add_text(slide, 80, 200 + title_h + 40, 1120, subtitle, 20,
                 color=NAVY_E6, line_height=1.6, height_px=60)

    # 日付＋著者（右下右揃え・条項9・幅を広く）
    add_text(slide, 480, 560, 720, f'{date}　{author}', 20, bold=True,
             color=WHITE, align=PP_ALIGN.RIGHT, height_px=70, line_height=1.4)

    # 白フッター背景（表紙はページ番号を省略）
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 661, CANVAS_W_PX, 59, fill=WHITE)
    add_text(slide, 40, 682, 400, '紺＆クリーン スライド作成', 14, bold=True,
             color=NAVY)

    return slide


# =====================================================================
# Layout 2: アジェンダ（add_agenda）
# =====================================================================
def add_agenda(prs, items, page_num=2, total=10,
               title='Agenda', lead='本資料の構成は以下のとおりです。'):
    """
    アジェンダスライド。

    Args:
        items: list of str or dict
            str: 章タイトルのみ
            dict: {'title': str, 'desc': str}
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Agenda')
    _add_lead(slide, lead)

    # 項目を縦に並べる（最大8項目まで）
    n = len(items)
    top_start = 170
    available_h = 460
    item_h = max(60, min(80, available_h // max(n, 1)))
    num_size = 32 if n >= 7 else 36
    title_size = 20 if n >= 7 else 22

    for i, item in enumerate(items):
        if isinstance(item, str):
            item = {'title': item, 'desc': ''}
        y = top_start + i * item_h

        # 番号（紺・大きく）
        add_text(slide, 40, y, 80, f'{i+1:02d}', num_size, bold=True,
                 color=NAVY, line_height=1.0, height_px=num_size + 4)

        # タイトル
        add_text(slide, 130, y + 4, 600, item['title'], title_size, bold=True,
                 color=NAVY, height_px=title_size + 8)

        # 説明（任意、タイトルの右側に配置して高さを節約）
        if item.get('desc'):
            add_text(slide, 740, y + 10, 500, item['desc'], 14,
                     color=SUB_TEXT, line_height=1.5, height_px=20)

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 3: 課題整理（add_issue_summary）
# =====================================================================
def add_issue_summary(prs, title, cards, page_num=3, total=10,
                      lead='現状の主な課題を整理します。',
                      conclusion=None):
    """
    課題整理スライド。3カード横並び。

    Args:
        cards: list of dict, 各 dict は {'no': str, 'heading': str, 'body': str}
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Issues')
    _add_lead(slide, lead)

    n = len(cards)
    card_w = (1200 - 15 * (n - 1)) // n  # カード間 15px
    card_top = 160
    card_h = 390 if conclusion else 460

    # ヘッダ帯（紺）の高さ
    HEADER_BAND_H = 88

    for i, card in enumerate(cards):
        left = 40 + i * (card_w + 15)
        # ---- カード背景（白＋薄ボーダー） ----
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  left, card_top, card_w, card_h,
                  fill=WHITE, line=BORDER_GRAY, line_width_pt=1)

        # ---- 上端ヘッダ帯（紺ベタ・大型ナンバー） ----
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  left + 1, card_top + 1,
                  card_w - 2, HEADER_BAND_H, fill=NAVY)
        # ナンバー（左端、大型）
        no_text = card.get('no', f'{i + 1:02d}')
        add_text(slide, left + 22, card_top + 14, 90, no_text,
                 40, bold=True, color=WHITE, line_height=1.0,
                 height_px=52)
        # 縦罫線（ナンバーと見出しの区切り）
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  left + 110, card_top + 22, 1, HEADER_BAND_H - 44,
                  fill=NAVY_LIGHT)
        # 見出し（ヘッダ帯内、白文字）
        add_text(slide, left + 124, card_top + 20,
                 card_w - 140, card.get('heading', ''),
                 18, bold=True, color=WHITE,
                 height_px=HEADER_BAND_H - 32, line_height=1.35)

        # ---- 本文領域（カード白地、縦中央寄せ） ----
        body_top = card_top + HEADER_BAND_H + 18
        body_bottom = card_top + card_h - 18
        body_h = body_bottom - body_top
        body = card.get('body', '')
        body_paragraphs = [
            {'text': '● ' + line.strip(), 'size': 15,
             'color': TEXT, 'line_height': 1.6}
            for line in body.split('\n') if line.strip()
        ]
        add_paragraph_box(slide, left + 22, body_top, card_w - 40,
                           body_paragraphs, height_px=body_h,
                           line_height=1.6)

    if conclusion:
        _add_conclusion_band(slide, 560, conclusion, label='POINT')

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 4: 優先度マトリクス（add_priority_matrix）
# =====================================================================
def add_priority_matrix(prs, title, items, page_num=4, total=10,
                        x_label_low='低', x_label_high='高',
                        y_label_low='低', y_label_high='高',
                        x_axis_name='重要度', y_axis_name='緊急度',
                        lead='4象限で施策を整理します。'):
    """
    優先度マトリクス。

    Args:
        items: list of dict, 各 {'label': str, 'x': float, 'y': float, 'priority': 'S'|'A'|'B'|'C'}
               x, y は 0.0〜1.0 で象限内の相対位置
               priority: S/A=赤, B/C=紺
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Priority Matrix')
    _add_lead(slide, lead)

    # マトリクス領域 (left=220, top=190, w=880, h=420)
    # リード文(y=90〜150)・Y軸名(mx_t-36)が重ならないよう mx_t を下にずらす
    mx_l, mx_t, mx_w, mx_h = 220, 190, 880, 420

    # 4象限の背景（右上が最優先）
    qw, qh = mx_w // 2, mx_h // 2
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l + qw, mx_t, qw, qh,
              fill=RGBColor(0xFF, 0xE8, 0xE8))  # 右上 = 最優先（薄赤）
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l, mx_t, qw, qh,
              fill=LIGHT_GRAY)  # 左上
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l, mx_t + qh, qw, qh,
              fill=STRIPE)  # 左下
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l + qw, mx_t + qh, qw, qh,
              fill=LIGHT_GRAY)  # 右下

    # 枠線
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l, mx_t, mx_w, mx_h,
              fill=None, line=BORDER_GRAY, line_width_pt=1)
    # 十字線
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l + qw, mx_t, 1, mx_h,
              fill=BORDER_GRAY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l, mx_t + qh, mx_w, 1,
              fill=BORDER_GRAY)

    # 軸ラベル（高/低）
    add_text(slide, mx_l - 60, mx_t - 8, 60, y_label_high, 14, bold=True,
             color=NAVY, align=PP_ALIGN.RIGHT, height_px=22)
    add_text(slide, mx_l - 60, mx_t + mx_h - 20, 60, y_label_low, 14,
             bold=True, color=NAVY, align=PP_ALIGN.RIGHT, height_px=22)
    add_text(slide, mx_l - 8, mx_t + mx_h + 8, 60, x_label_low, 14,
             bold=True, color=NAVY, height_px=22)
    add_text(slide, mx_l + mx_w - 30, mx_t + mx_h + 8, 60, x_label_high, 14,
             bold=True, color=NAVY, align=PP_ALIGN.RIGHT, height_px=22)

    # 軸の名称（X軸名は下に、十分な幅でセンタリング）
    x_axis_text = f'{x_axis_name} →'
    add_text(slide, mx_l, mx_t + mx_h + 36, mx_w,
             x_axis_text, 14, color=SUB_TEXT,
             align=PP_ALIGN.CENTER, height_px=22)
    # Y軸名：マトリクス左外、横長領域に十分な幅で表示（折返し防止）
    # 14pt全角約19px、全角6文字「↑ インパクト」→ 114px必要 → 幅180pxで余裕
    add_text(slide, mx_l - 200, mx_t + mx_h // 2 - 12, 190,
             f'↑ {y_axis_name}', 14, color=SUB_TEXT,
             align=PP_ALIGN.RIGHT, height_px=22)

    # プロット
    # ラベル領域の最大幅を計算（マトリクス右端を超えないよう動的調整）
    for i, it in enumerate(items):
        cx = mx_l + int(it['x'] * mx_w)
        cy = mx_t + int((1 - it['y']) * mx_h)
        d = 32
        priority = it.get('priority', 'B')
        fill = RED if priority in ('S', 'A') else NAVY
        add_shape(slide, MSO_SHAPE.OVAL, cx - d // 2, cy - d // 2, d, d,
                  fill=fill)
        # 番号（円の中・白）
        add_text(slide, cx - d // 2, cy - 14, d, f'{i+1:02d}', 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, height_px=22)
        # ラベル位置の自動判定：
        #  - 円の右側に十分な余白があれば右配置
        #  - 右余白が少ない（マトリクス右端まで180px未満）なら左配置
        label_w = 200
        label_left_default = cx + d // 2 + 6
        # マトリクス右端 mx_l + mx_w を超えないか判定
        if label_left_default + label_w > mx_l + mx_w:
            # 円の左側にラベル配置（右寄せ）
            label_left = cx - d // 2 - 6 - label_w
            label_align = PP_ALIGN.RIGHT
        else:
            label_left = label_left_default
            label_align = PP_ALIGN.LEFT
        add_text(slide, label_left, cy - 12, label_w, it['label'], 14,
                 bold=True, color=TEXT, align=label_align, height_px=22)

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 5: OK・NG例の対比（add_ok_ng_pair）
# =====================================================================
def add_ok_ng_pair(prs, title, ng, ok, page_num=5, total=10,
                   lead='良い例と悪い例を対比します。',
                   conclusion=None):
    """
    OK/NG対比カード（スタイリッシュ版）。

    レイアウト：
      ┌─────────────────────┐
      │ ▌NG ✕  見出し       │ ← 上端帯（色帯）
      │ ─────────────────── │
      │ 本文テキスト         │ ← 中央領域
      │                     │
      │ ─────────────────── │
      │ ⚠ キャプション帯    │ ← 下端帯（淡色）
      └─────────────────────┘

    Args:
        ng: dict {'heading': str, 'body': str, 'caption': str}
        ok: dict {'heading': str, 'body': str, 'caption': str}
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'OK / NG')
    _add_lead(slide, lead)

    col_w = 580
    col_h = 380 if conclusion else 460
    card_top = 160

    # 視覚的に間延びを防ぐため、上端ヘッダ帯／中央本文／下端キャプション帯を
    # それぞれ独立した「帯」として配置する
    HEADER_BAND_H = 56   # 上端帯（バッジ+見出し）
    CAPTION_BAND_H = 44  # 下端帯（一言サマリ）

    for i, (col, badge_text, badge_color, label_color, accent_bg) in enumerate(
        [(ng, 'NG ✕', RGBColor(0xDC, 0x35, 0x45), RED,
          RGBColor(0xFD, 0xEC, 0xEC)),                  # NG淡赤帯
         (ok, 'OK ✓', RGBColor(0x10, 0x88, 0x4A),
          RGBColor(0x10, 0x88, 0x4A),
          RGBColor(0xE6, 0xF4, 0xEA))]                  # OK淡緑帯
    ):
        left = 40 + i * (col_w + 40)
        # ---- カード背景（薄いボーダー） ----
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top,
                  col_w, col_h, fill=WHITE,
                  line=BORDER_GRAY, line_width_pt=1)

        # ---- 上端ヘッダ帯（色付き） ----
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + 1, card_top + 1,
                  col_w - 2, HEADER_BAND_H,
                  fill=accent_bg, line=accent_bg)
        # 色帯（左4px、強アクセント）
        add_shape(slide, MSO_SHAPE.RECTANGLE, left, card_top,
                  6, HEADER_BAND_H + 1, fill=badge_color)
        # バッジ
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  left + 20, card_top + 12, 96, 32,
                  fill=badge_color)
        add_text(slide, left + 20, card_top + 16, 96, badge_text, 16,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=24)
        # 見出し（バッジの右）
        add_text(slide, left + 130, card_top + 16, col_w - 150,
                 col.get('heading', ''), 18, bold=True,
                 color=label_color, height_px=28)

        # ---- 中央本文領域（縦中央寄せ） ----
        body_top = card_top + HEADER_BAND_H + 16
        body_bottom = card_top + col_h - CAPTION_BAND_H - 10
        body_h = body_bottom - body_top
        body = col.get('body', '')
        body_paragraphs = [{'text': line, 'size': 16, 'color': TEXT,
                            'line_height': 1.6}
                            for line in body.split('\n') if line.strip()]
        if body_paragraphs:
            add_paragraph_box(slide, left + 24, body_top, col_w - 48,
                               body_paragraphs, height_px=body_h,
                               line_height=1.6)

        # ---- 下端キャプション帯（淡色背景＋一言サマリ） ----
        if col.get('caption'):
            caption_top = card_top + col_h - CAPTION_BAND_H
            # 帯背景
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      left + 1, caption_top, col_w - 2, CAPTION_BAND_H - 1,
                      fill=accent_bg, line=accent_bg)
            # 左マーカー（▶）
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      left + 1, caption_top, 4, CAPTION_BAND_H - 1,
                      fill=badge_color)
            add_text(slide, left + 20, caption_top + 11,
                     col_w - 40, col['caption'], 14, bold=True,
                     color=label_color, height_px=22)

    if conclusion:
        _add_conclusion_band(slide, 558, conclusion, label='POINT')

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 6: 施策一覧表（add_action_table）
# =====================================================================
def add_action_table(prs, title, columns, rows, page_num=6, total=10,
                     lead='主要な施策を一覧化します。',
                     source=None, emphasize=None,
                     col_widths=None):
    """
    表組スライド（条項4準拠）。

    Args:
        columns: list of str（ヘッダ）
        rows: list of list（各行のセル値）
        emphasize: list of (row_idx, col_idx, original_text, emphasis_text)
                   セル内で original_text を emphasis_text に置換し、赤で大きく表示
        col_widths: list of int（各列の幅%。合計100。Noneなら均等）
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Action List')
    _add_lead(slide, lead)

    n_cols = len(columns)
    n_rows = len(rows) + 1  # ヘッダ含む

    table_left, table_top = 40, 160
    table_w = 1200
    row_h = 38

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, px(table_left), px(table_top),
        px(table_w), px(row_h * n_rows)
    )
    table = table_shape.table

    # 列幅
    if col_widths is None:
        col_widths = [100 // n_cols] * n_cols
    for i, w in enumerate(col_widths):
        table.columns[i].width = px(int(table_w * w / 100))

    # ヘッダ行
    for c, col_name in enumerate(columns):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ''
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        set_run(run, col_name, 16, bold=True, color=WHITE)
        cell.margin_left = Pt(8)
        cell.margin_right = Pt(8)
        cell.margin_top = Pt(6)
        cell.margin_bottom = Pt(6)

    # データ行
    for r, row in enumerate(rows, start=1):
        bg = WHITE if r % 2 == 1 else STRIPE
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.text = ''
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c > 0 else PP_ALIGN.CENTER

            # 強調指定があれば run を分割
            emphasis_for_cell = None
            if emphasize:
                for em in emphasize:
                    if em[0] == r - 1 and em[1] == c:
                        emphasis_for_cell = em
                        break

            if emphasis_for_cell:
                _, _, original, target = emphasis_for_cell
                parts = str(val).split(target)
                for k, part in enumerate(parts):
                    if part:
                        run = para.add_run()
                        set_run(run, part, 14, color=TEXT)
                    if k < len(parts) - 1:
                        emp_run = para.add_run()
                        set_run(emp_run, target, 18, bold=True, color=RED)
            else:
                run = para.add_run()
                set_run(run, str(val), 14, color=TEXT)

            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(6)
            cell.margin_bottom = Pt(6)

    # 表終端の y を実測（条項11準拠）
    table_end_y = table_top + row_h * n_rows

    if source:
        # 出所注記（フッター帯の上、表終端 + 30px）
        add_text(slide, 40, table_end_y + 30, 1200, source, 14,
                 color=SUB_TEXT, line_height=1.5)

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 7: KPI カード（add_kpi_card）
# =====================================================================
def add_kpi_card(prs, title, kpis, page_num=7, total=10,
                 lead='主要KPIの進捗を確認します。',
                 conclusion=None):
    """
    KPIカード（3〜4個横並び）。

    Args:
        kpis: list of dict
              {'label': str, 'value': str, 'unit': str, 'desc': str, 'color': 'red'|'navy'}
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Key Metrics')
    _add_lead(slide, lead)

    n = len(kpis)
    gap = 20
    card_w = (1200 - gap * (n - 1)) // n
    card_h = 280
    top = 160

    for i, kpi in enumerate(kpis):
        left = 40 + i * (card_w + gap)
        is_red = kpi.get('color', 'red') == 'red'
        value_color = RED if is_red else NAVY

        # カード（白＋紺枠）
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  left, top, card_w, card_h,
                  fill=WHITE, line=NAVY, line_width_pt=2)
        # ラベル（上部）
        add_text(slide, left + 24, top + 26, card_w - 48,
                 kpi.get('label', ''), 14, bold=True, color=NAVY,
                 letter_spacing=2)
        # 大きな数値（条項13: サイズ拡大＋色変更）
        value_size = 48 if len(kpi.get('value', '')) <= 6 else 36
        add_text(slide, left + 24, top + 80, card_w - 48,
                 kpi.get('value', '') + kpi.get('unit', ''),
                 value_size, bold=True, color=value_color,
                 line_height=1.0, height_px=70)
        # 説明（複数行対応：改行 \n を段落として保持）
        desc = kpi.get('desc', '')
        desc_paragraphs = [
            {'text': line.strip(), 'size': 14, 'color': SUB_TEXT,
             'line_height': 1.55}
            for line in desc.split('\n') if line.strip()
        ] or [{'text': '', 'size': 14, 'color': SUB_TEXT,
               'line_height': 1.55}]
        add_paragraph_box(slide, left + 24, top + 170, card_w - 48,
                           desc_paragraphs, height_px=90,
                           line_height=1.55)

    if conclusion:
        _add_conclusion_band(slide, 470, conclusion, label='SUMMARY')

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 8: スケジュール／ガント風（add_schedule_gantt）
# =====================================================================
def add_schedule_gantt(prs, title, months, tasks, page_num=8, total=10,
                       lead='実行スケジュールは以下のとおりです。'):
    """
    ガント風スケジュール。

    Args:
        months: list of str（例: ['4月','5月','6月','7月','8月','9月']）
        tasks: list of dict {'name': str, 'start': int, 'end': int, 'milestone': bool}
               start/end は months のインデックス（0始まり、end は inclusive）
               milestone=True ならバーではなく◆マーク
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Schedule')
    _add_lead(slide, lead)

    chart_left = 280
    chart_top = 170
    chart_w = 960
    n_months = len(months)
    col_w = chart_w // n_months

    # 月ヘッダ
    for i, m in enumerate(months):
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  chart_left + i * col_w, chart_top, col_w, 36,
                  fill=NAVY)
        add_text(slide, chart_left + i * col_w, chart_top + 8, col_w, m,
                 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # タスク行
    n_tasks = len(tasks)
    row_h = min(48, (560 - chart_top - 50) // max(n_tasks, 1))

    for t_idx, t in enumerate(tasks):
        y = chart_top + 40 + t_idx * row_h
        # タスク名
        add_text(slide, 40, y + 8, 230, t['name'], 16, bold=True,
                 color=TEXT, height_px=row_h - 4)
        # 列の縞模様
        for i in range(n_months):
            stripe_bg = STRIPE if t_idx % 2 == 1 else WHITE
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      chart_left + i * col_w, y, col_w, row_h - 4,
                      fill=stripe_bg)
        # バー or マイルストーン
        if t.get('milestone'):
            mx = chart_left + t['start'] * col_w + col_w // 2 - 12
            my = y + (row_h - 4) // 2 - 12
            add_shape(slide, MSO_SHAPE.DIAMOND, mx, my, 24, 24,
                      fill=RED)
        else:
            bar_left = chart_left + t['start'] * col_w + 4
            bar_right = chart_left + (t['end'] + 1) * col_w - 4
            bar_w = bar_right - bar_left
            bar_h = row_h - 18
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      bar_left, y + 6, bar_w, bar_h,
                      fill=NAVY)

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 9: 運用フロー比較（add_flow_compare）
# =====================================================================
def add_flow_compare(prs, title, before_steps, after_steps,
                     page_num=9, total=10,
                     lead='現状フローと改善後フローを対比します。',
                     conclusion=None):
    """
    Before/After フロー比較。

    Args:
        before_steps: list of str（現状のステップ名）
        after_steps: list of str（改善後のステップ名）
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Flow Compare')
    _add_lead(slide, lead)

    def _draw_flow(top, label, steps, badge_color, card_fill, card_line):
        # バッジ
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 40, top + 20, 100, 36,
                  fill=badge_color)
        add_text(slide, 40, top + 26, 100, label, 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)

        # ステップカード（横並び）
        n = len(steps)
        avail_w = 1080
        card_w = min(180, (avail_w - 30 * (n - 1)) // n)
        gap = (avail_w - card_w * n) // max(n - 1, 1) if n > 1 else 0
        start_x = 160

        for i, step in enumerate(steps):
            x = start_x + i * (card_w + gap)
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_w, 76,
                      fill=card_fill, line=card_line, line_width_pt=2)
            add_text(slide, x + 10, top + 22, card_w - 20, step, 15,
                     bold=True, color=TEXT, align=PP_ALIGN.CENTER,
                     height_px=40, line_height=1.4)
            # 矢印
            if i < n - 1:
                arrow_x = x + card_w + gap // 2 - 10
                add_text(slide, arrow_x, top + 28, 20, '▸', 22,
                         bold=True, color=RED, align=PP_ALIGN.CENTER)

    _draw_flow(170, 'BEFORE', before_steps,
               RGBColor(0x88, 0x88, 0x88), LIGHT_GRAY, BORDER_GRAY)
    _draw_flow(320, 'AFTER', after_steps, NAVY, WHITE, NAVY)

    if conclusion:
        _add_conclusion_band(slide, 470, conclusion, label='IMPROVEMENT')

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 10: クロージング（add_closing）
# =====================================================================
def add_closing(prs, message='Thank you.', next_step='',
                contact='', page_num=10, total=10):
    """
    クロージングスライド。

    Args:
        message: 大きく表示するメッセージ
        next_step: 次のアクション説明
        contact: 問い合わせ先
    """
    slide = _blank_slide(prs)

    # ナビー背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W_PX, 660, fill=NAVY)

    # メッセージ
    add_text(slide, 80, 200, 1120, message, 56, bold=True, color=WHITE,
             line_height=1.2, height_px=80)

    # 赤ライン
    add_shape(slide, MSO_SHAPE.RECTANGLE, 80, 300, 480, 3, fill=RED)

    # ネクストステップ
    if next_step:
        add_text(slide, 80, 340, 1120, next_step, 22, color=NAVY_E6,
                 line_height=1.6, height_px=120)

    # 問い合わせ先（右下右揃え）
    if contact:
        add_text(slide, 740, 560, 460, contact, 16, color=WHITE,
                 align=PP_ALIGN.RIGHT, line_height=1.6, height_px=60)

    # 白フッター背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 661, CANVAS_W_PX, 59, fill=WHITE)
    add_text(slide, 40, 682, 400, '紺＆クリーン スライド作成', 14, bold=True,
             color=NAVY)
    add_text(slide, 1140, 685, 120, f'{page_num} / {total}', 14,
             color=PAGE_NUM, align=PP_ALIGN.RIGHT)

    return slide


# =====================================================================
# C-1 / Layout 11: UI診断スコアカード概要（add_scorecard_overview）
# =====================================================================
def add_scorecard_overview(prs, diagnosis, page_num=1, total=3,
                            author='UI/UX診断 by GPTs'):
    """
    UI診断スコアカード スライド1（概要）。

    総合スコア（44〜56pt）＋ランク＋サービス情報＋強み（箇条書き）＋
    最優先課題（箇条書き）を1枚に集約。

    Args:
        diagnosis: dict
            - service_name (str): 診断対象サービス名
            - url (str): URL
            - input_type (str): "URL入力" / "スクショ入力"
            - total_score (int): 50点満点
            - rank (str): S/A/B/C/D
            - rank_label (str): "優秀" など
            - strengths (list of str): 2〜4項目
            - priority_issues (list of str): 1〜3項目
    """
    slide = _blank_slide(prs)
    _add_header(slide, 'UI診断スコアカード', 'UI Diagnosis Scorecard')

    # ── 上段左：サービス情報（left=40, top=88, width=720）
    add_text(slide, 40, 90, 720, '診断対象', 14, bold=True,
             color=NAVY, letter_spacing=2)
    add_text(slide, 40, 114, 720, diagnosis.get('service_name', ''), 22,
             bold=True, color=TEXT, height_px=34)
    # URL とinput_type を1行で
    url = diagnosis.get('url', '')
    input_type = diagnosis.get('input_type', '')
    meta = f'{url}　／　{input_type}' if url else input_type
    add_text(slide, 40, 150, 720, meta, 14, color=SUB_TEXT, height_px=22)

    # ── 上段右：総合スコア＋ランクカード（left=800, top=88, width=440, height=130）
    score = diagnosis.get('total_score', 0)
    rank = diagnosis.get('rank', '-')
    rank_label = diagnosis.get('rank_label', '')
    # ランクに応じた色
    rank_color = NAVY if rank in ('S', 'A', 'B') else RED

    # スコアカード背景
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 800, 88, 440, 130,
              fill=LIGHT_GRAY)
    # 左帯
    add_shape(slide, MSO_SHAPE.RECTANGLE, 800, 88, 6, 130, fill=rank_color)

    # ラベル
    add_text(slide, 822, 100, 200, 'TOTAL SCORE', 14, bold=True,
             color=NAVY, letter_spacing=2)
    # 総合スコア（条項13：核心数値、48ptで強調）
    add_text(slide, 822, 124, 280,
             f'{score}', 48, bold=True, color=rank_color,
             line_height=1.0, height_px=64)
    # 「/50」を小さく
    add_text(slide, 822 + 100, 162, 80, '／ 50', 18, bold=True,
             color=SUB_TEXT, height_px=24)

    # ランクバッジ（右半分）
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1060, 110, 160, 86,
              fill=rank_color)
    add_text(slide, 1060, 118, 160, 'RANK', 14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, letter_spacing=2)
    add_text(slide, 1060, 138, 160, rank, 40, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, line_height=1.0,
             height_px=52)
    if rank_label:
        add_text(slide, 800, 226, 440, rank_label, 16, color=SUB_TEXT,
                 align=PP_ALIGN.CENTER)

    # ── 下段左：強み（left=40, top=260, width=590）
    # 見出し帯（紺ベタ＋白文字）でスタイリッシュに
    SECTION_TOP = 260
    SECTION_W = 590
    LABEL_H = 36
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, SECTION_TOP, SECTION_W,
              LABEL_H, fill=NAVY)
    add_text(slide, 56, SECTION_TOP + 8, SECTION_W - 32,
             '◎ 強み （Strengths）', 16, bold=True, color=WHITE,
             height_px=24, letter_spacing=2)

    # 各項目をカード化して並べる（最大3項目想定）
    # カード幅 590px、本文width = 590 - 14(左padding) - 28(badge) - 8(gap) - 14(右padding) = 526px
    # 14pt全角約19px → 1行=27.6全角文字 → 30文字までの本文を1行表示できる
    strengths = diagnosis.get('strengths', [])[:3]
    item_top = SECTION_TOP + LABEL_H + 10
    item_h = 104
    item_gap = 10
    BADGE_D = 28  # バッジ直径（小さくして本文領域を確保）
    BADGE_X_OFFSET = 14
    BODY_LEFT_OFFSET = BADGE_X_OFFSET + BADGE_D + 10  # 52px
    BODY_RIGHT_PADDING = 14

    def _draw_item(left_x, y, idx, s, badge_color):
        # カード背景
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_x, y,
                  SECTION_W, item_h, fill=LIGHT_GRAY)
        # 番号バッジ（サークル、縦中央）
        bx = left_x + BADGE_X_OFFSET
        by = y + (item_h - BADGE_D) // 2
        add_shape(slide, MSO_SHAPE.OVAL, bx, by, BADGE_D, BADGE_D,
                  fill=badge_color)
        add_text(slide, bx, by + 4, BADGE_D, f'{idx + 1:02d}', 14,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=22)
        # 本文（縦中央寄せ）
        body_left = left_x + BODY_LEFT_OFFSET
        body_w = SECTION_W - BODY_LEFT_OFFSET - BODY_RIGHT_PADDING
        add_paragraph_box(slide, body_left, y + 18, body_w,
            [{'text': s, 'size': 14, 'color': TEXT, 'line_height': 1.6}],
            height_px=item_h - 32, line_height=1.6)

    for idx, s in enumerate(strengths):
        y = item_top + idx * (item_h + item_gap)
        _draw_item(40, y, idx, s, NAVY)

    # ── 下段右：最優先課題（left=650, top=260, width=590）
    add_shape(slide, MSO_SHAPE.RECTANGLE, 650, SECTION_TOP, SECTION_W,
              LABEL_H, fill=RED)
    add_text(slide, 666, SECTION_TOP + 8, SECTION_W - 32,
             '⚠ 最優先課題 （Priority Issues）', 16, bold=True,
             color=WHITE, height_px=24, letter_spacing=2)

    issues = diagnosis.get('priority_issues', [])[:3]
    for idx, s in enumerate(issues):
        y = item_top + idx * (item_h + item_gap)
        _draw_item(650, y, idx, s, RED)

    _add_footer(slide, page_num, total, author=author)
    return slide


# =====================================================================
# C-1 / Layout 12: UI診断 10項目スコア表（add_scorecard_table）
# =====================================================================
def add_scorecard_table(prs, diagnosis, page_num=2, total=3,
                         author='UI/UX診断 by GPTs',
                         high_threshold=0.75, low_threshold=0.50):
    """
    UI診断スコアカード スライド2（10項目スコア一覧）。

    スコア列のみ色分け：
      - スコア / max >= high_threshold → 紺（高評価）
      - スコア / max <= low_threshold → 赤（低評価）
      - その他 → 通常テキスト色

    Args:
        diagnosis: scores キーに list of {category, score, max, comment}
        high_threshold: 高評価とみなす比率（既定 0.75）
        low_threshold: 低評価とみなす比率（既定 0.50）
    """
    slide = _blank_slide(prs)
    _add_header(slide, 'UI診断スコアカード', '10-Item Score Detail')
    _add_lead(slide, f'{diagnosis.get("service_name", "")} の10項目評価詳細です。'
                       'スコア列は高評価=紺、低評価=赤で色分けしています。')

    scores = diagnosis.get('scores', [])
    columns = ['No', '評価項目', 'スコア', 'コメント']
    col_widths_pct = [8, 28, 14, 50]

    # 凡例帯（高さ36px）を表の下に確実に置くため、表の縦サイズを管理
    table_left, table_top = 40, 160
    table_w = 1200
    n_rows = len(scores) + 1  # ヘッダ含む
    # フッター(660)・余白(12px)・凡例帯(36px)・余白(12px) を引いた領域を使用
    # 表下端 = 660 - 12 - 36 - 12 = 600
    table_max_bottom = 600
    table_avail_h = table_max_bottom - table_top
    # ヘッダ行はやや高め、データ行は均等
    header_h = 38
    data_avail = table_avail_h - header_h
    row_h = max(28, min(36, data_avail // max(1, len(scores))))

    table_shape = slide.shapes.add_table(
        n_rows, len(columns),
        px(table_left), px(table_top),
        px(table_w), px(header_h + row_h * len(scores))
    )
    table = table_shape.table
    # 行高を明示設定（PowerPoint側の自動拡大を抑制）
    table.rows[0].height = px(header_h)
    for r in range(1, n_rows):
        table.rows[r].height = px(row_h)

    # 列幅
    for i, w in enumerate(col_widths_pct):
        table.columns[i].width = px(int(table_w * w / 100))

    # ヘッダ行
    for c, col_name in enumerate(columns):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ''
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        set_run(run, col_name, 15, bold=True, color=WHITE)
        cell.margin_left = Pt(8)
        cell.margin_right = Pt(8)
        cell.margin_top = Pt(6)
        cell.margin_bottom = Pt(6)

    # データ行
    for r, item in enumerate(scores, start=1):
        bg = WHITE if r % 2 == 1 else STRIPE
        score_val = item.get('score', 0)
        max_val = item.get('max', 5)
        ratio = score_val / max_val if max_val else 0
        if ratio >= high_threshold:
            score_color = NAVY
            score_bold = True
        elif ratio <= low_threshold:
            score_color = RED
            score_bold = True
        else:
            score_color = TEXT
            score_bold = False

        # No
        _fill_cell(table.cell(r, 0), str(r), 14, color=TEXT,
                    bg=bg, align=PP_ALIGN.CENTER, margin_v=3, no_wrap=True)
        # 評価項目
        # [FB対応 2026-06-30] scores の `name` キーを優先（schema統一）
        _fill_cell(table.cell(r, 1),
                    item.get('name') or item.get('category', ''), 14,
                    bold=True, color=TEXT, bg=bg, align=PP_ALIGN.LEFT,
                    margin_v=3, no_wrap=True)
        # スコア（色分け＋大きめ）
        score_text = f'{score_val} / {max_val}'
        _fill_cell(table.cell(r, 2), score_text, 16, bold=score_bold,
                    color=score_color, bg=bg, align=PP_ALIGN.CENTER,
                    margin_v=3, no_wrap=True)
        # コメント（長文は省略して1行で確実に収める）
        comment_raw = item.get('comment', '')
        comment = comment_raw
        total_w = 0.0
        for idx, ch in enumerate(comment_raw):
            total_w += 0.5 if ord(ch) < 128 else 1.0
            if total_w > 28:
                comment = comment_raw[:idx].rstrip() + '…'
                break
        _fill_cell(table.cell(r, 3), comment, 14,
                    color=TEXT, bg=bg, align=PP_ALIGN.LEFT,
                    margin_v=3, no_wrap=True)

    # 凡例（表の下に色チップ付きで配置）
    table_end_y = table_top + header_h + row_h * len(scores)
    legend_y = table_end_y + 12
    # 凡例帯背景（淡グレー、高さ36px、14pt下限厳守）
    legend_h = 36
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
              40, legend_y, 1200, legend_h,
              fill=LIGHT_GRAY, line=BORDER_GRAY, line_width_pt=0.5)
    # 凡例ラベル「凡例」
    add_text(slide, 56, legend_y + 8, 60, '凡例', 14,
             bold=True, color=SUB_TEXT, height_px=22)
    # 紺チップ＋説明
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              120, legend_y + 12, 16, 12, fill=NAVY)
    add_text(slide, 142, legend_y + 8, 360,
             f'紺＝高評価（達成率 {int(high_threshold*100)}% 以上）',
             14, color=SUB_TEXT, height_px=22)
    # 赤チップ＋説明
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              520, legend_y + 12, 16, 12, fill=RED)
    add_text(slide, 542, legend_y + 8, 360,
             f'赤＝低評価（達成率 {int(low_threshold*100)}% 以下）',
             14, color=SUB_TEXT, height_px=22)
    # 通常チップ＋説明
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              920, legend_y + 12, 16, 12, fill=BORDER_GRAY)
    add_text(slide, 942, legend_y + 8, 260,
             'グレー＝中位', 14, color=SUB_TEXT, height_px=22)

    _add_footer(slide, page_num, total, author=author)
    return slide


def _truncate_full(s, max_w):
    """文字列を全角換算 max_w で切り、超えたら末尾を '…' に置換。

    全角=1.0、半角(ASCII)=0.5 でカウント。
    """
    total = 0.0
    for idx, ch in enumerate(s):
        total += 0.5 if ord(ch) < 128 else 1.0
        if total > max_w:
            return s[:idx].rstrip() + '…'
    return s


def _add_multi_run_box(slide, left_px, top_px, width_px, height_px,
                        paragraphs, *, line_height=None,
                        space_after_pt=None, space_before_pt=None,
                        vertical_anchor=None):
    """1テキストボックス内に複数段落＋段落内の複数runを配置するヘルパ。

    1コンテキスト1ブロック原則の核となる関数。
    タイトル+サブタイトル、ラベル+値、ヘッダ+本文など、論理的に
    1つの文脈に属するテキスト群を「1つのテキストボックス」に集約する。

    paragraphs: list of dict
      各dictは段落を表す。以下のキーをサポート：
        - text (str): 単一run。size, bold, color, line_height を併用可。
        - runs (list of dict): 段落内に複数run。各runは {text, size, bold, color}。
        - line_height (float): 段落の line_spacing（＝段落**内部**の折返し行間）
        - space_after_pt (float): 段落末尾に追加する余白（pt）。項目間の分離余白に使う。
        - space_before_pt (float): 段落先頭に追加する余白（pt）。
        - align (PP_ALIGN): 段落の整列

    Parameters
    ----------
    line_height : float, optional
        全段落に適用する既定の line_spacing（各段落dictで上書き可）。
        **段落内の折返し行間** を制御する。項目間の余白ではない。
    space_after_pt : float, optional
        全段落に適用する既定の段落末尾余白（pt）。**リスト項目間の分離感** に使う。
        [v15/2026-07-12] リスト表示（強み・課題Top3等）で
        line_height と項目間余白を分離制御するために新設。
    space_before_pt : float, optional
        全段落に適用する既定の段落先頭余白（pt）。通常は使わない。

    Notes
    -----
    - line_height と space_after_pt を分離する設計思想：
      折返し行間（line_height）は詰めて「1文章の連続感」を担保し、
      項目間余白（space_after_pt）で「別項目である分離感」を明示する。
      これによりリスト表示の可読性が根本改善する（v15対応）。
    """
    box = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    # [v15/2026-07-12] 縦中央寄せ（総評帯など「1行運用でボックス高さより
    # テキストが小さい」場合に、上下余白の不均衡を解消するため）
    if vertical_anchor is not None:
        try:
            tf.vertical_anchor = vertical_anchor
        except Exception:
            pass

    for i, p in enumerate(paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        lh = p.get('line_height', line_height)
        if lh:
            para.line_spacing = lh
        # [v15] 段落間余白の分離制御
        sa = p.get('space_after_pt', space_after_pt)
        if sa is not None and sa != 0:
            para.space_after = Pt(sa)
        sb = p.get('space_before_pt', space_before_pt)
        if sb is not None and sb != 0:
            para.space_before = Pt(sb)
        if 'runs' in p:
            for r_def in p['runs']:
                run = para.add_run()
                set_run(run, r_def.get('text', ''),
                        r_def.get('size', 14),
                        bold=r_def.get('bold', False),
                        color=r_def.get('color', TEXT))
        else:
            run = para.add_run()
            set_run(run, p.get('text', ''),
                    p.get('size', 14),
                    bold=p.get('bold', False),
                    color=p.get('color', TEXT))
    return box


def _fill_cell(cell, text, size_pt, *, bold=False, color=TEXT, bg=WHITE,
                align=PP_ALIGN.LEFT, margin_v=6, no_wrap=False):
    """テーブルセルを単一スタイルで埋めるヘルパー

    Args:
        margin_v: セル上下マージン(pt)。行高をタイトに保ちたいときは小さく。
        no_wrap: True で word_wrap を OFF（長文を1行で描画、はみ出しはクリップ）
    """
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.text = ''
    tf = cell.text_frame
    if no_wrap:
        tf.word_wrap = False
    # 行高の自動拡大を抑制（LibreOffice 描画対策）
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    # 縦中央寄せ（行高に余裕がない場合でも文字が枠に収まる）
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    para = tf.paragraphs[0]
    para.alignment = align
    # 段落 line_spacing を pt で明示固定（LibreOffice の行高拡大を抑制）
    para.line_spacing = Pt(size_pt + 2)
    para.space_before = Pt(0)
    para.space_after = Pt(0)
    run = para.add_run()
    set_run(run, str(text), size_pt, bold=bold, color=color)
    cell.margin_left = Pt(8)
    cell.margin_right = Pt(8)
    cell.margin_top = Pt(margin_v)
    cell.margin_bottom = Pt(margin_v)


# =====================================================================
# C-1 / Layout 13: UI診断 結論（add_scorecard_conclusion）
# =====================================================================
def add_scorecard_conclusion(prs, diagnosis, page_num=3, total=3,
                              author='UI/UX診断 by GPTs',
                              cta='詳細な改善提案は別途お送りします。'):
    """
    UI診断スコアカード スライド3（結論）。

    結論テキストを大きく表示＋CONCLUSION結論帯＋CTA。

    Args:
        diagnosis: conclusion キー（総評3〜5文）
        cta: 結論帯に表示する次のアクション文
    """
    slide = _blank_slide(prs)
    _add_header(slide, 'UI診断スコアカード', 'Conclusion')

    # ラベル
    add_text(slide, 40, 90, 1200, '総評', 14, bold=True,
             color=NAVY, letter_spacing=2, height_px=22)

    # サービス名 + ランク要約
    score = diagnosis.get('total_score', 0)
    rank = diagnosis.get('rank', '-')
    rank_label = diagnosis.get('rank_label', '')
    summary_line = (f'{diagnosis.get("service_name", "")} の総合評価は '
                     f'{score}/50点・ランク {rank}（{rank_label}）')
    add_text(slide, 40, 116, 1200, summary_line, 18, bold=True,
             color=TEXT, height_px=28)

    # 区切り赤ライン
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, 156, 200, 3, fill=RED)

    # 結論テキスト（18〜20pt）
    conclusion_text = diagnosis.get('conclusion', '')
    # 改行（句点で改行を入れて読みやすく）
    paragraphs = [p.strip() for p in conclusion_text.split('\n') if p.strip()]
    if not paragraphs:
        paragraphs = [conclusion_text]

    paragraph_specs = []
    for p in paragraphs:
        paragraph_specs.append({
            'text': p, 'size': 19, 'color': TEXT, 'line_height': 1.8,
        })

    add_paragraph_box(slide, 40, 180, 1200, paragraph_specs,
                      height_px=340, line_height=1.8)

    # 結論帯（CONCLUSION）
    _add_conclusion_band(slide, 540, cta, label='NEXT STEP')

    _add_footer(slide, page_num, total, author=author)
    return slide


# =====================================================================
# C-2 / Layout 14: 改善提案 一覧（add_proposal_overview）
# =====================================================================
def add_proposal_overview(prs, proposals_data, page_num=1, total=3,
                            author='UI/UX診断 by GPTs'):
    """
    改善提案 スライド1（一覧表）。

    add_action_table を流用したPPTX表組。
    優先度列は S/A=赤、B=紺 で色分け。

    Args:
        proposals_data: dict
            - service_name (str)
            - proposals (list of dict): {no, title, category, priority, ...}
    """
    slide = _blank_slide(prs)
    _add_header(slide, '改善提案リスト', 'Improvement Proposals')
    _add_lead(slide, f'{proposals_data.get("service_name", "")} に対する改善提案の一覧です。'
                       '詳細は次スライド以降に記載しています。')

    proposals = proposals_data.get('proposals', [])
    columns = ['No', '改善タイトル', '対象カテゴリ', '優先度']
    col_widths_pct = [8, 55, 22, 15]

    table_left, table_top = 40, 160
    table_w = 1200
    n_rows = len(proposals) + 1
    row_h = 44

    table_shape = slide.shapes.add_table(
        n_rows, len(columns),
        px(table_left), px(table_top),
        px(table_w), px(row_h * n_rows)
    )
    table = table_shape.table

    for i, w in enumerate(col_widths_pct):
        table.columns[i].width = px(int(table_w * w / 100))

    # ヘッダ
    for c, col_name in enumerate(columns):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ''
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        set_run(run, col_name, 16, bold=True, color=WHITE)
        cell.margin_left = Pt(8); cell.margin_right = Pt(8)
        cell.margin_top = Pt(8); cell.margin_bottom = Pt(8)

    # データ行
    for r, p in enumerate(proposals, start=1):
        bg = WHITE if r % 2 == 1 else STRIPE
        priority = p.get('priority', 'B')
        priority_color = RED if priority in ('S', 'A') else NAVY

        # No
        _fill_cell(table.cell(r, 0), str(p.get('no', r)), 16,
                    bold=True, color=NAVY, bg=bg, align=PP_ALIGN.CENTER)
        # タイトル
        _fill_cell(table.cell(r, 1), p.get('title', ''), 15,
                    bold=True, color=TEXT, bg=bg, align=PP_ALIGN.LEFT)
        # カテゴリ
        _fill_cell(table.cell(r, 2), p.get('category', ''), 14,
                    color=SUB_TEXT, bg=bg, align=PP_ALIGN.CENTER)
        # 優先度（バッジ風：背景色付き）
        cell = table.cell(r, 3)
        cell.fill.solid()
        cell.fill.fore_color.rgb = priority_color
        cell.text = ''
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        set_run(run, priority, 18, bold=True, color=WHITE)
        cell.margin_left = Pt(8); cell.margin_right = Pt(8)
        cell.margin_top = Pt(8); cell.margin_bottom = Pt(8)

    _add_footer(slide, page_num, total, author=author)
    return slide


# =====================================================================
# C-2 / Layout 15: 改善提案 詳細（add_proposal_detail）
# =====================================================================
def add_proposal_detail(prs, proposal, page_num=2, total=3,
                         author='UI/UX診断 by GPTs',
                         summary=None,
                         force_two_column=False):
    """
    改善提案 詳細スライド。

    **デフォルト：1スライド1提案（全幅1200pxカード）**
      - 「現状」「アクション」「期待効果」の各セクション本文を
        省略なしでフル表示
      - 原文の改行（\\n）はそのまま段落として尊重

    Args:
        proposal: dict（推奨）または list of dict（互換）
            dict 形式: {no, title, category, priority,
                        current, action, expected}
            list 形式: 旧API互換。1件なら1スライド表示、
                       2件 + force_two_column=True なら2列表示。
        summary: 結論帯に表示する総括コメント（最終スライドに付与）
        force_two_column: True のとき、list で2件渡された場合に
                          従来の2列レイアウトを強制（後方互換用）。
                          デフォルト False（1件1スライド推奨）。
    """
    # ---- 引数正規化：dict も list も受け付ける ----
    if isinstance(proposal, dict):
        proposals_slice = [proposal]
    elif isinstance(proposal, (list, tuple)):
        proposals_slice = list(proposal)
    else:
        raise TypeError('proposal は dict または list を渡してください')

    if len(proposals_slice) == 0:
        raise ValueError('proposal が空です')
    if len(proposals_slice) > 2:
        raise ValueError(
            'add_proposal_detail は1スライド最大2件まで。'
            '3件以上は分割して複数回呼び出してください。')

    # 2件渡されたとき：force_two_column=True なら2列、未指定なら警告
    n = len(proposals_slice)
    if n == 2 and not force_two_column:
        raise ValueError(
            '2件渡されましたが、1スライド1提案がデフォルトです。'
            '2列表示が必要なら force_two_column=True を指定してください。')

    slide = _blank_slide(prs)
    _add_header(slide, '改善提案リスト', 'Proposal Detail')

    # ---- リード文 ----
    nums = [str(p.get('no', i + 1)) for i, p in enumerate(proposals_slice)]
    if n == 1:
        lead_text = (f'提案 #{nums[0]} ── '
                     f'現状 → アクション → 期待効果 の3段で整理。')
    else:
        lead_text = (f'提案 #{nums[0]} / #{nums[1]} ── '
                     f'現状 → アクション → 期待効果 の3段で整理。')
    add_text(slide, 40, 90, 1200, lead_text, 16,
             color=SUB_TEXT, height_px=24)

    # ---- カード寸法 ----
    card_top = 130
    # 結論帯ありなら top=568, height=82 → 下端650
    # 結論帯なしならフッター（top=660）まで使用
    card_max_bottom = 552 if summary else 640
    card_h = card_max_bottom - card_top
    card_w = (1200 - 30) // 2 if n == 2 else 1200

    # フォントサイズ（2列時は少し小さく；14pt下限厳守）
    title_size = 17 if n == 2 else 22
    body_size = 14
    label_size = 14
    no_size = 26 if n == 2 else 36

    # カード内領域配分
    PAD_TOP = 14
    HEADER_H = 32 if n == 2 else 40
    # タイトル：1列なら最大3行折返しOK、2列なら1行に制限（…で省略）
    TITLE_H = 30 if n == 2 else 72
    CATEGORY_H = 22
    DIVIDER_GAP = 8
    SECTION_LABEL_H = 22
    SECTION_BODY_LINE_H = int(body_size * 1.6)  # 22px / 1行
    SECTION_GAP = 12  # セクション間の明示的余白

    # ---- タイトル省略ヘルパ（2列のみ使用） ----
    def _truncate_title(s, max_w):
        total = 0.0
        for idx, ch in enumerate(s):
            add = 0.5 if ord(ch) < 128 else 1.0
            if total + add > max_w:
                return s[:idx].rstrip() + '…'
            total += add
        return s

    for i, p in enumerate(proposals_slice):
        left = 40 + i * (card_w + 30) if n == 2 else 40
        priority = p.get('priority', 'B')
        priority_color = RED if priority in ('S', 'A') else NAVY

        # カード背景＋左帯
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top,
                  card_w, card_h, fill=LIGHT_GRAY)
        add_shape(slide, MSO_SHAPE.RECTANGLE, left, card_top, 6, card_h,
                  fill=priority_color)

        # 番号
        add_text(slide, left + 22, card_top + PAD_TOP, 200,
                 f'#{p.get("no", "")}', no_size, bold=True,
                 color=priority_color, line_height=1.0,
                 height_px=no_size + 4)
        # 優先度バッジ（右上）
        badge_w = 110 if n == 1 else 90
        badge_h = 34 if n == 1 else 30
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  left + card_w - badge_w - 16, card_top + PAD_TOP + 2,
                  badge_w, badge_h, fill=priority_color)
        add_text(slide, left + card_w - badge_w - 16,
                 card_top + PAD_TOP + (10 if n == 1 else 8),
                 badge_w, f'優先度 {priority}',
                 16 if n == 1 else 14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, height_px=20)

        # タイトル
        title_y = card_top + PAD_TOP + HEADER_H + 4
        raw_title = p.get('title', '')
        display_title = _truncate_title(raw_title, max_w=14) if n == 2 else raw_title
        add_text(slide, left + 22, title_y, card_w - 44,
                 display_title, title_size, bold=True, color=NAVY,
                 height_px=TITLE_H, line_height=1.3)

        # カテゴリ
        cat_y = title_y + TITLE_H + 2
        add_text(slide, left + 22, cat_y, card_w - 44,
                 p.get('category', ''), 14,
                 color=SUB_TEXT, height_px=CATEGORY_H)

        # 区切り線
        divider_y = cat_y + CATEGORY_H + DIVIDER_GAP
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  left + 22, divider_y, card_w - 44, 1, fill=BORDER_GRAY)

        # 3段構成
        sections = [
            ('現状',     p.get('current', ''),  SUB_TEXT),
            ('アクション', p.get('action', ''),  NAVY),
            ('期待効果', p.get('expected', ''), RED),
        ]
        section_top_start = divider_y + DIVIDER_GAP + 4
        section_avail = (card_h - (section_top_start - card_top)
                         - 14 - SECTION_GAP * 2)
        section_h = section_avail // 3

        if n == 1:
            # ---- 1列モード（フル幅）：省略なし、原文改行を尊重 ----
            for s_idx, (label, body, label_color) in enumerate(sections):
                sy = section_top_start + s_idx * (section_h + SECTION_GAP)
                # ラベル
                add_text(slide, left + 22, sy, 240, label, label_size,
                         bold=True, color=label_color, letter_spacing=1,
                         height_px=SECTION_LABEL_H)
                # 本文（原文改行を段落として保持）
                body_top = sy + SECTION_LABEL_H + 4
                body_h = section_h - SECTION_LABEL_H - 6
                body_paragraphs = [
                    {'text': line.strip(),
                     'size': body_size,
                     'color': TEXT,
                     'line_height': 1.55}
                    for line in (body or '').split('\n') if line.strip()
                ] or [{'text': '', 'size': body_size,
                       'color': TEXT, 'line_height': 1.55}]
                add_paragraph_box(slide, left + 22, body_top,
                                  card_w - 44,
                                  body_paragraphs,
                                  height_px=body_h, line_height=1.55)
        else:
            # ---- 2列モード（互換）：文字数制限で省略 ----
            body_inner_w = card_w - 44
            approx_fullwidth_per_line = max(7, body_inner_w // 22)
            raw_lines = max(1, (section_h - SECTION_LABEL_H - 6)
                            // SECTION_BODY_LINE_H)
            safety_factor = 0.82
            max_fullwidth_chars = max(6, int(approx_fullwidth_per_line
                                              * raw_lines * safety_factor))

            def _truncate_to_fullwidth(s, max_w):
                total = 0.0
                for idx, ch in enumerate(s):
                    add = 0.5 if ord(ch) < 128 else 1.0
                    if total + add > max_w:
                        return s[:idx].rstrip() + '…'
                    total += add
                return s

            for s_idx, (label, body, label_color) in enumerate(sections):
                sy = section_top_start + s_idx * (section_h + SECTION_GAP)
                add_text(slide, left + 22, sy, 200, label, label_size,
                         bold=True, color=label_color, letter_spacing=1,
                         height_px=SECTION_LABEL_H)
                body_top = sy + SECTION_LABEL_H + 2
                body_h = section_h - SECTION_LABEL_H - 4
                flat_body = (body or '').replace('\n', ' ')
                flat_body = _truncate_to_fullwidth(flat_body,
                                                   max_fullwidth_chars)
                add_paragraph_box(slide, left + 22, body_top,
                                  card_w - 44,
                                  [{'text': flat_body, 'size': body_size,
                                    'color': TEXT, 'line_height': 1.55}],
                                  height_px=body_h, line_height=1.55)

    # 総括（最終スライドのみ）
    if summary:
        _add_conclusion_band(slide, 568, summary,
                             label='SUMMARY', height_px=82)

    _add_footer(slide, page_num, total, author=author)
    return slide


# =====================================================================
# C-1 / Layout 16: UI診断スコアカード（2枚構成）
# （add_scorecard_onepager）
#   1/2 サマリ：上段カード + 強み/最優先課題（2列） + 結論
#   2/2 詳細スコア表：10項目を横幅1200px全面活用、一言所見は全文表示
# =====================================================================
def _scorecard_title(slide, slide_no, subtitle):
    """両スライド共通のタイトル帯描画（左赤縦帯10px + タイトル + サブタイトル）

    [ビジュアル強化 2026-06-29]
    - 左赤縦帯：6px → 10px に拡幅して視認性向上
    - タイトル位置調整：60pxへシフト（縦帯余白）
    [FB対応 2026-06-30]
    - タイトル帯高さ：56 → 64px に拡大
    - サブタイトル top：+8px下げて可読性向上
    """
    TITLE_TOP = 16
    TITLE_H = 64
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, TITLE_TOP, 10, TITLE_H,
              fill=RED)
    add_text(slide, 60, TITLE_TOP + 4, 1140,
             f'【{slide_no}】UI診断スコアカード', 24, bold=True,
             color=NAVY, height_px=34)
    add_text(slide, 60, TITLE_TOP + 42, 1140, subtitle,
             14, color=SUB_TEXT, height_px=22)


def _scorecard_info_cards(slide, diagnosis):
    """両スライド共通：上段4分割カード（診断対象/入力種別/総合スコア/ランク）

    [ビジュアル強化 2026-06-29]
    - 総合スコア値：26pt → 32pt に拡大
    - ランク値：26pt → 32pt に拡大
    - 左色帯：4px → 6px に拡幅
    - カード背景：NAVY 5%透過の極薄背景で奥行き演出
    - ランクをバッジ風に強調（数字大型 + ラベル小型）
    [FB対応 2026-06-30]
    - INFO_TOP：82 → 90px（タイトル帯64px拡大に追従、+8pxシフト）
    """
    INFO_TOP = 90
    INFO_H = 80
    card_w = (1200 - 12 * 3) // 4  # 291px

    score_val = diagnosis.get('total_score', 0)
    rank = diagnosis.get('rank', '-')
    rank_label = diagnosis.get('rank_label', '')
    rank_color = NAVY if rank in ('S', 'A', 'B') else RED

    # [FB対応 2026-06-30] 左色帯を全4枚NAVYに統一
    # 装飾目的の色分けは廃止。スコア値・ランク値は内部の数字フォント色で
    # 個別強調（rank_color = NAVY or RED）するため左色帯の機能的役割は薄い。
    info_cards = [
        {'bar': NAVY, 'paragraphs': [
            {'text': '診断対象', 'size': 14, 'bold': True,
             'color': SUB_TEXT, 'line_height': 1.0},
            {'text': diagnosis.get('service_name', '') or '-',
             'size': 18, 'bold': True, 'color': TEXT, 'line_height': 1.2},
        ]},
        {'bar': NAVY, 'paragraphs': [
            {'text': '入力種別', 'size': 14, 'bold': True,
             'color': SUB_TEXT, 'line_height': 1.0},
            {'text': diagnosis.get('input_type', '') or '-',
             'size': 16, 'bold': False, 'color': TEXT, 'line_height': 1.2},
        ]},
        {'bar': NAVY, 'paragraphs': [
            {'text': '総合スコア', 'size': 14, 'bold': True,
             'color': SUB_TEXT, 'line_height': 1.0},
            {'runs': [
                {'text': str(score_val), 'size': 32, 'bold': True,
                 'color': rank_color},
                {'text': '  / 50', 'size': 14, 'bold': False,
                 'color': SUB_TEXT},
            ], 'line_height': 1.0},
        ]},
        {'bar': NAVY, 'paragraphs': [
            {'text': 'ランク', 'size': 14, 'bold': True,
             'color': SUB_TEXT, 'line_height': 1.0},
            {'runs': [
                {'text': rank, 'size': 32, 'bold': True, 'color': rank_color},
                {'text': '  ' + (rank_label or ''),
                 'size': 14, 'bold': False, 'color': SUB_TEXT},
            ], 'line_height': 1.0},
        ]},
    ]
    for i, card in enumerate(info_cards):
        x = 40 + i * (card_w + 12)
        # カード本体：白＋極薄NAVYの2層構造で奥行き演出
        _add_bg_frame(slide, x, INFO_TOP, card_w, INFO_H,
                      fill=WHITE,
                      line=BORDER_GRAY, line_width_pt=1)
        _add_bg_frame(slide, x, INFO_TOP, card_w, INFO_H,
                      fill=NAVY, fill_alpha=4000)
        # 左色帯：4px → 6px に拡幅
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, INFO_TOP, 6, INFO_H,
                  fill=card['bar'])
        _add_multi_run_box(slide, x + 18, INFO_TOP + 10,
                           card_w - 26, INFO_H - 14, card['paragraphs'])


def add_scorecard_onepager(prs, diagnosis, page_num=1, total=2,
                            author='UI/UX診断 by GPTs',
                            high_threshold=0.75, low_threshold=0.50,
                            slide_no='1'):
    """
    UI診断スコアカードを2枚に分割して出力。

    1/2 サマリ：
      - タイトル帯
      - 上段4分割カード（診断対象/入力種別/総合スコア/ランク）
      - 強み / 最優先課題（横並び2列、各3項目を line_height 1.7 で余裕表示）
      - 結論帯
    2/2 詳細スコア表：
      - タイトル帯
      - 上段4分割カード（同上：コンテキストを引継ぎ）
      - 10項目スコア表（横幅1200px全面、一言所見は省略なしで折返し表示）

    Returns: (slide_summary, slide_detail) tuple
    """
    # ==============================================
    # 文字数バリデーション（design_system.md §1.3 準拠）
    # ==============================================
    validate_length(diagnosis.get('service_name'), 'service_name',
                    'C-1 サービス名 (diagnosis.service_name)')
    for i, sc in enumerate(diagnosis.get('scores', [])[:10], start=1):
        if isinstance(sc, dict):
            validate_length(sc.get('comment'), 'c1_comment',
                            f'C-1 一言所見 #{i} (scores[{i-1}].comment / 項目: {sc.get("name","?")})')
    for i, st in enumerate(diagnosis.get('strengths', [])[:3], start=1):
        validate_length(st, 'c1_strength',
                        f'C-1 強み #{i} (strengths[{i-1}])')
    for i, iss in enumerate(diagnosis.get('priority_issues', [])[:3], start=1):
        validate_length(iss, 'c1_issue',
                        f'C-1 最優先課題 #{i} (priority_issues[{i-1}])')
    validate_length(diagnosis.get('conclusion'), 'c1_conclusion',
                    'C-1 結論 (diagnosis.conclusion)')

    # ==============================================
    # スライド 1/2 : サマリ
    # ==============================================
    slide = _blank_slide(prs)
    _scorecard_title(slide, slide_no,
                     'サマリ：診断対象・スコア・強み・最優先課題・結論')
    _scorecard_info_cards(slide, diagnosis)

    # ---- 強み と 最優先課題（横並び 2分割、余裕ある行高） ----
    # [ビジュアル強化 2026-06-29]
    # 情報カード下端=82+80=162 → 余白20px → LOWER_TOP=182
    # 結論帯top=608 まで → LOWER_H=426（ヘッダ36 + 本文390px）
    # 本文 line_height=1.6、16pt実描画~28px×折返し2行×3項目 = 168px、余裕あり
    # [FB対応 2026-06-30] タイトル帯64px化に追従、+8pxシフト
    # 情報カード下端=90+80=170 → 余白20px → LOWER_TOP=190
    # [FB対応 2026-07-02] 結論帯62px化に伴い top 614→594、下段下端を
    # 586 (=594-8余白) に揃えるため LOWER_H を 416→396 に縮小。
    FULL_W = 1200
    LOWER_TOP = 190
    LOWER_H = 396  # 606→586 に短縮（結論帯top=594 の直前に安全余白8px）
    GAP = 16
    HALF_W = (FULL_W - GAP) // 2  # 592px
    HEADER_BAR_H = 36  # 32→36：余白拡大
    LEFT_X = 40
    RIGHT_X = 40 + HALF_W + GAP

    # 強み（左半分） - NAVYヘッダ
    add_shape(slide, MSO_SHAPE.RECTANGLE, LEFT_X, LOWER_TOP, HALF_W,
              HEADER_BAR_H, fill=NAVY)
    add_text(slide, LEFT_X + 18, LOWER_TOP, HALF_W - 36,
             '◎ 強み（活かすべき点）', 16, bold=True, color=WHITE,
             height_px=HEADER_BAR_H, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    str_body_top = LOWER_TOP + HEADER_BAR_H
    str_body_h = LOWER_H - HEADER_BAR_H
    # 本文背景：LIGHT_GRAY → 白＋NAVY 3%極薄で上品さ演出
    add_shape(slide, MSO_SHAPE.RECTANGLE, LEFT_X, str_body_top,
              HALF_W, str_body_h, fill=WHITE,
              line=BORDER_GRAY, line_width_pt=0.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, LEFT_X, str_body_top,
              HALF_W, str_body_h, fill=NAVY, fill_alpha=3000)
    # [v15/2026-07-12] リスト内折返しの連続感と項目間の分離感を分離制御：
    # line_height 1.6→1.2（折返し行間タイト）＋ space_after_pt=8pt（項目間余白）
    strengths = diagnosis.get('strengths', [])[:3]
    str_paragraphs = []
    for s in strengths:
        str_paragraphs.append({
            'text': '● ' + s, 'size': 16, 'bold': False,
            'color': TEXT, 'line_height': 1.2, 'space_after_pt': 8,
        })
    _add_multi_run_box(slide, LEFT_X + 20, str_body_top + 18,
                       HALF_W - 40, str_body_h - 32,
                       str_paragraphs)

    # 最優先課題（右半分） - REDヘッダ
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, LOWER_TOP, HALF_W,
              HEADER_BAR_H, fill=RED)
    add_text(slide, RIGHT_X + 18, LOWER_TOP, HALF_W - 36,
             '⚠ 最優先で直すべき点', 16, bold=True, color=WHITE,
             height_px=HEADER_BAR_H, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    issue_body_top = LOWER_TOP + HEADER_BAR_H
    issue_body_h = LOWER_H - HEADER_BAR_H
    # 課題側はRED 3%極薄で警告色トーンを醸成
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, issue_body_top,
              HALF_W, issue_body_h, fill=WHITE,
              line=BORDER_GRAY, line_width_pt=0.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, issue_body_top,
              HALF_W, issue_body_h, fill=RED, fill_alpha=3000)
    # [v15/2026-07-12] リスト内折返しの連続感と項目間の分離感を分離制御
    issues = diagnosis.get('priority_issues', [])[:3]
    issue_paragraphs = []
    for i, s in enumerate(issues):
        issue_paragraphs.append({
            'text': f'{i + 1}. ' + s, 'size': 16, 'bold': False,
            'color': TEXT, 'line_height': 1.2, 'space_after_pt': 8,
        })
    _add_multi_run_box(slide, RIGHT_X + 20, issue_body_top + 18,
                       HALF_W - 40, issue_body_h - 32,
                       issue_paragraphs)

    # 結論帯（最下部）─ [ビジュアル強化] ORANGE系へ変更
    # ORANGE基調＋NAVYラベルで「行動喚起」の視覚的トーン演出
    # [FB対応 2026-07-02] 文字切れ根絶のため帯高さ 38→62px（2行対応）、
    #   最大120字＋word_wrap=True、top を 614→594 に上シフトしフッター境界
    #   （660）まで6pxの安全余白確保。ラベル箱は帯全高で垂直センター視認性向上。
    CONCL_TOP = 594
    CONCL_H = 62  # 帯本体高さ（2行対応・+24px）
    LABEL_W = 88  # ラベル箱幅
    GAP = 40  # ラベル⇔本文間隔（案B推奨値）
    conclusion_text = diagnosis.get('conclusion', '')
    if conclusion_text:
        flat = conclusion_text.replace('\n', ' ')
        # 本体：ORANGE薄塗背景＋ORANGE枠
        add_shape(slide, MSO_SHAPE.RECTANGLE, 40, CONCL_TOP, 1200,
                  CONCL_H, fill=WHITE,
                  line=ORANGE, line_width_pt=1)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 40, CONCL_TOP, 1200,
                  CONCL_H, fill=ORANGE, fill_alpha=10000)
        # ラベル部：ORANGEベタ、帯全高で垂直センター（視認性優先）
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  40, CONCL_TOP, LABEL_W, CONCL_H,
                  fill=ORANGE)
        add_text(slide, 40, CONCL_TOP, LABEL_W, '結論', 14,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=CONCL_H, letter_spacing=2,
                 vertical_anchor=MSO_ANCHOR.MIDDLE)
        # 120字上限で切り詰め（設計方針：GPTs側で120字以内に生成させる）
        flat_disp = _truncate_full(flat, 120)
        # 本文 left = 40 + LABEL_W + GAP = 168、幅 = 40 + 1200 - 168 = 1072
        body_left = 40 + LABEL_W + GAP
        body_w = 40 + 1200 - body_left
        body_box = add_text(slide, body_left, CONCL_TOP, body_w, flat_disp, 14,
                            bold=True, color=NAVY, height_px=CONCL_H,
                            line_height=1.35,
                            vertical_anchor=MSO_ANCHOR.MIDDLE)
        # word_wrap=True で2行折り返しを許容（文字切れ根絶）
        body_box.text_frame.word_wrap = True

    _add_footer(slide, page_num, total, author=author)
    slide_summary = slide

    # ==============================================
    # スライド 2/2 : 詳細スコア表
    # ==============================================
    slide = _blank_slide(prs)
    _scorecard_title(slide, slide_no,
                     '10項目評価：各項目のスコアと一言所見（全文表示）')
    _scorecard_info_cards(slide, diagnosis)

    # 表セクション（横幅1200px全面、一言所見は省略なしで折返し）
    # 情報カード下端=160 → 余白20px → TABLE_TOP=180
    # フッター帯top=660 まで → 表領域=460px → ヘッダ28 + 10行×43px = 458px
    FULL_W = 1200
    TABLE_TOP = 180
    header_h = 28
    n_items = len(diagnosis.get('scores', [])[:10])
    # 残り432px ÷ 10行 = 43px/行（一言所見の2行折返しに十分）
    row_h = 30  # 設定値30px、LibreOffice実描画で約40-43px
    # No:項目:スコア:一言所見 = 5:17:7:71 → 60:204:84:852px
    # 「ファーストビュー」など全角8文字を1行で収めるため項目幅を拡大
    col_widths_pct_local = [5, 17, 7, 71]
    columns_local = ['No', '項目', 'スコア', '一言所見']
    n_rows = n_items + 1

    table_shape = slide.shapes.add_table(
        n_rows, len(columns_local),
        px(40), px(TABLE_TOP),
        px(FULL_W), px(header_h + row_h * n_items))
    table = table_shape.table
    for i, wpct in enumerate(col_widths_pct_local):
        table.columns[i].width = px(int(FULL_W * wpct / 100))
    table.rows[0].height = px(header_h)
    for r in range(1, n_rows):
        table.rows[r].height = px(row_h)

    # ヘッダ行
    header_aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT,
                     PP_ALIGN.CENTER, PP_ALIGN.LEFT]
    for c, (col_name, a) in enumerate(zip(columns_local, header_aligns)):
        _fill_cell(table.cell(0, c), col_name, 14, bold=True,
                   color=WHITE, bg=NAVY, align=a,
                   margin_v=2, no_wrap=True)

    # データ行
    scores = diagnosis.get('scores', [])[:10]
    for r, item in enumerate(scores, start=1):
        bg = WHITE if r % 2 == 1 else STRIPE
        sv = item.get('score', 0)
        mv = item.get('max', 5)
        ratio = sv / mv if mv else 0
        if ratio >= high_threshold:
            score_color = NAVY; score_bold = True
        elif ratio <= low_threshold:
            score_color = RED; score_bold = True
        else:
            score_color = TEXT; score_bold = False

        _fill_cell(table.cell(r, 0), str(r), 14, color=TEXT,
                   bg=bg, align=PP_ALIGN.CENTER,
                   margin_v=2, no_wrap=True)
        # [FB対応 2026-06-30] scores の `name` キーを優先（schema統一）
        cat_raw = item.get('name') or item.get('category', '')
        cat_disp = _truncate_full(cat_raw, 10)
        _fill_cell(table.cell(r, 1), cat_disp, 14, bold=True,
                   color=TEXT, bg=bg, align=PP_ALIGN.LEFT,
                   margin_v=2, no_wrap=True)
        score_text = f'{sv}/{mv}'
        _fill_cell(table.cell(r, 2), score_text, 14, bold=score_bold,
                   color=score_color, bg=bg, align=PP_ALIGN.CENTER,
                   margin_v=2, no_wrap=True)
        # 一言所見：全文表示（折返し許可）
        comment_raw = item.get('comment', '')
        _fill_cell(table.cell(r, 3), comment_raw, 14, color=TEXT,
                   bg=bg, align=PP_ALIGN.LEFT,
                   margin_v=2, no_wrap=False)

    _add_footer(slide, page_num + 1, total, author=author)
    slide_detail = slide

    return (slide_summary, slide_detail)


# =====================================================================
# C-2 / Layout 17: 改善提案リスト（2枚構成）
# （add_proposal_onepager）
#   1/2 前半提案：5件中の前半3件をカード形式で表示
#   2/2 後半提案：残り2件のカード + POINT帯（最優先の一手まとめ）
# =====================================================================
def _proposal_legend(slide, title_top):
    """両スライド共通：右上の優先度凡例（高/中/低バッジ）"""
    LEGEND_TOP = title_top + 14
    add_text(slide, 880, LEGEND_TOP, 80, '優先度', 14, bold=True,
             color=SUB_TEXT, align=PP_ALIGN.RIGHT, height_px=22)
    PRIO_RED = RGBColor(0xD0, 0x02, 0x1B)
    PRIO_ORANGE = RGBColor(0xE8, 0x8B, 0x1F)
    PRIO_GRAY = RGBColor(0x88, 0x88, 0x88)
    for i, (lbl, c) in enumerate([('高', PRIO_RED),
                                   ('中', PRIO_ORANGE),
                                   ('低', PRIO_GRAY)]):
        bx = 970 + i * 80
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, LEGEND_TOP,
                  70, 28, fill=c)
        add_text(slide, bx, LEGEND_TOP + 3, 70, lbl, 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, height_px=22)


def _proposal_title(slide, slide_no, subtitle):
    """両スライド共通：タイトル帯

    [ビジュアル強化 2026-06-29]
    - 左赤縦帯：6px → 10px に拡幅
    [FB対応 2026-06-30]
    - タイトル帯高さ：56 → 64px に拡大（サブタイトル収納問題解消）
    - タイトル段落の line_height 拡大で視覚的間隔確保
    - サブタイトル space_before 相当の余白追加
    """
    TITLE_TOP = 16
    TITLE_H = 64
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, TITLE_TOP, 10, TITLE_H,
              fill=RED)
    # タイトル（24pt）とサブタイトル（14pt）を別々の add_text で配置
    # → 1ボックスにまとめる方式だと line_height でしか間隔調整できないため
    add_text(slide, 60, TITLE_TOP + 4, 820,
             f'【{slide_no}】改善提案リスト', 24, bold=True,
             color=NAVY, height_px=34)
    add_text(slide, 60, TITLE_TOP + 42, 820, subtitle,
             14, color=SUB_TEXT, height_px=22)
    _proposal_legend(slide, TITLE_TOP)


def _proposal_card(slide, p, y, card_h, prio_color, prio_label):
    """1提案カードを描画（共通ヘルパ）。

    [ビジュアル強化 2026-06-29]
    - カード背景：白＋NAVY 3%極薄の2層構造で奥行き演出
    - 左色帯：6px → 8px に拡幅
    - 番号サークル：30→34px、文字16pt化
    - 「改善後」帯：淡赤 → ORANGE薄塗（行動喚起トーン）
    """
    # カード背景：白＋NAVY 3%極薄
    _add_bg_frame(slide, 40, y, 1200, card_h,
                  fill=WHITE, line=BORDER_GRAY, line_width_pt=1)
    _add_bg_frame(slide, 40, y, 1200, card_h,
                  fill=NAVY, fill_alpha=3000)
    # 左色帯：6 → 8px に拡幅
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, y, 8, card_h,
              fill=prio_color)

    # ヘッダ行（番号サークル + タイトル + 優先度バッジ + 工数）
    head_top = y + 10
    # 番号サークル：30→34px、文字16pt化
    add_shape(slide, MSO_SHAPE.OVAL, 58, head_top, 34, 34,
              fill=prio_color)
    add_text(slide, 58, head_top + 4, 34,
             str(p.get('no', '')), 16, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, height_px=26,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    raw_title = p.get('title', '')
    t_disp = _truncate_full(raw_title, 36)
    add_text(slide, 102, head_top + 2, 770, t_disp, 18,
             bold=True, color=NAVY, height_px=32)
    # 優先度バッジ：中央揃え強化（vertical_anchor=MIDDLE）
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 980, head_top + 1,
              70, 30, fill=prio_color)
    add_text(slide, 980, head_top + 1, 70, prio_label, 14,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             height_px=30, vertical_anchor=MSO_ANCHOR.MIDDLE)
    effort = p.get('effort', '中')
    add_text(slide, 1060, head_top + 5, 170,
             f'工数：{effort}', 14, color=SUB_TEXT, height_px=22)

    # 区切り線
    add_shape(slide, MSO_SHAPE.RECTANGLE, 56, head_top + 36,
              1170, 1, fill=BORDER_GRAY)

    # 改善箇所 / 課題（2行表示・各1ブロック）
    BODY_TOP = head_top + 44
    target_area = p.get('target_area', '')
    issue = p.get('issue', '')
    _add_multi_run_box(slide, 56, BODY_TOP, 568, 48, [
        {'runs': [
            {'text': '改善箇所　', 'size': 14, 'bold': True, 'color': NAVY},
            {'text': target_area, 'size': 14, 'bold': False, 'color': TEXT},
        ], 'line_height': 1.5},
    ])
    _add_multi_run_box(slide, 640, BODY_TOP, 584, 48, [
        {'runs': [
            {'text': '課題　', 'size': 14, 'bold': True, 'color': NAVY},
            {'text': issue, 'size': 14, 'bold': False, 'color': TEXT},
        ], 'line_height': 1.5},
    ])

    # 現状 / 改善後（2行表示・色違い背景帯）
    # [ビジュアル強化 2026-06-29]
    # - 現状：淡グレー（変更なし、現状を中立的に表現）
    # - 改善後：淡赤 → ORANGE 12%薄塗（行動喚起トーン）、ラベルORANGE化
    bottom_top = BODY_TOP + 52
    # 現状（左、淡グレー帯）
    add_shape(slide, MSO_SHAPE.RECTANGLE, 56, bottom_top, 568, 48,
              fill=LIGHT_GRAY, line=BORDER_GRAY, line_width_pt=0.5)
    _add_multi_run_box(slide, 64, bottom_top + 4, 552, 40, [
        {'runs': [
            {'text': '現状　', 'size': 14, 'bold': True, 'color': SUB_TEXT},
            {'text': p.get('before', ''), 'size': 14,
             'bold': False, 'color': TEXT},
        ], 'line_height': 1.5},
    ])
    # 改善後（右、ORANGE 12%薄塗 帯）
    # [FB対応 2026-07-02] 色ルール明確化：
    #   ORANGE は「行動喚起帯（結論・POINT・改善方向）」専用とし、
    #   ここの識別枠は NAVY_LIGHT に統一。塗りの ORANGE 12% は
    #   Before との対比を出す視覚要素として維持（枠線のみ変更）。
    add_shape(slide, MSO_SHAPE.RECTANGLE, 640, bottom_top, 584, 48,
              fill=WHITE, line=NAVY_LIGHT, line_width_pt=0.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 640, bottom_top, 584, 48,
              fill=ORANGE, fill_alpha=12000)
    _add_multi_run_box(slide, 648, bottom_top + 4, 568, 40, [
        {'runs': [
            {'text': '改善後　', 'size': 14, 'bold': True, 'color': ORANGE},
            {'text': p.get('after', ''), 'size': 14,
             'bold': False, 'color': TEXT},
        ], 'line_height': 1.5},
    ])


def _proposal_point_band(slide, summary, point_band_top=594, point_band_h=62):
    """POINT帯描画（スライド最下部）

    [ビジュアル強化 2026-06-29]
    - 本体：NAVYベタ → ORANGE薄塗（行動喚起トーン）
    - ラベル：RED → ORANGEベタ
    - 本文：WHITE → NAVY太字（コントラスト確保＋ブランド統一）
    [FB対応 2026-06-30]
    - ラベル⇔本文間隔 12→24px拡大
    [FB対応 2026-07-01 案B]
    - ラベル箱を上下2pxずつ縮めて立体感を演出
    - ラベル⇔本文間隔 24→40pxに拡大
    [FB対応 2026-07-02]
    - 帯高さ 32→62px（2行対応）、top 620→594 に上シフト
    - 文字上限 56→120字＋word_wrap=True（文字切れ根絶）
    - ラベル箱は帯全高で垂直センター視認性向上
    """
    if not summary:
        return
    flat = summary.replace('\n', ' ')
    LABEL_W = 90
    GAP = 40
    # 本体：ORANGE薄塗背景＋ORANGE枠
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, point_band_top,
              1200, point_band_h, fill=WHITE,
              line=ORANGE, line_width_pt=1)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, point_band_top,
              1200, point_band_h, fill=ORANGE, fill_alpha=10000)
    # ラベル部：ORANGEベタ、帯全高で垂直センター（視認性優先）
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              40, point_band_top, LABEL_W, point_band_h,
              fill=ORANGE)
    add_text(slide, 40, point_band_top, LABEL_W, 'POINT', 14,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             height_px=point_band_h, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    # 120字上限で切り詰め（GPTs側で120字以内に生成させる）
    flat_disp = _truncate_full(flat, 120)
    # [FB対応 2026-07-16] 句点孤立防止：末尾の「。」直前に word joiner
    # (U+2060) を挿入し、直前文字と結束させる。これにより「。」だけが
    # 2行目に取り残される日本語組版として不自然な折返しを回避する。
    # 副作用：word joiner は zero-width の非表示文字のため見た目に影響なし。
    if flat_disp.endswith('。'):
        flat_disp = flat_disp[:-1] + '\u2060。'
    # 文中の句点も同様に処理（「〜こと。〜」のように文中で句点＋続きがある場合）
    # は稀なので、末尾のみ対応で実運用の大半をカバー
    # 本文 left = 40 + LABEL_W + GAP = 170、幅 = 1200 + 40 - 170 = 1070
    body_left = 40 + LABEL_W + GAP
    body_w = 40 + 1200 - body_left
    body_box = add_text(slide, body_left, point_band_top, body_w, flat_disp, 14,
                        bold=True, color=NAVY, height_px=point_band_h,
                        line_height=1.35,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)
    # word_wrap=True で2行折り返しを許容（文字切れ根絶）
    body_box.text_frame.word_wrap = True


def add_proposal_onepager(prs, proposals_data, page_num=1, total=None,
                           author='UI/UX診断 by GPTs',
                           slide_no='2'):
    """
    改善提案リストを件数に応じて自動的に1or2枚で出力。

    - 4件以上：2枚構成（1/2 前半3件 / 2/2 後半残り + POINT帯）
    - 3件以下：1枚構成（全件 + POINT帯を1枚に集約）

    Args:
        proposals_data: dict
            service_name (str), proposals (list of max 5 dicts):
              no, title, priority ('高'/'中'/'低' or 'S'/'A'/'B'/'C'),
              effort ('小'/'中'/'大'), target_area, issue, before, after,
              target_score_item (optional)
            summary (str): POINT帯本文
        total: 全体総ページ数（None指定時は提案件数から自動算出）
    Returns:
        slide1 (3件以下) or (slide1, slide2) tuple (4件以上)
    """
    PRIO_RED = RGBColor(0xD0, 0x02, 0x1B)
    PRIO_ORANGE = RGBColor(0xE8, 0x8B, 0x1F)
    PRIO_GRAY = RGBColor(0x88, 0x88, 0x88)

    def _prio_color(p):
        if p in ('S', 'A', '高'):
            return PRIO_RED
        if p in ('B', '中'):
            return PRIO_ORANGE
        return PRIO_GRAY

    prio_label_map = {'S': '高', 'A': '高', 'B': '中', 'C': '低',
                      '高': '高', '中': '中', '低': '低'}

    # ==============================================
    # 文字数バリデーション（design_system.md §1.4 準拠）
    # ==============================================
    validate_length(proposals_data.get('service_name'), 'service_name',
                    'C-2 サービス名 (proposals_data.service_name)')
    for i, p in enumerate(proposals_data.get('proposals', [])[:5], start=1):
        if isinstance(p, dict):
            validate_length(p.get('title'), 'c2_title',
                            f'C-2 提案#{i} タイトル (proposals[{i-1}].title)')
            # POINT本文：issue / before / after / target_area などを連結したもの。
            # 各フィールド単体ではなく point があれば優先、なければ issue を上限内に。
            point_text = p.get('point') or p.get('issue') or ''
            validate_length(point_text, 'c2_point',
                            f'C-2 提案#{i} POINT (proposals[{i-1}].point|issue)')
            if p.get('priority') is not None:
                validate_length(str(p.get('priority')), 'c2_priority',
                                f'C-2 提案#{i} 優先度 (proposals[{i-1}].priority)')
            if p.get('category') is not None:
                validate_length(p.get('category'), 'c2_category',
                                f'C-2 提案#{i} カテゴリ (proposals[{i-1}].category)')

    proposals = proposals_data.get('proposals', [])[:5]
    n = len(proposals)
    summary = proposals_data.get('summary', '')

    # 全体ページ数を自動決定（1スライド or 2スライド）
    is_split = n >= 4
    auto_total = 2 if is_split else 1
    if total is None:
        total = auto_total

    # [FB対応 2026-06-30] タイトル帯64px化に追従、92→100に+8pxシフト
    card_top = 100
    card_gap = 10

    if not is_split:
        # ===============================================
        # 1スライド構成（3件以下：全件 + POINT帯を1枚に集約）
        # ===============================================
        slide = _blank_slide(prs)
        _proposal_title(slide, slide_no,
                        f'優先度の高い改善アクションを整理（全{n}件）')

        # POINT帯は最下部、カード領域はその上
        # [FB対応 2026-07-02] POINT帯を 62px に拡大し top を 620→594 に上シフト。
        # AREA_BOTTOM も 594-8=586 に縮小（カード高さは件数に応じ自動均等割）。
        POINT_BAND_TOP = 594
        AREA_BOTTOM = POINT_BAND_TOP - 8  # 586
        avail = AREA_BOTTOM - card_top
        if n > 0:
            card_h = (avail - card_gap * (n - 1)) // n
        else:
            card_h = 200

        for i, p in enumerate(proposals):
            y = card_top + i * (card_h + card_gap)
            priority = p.get('priority', 'B')
            _proposal_card(slide, p, y, card_h,
                           _prio_color(priority),
                           prio_label_map.get(priority, '中'))

        _proposal_point_band(slide, summary, point_band_top=POINT_BAND_TOP)
        _add_footer(slide, page_num, total, author=author)
        return slide

    # ===============================================
    # 2スライド構成（4件以上：前半3件 + 後半残り+POINT）
    # ===============================================
    first_n = 3
    proposals_a = proposals[:first_n]
    proposals_b = proposals[first_n:]

    # --- スライド 1/2 : 前半提案 ---
    slide = _blank_slide(prs)
    _proposal_title(slide, slide_no,
                    f'優先度の高い改善アクションを整理（1/2 前半 {first_n}件）')

    AREA_BOTTOM_A = 658  # フッター帯top=660 直前
    avail_a = AREA_BOTTOM_A - card_top
    card_h_a = (avail_a - card_gap * (first_n - 1)) // first_n

    for i, p in enumerate(proposals_a):
        y = card_top + i * (card_h_a + card_gap)
        priority = p.get('priority', 'B')
        _proposal_card(slide, p, y, card_h_a,
                       _prio_color(priority),
                       prio_label_map.get(priority, '中'))

    _add_footer(slide, page_num, total, author=author)
    slide1 = slide

    # --- スライド 2/2 : 後半提案 + POINT帯 ---
    slide = _blank_slide(prs)
    second_n = len(proposals_b)
    _proposal_title(slide, slide_no,
                    f'優先度の高い改善アクションを整理（2/2 後半 {second_n}件・POINT）')

    # [FB対応 2026-07-02] POINT帯を 62px に拡大し top を 620→594 に上シフト。
    # [FB対応 2026-07-16] 2枚目のカード高さを1枚目と統一し上詰め配置。
    #   従来：件数で領域を均等割 → 2件時カードが肥大化し下部余白が不自然に大
    #   改善：card_h_b = card_h_a（1枚目と同じ）で上詰め、下部余白は自然な
    #         「資料的余白」として機能させる。POINT帯位置は594px据え置き。
    POINT_BAND_TOP = 594
    AREA_BOTTOM_B = POINT_BAND_TOP - 8  # 586（現状は参考値、上詰めのため未使用）
    card_h_b = card_h_a  # 1枚目と統一（上詰め・下余白自然化）

    for i, p in enumerate(proposals_b):
        y = card_top + i * (card_h_b + card_gap)
        priority = p.get('priority', 'B')
        _proposal_card(slide, p, y, card_h_b,
                       _prio_color(priority),
                       prio_label_map.get(priority, '中'))

    _proposal_point_band(slide, summary, point_band_top=POINT_BAND_TOP)
    _add_footer(slide, page_num + 1, total, author=author)
    slide2 = slide

    return (slide1, slide2)


# =====================================================================
# C-3 / Layout 18: ビジュアル診断ボード（3枚構成）
# （add_visual_board）— フェーズ1：レーダーは表形式代替
#   1/3 LP構造マップ + 総評・最重要課題 + 6ステップ行動フロー
#   2/3 スコア視覚化（10項目表＋強み/課題コンパクトカード）
#   3/3 Before/After Top3 ハイライト（赤帯Before／オレンジ帯After）
# =====================================================================
def _visual_title(slide, slide_no, sub_title):
    """C-3共通：タイトル帯（左赤縦帯 + タイトル + サブタイトル）

    [ビジュアル強化 2026-06-29]
    - 左赤縦帯：6px → 10px に拡幅
    [FB対応 2026-06-30]
    - タイトル帯高さ：56 → 64px に拡大
    - サブタイトル top：+8px下げて可読性向上
    """
    TITLE_TOP = 16
    TITLE_H = 64
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, TITLE_TOP, 10, TITLE_H, fill=RED)
    add_text(slide, 60, TITLE_TOP + 4, 1140,
             f'【{slide_no}】ビジュアル診断ボード', 24, bold=True,
             color=NAVY, height_px=34)
    add_text(slide, 60, TITLE_TOP + 42, 1140, sub_title,
             14, color=SUB_TEXT, height_px=22)


def _visual_meta_header(slide, visual_data, top=90):
    """C-3共通：目的/対象/診断日/スコアの4分割メタカード

    [ビジュアル強化 2026-06-29]
    - 左色帯：4px → 6px に拡幅
    - カード背景：白＋NAVY 4%極薄の2層構造で奥行き演出
    [FB対応 2026-06-30]
    - top デフォルト 82 → 90px（タイトル帯64px化に追従）
    """
    INFO_H = 64
    card_w = (1200 - 12 * 3) // 4  # 291px
    score_val = visual_data.get('total_score', 0)
    rank = visual_data.get('rank', '-')
    rank_label = visual_data.get('rank_label', '')

    # [FB対応 2026-07-04] 全4カード NAVY 統一（装飾目的の色分けは廃止）
    # [FB対応 2026-07-16 v15.4] 総合スコアカードのみ訴求強化：
    #   スコア値を14pt→26pt に、ランク値を14pt→32pt に大型化。
    #   販売資産として最も訴求すべきポイントを視覚的に主役化する。
    rank_color = NAVY if rank in ('S', 'A', 'B') else RED
    meta_cards = [
        {'bar': NAVY, 'label': '目的',
         'value': visual_data.get('purpose', '-'), 'size': 14},
        {'bar': NAVY, 'label': '対象',
         'value': visual_data.get('target', '-'), 'size': 14},
        {'bar': NAVY, 'label': '診断日',
         'value': visual_data.get('diagnosis_date', '-'), 'size': 14},
        # 総合スコアは runs で「スコア + / 50 + ランク」を大型化・色分け
        {'bar': NAVY, 'label': '総合スコア',
         'is_score_card': True},
    ]

    for i, card in enumerate(meta_cards):
        x = 40 + i * (card_w + 12)
        # カード本体：白＋極薄NAVYの2層構造
        _add_bg_frame(slide, x, top, card_w, INFO_H,
                      fill=WHITE, line=BORDER_GRAY, line_width_pt=1)
        _add_bg_frame(slide, x, top, card_w, INFO_H,
                      fill=NAVY, fill_alpha=4000)
        # 左色帯：4 → 6px に拡幅
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, top, 6, INFO_H,
                  fill=card['bar'])
        # 総合スコアカードのみ、C-1側と同じ大型化仕様を適用
        if card.get('is_score_card'):
            _add_multi_run_box(slide, x + 18, top + 8, card_w - 26,
                               INFO_H - 14, [
                {'text': card['label'], 'size': 14, 'bold': True,
                 'color': SUB_TEXT, 'line_height': 1.0},
                {'runs': [
                    {'text': str(score_val), 'size': 26, 'bold': True,
                     'color': rank_color},
                    {'text': '  / 50    ', 'size': 14, 'bold': False,
                     'color': SUB_TEXT},
                    {'text': rank, 'size': 26, 'bold': True,
                     'color': rank_color},
                ], 'line_height': 1.0},
            ])
        else:
            _add_multi_run_box(slide, x + 18, top + 8, card_w - 26,
                               INFO_H - 14, [
                {'text': card['label'], 'size': 14, 'bold': True,
                 'color': SUB_TEXT, 'line_height': 1.0},
                {'text': card['value'],
                 'size': card.get('size', 16),
                 'bold': card.get('value_bold', False),
                 'color': card.get('value_color', TEXT),
                 'line_height': 1.1},
            ])


def _visual_direction_footer(slide, direction, top=594):
    """C-3共通：フッター上の改善方向性帯

    [ビジュアル強化 2026-06-29]
    - 帯高さ：28 → 36px に拡大（視認性向上）
    - ORANGE_LIGHT → ORANGE 12%透過に変更
    - 本文 14pt → 16pt 太字 NAVY色（メリハリ強化）
    [FB対応 2026-06-30]
    - ラベル⇔本文間隔 20→28px拡大
    [FB対応 2026-07-01 案B]
    - ラベル箱を上下2pxずつ縮めて立体感を演出
    - ラベル⇔本文間隔 28→40pxに拡大
    [FB対応 2026-07-02]
    - 帯高さ 36→62px（2行対応）、top を 594 へ統一
    - 文字上限 44→120字＋word_wrap=True（文字切れ根絶）
    - 本文14pt太字に変更（2行構成時の可読性優先）
    - ラベル箱は帯全高で垂直センター視認性向上
    - 後方互換：top=620 指定時も 594 として扱う（旧呼出しの吸収）
    """
    if not direction:
        return
    # 後方互換：旧デフォルト620/612を渡された場合も新top 594 に強制上書き
    BAND_TOP = 594 if top in (620, 612) else top
    BAND_H = 62
    LABEL_W = 110
    GAP = 40
    # 本体：ORANGE薄塗＋枠
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, BAND_TOP, 1200, BAND_H,
              fill=WHITE, line=ORANGE, line_width_pt=1)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, BAND_TOP, 1200, BAND_H,
              fill=ORANGE, fill_alpha=12000)
    # ラベル部：ORANGEベタ、帯全高で垂直センター（視認性優先）
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              40, BAND_TOP, LABEL_W, BAND_H,
              fill=ORANGE)
    add_text(slide, 40, BAND_TOP, LABEL_W, '改善方向', 14,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             height_px=BAND_H, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    flat = direction.replace('\n', ' ')
    # 120字上限で切り詰め（GPTs側で120字以内に生成させる）
    flat = _truncate_full(flat, 120)
    # 本文 left = 40 + LABEL_W + GAP = 190、幅 = 40 + 1200 - 190 = 1050
    body_left = 40 + LABEL_W + GAP
    body_w = 40 + 1200 - body_left
    body_box = add_text(slide, body_left, BAND_TOP, body_w, flat, 14,
                        bold=True, color=NAVY, height_px=BAND_H,
                        line_height=1.35,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)
    # word_wrap=True で2行折り返しを許容（文字切れ根絶）
    body_box.text_frame.word_wrap = True


def _normalize_flow_status(s):
    """flow_steps.status の表記揺れを正規化"""
    if s in ('✓', 'OK', True, 'success', '成功'):
        return '✓'
    if s in ('✕', '×', 'NG', False, 'fail', '離脱'):
        return '✕'
    return '？'


def add_visual_board(prs, visual_data, page_num=1, total=3, slide_no='3',
                     author='UI/UX診断 by GPTs'):
    """
    ビジュアル診断ボードを3枚構成で出力。

    Args:
        visual_data: dict （詳細は visual_data_schema.md 参照）
        page_num: 開始ページ番号（個別出力時=1、統合版では既存ページ数+1）
        total: 全体ページ数（個別=3、統合版では C-1+C-2+3 を渡す）
        slide_no: 章番号（既定 '3'）
    Returns:
        (slide1, slide2, slide3) tuple
    """
    # ==============================================
    # 文字数バリデーション（design_system.md §1.5 準拠）
    # ==============================================
    validate_length(visual_data.get('service_name'), 'service_name',
                    'C-3 サービス名 (visual_data.service_name)')
    validate_length(visual_data.get('summary'), 'c3_summary',
                    'C-3 総評 (visual_data.summary)')
    validate_length(visual_data.get('top_issue'), 'c3_top_issue',
                    'C-3 最優先課題 (visual_data.top_issue)')
    validate_length(visual_data.get('direction'), 'c3_direction',
                    'C-3 改善方向 (visual_data.direction)')
    for i, sec in enumerate(visual_data.get('sections', []), start=1):
        if isinstance(sec, dict):
            validate_length(sec.get('label'), 'c3_section_label',
                            f'C-3 LP構造マップ#{i} (sections[{i-1}].label)')
    for i, step in enumerate(visual_data.get('flow_steps', [])[:6], start=1):
        if isinstance(step, dict):
            # ステップ全体（label + note）は1ステップ枠（25文字以内）
            combined = (step.get('label', '') + ' ' + step.get('note', '')).strip()
            validate_length(combined, 'c3_flow_step',
                            f'C-3 行動フロー#{i} (flow_steps[{i-1}].label+note)')
    for i, hl in enumerate(visual_data.get('highlights', [])[:3], start=1):
        if isinstance(hl, dict):
            validate_length(hl.get('title'), 'c3_highlight_title',
                            f'C-3 ハイライト#{i} タイトル (highlights[{i-1}].title)')
            validate_length(hl.get('before'), 'c3_before',
                            f'C-3 ハイライト#{i} Before (highlights[{i-1}].before)')
            validate_length(hl.get('after'), 'c3_after',
                            f'C-3 ハイライト#{i} After (highlights[{i-1}].after)')

    # ==============================================================
    # スライド 1/3 : LP構造マップ + 総評・最重要課題 + 行動フロー
    # ==============================================================
    slide = _blank_slide(prs)
    _visual_title(slide, slide_no,
                  '構造マップ・総評・最重要課題・ユーザー行動フロー')
    _visual_meta_header(slide, visual_data, top=90)

    # ---- 左：LP構造マップ ----
    # [FB対応 2026-07-03] 番号サークルの色（NAVY=良好 / RED=要改善）が
    # 凡例なしで意味不明だったため、ヘッダ帯の直下にコンパクトな凡例チップを追加。
    # ヘッダタイトルは幅を切って右側に凡例スペースを確保。
    LEFT_X = 40
    LEFT_W = 380
    MAP_TOP = 158
    MAP_HEADER_H = 28
    add_shape(slide, MSO_SHAPE.RECTANGLE, LEFT_X, MAP_TOP, LEFT_W,
              MAP_HEADER_H, fill=NAVY)
    # タイトル短縮（読み進み順 → 順）で凡例と共存
    add_text(slide, LEFT_X + 16, MAP_TOP, 220,
             '◤ LP構造マップ', 14, bold=True, color=WHITE,
             height_px=MAP_HEADER_H, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    # 凡例チップ（ヘッダ右端に配置）：●良好 / ●要改善
    LEGEND_Y = MAP_TOP + 6
    LEGEND_H = 16
    # 良好チップ（NAVY_LIGHT の丸 + "良好"）
    add_shape(slide, MSO_SHAPE.OVAL, LEFT_X + LEFT_W - 148, LEGEND_Y + 3,
              10, 10, fill=NAVY_LIGHT)
    add_text(slide, LEFT_X + LEFT_W - 132, LEGEND_Y, 40, '良好',
             14, bold=False, color=WHITE, height_px=LEGEND_H,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    # 要改善チップ（RED の丸 + "要改善"）
    add_shape(slide, MSO_SHAPE.OVAL, LEFT_X + LEFT_W - 78, LEGEND_Y + 3,
              10, 10, fill=RED)
    add_text(slide, LEFT_X + LEFT_W - 62, LEGEND_Y, 60, '要改善',
             14, bold=False, color=WHITE, height_px=LEGEND_H,
             vertical_anchor=MSO_ANCHOR.MIDDLE)

    sections = visual_data.get('sections', [])[:9]
    # [FB対応 2026-07-20 v15.6] 空値フォールバック実装
    # sectionsが空の場合、代替テキストを1件だけ描画
    if not sections:
        sections = [{'no': 1, 'label': 'データ未取得',
                     'desc': 'LP構造データが渡されませんでした',
                     'has_issue': False}]
    n_sec = len(sections)
    MAP_BODY_TOP = MAP_TOP + MAP_HEADER_H
    MAP_BODY_H = 276 - MAP_HEADER_H  # =248 (158+28+248=434, 行動フロー top=444 直前)
    add_shape(slide, MSO_SHAPE.RECTANGLE, LEFT_X, MAP_BODY_TOP, LEFT_W,
              MAP_BODY_H, fill=LIGHT_GRAY,
              line=BORDER_GRAY, line_width_pt=0.5)

    if n_sec > 0:
        # 行高は最小26pxを確保し、name｜desc を1行 runs 構造で表示
        row_h = max(26, MAP_BODY_H // n_sec)
        for i, sec in enumerate(sections):
            y = MAP_BODY_TOP + 4 + i * row_h
            # [FB対応 2026-06-30] has_issue 判定を status からも派生
            # status='!' or '✕' → 課題あり、'✓' → 課題なし
            status_val = sec.get('status', '')
            has_issue = (
                sec.get('has_issue', False)
                or status_val in ('!', '✕', '×', 'NG', 'fail', '離脱')
            )
            # 番号サークル（課題ありは赤、なしは紺）
            circle_color = RED if has_issue else NAVY
            add_shape(slide, MSO_SHAPE.OVAL, LEFT_X + 8, y, 22, 22,
                      fill=circle_color)
            add_text(slide, LEFT_X + 8, y, 22,
                     str(sec.get('no', i + 1)), 14, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER,
                     height_px=22, vertical_anchor=MSO_ANCHOR.MIDDLE)
            # 名前｜説明（1ブロック内に runs で混在、1行で表示）
            # [FB対応 2026-06-30] スキーマ不整合バグ修正
            # visual_data_schema.md と test データは `label` を使うため、
            # まず `label` を優先し、後方互換で `name` も読む。
            # [FB対応 2026-07-03] status 記号（✓/!/○）は凡例なしで意味不明のため、
            # 状態はサークル色（NAVY=良好/RED=要改善）で伝える設計に変更。
            # desc フィールドがあればそれを表示、なければ空文字（記号は使わない）。
            # [FB対応 2026-07-08 v12] LP構造マップで name=「FV/メインビジュアル」(10字)
            # が 8字上限で「FV/メインビジュ…」に末尾途切れを発生。
            # 実描画基準（14pt × 幅約300px, LEFT_W - 48）で 8→ 12字へ緩和。
            name = _truncate_full(
                sec.get('label') or sec.get('name', ''), 12)
            desc_raw = sec.get('desc', '')  # status は表示から除外
            desc = _truncate_full(desc_raw, 16)
            _add_multi_run_box(slide, LEFT_X + 38, y, LEFT_W - 48, row_h - 2, [
                {'runs': [
                    {'text': name, 'size': 14, 'bold': True, 'color': TEXT},
                    {'text': f'  {desc}' if desc else '',
                     'size': 14, 'bold': False, 'color': SUB_TEXT},
                ], 'line_height': 1.1},
            ])

    # ---- 中央上：総評 ----
    # [FB対応 2026-07-05 v10] v9で SUMMARY_H=66 に縮めた結果、
    # 内部領域 66-36=30px に対し summary_text=40字は 14pt×1.35≈
    # 26px×2行=52px 必要で 22px はみ出し→ ISSUE_TOP(234) に衝突。
    # 根本対策：
    #   1) summary上限を 40→ 24字 (1行収容確定)
    #   2) SUMMARY_H を 66→ 76（+10px）で内部領域 40px を確保
    #      → 1行(26px)+上下余白の安全値
    #   3) ISSUE_H を 200→ 190（-10px）で縦合計を維持
    #      SUMMARY_END = 158+76 = 234（旧ISSUE_TOPと同値）
    #      ISSUE_TOP = 234+10 = 244、ISSUE_END = 244+190 = 434（同一）
    #      内部領域 190-28-14 = 148px → 22字×3件×26px = 78px 安全
    # [FB対応 2026-07-09 v14・1行運用確定] 大原則：総評は"1行運用"を絶対条件とする。
    # 2行折返しが起きたら被り事故が発生するため、コード側で「1行に確実収容できる
    # 最大文字数」を上限に設定し、Instructions側でAI生成コピーがその範囲内で
    # 書くよう指示する。折返しを許容するのではなく発生させない設計に統一。
    #
    # 実測結果（14pt Meiryo × width 352px, 2026-07-09）：
    #   全角 17字（換算17.0）→ 1行内収容 ✅ 限界値
    #   全角 18字（換算18.0）→ 2行折返し ❌
    #   半角混 18字（換算16.0）→ 1行内収容 ✅ 安全
    #   半角混 19字（換算17.5）→ 2行折返し ❌
    # ⇒ 全角換算 17 が明確な境界値。_truncate_full(17) が正解。
    #
    # 座標最適化（v13 SUMMARY_H=84 は2行想定で過剰。1行分に最適化）：
    #   SUMMARY_H 84→ 62（ヘッダ28px + 本文帯34px = 1行安全収容）
    #   ISSUE_TOPギャップ +14→ +12（元設計値）
    #   ISSUE_H 178→ 200（削減分22を吸収、Top3内部余白拡大＝可読性向上）
    #   下端 158+62+12+200 = 432（右カラム下端434とほぼ同値・±2px許容）
    # [v15/2026-07-12] 総評本文帯の天地余白拡大（62→68）：
    # v14 では本文帯 34px、テキスト vertical_anchor=MIDDLE でも
    # 「詰まっている感」が残っていた。本文帯を 40px に拡大して天地に
    # 各6px の余裕を確保する。下流連動：ISSUE_TOP=238、ISSUE_END=438
    # （行動フロー帯 top=444 との安全余白 6px）。
    # [FB対応 2026-07-16 v15.5・入江さん提案] 3カラム再配分：
    # 中央（総評・最重要課題）を主役化するため、右（強み）から80px移譲。
    # 強み：3点固定・28字上限で折返し許容の設計 → 336px幅でも表示可
    # 中央：内部描画幅 380→460、20字（380px）に52pxの安全マージン確保
    CENTER_X = 432
    CENTER_W = 460  # 380 → 460（+80px拡張）
    SUMMARY_TOP = 158
    SUMMARY_H = 68  # v14=62 → v15=68（本文帯 34→40px, 天地余白 +6px）
    add_shape(slide, MSO_SHAPE.RECTANGLE, CENTER_X, SUMMARY_TOP, CENTER_W,
              28, fill=ORANGE)
    add_text(slide, CENTER_X + 16, SUMMARY_TOP, CENTER_W - 32,
             '◤ 総評', 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, CENTER_X, SUMMARY_TOP + 28, CENTER_W,
              SUMMARY_H - 28, fill=ORANGE_LIGHT,
              line=BORDER_GRAY, line_width_pt=0.5)
    # [v14] 全角換算17が実測境界値（全角17字ギリギリ1行、18字で折返し発生）。
    # Instructions側でも「総評は17字以内・1行運用」をルール化。
    # [v15/2026-07-12] 本文帯（34px）内でテキスト縦位置が上寄り（上6px/下2px）
    # だったため、テキストボックスを本文帯フル高（34px）で作成し
    # vertical_anchor=MIDDLE で縦中央寄せ。上下余白の視覚的不均衡を解消。
    summary_text = _truncate_full(visual_data.get('summary', ''), 17)
    _add_multi_run_box(slide, CENTER_X + 14, SUMMARY_TOP + 28,
                       CENTER_W - 28, SUMMARY_H - 28, [
        {'text': summary_text, 'size': 14, 'bold': False,
         'color': TEXT, 'line_height': 1.35},
    ], vertical_anchor=MSO_ANCHOR.MIDDLE)

    # ---- 中央下：最重要課題（Top3 or Top4） ----
    # [FB対応 2026-07-03] Top3レイアウト崩れ修正：line_height 1.35 + 22字切詰
    # [FB対応 2026-07-04 v9・案C] ISSUE_H 172→200（内部158px確保）
    # [FB対応 2026-07-16] 22字境界LibreOffice折返し回避で20字に厳格化
    # [FB対応 2026-07-16 v15.3] Top3/Top4動的対応：
    #   - Instructions v3.2 で C=4項目化に伴い、GPTs側が独自書換して
    #     ヘッダーが黒文字化する副作用が実測で確認された
    #   - コード側で件数判定→ヘッダー文字列自動生成でGPTs書換余地を排除
    #   - 4項目時は ISSUE_H を 200→240 に拡張し尾切れを根本解消
    # top_issues (複数) と top_issue (単数) の両方を扱える後方互換実装
    top_issues_all = (
        visual_data.get('top_issues')
        or ([visual_data.get('top_issue')] if visual_data.get('top_issue') else [])
        or visual_data.get('priority_issues', [])
    )
    # ランクに応じて表示件数を決定（B=3, C=4, D=5 は将来対応）
    # 現状は Top3/Top4 の2パターンをサポート
    n_issues = min(len(top_issues_all), 4)  # 最大4項目まで
    top_issues = top_issues_all[:n_issues]
    # ヘッダー文字列をコード側で動的生成（GPTs書換によるカラー崩壊防止）
    issue_header_text = f'⚠ 最重要課題（Top{n_issues}）'
    # [FB対応 2026-07-16 v15.5] 中央カラム幅拡張（380→460）により20字が
    # 1行完結する境界（内部432px確保）。ISSUE_H は元の200pxに復帰。
    # 3項目：26px×3 + space8pt×2 = 約94px（内部158pxに余裕64px）
    # 4項目：26px×4 + space2pt×3 = 約112px（内部158pxに余裕46px）
    ISSUE_H = 200
    space_after = 8 if n_issues <= 3 else 2
    ISSUE_TOP = SUMMARY_TOP + SUMMARY_H + 12
    add_shape(slide, MSO_SHAPE.RECTANGLE, CENTER_X, ISSUE_TOP, CENTER_W,
              28, fill=RED)
    add_text(slide, CENTER_X + 16, ISSUE_TOP, CENTER_W - 32,
             issue_header_text, 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, CENTER_X, ISSUE_TOP + 28, CENTER_W,
              ISSUE_H - 28, fill=RED_LIGHT,
              line=BORDER_GRAY, line_width_pt=0.5)
    # [FB対応 2026-07-16 v15.5] 中央カラム460px化で内部432px確保。
    # 20字（380px）が1行完結する境界。安全マージン+52pxで折返し防止。
    issue_paragraphs = [
        {'text': f'{i + 1}. ' + _truncate_full(s, 20),
         'size': 14, 'bold': False, 'color': TEXT,
         'line_height': 1.2, 'space_after_pt': space_after}
        for i, s in enumerate(top_issues)
    ]
    _add_multi_run_box(slide, CENTER_X + 14, ISSUE_TOP + 34,
                       CENTER_W - 28, ISSUE_H - 42, issue_paragraphs)

    # ---- 右：強み3点 ----
    # [FB対応 2026-07-16 v15.5・入江さん提案] 右カラム縮小・80px中央へ移譲。
    # 強み3点固定・28字上限の設計上、内部304pxで折返し許容表示は成立。
    RIGHT_X = 904  # 824 → 904（80px右シフト）
    RIGHT_W = 336  # 416 → 336（-80px縮小）
    STR_TOP = 158
    STR_H = 276  # 元の値に復帰（ISSUE_H拡張が不要になったため）
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, STR_TOP, RIGHT_W,
              28, fill=NAVY)
    add_text(slide, RIGHT_X + 16, STR_TOP, RIGHT_W - 32,
             '◎ 強み（活かすべき点）', 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, STR_TOP + 28, RIGHT_W,
              STR_H - 28, fill=LIGHT_GRAY,
              line=BORDER_GRAY, line_width_pt=0.5)
    # [FB対応 2026-07-08 v12] 強み項目の「ファーストビュー設/計」のような
    # 語の途中での不自然な折返しを避けるため、実描画基準で 40→ 28字に厳格化。
    # 14pt Meiryo × width 384px（実領域 352px）の実測描画で 1行に収まる上限。
    # [v15/2026-07-12] 折返し発生時の「1文の連続感」と「項目間の分離感」を分離制御：
    #   - line_height 1.7→1.2（折返し行間を詰める＝連続感）
    #   - space_after_pt=8pt を項目間に追加（項目間の分離感）
    # これにより「明確な CTA ボタンとファーストビュー設計」のような28字文が
    # 折返しても、2行目が「同一文の続き」と視認できる。
    strengths = visual_data.get('strengths', [])[:3]
    str_paragraphs = [
        {'text': '● ' + _truncate_full(s, 28),
         'size': 14, 'bold': False, 'color': TEXT,
         'line_height': 1.2, 'space_after_pt': 8}
        for s in strengths
    ]
    _add_multi_run_box(slide, RIGHT_X + 16, STR_TOP + 36,
                       RIGHT_W - 32, STR_H - 44, str_paragraphs)

    # ---- 下：行動フロー6ステップ ----
    # [ビジュアル強化 2026-06-29] ステップ円36→44px、ステータス16→18pt
    # [FB対応 2026-07-02] direction帯62px化に伴い FLOW_H 156→146
    # [FB対応 2026-07-16 v15.5] 中央カラム拡張で情報表示問題を解決したため
    # FLOW_TOP/FLOW_H は元の値（444/146）に復帰。レイアウト整合性を保持。
    FLOW_TOP = 444
    FLOW_H = 146
    FLOW_W = 1200
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, FLOW_TOP, FLOW_W, 28, fill=NAVY)
    add_text(slide, 56, FLOW_TOP, 500,
             '◤ ユーザー行動フローと課題', 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    # [FB対応 2026-07-02] フロー総括をヘッダ右端に移設（従来はステップnoteと
    # 重なっていたため）。40字上限は据置き。GOLDでナビー帯内の視認性確保。
    flow_summary_early = visual_data.get('flow_summary', '')
    if flow_summary_early:
        add_text(slide, 560, FLOW_TOP, FLOW_W - 560 - 16,
                 '→ ' + _truncate_full(flow_summary_early, 40), 14, bold=True,
                 color=GOLD, align=PP_ALIGN.RIGHT,
                 height_px=28, vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, FLOW_TOP + 28, FLOW_W,
              FLOW_H - 28, fill=WHITE,
              line=BORDER_GRAY, line_width_pt=0.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, FLOW_TOP + 28, FLOW_W,
              FLOW_H - 28, fill=NAVY, fill_alpha=3000)

    flow_steps = visual_data.get('flow_steps', [])[:6]
    n_flow = len(flow_steps)
    if n_flow > 0:
        step_gap = 8
        step_w = (FLOW_W - 20 - step_gap * (n_flow - 1) - 60) // n_flow
        arrow_w = 60 // max(1, n_flow - 1) if n_flow > 1 else 0
        step_top = FLOW_TOP + 36
        STEP_D = 44  # 36→44に拡大
        for i, step in enumerate(flow_steps):
            x = 50 + i * (step_w + step_gap + arrow_w)
            status = _normalize_flow_status(step.get('status', '✓'))
            step_color = NAVY if status == '✓' else RED
            # ステップ円
            add_shape(slide, MSO_SHAPE.OVAL, x + step_w // 2 - STEP_D // 2,
                      step_top, STEP_D, STEP_D, fill=step_color)
            add_text(slide, x + step_w // 2 - STEP_D // 2, step_top,
                     STEP_D, status, 18, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, height_px=STEP_D,
                     vertical_anchor=MSO_ANCHOR.MIDDLE)
            # ラベル + note（1ブロック2段落）
            # ステップ円44pxに合わせてラベル開始位置を+8px調整
            label = _truncate_full(step.get('label', ''), 6)
            note = _truncate_full(step.get('note', ''), 14)
            _add_multi_run_box(slide, x, step_top + STEP_D + 4, step_w, 56, [
                {'text': label, 'size': 14, 'bold': True,
                 'color': TEXT, 'align': PP_ALIGN.CENTER,
                 'line_height': 1.1},
                {'text': note, 'size': 14, 'bold': False,
                 'color': SUB_TEXT, 'align': PP_ALIGN.CENTER,
                 'line_height': 1.2},
            ])
            # 矢印（最終ステップ以外）
            if i < n_flow - 1:
                arrow_x = x + step_w + step_gap
                add_text(slide, arrow_x, step_top + 8, arrow_w, '→',
                         20, bold=True, color=ORANGE,
                         align=PP_ALIGN.CENTER, height_px=STEP_D,
                         vertical_anchor=MSO_ANCHOR.MIDDLE)

    # フロー総括はヘッダ右端に移設済み（上のフロー帯ヘッダ内 GOLD 表示）。
    # 下段位置はステップ note と重なっていたため廃止。

    # 改善方向フッター帯
    _visual_direction_footer(slide, visual_data.get('direction', ''))
    _add_footer(slide, page_num, total, author=author)
    slide1 = slide

    # ==============================================================
    # スライド 2/3 : スコア視覚化（10項目表 + 強み/課題コンパクト）
    # ==============================================================
    slide = _blank_slide(prs)
    _visual_title(slide, slide_no,
                  'スコア視覚化：10項目評価と強み・最優先課題の俯瞰')
    _visual_meta_header(slide, visual_data, top=90)

    # ---- 左：10項目スコア表（フェーズ1=表形式、フェーズ2でレーダー差し替え） ----
    SCORE_TOP = 158
    SCORE_W = 580
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, SCORE_TOP, SCORE_W, 28, fill=NAVY)
    add_text(slide, 56, SCORE_TOP, SCORE_W - 32,
             '◤ 10項目スコア（満点5・ギャップ可視化）',
             14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)

    scores = visual_data.get('scores', [])[:10]
    # [FB対応 2026-07-20 v15.6] 空値フォールバック実装
    # scoresが空の場合、代替表示を提供（バーは描かず、警告文を1行）
    n_scores = len(scores)
    SCORE_BODY_TOP = SCORE_TOP + 28
    # [FB対応 2026-07-03] 改善方向帯top=594 と18px重複していたため
    # SCORE_BODY_H を 426→400 に短縮。end=158+28+400=586（帯top594との
    # 間に8px安全余白確保）。
    SCORE_BODY_H = 400
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, SCORE_BODY_TOP, SCORE_W,
              SCORE_BODY_H, fill=WHITE,
              line=BORDER_GRAY, line_width_pt=0.5)

    if n_scores > 0:
        row_h = min(42, SCORE_BODY_H // n_scores)
        for i, item in enumerate(scores):
            y = SCORE_BODY_TOP + 6 + i * row_h
            sv = item.get('score', '－')
            mv = item.get('max', 5)
            is_na = (sv == '－' or sv == '－/5')
            ratio = (sv / mv) if (not is_na and mv) else 0
            # スコア色分け：2点以下=赤、4点以上=紺、それ以外=グレー
            if is_na:
                bar_color = BORDER_GRAY
                score_color = SUB_TEXT
                score_bold = False
            elif sv <= 2:
                bar_color = RED
                score_color = RED
                score_bold = True
            elif sv >= 4:
                bar_color = NAVY
                score_color = NAVY
                score_bold = True
            else:
                bar_color = ORANGE
                score_color = TEXT
                score_bold = False
            # 項目名
            # [FB対応 2026-06-30] スキーマ不整合バグ修正
            # scores 配列は C-1 と同じ `name` キーを使う（schema 統一）
            # `category` は後方互換のため最終フォールバック
            cat_disp = _truncate_full(
                item.get('name') or item.get('category', ''), 10)
            add_text(slide, 56, y + 8, 180, cat_disp, 14, bold=True,
                     color=TEXT, height_px=22)
            # ★バー視覚化（max=5、塗りつぶし数=sv）
            bar_x = 240
            bar_w = 28
            for star_i in range(5):
                star_color = bar_color if (not is_na and star_i < sv) else BORDER_GRAY
                add_shape(slide, MSO_SHAPE.RECTANGLE,
                          bar_x + star_i * (bar_w + 4), y + 12,
                          bar_w, 16, fill=star_color)
            # スコア数値
            score_text = '－/5' if is_na else f'{sv}/5'
            add_text(slide, bar_x + 5 * (bar_w + 4) + 8, y + 8, 50, score_text,
                     14, bold=score_bold, color=score_color, height_px=22)
    else:
        # [FB対応 2026-07-20 v15.6] scoresが空の時のフォールバック描画
        add_text(slide, 56, SCORE_BODY_TOP + 20, SCORE_W - 32,
                 '⚠ 10項目スコアデータが渡されていません', 14, bold=True,
                 color=RED, height_px=22)
        add_text(slide, 56, SCORE_BODY_TOP + 50, SCORE_W - 32,
                 'visual_dataに`scores`キー（C-1のscoresと同一10件）を必ず含めてください',
                 12, color=SUB_TEXT, height_px=40)

    # ---- 右上：強み3点（コンパクト） ----
    RIGHT_X = 640
    RIGHT_W = 600
    STR_TOP = 158
    STR_H = 218
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, STR_TOP, RIGHT_W,
              28, fill=NAVY)
    add_text(slide, RIGHT_X + 16, STR_TOP, RIGHT_W - 32,
             '◎ 強み（活かすべき点）', 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, STR_TOP + 28, RIGHT_W,
              STR_H - 28, fill=LIGHT_GRAY,
              line=BORDER_GRAY, line_width_pt=0.5)
    # [v15/2026-07-12] リスト内折返しの連続感と項目間の分離感を分離制御
    strengths = visual_data.get('strengths', [])[:3]
    # [FB対応 2026-07-20 v15.6] 空値フォールバック実装
    if not strengths:
        strengths = ['強みデータが渡されませんでした（visual_dataにstrengthsキー必須）']
    str_paragraphs = [
        {'text': '● ' + _truncate_full(s, 40),
         'size': 14, 'bold': False, 'color': TEXT,
         'line_height': 1.2, 'space_after_pt': 8}
        for s in strengths
    ]
    _add_multi_run_box(slide, RIGHT_X + 16, STR_TOP + 36,
                       RIGHT_W - 32, STR_H - 44, str_paragraphs)

    # ---- 右下：最優先課題3点（コンパクト） ----
    ISSUE_TOP = STR_TOP + STR_H + 12
    # [FB対応 2026-07-02] direction帯62px化(top=594)に伴い ISSUE_H を 220→198
    # に縮小し下端を 586 に (安全余白8px)。
    #   ISSUE_TOP = STR_TOP(158) + STR_H(218) + 12 = 388
    #   ISSUE_TOP + ISSUE_H = 388 + 198 = 586 < 594（direction帯 top）
    ISSUE_H = 198
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, ISSUE_TOP, RIGHT_W,
              28, fill=RED)
    add_text(slide, RIGHT_X + 16, ISSUE_TOP, RIGHT_W - 32,
             '⚠ 最優先で直すべき点', 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, ISSUE_TOP + 28, RIGHT_W,
              ISSUE_H - 28, fill=RED_LIGHT,
              line=BORDER_GRAY, line_width_pt=0.5)
    # [v15/2026-07-12] リスト内折返しの連続感と項目間の分離感を分離制御
    priority_issues = visual_data.get('priority_issues', [])[:3]
    # [FB対応 2026-07-20 v15.6] 空値フォールバック実装
    if not priority_issues:
        priority_issues = ['最優先課題データが渡されませんでした（visual_dataにpriority_issuesキー必須）']
    issue_paragraphs = [
        {'text': f'{i + 1}. ' + _truncate_full(s, 50),
         'size': 14, 'bold': False, 'color': TEXT,
         'line_height': 1.2, 'space_after_pt': 8}
        for i, s in enumerate(priority_issues)
    ]
    _add_multi_run_box(slide, RIGHT_X + 16, ISSUE_TOP + 36,
                       RIGHT_W - 32, ISSUE_H - 44, issue_paragraphs)

    # 改善方向フッター帯
    _visual_direction_footer(slide, visual_data.get('direction', ''))
    _add_footer(slide, page_num + 1, total, author=author)
    slide2 = slide

    # ==============================================================
    # スライド 3/3 : Before/After Top3 ハイライト
    # ==============================================================
    slide = _blank_slide(prs)
    _visual_title(slide, slide_no,
                  'Before/After ハイライト：優先度の高い改善Top3')

    highlights = visual_data.get('highlights', [])[:3]
    n_h = len(highlights)

    HL_TOP = 92
    HL_GAP = 12
    # [FB対応 2026-07-02] direction帯62px化(top=594)に伴い HL_AREA_H を
    # 520→494 に縮小し下端を 586 に (安全余白8px)。92 + 494 = 586。
    HL_AREA_H = 494
    hl_h = (HL_AREA_H - HL_GAP * max(0, n_h - 1)) // max(1, n_h)

    prio_color_map = {'高': PRIO_RED, '中': PRIO_ORANGE, '低': PRIO_GRAY,
                      'S': PRIO_RED, 'A': PRIO_RED,
                      'B': PRIO_ORANGE, 'C': PRIO_GRAY}

    for i, h in enumerate(highlights):
        y = HL_TOP + i * (hl_h + HL_GAP)
        priority = h.get('priority', '中')
        prio_col = prio_color_map.get(priority, PRIO_ORANGE)
        effort = h.get('effort', '中')

        # カード背景 + 左色帯
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 40, y, 1200,
                  hl_h, fill=WHITE,
                  line=BORDER_GRAY, line_width_pt=1)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 40, y, 6, hl_h, fill=prio_col)

        # ヘッダ：番号 + タイトル + 優先度バッジ + 工数
        head_top = y + 10
        add_shape(slide, MSO_SHAPE.OVAL, 60, head_top, 32, 32, fill=prio_col)
        add_text(slide, 60, head_top + 4, 32, str(h.get('no', i + 1)),
                 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=24, vertical_anchor=MSO_ANCHOR.MIDDLE)
        title = _truncate_full(h.get('title', ''), 30)
        # [FB対応 2026-07-16 v15.3] 優先度x=820→780 に伴い、タイトル幅を
        # 700→660 に縮小（右端 762、優先度左端 780 との安全余白 18px）
        add_text(slide, 102, head_top + 4, 660, title, 18,
                 bold=True, color=NAVY, height_px=28)
        # [FB対応 2026-07-16 v15.3] 優先度・工数・箇所のx座標配分を再設計。
        # 従来：優先度820/工数900/箇所1040-1220（箇所幅180pxで日本語折返し多発）
        # 改善：優先度780/工数860/箇所960-1220（箇所幅260pxに拡張）
        # 優先度バッジ
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 780, head_top + 2,
                  70, 28, fill=prio_col)
        add_text(slide, 780, head_top + 5, 70, priority, 14,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=22)
        # 工数（幅130→90に縮小、コンパクト化）
        add_text(slide, 860, head_top + 5, 90, f'工数：{effort}', 14,
                 color=SUB_TEXT, height_px=22)
        # 改善箇所（幅180→260に拡張、日本語テキストの折返し防止）
        target_area = _truncate_full(h.get('target_area', ''), 20)
        add_text(slide, 960, head_top + 5, 260,
                 f'箇所：{target_area}', 14, color=SUB_TEXT, height_px=22)

        # 区切り線
        add_shape(slide, MSO_SHAPE.RECTANGLE, 60, head_top + 38, 1170, 1,
                  fill=BORDER_GRAY)

        # Before/After 2カラム
        body_top = head_top + 50
        body_h = hl_h - 60
        col_w = 568

        # Before（赤帯、letter_spacing削除でBEFORE全文表示）
        add_shape(slide, MSO_SHAPE.RECTANGLE, 60, body_top, 90, 28, fill=RED)
        add_text(slide, 60, body_top, 90, 'BEFORE', 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=28,
                 vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 150, body_top,
                  col_w - 90, 28, fill=RED_LIGHT)
        before_text = _truncate_full(h.get('before', ''), 38)
        add_text(slide, 158, body_top, col_w - 106, before_text, 14,
                 color=TEXT, height_px=28,
                 vertical_anchor=MSO_ANCHOR.MIDDLE)

        # After（オレンジ帯、letter_spacing削除）
        add_shape(slide, MSO_SHAPE.RECTANGLE, 644, body_top, 90, 28,
                  fill=ORANGE)
        add_text(slide, 644, body_top, 90, 'AFTER', 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=28,
                 vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 734, body_top,
                  col_w - 90, 28, fill=ORANGE_LIGHT)
        after_text = _truncate_full(h.get('after', ''), 38)
        add_text(slide, 742, body_top, col_w - 106, after_text, 14,
                 color=TEXT, height_px=28,
                 vertical_anchor=MSO_ANCHOR.MIDDLE)

    # 改善方向フッター帯
    _visual_direction_footer(slide, visual_data.get('direction', ''))
    _add_footer(slide, page_num + 2, total, author=author)
    slide3 = slide

    return (slide1, slide2, slide3)


# =====================================================================
# build_full_report：C-1 + C-2 + C-3 を1ファイルに統合
# =====================================================================
def build_full_report(diagnosis, proposals, visual_data,
                      author='UI/UX診断 by GPTs'):
    """
    UI診断統合レポートを生成。C-1（2枚）+ C-2（1-2枚）+ C-3（3枚）を
    1つの Presentation に格納して返す。

    Args:
        diagnosis: dict（add_scorecard_onepager と同じスキーマ）
        proposals: dict（add_proposal_onepager と同じスキーマ）
        visual_data: dict（visual_data_schema.md 参照）
        author: フッター著者名

    Returns:
        prs: Presentation オブジェクト（.save() してダウンロード可能）
    """
    prs = create_presentation()

    # 提案件数からC-2のスライド数を推定
    n_proposals = len(proposals.get('proposals', []))
    c2_slides = 2 if n_proposals >= 4 else 1

    # 全体ページ数 = C-1(2) + C-2(1or2) + C-3(3)
    total = 2 + c2_slides + 3

    # C-1: スコアカード (2枚)
    add_scorecard_onepager(prs, diagnosis,
                            page_num=1, total=total,
                            author=author, slide_no='1')

    # C-2: 改善提案リスト (1 or 2枚)
    add_proposal_onepager(prs, proposals,
                          page_num=3, total=total,
                          author=author, slide_no='2')

    # C-3: ビジュアル診断ボード (3枚)
    add_visual_board(prs, visual_data,
                     page_num=3 + c2_slides, total=total,
                     slide_no='3', author=author)

    return prs


# =====================================================================
# 動作確認用（直接実行時）
# =====================================================================
if __name__ == '__main__':
    prs = create_presentation()

    add_cover(prs,
              title='Q2 販売戦略レビュー',
              date='2026年7月15日',
              author='○○株式会社 営業企画部',
              subtitle='上期実績の振り返りと下期に向けた重点施策')

    add_agenda(prs, items=[
        {'title': '市場環境の整理', 'desc': 'マクロ・競合・顧客動向'},
        {'title': '現状の課題', 'desc': '上期で見えた3つの論点'},
        {'title': '優先施策', 'desc': '4象限で整理した優先順位'},
        {'title': '実行スケジュール', 'desc': '下期6か月のロードマップ'},
    ], page_num=2, total=4)

    add_issue_summary(prs,
        title='現状の3つの課題',
        cards=[
            {'no': '01', 'heading': 'リード獲得の伸び悩み',
             'body': '主要チャネルで前年比横ばい。\n新規流入経路の追加検討が必要。'},
            {'no': '02', 'heading': '商談化率の低下',
             'body': 'リードから商談への転換が\n前年比10pt低下。\nスコアリング再設計が必要。'},
            {'no': '03', 'heading': '受注後の解約率増',
             'body': '導入後3か月の解約が増加。\nオンボーディング強化が急務。'},
        ],
        page_num=3, total=4,
        lead='上期実績の分析から、下期に向けた3つの課題を整理しました。',
        conclusion='3課題は独立ではなく「リード→商談→受注→定着」の流れで連動している。')

    add_closing(prs,
                message='Thank you.',
                next_step='次回ミーティング：7月22日（火）\n各施策の責任者から進捗報告をお願いします。',
                contact='問い合わせ：営業企画部 山田',
                page_num=4, total=4)

    prs.save('test_output.pptx')
    print('Generated: test_output.pptx')

    # =====================================================================
    # C-1 / C-2 動作確認サンプル
    # =====================================================================
    # ── C-1: UI診断スコアカード（3枚） ──
    diagnosis_sample = {
        'service_name': 'サンプルECサイト',
        'url': 'https://example.com',
        'input_type': 'URL入力',
        'total_score': 32,
        'rank': 'B',
        'rank_label': '標準的（改善余地あり）',
        'scores': [
            {'category': 'ファーストビュー訴求',  'score': 4.0, 'max': 5, 'comment': 'メインビジュアルは明快。サブコピーで価値訴求を補強したい。'},
            {'category': 'ナビゲーション設計',    'score': 3.5, 'max': 5, 'comment': 'グローバルナビは整理されているが、カテゴリ深度がやや深い。'},
            {'category': '視覚的階層',            'score': 4.5, 'max': 5, 'comment': '見出し階層が明確で読みやすい。'},
            {'category': '検索・絞り込み',        'score': 2.0, 'max': 5, 'comment': '絞り込み条件の保存・組み合わせができない。'},
            {'category': '商品詳細の情報量',      'score': 3.0, 'max': 5, 'comment': 'サイズ感や使用シーンの情報が不足している。'},
            {'category': 'カート/購入導線',       'score': 4.0, 'max': 5, 'comment': 'カート遷移はスムーズ。ゲスト購入も可能。'},
            {'category': 'フォーム使いやすさ',    'score': 2.5, 'max': 5, 'comment': '入力エラー表示が遅延し、修正箇所が分かりにくい。'},
            {'category': 'モバイル最適化',        'score': 3.5, 'max': 5, 'comment': 'タップ領域は十分。横スクロールが一部発生。'},
            {'category': '表示速度',              'score': 3.0, 'max': 5, 'comment': '初回読み込みが3.2秒。画像最適化が必要。'},
            {'category': 'アクセシビリティ',      'score': 2.0, 'max': 5, 'comment': 'コントラスト比不足の箇所が散見される。alt属性も欠落多い。'},
        ],
        'strengths': [
            'ファーストビューのメインビジュアルが明快で価値が伝わる',
            '見出し階層が整理されており情報の優先順位が分かりやすい',
            'カート→購入の導線がスムーズで離脱が起きにくい',
        ],
        'priority_issues': [
            '検索・絞り込みの組み合わせ条件が保存できず再操作の手間が大きい',
            'フォーム入力時のエラー表示が遅延し、入力者が修正箇所を特定しにくい',
            'アクセシビリティ（コントラスト・alt属性）の改善が必要',
        ],
        'conclusion': (
            '総合32/50点・ランクBで、基本的なUI設計は整っているが改善余地が大きい状態です。\n'
            '特に「検索・絞り込み」「フォーム」「アクセシビリティ」の3領域は緊急度が高く、'
            'いずれもCVR・離脱率に直接影響する論点です。\n'
            '一方、ファーストビュー訴求や視覚的階層は強みであり、これらを軸に他領域の改善を進めることで'
            '効率的なUX向上が見込めます。'
        ),
    }

    prs2 = create_presentation()
    add_scorecard_overview(prs2, diagnosis_sample, page_num=1, total=3)
    add_scorecard_table(prs2, diagnosis_sample, page_num=2, total=3)
    add_scorecard_conclusion(prs2, diagnosis_sample, page_num=3, total=3,
        cta='続けて C-2 で改善提案リストをご確認ください。')
    prs2.save('test_scorecard.pptx')
    print('Generated: test_scorecard.pptx')

    # ── C-2: 改善提案リスト（3枚） ──
    proposals_sample = {
        'service_name': 'サンプルECサイト',
        'proposals': [
            {'no': 1, 'title': '絞り込み条件の保存機能と複数条件の組み合わせ対応',
             'category': '検索・絞り込み', 'priority': 'S',
             'current': '条件選択後に他ページへ移動すると条件がリセットされ、再度入力が必要。\n複数カテゴリの同時選択にも対応していない。',
             'action': 'URLパラメータ／LocalStorageで条件を保持し、戻る操作でも復元する。\nチェックボックスUIに変更し、複数条件のAND選択を可能にする。',
             'expected': '再操作率を約40%削減、検索ページ滞在時間を-25%。\n結果として商品詳細到達率の向上が見込める。'},
            {'no': 2, 'title': 'フォーム入力エラーのリアルタイム表示',
             'category': 'フォーム使いやすさ', 'priority': 'S',
             'current': '送信ボタン押下後にまとめてエラーが表示され、修正箇所が分かりにくい。\nエラー文言も「入力が誤っています」と汎用的で具体性が低い。',
             'action': '各フィールドのblur時にバリデーションを実行し、エラー文言を「具体的に何がおかしいか」「どう直すか」の2点を含めて表示。',
             'expected': 'フォーム離脱率を-15pt、購入完了率を+5pt見込み。'},
            {'no': 3, 'title': 'アクセシビリティ最低基準（WCAG AA）への準拠',
             'category': 'アクセシビリティ', 'priority': 'A',
             'current': '主要テキストのコントラスト比が4.5未満の箇所が多数。\n商品画像のalt属性が空欄またはファイル名のままで、SR利用者が情報を取得できない。',
             'action': 'デザインシステムのカラートークンを見直し、AA基準を満たすよう調整。\nalt属性は「商品名＋色／柄」を必須項目化し、CMS入力時にバリデーションする。',
             'expected': '潜在顧客層（視覚障碍ユーザ・高齢者）への到達拡大。\nSEO面でも画像検索流入の改善が見込める。'},
            {'no': 4, 'title': '商品詳細ページの情報量強化',
             'category': '商品詳細', 'priority': 'A',
             'current': 'サイズ表記が「M」「L」のみで実寸が不明。使用シーン写真も少なく、購入後イメージが湧きにくい。',
             'action': 'サイズ実寸表（cm）と着用モデル身長を必須化。シーン別の使用写真を最低3枚追加する。',
             'expected': '商品詳細→カート遷移率を+8pt。返品率の低下も期待される。'},
        ],
        'summary': '優先度SとAの計4件を3か月以内に着手することで、ランクBからAへの引き上げが現実的。',
    }

    prs3 = create_presentation()
    # 一覧（1枚）＋ 詳細スライド4枚（1スライド1提案）＝ 計5枚
    proposals_list = proposals_sample['proposals']
    total_pages = 1 + len(proposals_list)  # = 5
    add_proposal_overview(prs3, proposals_sample,
                          page_num=1, total=total_pages)
    for idx, proposal in enumerate(proposals_list):
        is_last = (idx == len(proposals_list) - 1)
        add_proposal_detail(
            prs3, proposal,
            page_num=2 + idx,
            total=total_pages,
            summary=proposals_sample['summary'] if is_last else None,
        )
    prs3.save('test_proposals.pptx')
    print('Generated: test_proposals.pptx')

    # =====================================================================
    # C-1 onepager / C-2 onepager 動作確認サンプル
    # （添付参照デザイン: SpeakUpEnglish）
    # =====================================================================
    onepager_diagnosis = {
        'service_name': 'スピークアップ英会話',
        'input_type': 'スクリーンショット画像',
        'total_score': 30,
        'rank': 'B',
        'rank_label': '標準的（改善余地あり）',
        'scores': [
            {'category': 'ファーストビュー',   'score': 3, 'max': 5,
             'comment': '清潔感はあるが、体験価値と成果が一瞬で伝わりにくい'},
            {'category': 'キャッチコピー',     'score': 3, 'max': 5,
             'comment': '開始喚起は明快だが、選ばれる理由が弱い'},
            {'category': 'CTA設計',           'score': 1, 'max': 5,
             'comment': 'FV・中盤に主要CTAがなく行動導線が弱い'},
            {'category': '信頼性・権威性',     'score': 3, 'max': 5,
             'comment': '実績数値はあるが、根拠や証拠の厚みが不足'},
            {'category': 'フォーム設計',       'score': 1, 'max': 5,
             'comment': '問い合わせ導線はあるが入力前の不安解消が少ない'},
            {'category': 'レスポンシブ',       'score': 3, 'max': 5,
             'comment': 'PC表示は整っているが、SPでCTA固定が必要'},
            {'category': '読みやすさ',         'score': 5, 'max': 5,
             'comment': '余白・行間・色数が整理され読みやすい'},
            {'category': '情報設計',           'score': 3, 'max': 5,
             'comment': '特徴と料金は見やすいが、比較・成果・不安解消が不足'},
            {'category': 'ブランド一貫性',     'score': 5, 'max': 5,
             'comment': '青基調で教育サービスらしい安心感がある'},
            {'category': '表示速度・技術',     'score': 3, 'max': 5,
             'comment': '画像中心のため最適化余地がある'},
        ],
        'strengths': [
            '余白が広く、読みやすいレイアウトで安心感がある',
            '料金・特徴・講師情報が整理され、サービス概要を把握しやすい',
            '青基調の配色で教育サービスらしい信頼感が出ている',
        ],
        'priority_issues': [
            'ファーストビューに主要CTAがなく、次の行動が分かりにくい',
            '受講後の成果・具体的なベネフィットが弱く、比較検討時の決め手に欠ける',
            '受講生の声が短く、信頼材料としての説得力が不足している',
        ],
        'conclusion': '興味を持ったユーザーも、無料体験への問い合わせ前に離脱する可能性が高い。FV内CTAと成果訴求・信頼材料の3点を最優先で改善する。',
    }

    prs4 = create_presentation()
    add_scorecard_onepager(prs4, onepager_diagnosis,
                            page_num=1, total=1, slide_no='1')
    prs4.save('test_onepager_scorecard.pptx')
    print('Generated: test_onepager_scorecard.pptx')

    onepager_proposals = {
        'service_name': 'スピークアップ英会話',
        'proposals': [
            {'no': 1, 'title': 'ファーストビューに無料体験CTAを追加',
             'priority': '高', 'effort': '小',
             'target_area': 'ファーストビュー',
             'issue': '興味を持った直後に行動できない',
             'before': '見出しと説明文はあるが、主要ボタンが見当たらない',
             'after': '無料体験・料金確認の2ボタンをFV内に配置',
             'target_score_item': 'CTA設計'},
            {'no': 2, 'title': 'キャッチコピーを成果訴求型に変更',
             'priority': '高', 'effort': '小',
             'target_area': 'メインコピー',
             'issue': '学習後のメリットが弱く差別化しにくい',
             'before': '「英語、始めませんか？」で汎用的な印象',
             'after': '忙しい人でも続く、成果が見えるオンライン英会話として訴求',
             'target_score_item': 'キャッチコピー'},
            {'no': 3, 'title': '実績数値に根拠と補足を追加',
             'priority': '中', 'effort': '中',
             'target_area': '実績帯',
             'issue': '数値はあるが信頼の裏付けが不足',
             'before': '開校年・採用率・段階数・時間帯のみ表示',
             'after': '累計受講者数、継続率、満足度など検討材料を追加',
             'target_score_item': '信頼性・権威性'},
            {'no': 4, 'title': '受講生の声を具体化',
             'priority': '中', 'effort': '小',
             'target_area': '受講生の声',
             'issue': '短文のみで利用シーンや成果が伝わりにくい',
             'before': '「続けられています」など一言コメント中心',
             'after': '課題・受講理由・成果を含む3行レビューへ拡張',
             'target_score_item': '信頼性・権威性'},
            {'no': 5, 'title': '料金セクション下に比較とCTAを追加',
             'priority': '高', 'effort': '中',
             'target_area': '料金プラン',
             'issue': 'プラン選択後の次アクションが弱い',
             'before': '料金カードの提示で終わり、相談導線が離れている',
             'after': 'おすすめプラン表示と無料相談CTAを直下に配置',
             'target_score_item': 'CTA設計'},
        ],
        'summary': '最優先の一手：ファーストビューに無料体験CTAを追加し、直下に成果・安心材料を3点で提示する。',
    }

    prs5 = create_presentation()
    add_proposal_onepager(prs5, onepager_proposals,
                           page_num=1, total=1, slide_no='2')
    prs5.save('test_onepager_proposals.pptx')
    print('Generated: test_onepager_proposals.pptx')


# =====================================================================
# ▼▼▼ v17 追加ブロック（P1｜category / breakdown / comparison） ▼▼▼
# 既存の C-1〜C-3 描画ロジックには一切触れない（後方互換完全維持）
# =====================================================================

# ---------------------------------------------------------------------
# v17-0｜拡張パターン仕様（原本3属性に対する拡張層）
#   出典: v35_core_extended_pattern_definitions.md 拡張定義集約表
#   ※ 原本 DIAGRAM_PATTERNS は書き換えない（2層構造を維持）
# ---------------------------------------------------------------------
DIAGRAM_PATTERN_SPEC = {
    'category': {
        'min_elements': 3, 'max_elements': 6, 'requires_axes': False,
        'direction': 'grid', 'color_gradation': 'uniform_parallel',
        'grid_map': {3: (3, 1), 4: (2, 2), 5: (3, 2), 6: (3, 2)},
    },
    'breakdown': {
        'min_elements': 3, 'max_elements': 7, 'requires_axes': False,
        'direction': 'vertical', 'color_gradation': 'proportional',
    },
    'comparison': {
        'min_elements': 2, 'max_elements': 3, 'requires_axes': False,
        'direction': 'horizontal', 'color_gradation': 'discrete_contrast',
    },
}

# 描画領域（ヘッダ60px／フッター境界660pxの内側）
V17_AREA = {
    'left': 40, 'right': 1240, 'width': 1200,
    'title_top': 90, 'body_top': 138, 'body_bottom': 646,
    'gap': 16,
}

# 警告オーバーライドの閾値（原則②｜B-6 4.4節）
V17_WARNING_SCORE = 40

# カード高さの上限（px）
# v17-fix1（8/24 実機検証）: 上限230pxではラベル+スコア+説明1行=108pxに対し
# 充填率47%となり「カード内が空白だらけ」に見える事象を検出。
# 説明文2行（44px）を収容しつつ充填率60%以上を確保する値として186pxを採用。
V17_CARD_H_MAX = 186


# ---------------------------------------------------------------------
# v17-1｜色ユーティリティ
# ---------------------------------------------------------------------
def hex_to_rgb(hex_str):
    """'#0017C1' / '0017C1' -> RGBColor"""
    from pptx.dml.color import RGBColor
    s = str(hex_str).lstrip('#')
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def get_theme_palette(theme_id):
    """テーマIDから8色階調辞書（hex文字列）を返す。未知IDは SolidGray。"""
    if theme_id not in DIGITAL_AGENCY_PALETTE:
        theme_id = 'SolidGray'
    return DIGITAL_AGENCY_PALETTE[theme_id]


def select_theme_by_project_type(project_type, warning_flag=False):
    """プロジェクトタイプ＋警告フラグからテーマIDを決定（決定論的）。"""
    if warning_flag:
        return 'Red'
    return {
        'corporate': 'Blue',
        'ec':        'Orange',
        'lp':        'Orange',
        'webapp':    'Cyan',
        'media':     'Green',
    }.get(project_type, 'SolidGray')


def _relative_luminance(hex_str):
    """WCAG 相対輝度（0.0〜1.0）"""
    s = str(hex_str).lstrip('#')
    out = []
    for i in (0, 2, 4):
        v = int(s[i:i + 2], 16) / 255.0
        out.append(v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4)
    return 0.2126 * out[0] + 0.7152 * out[1] + 0.0722 * out[2]


def _text_color_on(bg_hex):
    """背景色に対して可読な文字色を返す（WHITE または TEXT）。

    v17 拡張｜拡張定義集約表には規定がないため本実装で追加した。
    根拠: テーマ Green/Cyan/Red/Orange の warning 色は '#CCCCCC'（明色）で、
    白文字を固定すると score<40 の警告セルが判読不能になる。
    輝度0.5を境に自動反転する（決定論的｜同一入力→同一出力）。
    """
    return WHITE if _relative_luminance(bg_hex) < 0.5 else TEXT


def _tier_fill(palette, score, base_key):
    """基準色 + 警告オーバーライド（原則②）。

    score が None の場合は警告判定をスキップする。
    """
    if score is not None and score < V17_WARNING_SCORE:
        return palette['warning']
    return palette[base_key]


# ---------------------------------------------------------------------
# v17-2｜共通描画ヘルパー
# ---------------------------------------------------------------------
def _v17_title(slide, title, palette):
    """パターン共通のタイトル行（ヘッダ帯とは別のスライド内見出し）"""
    add_text(slide, V17_AREA['left'], V17_AREA['title_top'], V17_AREA['width'],
             str(title), 22, bold=True, color=hex_to_rgb(palette['primary']),
             height_px=34)


def _v17_axis_label(slide, text, palette):
    """任意の軸ラベル（比較軸・時間軸等）。斜体は使わない（条項13）。"""
    add_text(slide, V17_AREA['left'], V17_AREA['title_top'] + 36,
             V17_AREA['width'], str(text), 14, bold=True,
             color=hex_to_rgb(palette['secondary']), height_px=22)


def _v17_card(slide, x, y, w, h, fill_hex, palette, radius_px=8):
    """カード図形（塗り＋枠線のみ｜条項7：図形に文字を入れない）"""
    return _add_bg_frame(slide, x, y, w, h,
                         fill=hex_to_rgb(fill_hex),
                         line=hex_to_rgb(palette['midtone']),
                         line_width_pt=1, radius_px=radius_px)


def _v17_rule(slide, x, y, w, h, hex_color):
    """罫線・コネクタ代替の矩形（原則④：基本図形のみ）"""
    from pptx.enum.shapes import MSO_SHAPE
    return add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, max(w, 1), max(h, 1),
                     fill=hex_to_rgb(hex_color))


def _v17_normalize(items, key_label='label'):
    """要素を dict 化して正規化（str 入力も許容）"""
    out = []
    for it in items or []:
        if isinstance(it, dict):
            d = dict(it)
        else:
            d = {key_label: str(it)}
        d.setdefault(key_label, '')
        out.append(d)
    return out


def _v17_report(pattern, drawn, fallback_from=None, notes=None):
    """描画結果レポート（回帰確認・ログ記録用）"""
    return {
        'pattern': pattern,
        'elements_drawn': drawn,
        'fallback_from': fallback_from,
        'notes': list(notes or []),
    }


# ---------------------------------------------------------------------
# v17-3｜category（分類）｜grid / uniform_parallel / 3〜6要素
# ---------------------------------------------------------------------
def draw_category(slide, palette, data, _fallback_from=None):
    """診断結果を並列カテゴリ（グリッド）で描画する。

    Args:
        slide   : python-pptx Slide
        palette : 8色階調辞書（get_theme_palette の戻り値）
        data    : {'title': str,
                   'categories': [{'label': str, 'score': int|None,
                                   'description': str}, ...]}
    Returns:
        dict : 描画レポート

    仕様: 拡張定義集約表 P1-2 行（min3 / max6 / grid / uniform_parallel）
    フォールバック本体のため例外は送出しない（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['category']
    notes = []
    items = _v17_normalize(data.get('categories', []))

    # 要素数の丸め込み（例外を投げない｜原則①）
    if len(items) > spec['max_elements']:
        notes.append('要素数 %d > max %d｜先頭 %d 件で描画'
                     % (len(items), spec['max_elements'], spec['max_elements']))
        items = items[:spec['max_elements']]
    if len(items) < spec['min_elements']:
        notes.append('要素数 %d < min %d｜そのまま描画（劣化描画で通す）'
                     % (len(items), spec['min_elements']))
    if not items:
        _v17_title(slide, data.get('title', ''), palette)
        notes.append('要素0件｜タイトルのみ描画')
        return _v17_report('category', 0, _fallback_from, notes)

    _v17_title(slide, data.get('title', ''), palette)

    n = len(items)
    cols, rows = spec['grid_map'].get(n, (min(n, 3), (n + 2) // 3))
    gap = V17_AREA['gap']
    avail_h = V17_AREA['body_bottom'] - V17_AREA['body_top']
    cell_w = (V17_AREA['width'] - gap * (cols - 1)) / float(cols)
    cell_h = min((avail_h - gap * (rows - 1)) / float(rows), V17_CARD_H_MAX)
    block_h = cell_h * rows + gap * (rows - 1)
    y0 = V17_AREA['body_top'] + (avail_h - block_h) / 2.0

    for i, it in enumerate(items):
        r, c = i // cols, i % cols
        in_row = min(cols, n - r * cols)
        row_w = cell_w * in_row + gap * (in_row - 1)
        x = V17_AREA['left'] + (V17_AREA['width'] - row_w) / 2.0 + c * (cell_w + gap)
        y = y0 + r * (cell_h + gap)

        # uniform_parallel: 全セル同一の基準色（並列性の担保）＋警告オーバーライド
        score = it.get('score')
        fill_hex = _tier_fill(palette, score, 'secondary')
        fg = _text_color_on(fill_hex)
        _v17_card(slide, x, y, cell_w, cell_h, fill_hex, palette)

        pad = 16
        add_text(slide, x + pad, y + 12, cell_w - pad * 2,
                 str(it.get('label', '')), 16, bold=True, color=fg, height_px=26)
        if score is not None:
            add_text(slide, x + pad, y + 42, cell_w - pad * 2,
                     '%s%%' % score, 26, bold=True, color=fg, height_px=40)
        desc_top = y + (86 if score is not None else 46)
        desc = str(it.get('description', ''))
        if desc:
            add_text(slide, x + pad, desc_top, cell_w - pad * 2, desc, 14,
                     color=fg, height_px=max(int(cell_h - (desc_top - y) - 12), 22),
                     line_height=1.4)

    return _v17_report('category', n, _fallback_from, notes)


# ---------------------------------------------------------------------
# v17-4｜breakdown（分解）｜vertical / proportional / 3〜7要素
# ---------------------------------------------------------------------
def _breakdown_tier(ratio):
    """構成比から色階調キーを決定（proportional｜決定論的）"""
    if ratio >= 0.30:
        return 'primary'
    if ratio >= 0.20:
        return 'secondary'
    if ratio >= 0.10:
        return 'midtone'
    return 'light'


def draw_breakdown(slide, palette, data):
    """全体を構成要素に分解して縦積みで描画する（構成比連動の色階調）。

    Args:
        data : {'title': str,
                'whole': {'label': str, 'value': int|float|None},
                'components': [{'label': str, 'value': int|float,
                                'score': int|None, 'note': str}, ...]}
    仕様: 拡張定義集約表 P1-3 行（min3 / max7 / vertical / proportional）
    要素数逸脱・データ不正時は draw_category へフォールバック（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['breakdown']
    notes = []
    comps = _v17_normalize(data.get('components', []))

    # 数値健全性チェック
    values = []
    bad = False
    for c in comps:
        try:
            v = float(c.get('value'))
        except (TypeError, ValueError):
            bad = True
            break
        if v < 0:
            bad = True
            break
        values.append(v)
    total_in = sum(values) if values else 0.0

    if (not (spec['min_elements'] <= len(comps) <= spec['max_elements'])
            or bad or total_in <= 0):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': c.get('label', ''),
                            'score': c.get('score'),
                            'description': str(c.get('note', ''))} for c in comps],
        }, _fallback_from='breakdown')
        rep['notes'].insert(0, 'breakdown 前提不成立（要素数=%d／数値不正=%s／合計=%s）｜'
                            'category へフォールバック' % (len(comps), bad, total_in))
        return rep

    _v17_title(slide, data.get('title', ''), palette)

    whole = data.get('whole') or {}
    whole_label = str(whole.get('label', '全体'))
    whole_value = whole.get('value')
    if whole_value is None:
        whole_value = total_in

    # 上段｜全体ボックス
    wh_h = 56
    wh_y = V17_AREA['body_top']
    _v17_card(slide, V17_AREA['left'], wh_y, V17_AREA['width'], wh_h,
              palette['primary'], palette)
    fg_w = _text_color_on(palette['primary'])
    add_text(slide, V17_AREA['left'] + 20, wh_y + 15, 700, whole_label, 18,
             bold=True, color=fg_w, height_px=26)
    add_text(slide, V17_AREA['left'] + V17_AREA['width'] - 320, wh_y + 13, 300,
             _fmt_num(whole_value), 20, bold=True, color=fg_w, height_px=28,
             align=_pp_right())

    # 下段｜構成要素（縦積み・幅を構成比に比例）
    n = len(comps)
    gap = 10
    top = wh_y + wh_h + 22
    avail = V17_AREA['body_bottom'] - top
    row_h = min((avail - gap * (n - 1)) / float(n), 62)
    label_w = 300
    bar_x = V17_AREA['left'] + label_w + 12
    bar_max_w = V17_AREA['width'] - label_w - 12 - 150

    # 左スパイン（親→子の接続｜基本図形のみ）
    spine_bottom = top + row_h * n + gap * (n - 1) - row_h / 2.0
    _v17_rule(slide, V17_AREA['left'] + 24, wh_y + wh_h,
              3, max(spine_bottom - (wh_y + wh_h), 1), palette['midtone'])

    for i, c in enumerate(comps):
        v = values[i]
        ratio = v / total_in
        y = top + i * (row_h + gap)
        tier = _breakdown_tier(ratio)
        fill_hex = _tier_fill(palette, c.get('score'), tier)
        fg = _text_color_on(fill_hex)

        # 横枝（スパイン→ラベル）
        _v17_rule(slide, V17_AREA['left'] + 24, y + row_h / 2.0 - 1,
                  28, 3, palette['midtone'])

        add_text(slide, V17_AREA['left'] + 58, y + (row_h - 24) / 2.0,
                 label_w - 58, str(c.get('label', '')), 16, bold=True,
                 color=hex_to_rgb(palette['primary']), height_px=26)

        bar_w = max(bar_max_w * ratio, 40)
        _v17_card(slide, bar_x, y, bar_w, row_h, fill_hex, palette, radius_px=6)

        # 値ラベル｜バー内に収まる場合は内側（白抜き/濃字）、
        # 収まらない場合はバー右外に出す（文字切れ防止｜条項5）
        val_text = '%s（%.1f%%）' % (_fmt_num(v), ratio * 100)
        need_w = _est_text_w(val_text, 14) * 1.2 + 28
        if bar_w >= need_w:
            add_text(slide, bar_x + 14, y + (row_h - 24) / 2.0, bar_w - 28,
                     val_text, 14, bold=True, color=fg, height_px=24)
        else:
            add_text(slide, bar_x + bar_w + 10, y + (row_h - 24) / 2.0,
                     _est_text_w(val_text, 14) * 1.2 + 20, val_text, 14,
                     bold=True, color=hex_to_rgb(palette['primary']),
                     height_px=24)

        note = str(c.get('note', ''))
        if note:
            add_text(slide, bar_x + bar_max_w + 16, y + (row_h - 22) / 2.0, 134,
                     note, 14, color=hex_to_rgb(palette['secondary']),
                     height_px=24)

    notes.append('構成比合計=%.1f%%（入力合計 %s を100%%として正規化）'
                 % (100.0, _fmt_num(total_in)))
    return _v17_report('breakdown', n, None, notes)


# ---------------------------------------------------------------------
# v17-5｜comparison（比較）｜horizontal / discrete_contrast / 2〜3要素
# ---------------------------------------------------------------------
def draw_comparison(slide, palette, data):
    """2〜3要素を並列カラムで対比描画する。

    Args:
        data : {'title': str,
                'comparison_axis': str|None,      # 任意（原本 requires_axes=False）
                'attribute_labels': [str, ...],
                'items': [{'label': str, 'score': int|None,
                           'attributes': {key: value}}, ...]}
    仕様: 拡張定義集約表 P1-4 行（min2 / max3 / horizontal / discrete_contrast）
    要素数逸脱時は draw_category へフォールバック（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['comparison']
    notes = []
    items = _v17_normalize(data.get('items', []))

    if not (spec['min_elements'] <= len(items) <= spec['max_elements']):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': it.get('label', ''),
                            'score': it.get('score'),
                            'description': ''} for it in items],
        }, _fallback_from='comparison')
        rep['notes'].insert(0, 'comparison 要素数 %d が %d〜%d の範囲外｜'
                            'category へフォールバック'
                            % (len(items), spec['min_elements'], spec['max_elements']))
        return rep

    _v17_title(slide, data.get('title', ''), palette)
    axis = data.get('comparison_axis')
    body_top = V17_AREA['body_top']
    if axis:
        _v17_axis_label(slide, '比較軸：%s' % axis, palette)
        body_top += 26

    # discrete_contrast: 中間色を挟まず index で離散的に割り当てる
    contrast_keys = {2: ['midtone', 'primary'],
                     3: ['midtone', 'secondary', 'primary']}[len(items)]

    n = len(items)
    gap = 24
    col_w = (V17_AREA['width'] - gap * (n - 1)) / float(n)
    col_h = V17_AREA['body_bottom'] - body_top
    head_h = 108
    attr_labels = list(data.get('attribute_labels') or [])

    for i, it in enumerate(items):
        x = V17_AREA['left'] + i * (col_w + gap)
        score = it.get('score')
        fill_hex = _tier_fill(palette, score, contrast_keys[i])
        fg = _text_color_on(fill_hex)

        # 本体カード（薄色）＋ヘッダ帯（濃色）
        _v17_card(slide, x, body_top, col_w, col_h, palette['lightest'], palette)
        _v17_card(slide, x, body_top, col_w, head_h, fill_hex, palette)

        add_text(slide, x + 18, body_top + 14, col_w - 36,
                 str(it.get('label', '')), 18, bold=True, color=fg, height_px=28)
        if score is not None:
            add_text(slide, x + 18, body_top + 48, col_w - 36, '%s%%' % score,
                     32, bold=True, color=fg, height_px=48)

        # 属性行
        attrs = it.get('attributes') or {}
        keys = attr_labels or list(attrs.keys())
        ay = body_top + head_h + 18
        row_h = 46
        for k in keys:
            if ay + row_h > body_top + col_h - 8:
                notes.append('カラム%d｜属性が領域を超過したため以降を省略' % (i + 1))
                break
            add_text(slide, x + 18, ay, col_w - 36, str(k), 14, bold=True,
                     color=hex_to_rgb(palette['secondary']), height_px=22)
            add_text(slide, x + 18, ay + 20, col_w - 36,
                     str(attrs.get(k, '－')), 16,
                     color=hex_to_rgb(palette['primary']), height_px=24)
            _v17_rule(slide, x + 18, ay + row_h - 2, col_w - 36, 1,
                      palette['light'])
            ay += row_h

    return _v17_report('comparison', n, None, notes)


# ---------------------------------------------------------------------
# v17-6｜共通ディスパッチャ
# ---------------------------------------------------------------------
def resolve_pattern(diagnosis_key):
    """診断カテゴリ → パターンキー（原本 DIAGNOSIS_TO_PATTERN を参照）。

    未定義キーは 'category' にフォールバック（原則①）。
    """
    return DIAGNOSIS_TO_PATTERN.get(diagnosis_key, 'category')


def draw_pattern(slide, pattern_key, palette, data):
    """パターンキーで描画関数を振り分ける。未実装キーは category に退避。"""
    table = {
        'category':   draw_category,
        'breakdown':  draw_breakdown,
        'comparison': draw_comparison,
    }
    fn = table.get(pattern_key)
    if fn is None:
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': data.get('categories', []) or [
                {'label': str(x.get('label', '')), 'score': x.get('score'),
                 'description': ''} for x in (data.get('items')
                                              or data.get('components') or [])],
        }, _fallback_from=pattern_key)
        rep['notes'].append('パターン "%s" は v17 P1 の対象外（P2/P3 で実装）'
                            % pattern_key)
        return rep
    return fn(slide, palette, data)


def _est_text_w(text, size_pt):
    """メイリオ想定の概算テキスト幅（px）。

    全角=size_pt*1.34px／半角=size_pt*0.70px で見積る。
    バー内に値ラベルが収まるかの判定に用いる（文字切れ防止）。
    """
    w = 0.0
    for ch in str(text):
        w += size_pt * (0.70 if ord(ch) < 0x2000 else 1.34)
    return w


def add_diagram_slide(prs, pattern_key, palette, data, page_num=1, total=1,
                      header_label=None, author='紺＆クリーン スライド作成'):
    """図解パターン1枚を「ヘッダ帯＋フッター＋パターン描画」で1スライド出力する。

    既存の add_* 関数群と同じ呼び出し規約（prs を受け取り slide を返す）に揃えた
    ラッパ。v17-fix2（8/24 実機検証）で検出した「_add_header と draw_* の
    両方がタイトルを描き、同一文言が2箇所に出る」事象を構造的に防ぐ。

    ヘッダ帯には data['title'] を出さず、原本 DIAGRAM_PATTERNS の `ja`
    （例: 'category' → '分類'）を既定ラベルとして表示する。
    スライド内見出しは draw_* 側が data['title'] を1箇所だけ描く。

    Args:
        prs         : python-pptx Presentation
        pattern_key : 'category' / 'breakdown' / 'comparison' 等
        palette     : get_theme_palette の戻り値
        data        : 各 draw_* のデータ辞書
        page_num    : ページ番号
        total       : 総ページ数
        header_label: ヘッダ帯の左側文言（None なら原本 `ja` を使用）
    Returns:
        (slide, report)
    """
    slide = _blank_slide(prs)
    if header_label is None:
        meta = DIAGRAM_PATTERNS.get(pattern_key) or {}
        header_label = meta.get('ja', '図解')
    _add_header(slide, header_label, str(pattern_key).upper())
    _add_footer(slide, page_num, total, author=author)
    report = draw_pattern(slide, pattern_key, palette, data)
    return slide, report


def _fmt_num(v):
    """数値の表示整形（整数はそのまま、小数は1桁）"""
    try:
        f = float(v)
    except (TypeError, ValueError):
        return str(v)
    return str(int(f)) if abs(f - int(f)) < 1e-9 else '%.1f' % f


def _pp_right():
    from pptx.enum.text import PP_ALIGN
    return PP_ALIGN.RIGHT

# =====================================================================
# ▲▲▲ v17 追加ブロック ここまで ▲▲▲
# =====================================================================


# =====================================================================
# ▼▼▼ v17 P2 追加ブロック（pyramid / sequence / framework） ▼▼▼
# P1ブロックのヘルパーを再利用する。既存 C-1〜C-3 には一切触れない。
# =====================================================================

# ---------------------------------------------------------------------
# v17-P2-0｜拡張パターン仕様（集約表 P2-1／P2-2／P2-3 行を正とする）
# ---------------------------------------------------------------------
DIAGRAM_PATTERN_SPEC.update({
    'pyramid': {
        'min_elements': 3, 'max_elements': 5, 'requires_axes': False,
        'direction': 'vertical', 'color_gradation': 'hierarchical',
        # 頂点30%→基層90%（原本 use「階層・優先順位を上下で表現」）
        'width_ratio_top': 0.30, 'width_ratio_bottom': 0.90,
    },
    'sequence': {
        'min_elements': 3, 'max_elements': 6, 'requires_axes': False,
        'direction': 'horizontal', 'color_gradation': 'progressive',
        'arrow_gap_px': 44,
    },
    'framework': {
        'min_elements': 4, 'max_elements': 9, 'requires_axes': True,
        'direction': 'grid', 'color_gradation': 'positional_quadrant',
        # 要素数→グリッド形状（決定論的｜原則③）
        'grid_map': {4: (2, 2), 5: (3, 2), 6: (3, 2),
                     7: (3, 3), 8: (3, 3), 9: (3, 3)},
        'axis_y_w': 72, 'axis_x_h': 46,
    },
})

# 階層色（hierarchical）｜段数別の色キー列（決定論的）
V17_PYRAMID_TIERS = {
    3: ['primary', 'midtone', 'light'],
    4: ['primary', 'secondary', 'midtone', 'light'],
    5: ['primary', 'secondary', 'midtone', 'light', 'lightest'],
}

# 進行色（progressive）｜段数別の色キー列
V17_SEQUENCE_TIERS = {
    3: ['primary', 'secondary', 'midtone'],
    4: ['primary', 'secondary', 'midtone', 'light'],
    5: ['primary', 'secondary', 'midtone', 'light', 'lightest'],
    6: ['primary', 'primary', 'secondary', 'midtone', 'light', 'lightest'],
}


# ---------------------------------------------------------------------
# v17-P2-1｜pyramid（ピラミッド）｜vertical / hierarchical / 3〜5段
# ---------------------------------------------------------------------
def draw_pyramid(slide, palette, data):
    """階層・優先順位を上下で表現する（頂点＝最重要）。

    Args:
        data : {'title': str,
                'levels': [{'label': str, 'score': int|None,
                            'description': str}, ...]}   # 3〜5・index0が頂点
    Returns:
        dict : 描画レポート

    仕様: 集約表 P2-1 行（min3／max5／vertical／hierarchical／原本shape=triangle）
    ⚠️ 台形は `add_freeform` を使わず、幅可変の矩形段で近似する（原則④）。
       `add_freeform` は python-pptx のバージョン依存があり、PowerPoint実機での
       再現性を担保できないため。段ごとの幅差で三角形の輪郭を表現する。
    範囲外は draw_category へフォールバック（原則①｜例外を投げない）。
    """
    spec = DIAGRAM_PATTERN_SPEC['pyramid']
    notes = []
    levels = _v17_normalize(data.get('levels', []))

    if not (spec['min_elements'] <= len(levels) <= spec['max_elements']):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': lv.get('label', ''), 'score': lv.get('score'),
                            'description': str(lv.get('description', ''))}
                           for lv in levels],
        }, _fallback_from='pyramid')
        rep['notes'].insert(0, 'pyramid 段数 %d が %d〜%d の範囲外｜category へフォールバック'
                            % (len(levels), spec['min_elements'], spec['max_elements']))
        return rep

    _v17_title(slide, data.get('title', ''), palette)

    n = len(levels)
    tiers = V17_PYRAMID_TIERS[n]
    gap = 8
    avail = V17_AREA['body_bottom'] - V17_AREA['body_top']
    band_h = min((avail - gap * (n - 1)) / float(n), 104)
    block_h = band_h * n + gap * (n - 1)
    y0 = V17_AREA['body_top'] + (avail - block_h) / 2.0
    cx = V17_AREA['left'] + V17_AREA['width'] / 2.0

    r_top, r_bot = spec['width_ratio_top'], spec['width_ratio_bottom']
    for i, lv in enumerate(levels):
        # 段の幅：頂点 r_top → 基層 r_bot を線形に配分
        ratio = r_top + (r_bot - r_top) * (i / float(n - 1))
        w = V17_AREA['width'] * ratio
        x = cx - w / 2.0
        y = y0 + i * (band_h + gap)

        score = lv.get('score')
        fill_hex = _tier_fill(palette, score, tiers[i])
        fg = _text_color_on(fill_hex)
        _v17_card(slide, x, y, w, band_h, fill_hex, palette, radius_px=6)

        # 段ラベル｜頂点ほど大きく（優先度の視覚的強調）
        label_size = 20 if i == 0 else (18 if i == 1 else 16)
        pad = 20
        add_text(slide, x + pad, y + 8, w - pad * 2, str(lv.get('label', '')),
                 label_size, bold=True, color=fg, height_px=int(label_size * 1.5),
                 align=_pp_center())
        sub = []
        if score is not None:
            sub.append('%s%%' % score)
        desc = str(lv.get('description', ''))
        if desc:
            sub.append(desc)
        if sub:
            add_text(slide, x + pad, y + 8 + int(label_size * 1.5), w - pad * 2,
                     '　'.join(sub), 14, color=fg,
                     height_px=max(int(band_h - int(label_size * 1.5) - 16), 22),
                     align=_pp_center())

        # 優先度の序列を左外に添える（1が最上位）
        add_text(slide, V17_AREA['left'], y + band_h / 2.0 - 13, 40,
                 str(i + 1), 16, bold=True,
                 color=hex_to_rgb(palette['secondary']), height_px=26,
                 align=_pp_center())

    return _v17_report('pyramid', n, None, notes)


# ---------------------------------------------------------------------
# v17-P2-2｜sequence（順序）｜horizontal / progressive / 3〜6ステップ
# ---------------------------------------------------------------------
def draw_sequence(slide, palette, data):
    """ステップ・時系列を左→右で表現する。

    Args:
        data : {'title': str,
                'steps': [{'label': str, 'score': int|None,
                           'description': str}, ...]}    # 3〜6
    仕様: 集約表 P2-2 行（min3／max6／horizontal 固定／progressive）
    ⚠️ 原本 use「ステップ・時系列を**左→右**」により direction は horizontal 固定。
       集約表も direction=horizontal と規定するため、縦方向は実装しない。
    範囲外は draw_category へフォールバック（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['sequence']
    notes = []
    steps = _v17_normalize(data.get('steps', []))

    if not (spec['min_elements'] <= len(steps) <= spec['max_elements']):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': st.get('label', ''), 'score': st.get('score'),
                            'description': str(st.get('description', ''))}
                           for st in steps],
        }, _fallback_from='sequence')
        rep['notes'].insert(0, 'sequence ステップ数 %d が %d〜%d の範囲外｜'
                            'category へフォールバック'
                            % (len(steps), spec['min_elements'], spec['max_elements']))
        return rep

    _v17_title(slide, data.get('title', ''), palette)

    n = len(steps)
    tiers = V17_SEQUENCE_TIERS[n]
    agap = spec['arrow_gap_px']
    card_w = (V17_AREA['width'] - agap * (n - 1)) / float(n)
    card_h = min(V17_AREA['body_bottom'] - V17_AREA['body_top'], 232)
    y = V17_AREA['body_top'] + \
        (V17_AREA['body_bottom'] - V17_AREA['body_top'] - card_h) / 2.0

    for i, st in enumerate(steps):
        x = V17_AREA['left'] + i * (card_w + agap)
        score = st.get('score')
        fill_hex = _tier_fill(palette, score, tiers[i])
        fg = _text_color_on(fill_hex)
        _v17_card(slide, x, y, card_w, card_h, fill_hex, palette)

        pad = 14
        add_text(slide, x + pad, y + 12, card_w - pad * 2, 'STEP %d' % (i + 1),
                 14, bold=True, color=fg, height_px=22)
        add_text(slide, x + pad, y + 40, card_w - pad * 2,
                 str(st.get('label', '')), 16, bold=True, color=fg,
                 height_px=52, line_height=1.3)
        if score is not None:
            add_text(slide, x + pad, y + 98, card_w - pad * 2, '%s%%' % score,
                     24, bold=True, color=fg, height_px=38)
        desc = str(st.get('description', ''))
        if desc:
            dtop = y + (142 if score is not None else 98)
            add_text(slide, x + pad, dtop, card_w - pad * 2, desc, 14, color=fg,
                     height_px=max(int(y + card_h - dtop - 10), 22),
                     line_height=1.4)

        # ステップ間の矢印（原則④｜python-pptx標準図形・adjustments 未使用）
        if i < n - 1:
            ax = x + card_w + 6
            aw = agap - 12
            ah = 26
            add_shape(slide, _mso_right_arrow(), ax, y + card_h / 2.0 - ah / 2.0,
                      aw, ah, fill=hex_to_rgb(palette['accent']))

    return _v17_report('sequence', n, None, notes)


# ---------------------------------------------------------------------
# v17-P2-3｜framework（フレームワーク）｜grid / positional_quadrant / 4〜9セル
#   ⚠️ 実装済11パターン中で唯一 requires_axes=True
# ---------------------------------------------------------------------
def _framework_tier(row, col, rows, cols):
    """セル位置から色キーを決定（positional_quadrant｜決定論的）。

    2x2 は原本 use「4象限マトリクス等」に従い象限別の意味付けを行う
    （右上＝最重要／左下＝最軽微）。それ以外は行位置ベース。
    """
    if (rows, cols) == (2, 2):
        return {(0, 1): 'primary', (0, 0): 'secondary',
                (1, 1): 'midtone', (1, 0): 'light'}.get((row, col), 'midtone')
    if row == 0:
        return 'primary'
    if row == rows - 1:
        return 'light'
    return 'secondary'


def draw_framework(slide, palette, data):
    """2軸マトリクス（4象限等）で構造を提示する。

    Args:
        data : {'title': str,
                'axis_x_label': str, 'axis_y_label': str,     # ⚠️ 必須
                'axis_x_low': str, 'axis_x_high': str,
                'axis_y_low': str, 'axis_y_high': str,
                'cells': [{'row': int, 'col': int, 'label': str,
                           'score': int|None, 'items': [str]}, ...]}  # 4〜9
    仕様: 集約表 P2-3 行（min4／max9／requires_axes=True／grid／positional_quadrant）
    ⚠️ 集約表は「4〜9の範囲」を規定するため、5・7・8セルも受け付け、
       要素数からグリッド形状を決定論的に自動選定する（実装記録の固定値縛りは採らない）。
    ⚠️ 軸ラベル欠落時も例外を投げず draw_category へフォールバックする（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['framework']
    notes = []
    cells = _v17_normalize(data.get('cells', []))
    n = len(cells)

    def _fallback(reason):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': c.get('label', ''), 'score': c.get('score'),
                            'description': ''} for c in cells],
        }, _fallback_from='framework')
        # ⚠️ v17.1.0 修正：ローカル notes ではなく戻り値の notes に記録する。
        # 旧実装はローカル notes に append した後 draw_category の新しい report を
        # 返していたため、フォールバック理由が呼び出し側に届かなかった。
        rep['notes'].insert(0, reason)
        return rep

    if not (spec['min_elements'] <= n <= spec['max_elements']):
        return _fallback('framework セル数 %d が %d〜%d の範囲外｜category へフォールバック'
                         % (n, spec['min_elements'], spec['max_elements']))

    ax_label = str(data.get('axis_x_label') or '').strip()
    ay_label = str(data.get('axis_y_label') or '').strip()
    if not ax_label or not ay_label:
        return _fallback('framework は requires_axes=True｜軸ラベル欠落のため '
                         'category へフォールバック（例外は送出しない）')

    _v17_title(slide, data.get('title', ''), palette)

    cols, rows = spec['grid_map'][n]
    notes.append('セル数 %d → グリッド %dx%d を自動選定' % (n, cols, rows))

    ay_w, ax_h = spec['axis_y_w'], spec['axis_x_h']
    gx = V17_AREA['left'] + ay_w
    gy = V17_AREA['body_top']
    gw = V17_AREA['width'] - ay_w
    gh = V17_AREA['body_bottom'] - gy - ax_h
    gap = 10
    cw = (gw - gap * (cols - 1)) / float(cols)
    ch = (gh - gap * (rows - 1)) / float(rows)

    # 軸ラベル（framework固有｜requires_axes=True）
    add_text(slide, gx, V17_AREA['body_bottom'] - ax_h + 10, gw,
             '%s ← %s → %s' % (data.get('axis_x_low', '低'), ax_label,
                               data.get('axis_x_high', '高')),
             14, bold=True, color=hex_to_rgb(palette['secondary']),
             height_px=24, align=_pp_center())
    add_paragraph_box(slide, V17_AREA['left'], gy + gh / 2.0 - 56, ay_w,
                      [{'text': str(data.get('axis_y_high', '高')), 'size': 14,
                        'bold': True, 'align': _pp_center()},
                       {'text': '↑', 'size': 14, 'align': _pp_center()},
                       {'text': ay_label, 'size': 14, 'bold': True,
                        'align': _pp_center()},
                       {'text': '↓', 'size': 14, 'align': _pp_center()},
                       {'text': str(data.get('axis_y_low', '低')), 'size': 14,
                        'bold': True, 'align': _pp_center()}],
                      height_px=112, default_color=hex_to_rgb(palette['secondary']),
                      line_height=1.15)

    # セル（row/col 指定がなければ左上から順に充填）
    used = set()
    for idx, c in enumerate(cells):
        r = c.get('row')
        col = c.get('col')
        if not isinstance(r, int) or not isinstance(col, int) \
                or not (0 <= r < rows and 0 <= col < cols) or (r, col) in used:
            r, col = idx // cols, idx % cols
        used.add((r, col))

        x = gx + col * (cw + gap)
        y = gy + r * (ch + gap)
        score = c.get('score')
        fill_hex = _tier_fill(palette, score, _framework_tier(r, col, rows, cols))
        fg = _text_color_on(fill_hex)
        _v17_card(slide, x, y, cw, ch, fill_hex, palette)

        pad = 14
        add_text(slide, x + pad, y + 10, cw - pad * 2, str(c.get('label', '')),
                 16, bold=True, color=fg, height_px=26)
        top = y + 38
        if score is not None:
            add_text(slide, x + pad, top, cw - pad * 2, '%s%%' % score,
                     20, bold=True, color=fg, height_px=32)
            top += 34
        items = [str(i) for i in (c.get('items') or [])][:3]
        if items and top + 22 <= y + ch - 8:
            add_paragraph_box(slide, x + pad, top, cw - pad * 2,
                              [{'text': '・' + t, 'size': 14} for t in items],
                              height_px=max(int(y + ch - top - 8), 22),
                              default_color=fg, line_height=1.35,
                              space_after_pt=2)

    return _v17_report('framework', n, None, notes)


# ---------------------------------------------------------------------
# v17-P2-4｜ディスパッチ表への登録（P1の draw_pattern を置き換える）
# ---------------------------------------------------------------------
def draw_pattern(slide, pattern_key, palette, data):
    """パターンキーで描画関数を振り分ける（P1 3種＋P2 3種＝6種に対応）。

    未実装キー（funnel / timeline / contrast / cycle / network / integration）は
    category に退避し、notes に理由を記録する（例外は送出しない）。
    """
    table = {
        'category':   draw_category,
        'breakdown':  draw_breakdown,
        'comparison': draw_comparison,
        'pyramid':    draw_pyramid,
        'sequence':   draw_sequence,
        'framework':  draw_framework,
    }
    fn = table.get(pattern_key)
    if fn is None:
        src = (data.get('categories') or data.get('items')
               or data.get('components') or data.get('levels')
               or data.get('steps') or data.get('cells') or [])
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': str(x.get('label', '')) if isinstance(x, dict) else str(x),
                            'score': x.get('score') if isinstance(x, dict) else None,
                            'description': ''} for x in src],
        }, _fallback_from=pattern_key)
        rep['notes'].append('パターン "%s" は v17 P3 で実装予定（P1/P2 の対象外）'
                            % pattern_key)
        return rep
    return fn(slide, palette, data)


def _pp_center():
    from pptx.enum.text import PP_ALIGN
    return PP_ALIGN.CENTER


def _mso_right_arrow():
    from pptx.enum.shapes import MSO_SHAPE
    return MSO_SHAPE.RIGHT_ARROW

# =====================================================================
# ▲▲▲ v17 P2 追加ブロック ここまで ▲▲▲
# =====================================================================


# =====================================================================
# ▼▼▼ v17 P3 追加ブロック ここから ▼▼▼
#     対象: funnel / cycle / contrast / timeline / network（5種）
#     配置: 03_pptx_builder.py の末尾（v17 P2 ブロックの直後）に追記
#     版数: 17.1.0 → 17.2.0（__version__ を更新すること）
#
#     設計根拠:
#       - 拡張定義集約表（v35_core_extended_pattern_definitions.md）P3-1〜P3-5 行
#       - 集約表と実装記録（8/13〜8/15）の照合結果: ⭐ 相違 0件
#       - 原本 DIAGRAM_PATTERNS は無改変。拡張層は DIAGRAM_PATTERN_SPEC に分離。
#
#     横断設計原則（v3.5コア）:
#       原則① 描画不能時も例外を投げず劣化描画（category フォールバック）で通す
#       原則② score < 40 は警告色でオーバーライドする
#       原則③ 同一入力から同一出力（決定論性｜乱数・力学配置を使わない）
#       原則④ 基本図形を優先し、環境依存図形を増やさない
#
#     ⚠️ cycle の図形選定について（8/25 実測にもとづく設計変更）
#       実装記録 8/14 判断④は MSO_SHAPE.BLOCK_ARC ＋ rotation を採ると
#       していたが、8/25 に実機同等環境で実測した結果:
#         (a) adjustments は「度」単位（既定値 [108.0, 0.0, 0.25]）であり、
#             実装記録が想定した 0.0〜1.0 比率ではない。
#             → 比率で指定すると極小弧になり視覚的に「線」になる。
#         (b) 正しく度で指定すれば分割弧は描けるが、⚠️ 弧の内部に
#             テキストを格納できないため、他10パターンが持つ
#             「ラベル＋スコア＋説明」の情報構造を表現できない。
#       よって本実装は角丸矩形の円周配置＋RIGHT_ARROW（基本図形のみ）で
#       円環構造を構成する。原本 shape='circle_arrow'（円環＋矢印）との
#       整合は矢印の実装によって維持される。原則④とも整合する。
#       （判断原理14「テストが通ることは実機で描画されることの代替に
#         ならない」を設計段階に先取り適用して得た結論）
# =====================================================================

# ---------------------------------------------------------------------
# v17-P3-0｜拡張仕様の登録（原本 DIAGRAM_PATTERNS は無改変）
# ---------------------------------------------------------------------
DIAGRAM_PATTERN_SPEC.update({
    'funnel': {
        'min_elements': 3, 'max_elements': 6, 'requires_axes': False,
        'direction': 'vertical', 'color_gradation': 'progressive_narrowing',
        # 上段90%→下段40%（原本 use「上から下へ絞り込むファネル型」）
        'width_ratio_top': 0.90, 'width_ratio_bottom': 0.40,
    },
    'cycle': {
        'min_elements': 3, 'max_elements': 6, 'requires_axes': False,
        'direction': 'clockwise', 'color_gradation': 'uniform_cyclic',
        # 12時起点・時計回り固定（集約表 P3-2 行）
        'start_angle_deg': -90.0,
    },
    'contrast': {
        'min_elements': 2, 'max_elements': 2, 'requires_axes': False,
        'direction': 'horizontal', 'color_gradation': 'polarized_contrast',
        'divider_w': 6,
    },
    'timeline': {
        'min_elements': 3, 'max_elements': 7, 'requires_axes': True,
        'direction': 'horizontal', 'color_gradation': 'progressive',
        'axis_h': 8, 'dot_d': 26,
    },
    'network': {
        'min_elements': 3, 'max_elements': 7, 'requires_axes': False,
        'direction': 'hierarchical_top_down',
        'color_gradation': 'depth_hierarchical',
        'node_w': 210, 'node_h': 66, 'edge_w': 3,
    },
})

# 段階減衰色（progressive_narrowing）｜funnel の段数別色キー列
V17_FUNNEL_TIERS = {
    3: ['primary', 'secondary', 'midtone'],
    4: ['primary', 'secondary', 'midtone', 'light'],
    5: ['primary', 'secondary', 'midtone', 'light', 'lightest'],
    6: ['primary', 'primary', 'secondary', 'midtone', 'light', 'lightest'],
}

# 均等循環色（uniform_cyclic）｜4色を循環させる
# ⚠️ 段階減衰は「終わりがある」誤読を生むため使わない（集約表 P3-2 行の根拠）
V17_CYCLE_TIERS = ['primary', 'secondary', 'midtone', 'light']

# 2極化色（polarized_contrast）｜中間色を排し落差を可視化
V17_CONTRAST_TIERS = ['lightest', 'primary']

# 深度別階層色（depth_hierarchical）｜network の depth→色キー
V17_NETWORK_TIERS = ['primary', 'secondary', 'midtone', 'light']


# ---------------------------------------------------------------------
# v17-P3-1｜funnel（絞り込み）｜vertical / progressive_narrowing / 3〜6段
# ---------------------------------------------------------------------
def draw_funnel(slide, palette, data):
    """段階的な絞り込みを上から下へ幅を狭めて表現する。

    Args:
        data : {'title': str,
                'stages': [{'label': str, 'score': int|None,
                            'description': str}, ...]}   # 3〜6・index0が最上段
    Returns:
        dict : 描画レポート

    仕様: 集約表 P3-1 行（min3／max6／vertical／progressive_narrowing／
          原本shape=trapezoid）
    ⚠️ 台形は MSO_SHAPE.TRAPEZOID を使わず、幅可変の矩形段で近似する。
       pyramid（P2-1）と同一方針（原則④）。pyramid は上が狭いが、
       funnel は上が広い（絞り込み方向が逆）。
    範囲外は draw_category へフォールバック（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['funnel']
    notes = []
    stages = _v17_normalize(data.get('stages', []))
    n = len(stages)

    if not (spec['min_elements'] <= n <= spec['max_elements']):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': s.get('label', ''), 'score': s.get('score'),
                            'description': str(s.get('description', ''))}
                           for s in stages],
        }, _fallback_from='funnel')
        rep['notes'].append(
            '段数 %d は funnel の範囲外（min %d／max %d）｜category へフォールバック'
            % (n, spec['min_elements'], spec['max_elements']))
        return rep

    _v17_title(slide, data.get('title', ''), palette)

    top = V17_AREA['body_top']
    avail_h = V17_AREA['body_bottom'] - top
    gap = 10
    band_h = int((avail_h - gap * (n - 1)) / n)
    band_h = min(band_h, 118)

    r_top = spec['width_ratio_top']
    r_bot = spec['width_ratio_bottom']
    tiers = V17_FUNNEL_TIERS.get(n, V17_FUNNEL_TIERS[6])
    cx = V17_AREA['left'] + V17_AREA['width'] // 2

    drawn = 0
    y = top
    for i, st in enumerate(stages):
        ratio = r_top + (r_bot - r_top) * (i / float(max(n - 1, 1)))
        w = int(V17_AREA['width'] * ratio)
        x = cx - w // 2
        score = st.get('score')
        fill_hex = _tier_fill(palette, score, tiers[min(i, len(tiers) - 1)])
        _v17_card(slide, x, y, w, band_h, fill_hex, palette)

        # 段番号（左端｜描画領域内に収める）
        num_x = max(V17_AREA['left'], x - 34)
        add_text(slide, num_x, y + band_h // 2 - 12, 30, str(i + 1), 14,
                 bold=True, color=hex_to_rgb(palette['midtone']), height_px=24)

        txt_color = _text_color_on(fill_hex)
        pad = 12
        add_text(slide, x + pad, y + 8, w - pad * 2, str(st.get('label', '')),
                 16, bold=True, color=txt_color, height_px=24,
                 align=_pp_center())
        line2 = []
        if score is not None:
            line2.append('%s%%' % score)
        desc = str(st.get('description', '')).strip()
        if desc:
            line2.append(desc)
        if line2 and band_h >= 56:
            add_text(slide, x + pad, y + 34, w - pad * 2, '　'.join(line2),
                     14, color=txt_color, height_px=band_h - 42,
                     align=_pp_center())
        drawn += 1
        y += band_h + gap

    notes.append('段数 %d｜幅比 %.2f→%.2f で絞り込みを表現' % (n, r_top, r_bot))
    return _v17_report('funnel', drawn, None, notes)


# ---------------------------------------------------------------------
# v17-P3-2｜cycle（循環）｜clockwise / uniform_cyclic / 3〜6段
# ---------------------------------------------------------------------
def _compute_cycle_positions(n, cx, cy, radius, box_w, box_h, start_deg):
    """円周上の n 個のボックス左上座標を決定論的に返す（原則③）。

    12時起点・時計回り。三角関数のみを使い乱数・力学配置は使わない。
    """
    import math
    step = 360.0 / n
    out = []
    for i in range(n):
        ang = start_deg + i * step
        rad = math.radians(ang)
        px = cx + radius * math.cos(rad) - box_w / 2.0
        py = cy + radius * math.sin(rad) - box_h / 2.0
        out.append((int(round(px)), int(round(py)), ang))
    return out


def _draw_cycle_arrow(slide, palette, cx, cy, radius, angle_deg):
    """円周の接線方向へ向く矢印を1本描く（基本図形 RIGHT_ARROW ＋ rotation）。

    ⚠️ BLOCK_ARC を使わない理由は本ブロック冒頭のコメントを参照。
    """
    import math
    aw, ah = 56, 20
    rad = math.radians(angle_deg)
    ax = cx + radius * math.cos(rad) - aw / 2.0
    ay = cy + radius * math.sin(rad) - ah / 2.0
    sh = add_shape(slide, _mso_right_arrow(), int(round(ax)), int(round(ay)),
                   aw, ah, fill=hex_to_rgb(palette['midtone']))
    try:
        sh.rotation = angle_deg + 90.0
    except Exception:
        pass          # rotation 非対応環境でも矢印自体は残す（原則①）
    return sh


def draw_cycle(slide, palette, data):
    """反復プロセスを円環で表現する（12時起点・時計回り固定）。

    Args:
        data : {'title': str,
                'cycle_name': str,                        # 円環中心の名称（任意）
                'phases': [{'label': str, 'score': int|None,
                            'description': str}, ...]}    # 3〜6
    Returns:
        dict : 描画レポート

    仕様: 集約表 P3-2 行（min3／max6／clockwise／uniform_cyclic／
          原本shape=circle_arrow）
    ⚠️ 色は段階減衰させず4色を循環させる。段階減衰は「終わりがある」
       誤読を生むため（集約表 P3-2 行の根拠）。
    範囲外は draw_category へフォールバック（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['cycle']
    notes = []
    phases = _v17_normalize(data.get('phases', []))
    n = len(phases)

    if not (spec['min_elements'] <= n <= spec['max_elements']):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': p.get('label', ''), 'score': p.get('score'),
                            'description': str(p.get('description', ''))}
                           for p in phases],
        }, _fallback_from='cycle')
        rep['notes'].append(
            '段数 %d は cycle の範囲外（min %d／max %d）｜category へフォールバック'
            % (n, spec['min_elements'], spec['max_elements']))
        return rep

    _v17_title(slide, data.get('title', ''), palette)

    top = V17_AREA['body_top']
    field_h = V17_AREA['body_bottom'] - top
    cx = V17_AREA['left'] + V17_AREA['width'] // 2
    cy = top + field_h // 2

    # v17-P3-fix1（8/25 自己目視）: 説明文が2行に折返してカード内が窮屈になる
    # 事象を検出。ボックス幅を広げ高さも確保する（充填率と可読性の両立）。
    box_w = 268 if n <= 4 else 212
    box_h = 100 if n <= 4 else 90
    radius = int(min(V17_AREA['width'] / 2 - box_w / 2 - 16,
                     field_h / 2 - box_h / 2 - 6))

    positions = _compute_cycle_positions(n, cx, cy, radius, box_w, box_h,
                                         spec['start_angle_deg'])

    # 矢印は各ボックスの中間角に置く（時計回りの進行を示す）
    step = 360.0 / n
    for i in range(n):
        mid = spec['start_angle_deg'] + (i + 0.5) * step
        _draw_cycle_arrow(slide, palette, cx, cy, int(radius * 0.62), mid)

    drawn = 0
    for i, (x, y, _ang) in enumerate(positions):
        ph = phases[i]
        score = ph.get('score')
        base_key = V17_CYCLE_TIERS[i % len(V17_CYCLE_TIERS)]
        fill_hex = _tier_fill(palette, score, base_key)
        _v17_card(slide, x, y, box_w, box_h, fill_hex, palette)

        txt_color = _text_color_on(fill_hex)
        pad = 10
        head = '%d. %s' % (i + 1, str(ph.get('label', '')))
        add_text(slide, x + pad, y + 7, box_w - pad * 2, head, 15, bold=True,
                 color=txt_color, height_px=22)
        sub = []
        if score is not None:
            sub.append('%s%%' % score)
        desc = str(ph.get('description', '')).strip()
        if desc:
            sub.append(desc)
        if sub:
            add_text(slide, x + pad, y + 31, box_w - pad * 2, '　'.join(sub),
                     14, color=txt_color, height_px=box_h - 38)
        drawn += 1

    # 判断⑤｜円環中心にサイクル名を配置
    cname = str(data.get('cycle_name', '')).strip()
    if cname:
        add_text(slide, cx - 110, cy - 16, 220, cname, 18, bold=True,
                 color=hex_to_rgb(palette['primary']), height_px=32,
                 align=_pp_center())

    notes.append('段数 %d｜12時起点・時計回り｜4色循環（uniform_cyclic）' % n)
    notes.append('円環は角丸矩形の円周配置＋RIGHT_ARROW で構成'
                 '（BLOCK_ARC 不使用｜8/25 実測にもとづく設計判断）')
    return _v17_report('cycle', drawn, None, notes)


# ---------------------------------------------------------------------
# v17-P3-3｜contrast（対比）｜horizontal / polarized_contrast / 2固定
# ---------------------------------------------------------------------
def draw_contrast(slide, palette, data):
    """対照的な2要素を左右分割で並列強調する。

    Args:
        data : {'title': str,
                'sides': [{'label': str, 'score': int|None,
                           'items': [str]}, ...]}     # ⚠️ 2固定
    Returns:
        dict : 描画レポート

    仕様: 集約表 P3-3 行（min2／max2 固定／horizontal／polarized_contrast／
          原本shape=split_screen）
    ⚠️ 中間色を排し2極化（薄→濃）する。目的は「推移」ではなく「落差」の
       可視化であるため（集約表 P3-3 行の根拠）。
    ⚠️ requires_axes=False のため軸ラベルは描かない（framework との差異）。
    2要素以外は draw_category へフォールバック（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['contrast']
    notes = []
    sides = _v17_normalize(data.get('sides', []))
    n = len(sides)

    if n != 2:
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': s.get('label', ''), 'score': s.get('score'),
                            'description': ''} for s in sides],
        }, _fallback_from='contrast')
        rep['notes'].append(
            '要素数 %d は contrast の範囲外（2固定）｜category へフォールバック' % n)
        return rep

    _v17_title(slide, data.get('title', ''), palette)

    top = V17_AREA['body_top']
    dv = spec['divider_w']
    col_w = (V17_AREA['width'] - dv - V17_AREA['gap'] * 2) // 2

    # v17-P3-fix2（8/25 自己目視）: items が少ないとカード下部が大きく空き、
    # 充填率が低下する事象を検出（P2 framework の V17_CARD_H_MAX と同趣旨）。
    # 内容量（ラベル+スコア+items 行数）から必要高を算出し上限で丸める。
    # 行高の実測値: 14pt×line_height1.6 ≒ 32px／折返し1行を許容して 44px 見込む
    max_items = max(len([t for t in (sd.get('items') or [])][:5]) for sd in sides)
    has_score = any(sd.get('score') is not None for sd in sides)
    need_h = (14 + 28) + (40 if has_score else 0) + max_items * 44 + 24
    # v17-P3-fix3（8/25 自己目視）: 内容追従だけだとカード下部ではなく
    # スライド下部に大きな空白が残り、他10パターンと重心が揃わない。
    # 描画領域の62%を下限として確保する（決定論的｜原則③）。
    avail_h = V17_AREA['body_bottom'] - top
    h = min(avail_h, max(need_h, int(avail_h * 0.62)))

    drawn = 0
    for i, sd in enumerate(sides):
        x = V17_AREA['left'] + i * (col_w + dv + V17_AREA['gap'] * 2)
        score = sd.get('score')
        fill_hex = _tier_fill(palette, score, V17_CONTRAST_TIERS[i])
        _v17_card(slide, x, top, col_w, h, fill_hex, palette)

        txt_color = _text_color_on(fill_hex)
        pad = 18
        add_text(slide, x + pad, top + 14, col_w - pad * 2,
                 str(sd.get('label', '')), 18, bold=True, color=txt_color,
                 height_px=28)
        cur = top + 48
        if score is not None:
            add_text(slide, x + pad, cur, col_w - pad * 2, '%s%%' % score, 24,
                     bold=True, color=txt_color, height_px=34)
            cur += 40
        items = [str(t) for t in (sd.get('items') or [])][:5]
        if items:
            add_paragraph_box(
                slide, x + pad, cur, col_w - pad * 2,
                [{'text': '・' + t, 'size': 14} for t in items],
                default_color=txt_color, default_size=14, line_height=1.6,
                height_px=max(top + h - cur - 14, len(items) * 44))
        drawn += 1

    # 中央の分割罫（split_screen の視覚的境界）
    dx = V17_AREA['left'] + col_w + V17_AREA['gap']
    _v17_rule(slide, dx, top, dv, h, palette['primary'])

    notes.append('2要素固定｜2極化色（%s → %s）で落差を可視化'
                 % (V17_CONTRAST_TIERS[0], V17_CONTRAST_TIERS[1]))
    return _v17_report('contrast', drawn, None, notes)


# ---------------------------------------------------------------------
# v17-P3-4｜timeline（時間軸）｜horizontal / progressive / 3〜7
#           ⚠️ requires_axes=True（統括厳守事項｜8/13）
# ---------------------------------------------------------------------
def draw_timeline(slide, palette, data):
    """期間別のマイルストーンを水平時間軸上に配置する。

    Args:
        data : {'title': str,
                'axis_label': str,                          # ⚠️ 必須
                'milestones': [{'label': str, 'axis': str,  # ⚠️ axis も必須
                                'score': int|None,
                                'description': str}, ...]}  # 3〜7
    Returns:
        dict : 描画レポート

    仕様: 集約表 P3-4 行（min3／max7／requires_axes=True／horizontal／
          progressive／原本shape=horizontal_bar）
    ⚠️ requires_axes=True は timeline の本質（統括厳守事項｜8/13）。
       時間軸ラベル（axis_label／各 milestone の axis）が欠落した場合は
       例外を投げず draw_category へフォールバックする（原則①）。
       framework と同一の姿勢。
    """
    spec = DIAGRAM_PATTERN_SPEC['timeline']
    notes = []
    ms = _v17_normalize(data.get('milestones', []))
    n = len(ms)

    def _fallback(reason):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': m.get('label', ''), 'score': m.get('score'),
                            'description': str(m.get('description', ''))}
                           for m in ms],
        }, _fallback_from='timeline')
        rep['notes'].append(reason)
        return rep

    if not (spec['min_elements'] <= n <= spec['max_elements']):
        return _fallback(
            'マイルストーン数 %d は timeline の範囲外（min %d／max %d）'
            '｜category へフォールバック'
            % (n, spec['min_elements'], spec['max_elements']))

    # ⚠️ requires_axes=True の検査（厳守事項）
    axis_label = str(data.get('axis_label') or '').strip()
    missing = [i + 1 for i, m in enumerate(ms)
               if not str(m.get('axis') or '').strip()]
    if not axis_label:
        return _fallback('timeline は requires_axes=True｜axis_label が未指定'
                         '｜category へフォールバック')
    if missing:
        return _fallback('timeline は requires_axes=True｜%s 番目の axis が未指定'
                         '｜category へフォールバック'
                         % '／'.join(str(i) for i in missing))

    _v17_title(slide, data.get('title', ''), palette)
    _v17_axis_label(slide, axis_label, palette)

    top = V17_AREA['body_top'] + 24
    field_h = V17_AREA['body_bottom'] - top
    axis_y = top + field_h // 2
    dot_d = spec['dot_d']

    # 時間軸本体（水平バー）
    _v17_rule(slide, V17_AREA['left'], axis_y - spec['axis_h'] // 2,
              V17_AREA['width'], spec['axis_h'], palette['midtone'])

    slot_w = V17_AREA['width'] // n
    card_w = min(slot_w - V17_AREA['gap'], 220)
    card_h = 104
    tiers = V17_SEQUENCE_TIERS.get(n, V17_SEQUENCE_TIERS[6])

    drawn = 0
    for i, m in enumerate(ms):
        slot_cx = V17_AREA['left'] + slot_w * i + slot_w // 2
        score = m.get('score')
        fill_hex = _tier_fill(palette, score, tiers[min(i, len(tiers) - 1)])

        # マイルストーンドット（軸上）
        add_shape(slide, _mso_oval(), slot_cx - dot_d // 2,
                  axis_y - dot_d // 2, dot_d, dot_d,
                  fill=hex_to_rgb(fill_hex))

        # 交互配置（上下）で重なりを避ける｜決定論的（原則③）
        above = (i % 2 == 0)
        cy = axis_y - dot_d // 2 - 14 - card_h if above else axis_y + dot_d // 2 + 32
        cx = slot_cx - card_w // 2
        cx = max(V17_AREA['left'], min(cx, V17_AREA['right'] - card_w))
        _v17_card(slide, cx, cy, card_w, card_h, fill_hex, palette)

        txt_color = _text_color_on(fill_hex)
        pad = 10
        add_text(slide, cx + pad, cy + 7, card_w - pad * 2,
                 str(m.get('label', '')), 15, bold=True, color=txt_color,
                 height_px=22)
        cur = cy + 31
        if score is not None:
            add_text(slide, cx + pad, cur, card_w - pad * 2, '%s%%' % score,
                     18, bold=True, color=txt_color, height_px=26)
            cur += 28
        desc = str(m.get('description', '')).strip()
        if desc and cur + 20 <= cy + card_h - 6:
            add_text(slide, cx + pad, cur, card_w - pad * 2, desc, 14,
                     color=txt_color, height_px=cy + card_h - cur - 6)

        # 時間軸ラベル（軸の反対側｜requires_axes=True の実体）
        lab_y = axis_y + dot_d // 2 + 6 if above else axis_y - dot_d // 2 - 26
        add_text(slide, slot_cx - slot_w // 2, lab_y, slot_w,
                 str(m.get('axis', '')), 14, bold=True,
                 color=hex_to_rgb(palette['secondary']), height_px=22,
                 align=_pp_center())
        drawn += 1

    notes.append('マイルストーン %d｜requires_axes=True 充足（軸ラベル＋各期間ラベル）' % n)
    notes.append('カードは軸に対し交互配置（決定論的｜重なり回避）')
    return _v17_report('timeline', drawn, None, notes)


# ---------------------------------------------------------------------
# v17-P3-5｜network（ネットワーク）｜階層型 / depth_hierarchical / 3〜7
# ---------------------------------------------------------------------
def _network_depths(nodes, edges):
    """各ノードの深度を決定論的に算出する（原則③）。

    ルート（入次数0）を深度0とし、エッジをたどって深度を確定する。
    循環がある場合も無限ループしないよう訪問済みを管理する。
    """
    ids = [str(nd.get('id') or nd.get('label') or i) for i, nd in enumerate(nodes)]
    idx = {k: i for i, k in enumerate(ids)}
    indeg = dict((k, 0) for k in ids)
    adj = dict((k, []) for k in ids)
    for e in edges:
        a = str(e.get('from', ''))
        b = str(e.get('to', ''))
        if a in idx and b in idx:
            adj[a].append(b)
            indeg[b] += 1
    roots = [k for k in ids if indeg[k] == 0] or [ids[0]]
    depth = dict((k, None) for k in ids)
    queue = [(r, 0) for r in roots]
    while queue:
        k, d = queue.pop(0)
        if depth[k] is not None and depth[k] <= d:
            continue
        depth[k] = d
        for nb in adj[k]:
            queue.append((nb, d + 1))
    for k in ids:
        if depth[k] is None:
            depth[k] = 0          # 孤立ノードは深度0に置く（劣化描画｜原則①）
    return ids, idx, depth


def _draw_network_edge(slide, palette, x1, y1, x2, y2, width_px):
    """ノード間のエッジを描く。

    ⚠️ add_connector（MSO_CONNECTOR）を第一選択とし、非対応環境では
       極細矩形の回転で代替する（原則①④｜cycle の教訓を適用）。
    """
    try:
        from pptx.enum.shapes import MSO_CONNECTOR
        from pptx.util import Pt
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px(x1), px(y1), px(x2), px(y2))
        conn.line.color.rgb = hex_to_rgb(palette['midtone'])
        conn.line.width = Pt(max(width_px * 0.75, 1))
        return conn
    except Exception:
        import math
        length = int(round(math.hypot(x2 - x1, y2 - y1))) or 1
        ang = math.degrees(math.atan2(y2 - y1, x2 - x1))
        mx = (x1 + x2) // 2 - length // 2
        my = (y1 + y2) // 2 - width_px // 2
        sh = _v17_rule(slide, mx, my, length, max(width_px, 2),
                       palette['midtone'])
        try:
            sh.rotation = ang
        except Exception:
            pass
        return sh


def draw_network(slide, palette, data):
    """ノード間の関係性を階層レイアウト（上→下）で描画する。

    Args:
        data : {'title': str,
                'nodes': [{'id': str, 'label': str,
                           'score': int|None}, ...],        # 3〜7
                'edges': [{'from': str, 'to': str}, ...]}
    Returns:
        dict : 描画レポート

    仕様: 集約表 P3-5 行（min3／max7／hierarchical_top_down／
          depth_hierarchical／原本shape=node_edge）
    ⚠️ 再現性担保のため力学配置（force-directed）は採らず、
       決定論的な階層レイアウトを採用（集約表 P3-5 行の根拠｜原則③）。
    ⚠️ フォールバックは3起動条件（統括指示｜8/12議題4）:
       (1) ノード数が範囲外
       (2) エッジ参照が不整合（存在しないノードIDを参照）
       (3) エッジが1本もない（node_edge の本質を満たさない）
    """
    spec = DIAGRAM_PATTERN_SPEC['network']
    notes = []
    nodes = _v17_normalize(data.get('nodes', []))
    edges = [e for e in (data.get('edges') or []) if isinstance(e, dict)]
    n = len(nodes)

    def _fallback(reason):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': nd.get('label', ''), 'score': nd.get('score'),
                            'description': ''} for nd in nodes],
        }, _fallback_from='network')
        rep['notes'].append(reason)
        return rep

    # 起動条件(1)｜ノード数
    if not (spec['min_elements'] <= n <= spec['max_elements']):
        return _fallback(
            'ノード数 %d は network の範囲外（min %d／max %d）'
            '｜category へフォールバック'
            % (n, spec['min_elements'], spec['max_elements']))

    ids, idx, depth = _network_depths(nodes, edges)

    # 起動条件(2)｜エッジ参照の整合
    bad = [(str(e.get('from', '')), str(e.get('to', '')))
           for e in edges
           if str(e.get('from', '')) not in idx or str(e.get('to', '')) not in idx]
    if bad:
        return _fallback('エッジ参照が不整合（%s）｜category へフォールバック'
                         % '／'.join('%s→%s' % b for b in bad[:3]))

    # 起動条件(3)｜エッジ0本
    if not edges:
        return _fallback('エッジが0本｜node_edge の本質を満たさない'
                         '｜category へフォールバック')

    _v17_title(slide, data.get('title', ''), palette)

    top = V17_AREA['body_top']
    field_h = V17_AREA['body_bottom'] - top
    nw, nh = spec['node_w'], spec['node_h']

    # 深度ごとにグルーピング（決定論的）
    groups = {}
    for k in ids:
        groups.setdefault(depth[k], []).append(k)
    max_depth = max(groups.keys())
    row_h = field_h / float(max_depth + 1)

    center = {}
    for d in sorted(groups.keys()):
        row = groups[d]
        cnt = len(row)
        span = V17_AREA['width'] / float(cnt)
        y = int(top + row_h * d + (row_h - nh) / 2.0)
        for j, k in enumerate(row):
            x = int(V17_AREA['left'] + span * j + (span - nw) / 2.0)
            x = max(V17_AREA['left'], min(x, V17_AREA['right'] - nw))
            center[k] = (x, y)

    # エッジを先に描く（ノードの下に来るように）
    for e in edges:
        a, b = str(e.get('from')), str(e.get('to'))
        ax, ay = center[a]
        bx, by = center[b]
        _draw_network_edge(slide, palette,
                           ax + nw // 2, ay + nh, bx + nw // 2, by,
                           spec['edge_w'])

    drawn = 0
    for k in ids:
        nd = nodes[idx[k]]
        x, y = center[k]
        score = nd.get('score')
        base_key = V17_NETWORK_TIERS[min(depth[k], len(V17_NETWORK_TIERS) - 1)]
        fill_hex = _tier_fill(palette, score, base_key)
        _v17_card(slide, x, y, nw, nh, fill_hex, palette)

        txt_color = _text_color_on(fill_hex)
        pad = 10
        add_text(slide, x + pad, y + 8, nw - pad * 2,
                 str(nd.get('label', '')), 15, bold=True, color=txt_color,
                 height_px=22)
        if score is not None:
            add_text(slide, x + pad, y + 32, nw - pad * 2, '%s%%' % score, 16,
                     bold=True, color=txt_color, height_px=24)
        drawn += 1

    notes.append('ノード %d／エッジ %d｜深度 0〜%d の階層レイアウト（決定論的）'
                 % (n, len(edges), max_depth))
    notes.append('エッジ交差の回避は v3.5 範囲外（統括承認｜8/12議題4）')
    return _v17_report('network', drawn, None, notes)


# ---------------------------------------------------------------------
# v17-P3-6｜ディスパッチの拡張（P1 3種＋P2 3種＋P3 5種＝11種）
# ---------------------------------------------------------------------
def draw_pattern(slide, pattern_key, palette, data):
    """パターンキーで描画関数を振り分ける（11種に対応）。

    ⚠️ v3.5コアの実装対象は11パターン。原本 DIAGNOSIS_TO_PATTERN に
       対応診断カテゴリを持たない 'integration' は構造的に対象外であり、
       category へ退避して notes に理由を記録する（例外は送出しない）。
    """
    table = {
        # P1
        'category':   draw_category,
        'breakdown':  draw_breakdown,
        'comparison': draw_comparison,
        # P2
        'pyramid':    draw_pyramid,
        'sequence':   draw_sequence,
        'framework':  draw_framework,
        # P3
        'funnel':     draw_funnel,
        'cycle':      draw_cycle,
        'contrast':   draw_contrast,
        'timeline':   draw_timeline,
        'network':    draw_network,
    }
    fn = table.get(pattern_key)
    if fn is None:
        src = (data.get('categories') or data.get('items')
               or data.get('components') or data.get('levels')
               or data.get('steps') or data.get('cells')
               or data.get('stages') or data.get('phases')
               or data.get('sides') or data.get('milestones')
               or data.get('nodes') or [])
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': str(x.get('label', '')) if isinstance(x, dict) else str(x),
                            'score': x.get('score') if isinstance(x, dict) else None,
                            'description': ''} for x in src],
        }, _fallback_from=pattern_key)
        if pattern_key == 'integration':
            rep['notes'].append(
                'パターン "integration" は原本 DIAGNOSIS_TO_PATTERN に対応診断'
                'カテゴリがないため v3.5 の実装対象外（構造的に対象外）')
        else:
            rep['notes'].append('パターン "%s" は未定義キー｜category へ退避'
                                % pattern_key)
        return rep
    return fn(slide, palette, data)


def _mso_oval():
    from pptx.enum.shapes import MSO_SHAPE
    return MSO_SHAPE.OVAL

# =====================================================================
# ▲▲▲ v17 P3 追加ブロック ここまで ▲▲▲
# =====================================================================
