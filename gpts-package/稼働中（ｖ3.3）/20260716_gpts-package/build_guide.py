"""
UI診断ディレクター 使い方ガイド 15ページ制作スクリプト
- ベース：03_pptx_builder.py のカラー・shape機能を流用
- 出力：PPTX（A4横：29.7cm × 21cm）
- Section 1-2 を最優先で制作（表紙 + はじめに2P + 基本操作3P + 出力物3P + 裏表紙 = 9P）
"""
from pptx import Presentation
from pptx.util import Emu, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ==========================================
# カラー定数（03_pptx_builder.py から流用）
# ==========================================
NAVY = RGBColor(0x1C, 0x36, 0x6C)
NAVY_LIGHT = RGBColor(0x4A, 0x63, 0x9E)
RED = RGBColor(0xD0, 0x02, 0x1B)
ORANGE = RGBColor(0xF9, 0x73, 0x16)
ORANGE_LIGHT = RGBColor(0xFF, 0xF2, 0xE8)
GOLD = RGBColor(0xE0, 0xB4, 0x2C)
TEXT = RGBColor(0x1E, 0x1E, 0x1E)
SUB_TEXT = RGBColor(0x5C, 0x5C, 0x5C)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
BORDER_GRAY = RGBColor(0xD0, 0xD0, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# キャンバス設定：A4横 = 29.7cm × 21cm
CANVAS_W_CM = 29.7
CANVAS_H_CM = 21.0
PX = 9525  # 1px = 9525 EMU（ただしCM単位も併用）


def create_guide_presentation():
    """A4横キャンバスのPPTX作成"""
    prs = Presentation()
    prs.slide_width = Cm(CANVAS_W_CM)
    prs.slide_height = Cm(CANVAS_H_CM)
    return prs


def add_blank_slide(prs):
    """空スライドを追加"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, x_cm, y_cm, w_cm, h_cm, fill=None, line=None, line_width_pt=None):
    """矩形shape追加（cm指定）"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm)
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is not None:
        shape.line.color.rgb = line
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, x_cm, y_cm, w_cm, h_cm, text, size_pt=14,
             bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font='Meiryo'):
    """テキストボックス追加"""
    tb = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(w_cm), Cm(h_cm))
    tf = tb.text_frame
    tf.margin_left = Cm(0.15)
    tf.margin_right = Cm(0.15)
    tf.margin_top = Cm(0.05)
    tf.margin_bottom = Cm(0.05)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, list):
        # 複数段落
        for i, t in enumerate(text):
            if i > 0:
                p = tf.add_paragraph()
                p.alignment = align
            run = p.add_run()
            run.text = t
            run.font.size = Pt(size_pt)
            run.font.bold = bold
            run.font.color.rgb = color
            run.font.name = font
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_page_header(slide, page_num, total, section_name):
    """全ページ共通のヘッダー帯"""
    # 上端にNAVYの帯
    add_rect(slide, 0, 0, CANVAS_W_CM, 1.2, fill=NAVY)
    # 商品名
    add_text(slide, 0.8, 0.15, 20, 0.9,
             'UI診断ディレクター 使い方ガイド', 14, bold=True, color=WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    # セクション名（右寄せ）
    add_text(slide, 15, 0.15, 12, 0.9,
             section_name, 12, color=WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # ページ番号（右下）
    add_text(slide, CANVAS_W_CM - 3, CANVAS_H_CM - 1, 2.5, 0.7,
             f'{page_num:02d} / {total:02d}', 11, color=SUB_TEXT,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_page_title(slide, title, subtitle=None):
    """ページ内のタイトル（左端赤縦帯 + 大タイトル）"""
    # 左端に赤縦帯
    add_rect(slide, 0.8, 1.8, 0.3, 1.5, fill=RED)
    # タイトル
    add_text(slide, 1.4, 1.7, 25, 1.2, title, 24, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, 1.4, 2.7, 25, 0.7, subtitle, 14, color=SUB_TEXT,
                 anchor=MSO_ANCHOR.TOP)


# ==========================================
# スライド1: 表紙
# ==========================================
def build_cover(prs):
    slide = add_blank_slide(prs)
    # 全面NAVY背景
    add_rect(slide, 0, 0, CANVAS_W_CM, CANVAS_H_CM, fill=NAVY)
    # 中央帯（強調）
    add_rect(slide, 0, 7, CANVAS_W_CM, 6, fill=WHITE)
    # 商品名（大）
    add_text(slide, 2, 7.5, 25.7, 2.5,
             'UI診断ディレクター', 44, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # サブタイトル
    add_text(slide, 2, 9.8, 25.7, 1.5,
             '使い方ガイド', 28, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # キャッチコピー
    add_text(slide, 2, 11.5, 25.7, 1,
             'URLを1本送るだけ。Webディレクター視点の10項目診断を、PPTX3ファイルで即納品。',
             13, color=SUB_TEXT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 下部：バージョン情報
    add_text(slide, 2, 18.5, 25.7, 1,
             'v1.0（2026年7月版）', 11, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ==========================================
# スライド2: 目次
# ==========================================
def build_toc(prs):
    slide = add_blank_slide(prs)
    add_page_header(slide, 0, 15, 'CONTENTS')
    add_page_title(slide, '目次', 'このガイドの全体像')

    # セクション別リスト
    sections = [
        ('はじめに', '01. このガイドについて\n02. UI診断ディレクターでできること', NAVY),
        ('Section 1：基本操作', '03. GPTsへのアクセス方法\n04. 診断対象URLの送り方\n05. 生成完了までの流れ', NAVY),
        ('Section 2：出力物の理解', '06. UIスコアカードの読み方\n07. 改善提案リストの活用方法\n08. ビジュアル診断ボードの使い方', NAVY),
        ('Section 3：応用と対話', '09. 生成後にAIと相談する方法\n10. 業種別の活用パターン\n11. クライアント提案への転用', SUB_TEXT),
        ('Section 4：トラブルシューティング', '12. URLアクセス失敗時の対処\n13. 出力が期待と異なる場合\n14. サポート窓口の使い方', SUB_TEXT),
        ('おわりに', '15. アップデート情報の受け取り方', SUB_TEXT),
    ]

    y = 4
    for section, items, color in sections:
        # セクション名（見出し帯）
        add_rect(slide, 1.5, y, 6, 0.8, fill=LIGHT_GRAY,
                 line=BORDER_GRAY, line_width_pt=0.5)
        add_text(slide, 1.5, y, 6, 0.8, section, 12, bold=True, color=color,
                 anchor=MSO_ANCHOR.MIDDLE)
        # 項目一覧
        add_text(slide, 8, y - 0.05, 20, 0.9 * (items.count('\n') + 1),
                 items, 11, color=color, anchor=MSO_ANCHOR.TOP)
        y += 0.9 * (items.count('\n') + 1) + 0.5


# ==========================================
# スライド3: 01. このガイドについて
# ==========================================
def build_p01(prs):
    slide = add_blank_slide(prs)
    add_page_header(slide, 1, 15, 'はじめに')
    add_page_title(slide, '01. このガイドについて',
                   'UI診断ディレクターを、迷わず使い始めるための道しるべ')

    # 本文
    body_top = 4.2
    add_text(slide, 1.5, body_top, 26, 1.5,
             'このガイドの目的', 14, bold=True, color=NAVY)
    add_text(slide, 1.5, body_top + 0.8, 26, 3,
             'UI診断ディレクター（GPTs）を購入いただき、ありがとうございます。\n'
             'このガイドは、購入後の初回利用から、応用的な活用まで、\n'
             '「迷わず・気持ちよく使える」ようにするための解説書です。',
             12, color=TEXT)

    add_text(slide, 1.5, body_top + 3.5, 26, 1,
             '対象読者', 14, bold=True, color=NAVY)
    add_text(slide, 1.5, body_top + 4.3, 26, 3,
             '・ご購入いただいた全ての方\n'
             '・サイト運営者、EC担当者、Webディレクター、制作会社の実務者',
             12, color=TEXT)

    add_text(slide, 1.5, body_top + 7, 26, 1,
             '構成マップ', 14, bold=True, color=NAVY)
    add_text(slide, 1.5, body_top + 7.8, 26, 4,
             '本ガイドは以下の5セクション・15ページで構成されています。\n\n'
             '・はじめに（2ページ）\n'
             '・Section 1：基本操作（3ページ）\n'
             '・Section 2：出力物の理解（3ページ）\n'
             '・Section 3：応用と対話（3ページ）\n'
             '・Section 4：トラブルシューティング（3ページ）\n'
             '・おわりに（1ページ）',
             12, color=TEXT)

    # 注記帯
    add_rect(slide, 1.5, 18.5, 26.7, 1.5, fill=ORANGE_LIGHT,
             line=ORANGE, line_width_pt=0.75)
    add_text(slide, 1.8, 18.5, 26.2, 1.5,
             '💡 困った時は、Section 4「トラブルシューティング」を先に読んでいただいて構いません。',
             12, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


# ==========================================
# スライド4: 02. UI診断ディレクターでできること
# ==========================================
def build_p02(prs):
    slide = add_blank_slide(prs)
    add_page_header(slide, 2, 15, 'はじめに')
    add_page_title(slide, '02. UI診断ディレクターでできること',
                   'URLを1本送るだけで、3ファイルのレポートを納品')

    # 3ファイル説明カード
    body_top = 4.5
    card_w = 8.2
    card_h = 7
    gap = 0.6

    files = [
        ('①', 'UIスコアカード', 'NAVY',
         '10項目 × 5段階評価と\n総合ランクを1枚に集約。\n\n経営層への現状共有に。'),
        ('②', '改善提案リスト', 'RED',
         '優先度・工数付きの\n改善提案を5件明示。\n\n実装計画の指示書に。'),
        ('③', 'ビジュアル診断ボード', 'ORANGE',
         '構造マップ・行動フロー・\nBefore/After を視覚化。\n\nプレゼン資料として。'),
    ]

    color_map = {'NAVY': NAVY, 'RED': RED, 'ORANGE': ORANGE}

    for i, (num, title, color_key, desc) in enumerate(files):
        x = 1.5 + i * (card_w + gap)
        col = color_map[color_key]
        # 上部色帯
        add_rect(slide, x, body_top, card_w, 1, fill=col)
        # 番号
        add_text(slide, x + 0.3, body_top + 0.1, 1.5, 0.8,
                 num, 22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # タイトル
        add_text(slide, x + 1.8, body_top + 0.1, card_w - 2, 0.8,
                 title, 14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # カード本体
        add_rect(slide, x, body_top + 1, card_w, card_h - 1, fill=WHITE,
                 line=BORDER_GRAY, line_width_pt=0.75)
        # 説明
        add_text(slide, x + 0.4, body_top + 1.3, card_w - 0.8, card_h - 1.6,
                 desc, 12, color=TEXT, anchor=MSO_ANCHOR.TOP)

    # 想定シーン帯
    scene_top = body_top + card_h + 1
    add_rect(slide, 1.5, scene_top, 26.7, 3, fill=LIGHT_GRAY,
             line=BORDER_GRAY, line_width_pt=0.5)
    add_text(slide, 1.8, scene_top + 0.2, 26.2, 0.8,
             '💼 想定される活用シーン', 13, bold=True, color=NAVY)
    add_text(slide, 1.8, scene_top + 1, 26.2, 2,
             '・社内改善会議での現状共有（スコアカードで一目で把握）\n'
             '・クライアントへの提案書（改善提案リストをそのまま流用）\n'
             '・実装チームへの指示書（優先度付きで作業順序を明確化）',
             12, color=TEXT, anchor=MSO_ANCHOR.TOP)


# ==========================================
# スライド5: 03. GPTsへのアクセス方法
# ==========================================
def build_p03(prs):
    slide = add_blank_slide(prs)
    add_page_header(slide, 3, 15, 'Section 1：基本操作')
    add_page_title(slide, '03. GPTsへのアクセス方法',
                   '購入後、初めてUI診断ディレクターを開くまで')

    # ステップリスト
    steps = [
        ('STEP 1', 'ChatGPT にログイン',
         'ブラウザで chat.openai.com を開くか、ChatGPT アプリを起動します。\n'
         'ChatGPT Plus（有料プラン）のアカウントが必要です。'),
        ('STEP 2', 'GPTs を開く',
         '左サイドバーの「GPTs を探す」または「マイGPT」からアクセスします。\n'
         'または、購入時にご案内するURLから直接開けます。'),
        ('STEP 3', 'UI診断ディレクターを起動',
         'GPTsの一覧から「UI診断ディレクター」を選択します。\n'
         'よく使うので、右上の「☆」でお気に入りに追加すると便利です。'),
    ]

    y = 4.3
    for label, title, desc in steps:
        # STEP バッジ
        add_rect(slide, 1.5, y, 3, 1, fill=NAVY)
        add_text(slide, 1.5, y, 3, 1, label, 13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # タイトル
        add_text(slide, 4.8, y, 20, 1, title, 14, bold=True, color=NAVY,
                 anchor=MSO_ANCHOR.MIDDLE)
        # 説明
        add_text(slide, 4.8, y + 1.1, 22, 2, desc, 12, color=TEXT,
                 anchor=MSO_ANCHOR.TOP)
        y += 3.7

    # 注記
    add_rect(slide, 1.5, 17.5, 26.7, 2, fill=ORANGE_LIGHT,
             line=ORANGE, line_width_pt=0.75)
    add_text(slide, 1.8, 17.5, 26.2, 2,
             '⚠ ChatGPT Plus未加入の方へ：\n'
             'UI診断ディレクターの利用には、ChatGPTの有料プラン（月額20USD）が必要です。\n'
             '購入前に、ChatGPT Plus加入をご確認ください（2026年7月時点の情報）。',
             11, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)


# ==========================================
# スライド6: 04. 診断対象URLの送り方
# ==========================================
def build_p04(prs):
    slide = add_blank_slide(prs)
    add_page_header(slide, 4, 15, 'Section 1：基本操作')
    add_page_title(slide, '04. 診断対象URLの送り方',
                   '初回でも迷わない、正しい診断リクエストの書き方')

    # 良い例／悪い例
    y = 4.3
    # 良い例カード
    add_rect(slide, 1.5, y, 12.8, 6.5, fill=WHITE,
             line=NAVY, line_width_pt=1.5)
    add_rect(slide, 1.5, y, 12.8, 1, fill=NAVY)
    add_text(slide, 1.5, y, 12.8, 1, '✓ 良い送り方', 14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 1.9, y + 1.3, 12, 5,
             '・https:// から始まる完全なURLで送る\n'
             '・診断したいページを1つに絞る（トップページ推奨）\n'
             '・シンプルな依頼文で送る\n\n'
             '例：\n'
             '「https://example.com を診断してください」\n\n'
             '・ログイン必須ページはスクショで送る',
             12, color=TEXT, anchor=MSO_ANCHOR.TOP)

    # 悪い例カード
    add_rect(slide, 15, y, 12.8, 6.5, fill=WHITE,
             line=RED, line_width_pt=1.5)
    add_rect(slide, 15, y, 12.8, 1, fill=RED)
    add_text(slide, 15, y, 12.8, 1, '✗ 避けたい送り方', 14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 15.4, y + 1.3, 12, 5,
             '・URLだけで指示なし\n'
             '・複数URLを同時に送る\n'
             '・「https://」を省略した表記\n\n'
             '例：\n'
             '「example.com とか shop.example と\n'
             '  あと corp.example もお願いします」\n\n'
             '・詳細な要望を長文で追加する',
             12, color=TEXT, anchor=MSO_ANCHOR.TOP)

    # コツ帯
    tips_top = 11.5
    add_rect(slide, 1.5, tips_top, 26.7, 3.5, fill=LIGHT_GRAY,
             line=BORDER_GRAY, line_width_pt=0.5)
    add_text(slide, 1.8, tips_top + 0.2, 26.2, 0.8,
             '💡 診断精度を上げるコツ', 13, bold=True, color=NAVY)
    add_text(slide, 1.8, tips_top + 1.1, 26.2, 2.5,
             '・トップページ推奨：サイト全体の設計思想が反映されているため、診断示唆が実務者に刺さる\n'
             '・SPA（React/Vue製のサイト）は動的読込のため、スクリーンショット併用がおすすめ\n'
             '・広告用LPは単一動線として最適な診断対象、無料相談CVを高めるヒントが得られる',
             11, color=TEXT, anchor=MSO_ANCHOR.TOP)


# ==========================================
# スライド7: 05. 生成完了までの流れ
# ==========================================
def build_p05(prs):
    slide = add_blank_slide(prs)
    add_page_header(slide, 5, 15, 'Section 1：基本操作')
    add_page_title(slide, '05. 生成完了までの流れ',
                   '3〜5分の待ち時間中に、何が起きているか')

    # フロー図
    y = 4.5
    steps = [
        ('①', '受付', 'URLまたはスクショを解析', NAVY),
        ('②', '診断', '10項目のスコアと\n強み・課題を判定', NAVY),
        ('③', '提案', '優先度付きの\n改善提案5件を生成', ORANGE),
        ('④', 'PPTX生成', 'Code Interpreter で\n3ファイル同時生成', ORANGE),
        ('⑤', '納品', 'ダウンロードリンク\nを並列表示', RED),
    ]

    step_w = 5
    gap = 0.4
    total_w = step_w * 5 + gap * 4
    start_x = (CANVAS_W_CM - total_w) / 2

    for i, (num, title, desc, color) in enumerate(steps):
        x = start_x + i * (step_w + gap)
        # ステップ円（上部）
        add_rect(slide, x + step_w/2 - 1, y, 2, 2, fill=color)
        add_text(slide, x + step_w/2 - 1, y, 2, 2, num, 24, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # タイトル
        add_text(slide, x, y + 2.3, step_w, 1, title, 14, bold=True,
                 color=NAVY, align=PP_ALIGN.CENTER)
        # 説明
        add_text(slide, x, y + 3.3, step_w, 2, desc, 11,
                 color=SUB_TEXT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        # 矢印（最後以外）
        if i < len(steps) - 1:
            arrow_x = x + step_w + gap/2 - 0.3
            add_text(slide, arrow_x, y + 0.5, 0.6, 1, '→', 22, bold=True,
                     color=ORANGE, align=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE)

    # 想定時間
    time_top = 11.5
    add_rect(slide, 1.5, time_top, 26.7, 1.2, fill=NAVY)
    add_text(slide, 1.5, time_top, 26.7, 1.2,
             '▶ 想定所要時間：3〜5分（サイトの規模により変動）', 14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 注記
    note_top = time_top + 1.7
    add_text(slide, 1.5, note_top, 26.7, 1,
             '待ち時間中の表示メッセージ', 13, bold=True, color=NAVY)
    add_text(slide, 1.5, note_top + 0.8, 26.7, 4,
             '・「診断を開始します」→ URL/スクショの受付完了\n'
             '・「診断中です」→ 10項目分析中\n'
             '・「レポートを生成中です」→ PPTX組み立て中\n'
             '・「診断が完了しました」→ ダウンロードリンク表示\n\n'
             '5分以上待っても応答がない場合は、Section 4-12 の対処法をご確認ください。',
             11, color=TEXT, anchor=MSO_ANCHOR.TOP)


# ==========================================
# スライド8: 06. UIスコアカードの読み方
# ==========================================
def build_p06(prs):
    slide = add_blank_slide(prs)
    add_page_header(slide, 6, 15, 'Section 2：出力物の理解')
    add_page_title(slide, '06. UIスコアカードの読み方',
                   '10項目診断の結果を、経営層と共有する1枚')

    # 3つのポイント
    points = [
        ('1', '10項目診断とは', NAVY,
         'ファーストビュー、キャッチコピー、CTA設計、信頼性、フォーム、レスポンシブ、\n'
         '読みやすさ、情報設計、ブランド、表示速度の10軸で診断します。\n'
         '各項目5段階評価、合計50点満点。'),
        ('2', '5段階評価の基準', NAVY,
         '5=業界トップクラス／4=良好／3=標準／2=改善必要／1=致命的\n'
         '「3」が実務における標準的な状態。「4以上」で維持、「2以下」は最優先対応。'),
        ('3', 'ランクと総合スコア', RED,
         'S（45-50）／A（38-44）／B（30-37）／C（20-29）／D（0-19）\n'
         '購入検討層に最も響くのはB〜Cの帯。改善余地を明示することが価値提供の中核。'),
    ]

    y = 4.3
    for num, title, color, desc in points:
        # 番号バッジ
        add_rect(slide, 1.5, y, 1.2, 1.5, fill=color)
        add_text(slide, 1.5, y, 1.2, 1.5, num, 22, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # タイトル
        add_text(slide, 3, y, 24, 0.9, title, 14, bold=True, color=color)
        # 説明
        add_text(slide, 3, y + 0.9, 24, 2.5, desc, 12, color=TEXT,
                 anchor=MSO_ANCHOR.TOP)
        y += 3.9

    # 活用の指針
    tip_top = y + 0.3
    add_rect(slide, 1.5, tip_top, 26.7, 2, fill=ORANGE_LIGHT,
             line=ORANGE, line_width_pt=0.75)
    add_text(slide, 1.8, tip_top + 0.2, 26.2, 0.8,
             '💡 活用の指針', 13, bold=True, color=NAVY)
    add_text(slide, 1.8, tip_top + 1, 26.2, 1,
             '「強み3項目」で自信のポイントを、「最優先課題3項目」で次のアクションを明示できます。\n'
             '経営層への5分プレゼンには、この1枚が最強の武器になります。',
             11, color=TEXT, anchor=MSO_ANCHOR.TOP)


# ==========================================
# スライド9: 07. 改善提案リストの活用方法
# ==========================================
def build_p07(prs):
    slide = add_blank_slide(prs)
    add_page_header(slide, 7, 15, 'Section 2：出力物の理解')
    add_page_title(slide, '07. 改善提案リストの活用方法',
                   '優先度に応じた実装計画の立て方')

    # 3つのカラム：優先度別
    y = 4.5
    cols = [
        ('高優先', '今週やること', RED,
         '成果影響度・改善容易度がともに高い項目。\n\n'
         '・FVでのCTA可視化\n'
         '・購入導線の主要ボタン設置\n'
         '・入力フォームの負荷軽減'),
        ('中優先', '来月やること', ORANGE,
         '成果には効くが、実装に一定のリソースが必要な項目。\n\n'
         '・料金プランの比較表化\n'
         '・レビュー配置の再設計\n'
         '・情報設計の再構築'),
        ('低優先', '四半期でやること', SUB_TEXT,
         '中長期の改善余地。基盤刷新と併せて検討する項目。\n\n'
         '・表示速度の最適化\n'
         '・レスポンシブ全面刷新\n'
         '・アクセシビリティ強化'),
    ]

    card_w = 8.7
    gap = 0.3
    for i, (level, timing, color, desc) in enumerate(cols):
        x = 1.5 + i * (card_w + gap)
        # 上部帯
        add_rect(slide, x, y, card_w, 1.5, fill=color)
        add_text(slide, x, y + 0.05, card_w, 0.7, level, 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, y + 0.75, card_w, 0.7, timing, 12,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 本体
        add_rect(slide, x, y + 1.5, card_w, 8, fill=WHITE,
                 line=BORDER_GRAY, line_width_pt=0.75)
        add_text(slide, x + 0.4, y + 1.8, card_w - 0.8, 7.5, desc, 12,
                 color=TEXT, anchor=MSO_ANCHOR.TOP)

    # Before/After構造の読み方
    note_top = y + 10.2
    add_rect(slide, 1.5, note_top, 26.7, 2, fill=LIGHT_GRAY,
             line=BORDER_GRAY, line_width_pt=0.5)
    add_text(slide, 1.8, note_top + 0.2, 26.2, 0.8,
             '📖 Before/After構造の読み方', 13, bold=True, color=NAVY)
    add_text(slide, 1.8, note_top + 1, 26.2, 1,
             '各提案には「現状（Before）」と「改善後（After）」が明示されています。\n'
             'この構造は、そのままクライアント提案書のスライドとして流用できます。',
             11, color=TEXT, anchor=MSO_ANCHOR.TOP)


# ==========================================
# スライド10: 08. ビジュアル診断ボードの使い方
# ==========================================
def build_p08(prs):
    slide = add_blank_slide(prs)
    add_page_header(slide, 8, 15, 'Section 2：出力物の理解')
    add_page_title(slide, '08. ビジュアル診断ボードの使い方',
                   'プレゼン資料としてそのまま使う3枚構成')

    # 3枚構成の説明
    slides_desc = [
        ('スライド 1', '構造マップ + 総評 + 最重要課題', NAVY,
         '・LP構造マップ（サイトのセクション別状態）\n'
         '・◤ 総評（改善方向の1行スローガン）\n'
         '・⚠ 最重要課題（Top3または4）\n'
         '・◎ 強み（活かすべき点3点）\n'
         '・ユーザー行動フロー6ステップ'),
        ('スライド 2', 'スコア視覚化 + 強み・課題', NAVY,
         '・10項目のバーチャート（強み・弱み一目視認）\n'
         '・強み3項目\n'
         '・最優先課題3項目\n'
         '・改善方向の再掲\n\n'
         '経営層への5分プレゼンで最も刺さる1枚です。'),
        ('スライド 3', 'Before/After ハイライトTop3', ORANGE,
         '・優先度Top3の改善案を、Before/After形式で並列表示\n'
         '・優先度バッジ・工数・箇所を明示\n'
         '・「何をどう変えるか」が即座に伝わる構造\n\n'
         '実装チームへの指示書としてそのまま流用可能。'),
    ]

    y = 4.3
    card_w = 8.7
    gap = 0.3
    for i, (label, title, color, desc) in enumerate(slides_desc):
        x = 1.5 + i * (card_w + gap)
        # 上部帯
        add_rect(slide, x, y, card_w, 1.7, fill=color)
        add_text(slide, x + 0.3, y + 0.1, card_w - 0.6, 0.7, label, 13,
                 bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + 0.3, y + 0.8, card_w - 0.6, 0.8, title, 12,
                 color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        # 本体
        add_rect(slide, x, y + 1.7, card_w, 8, fill=WHITE,
                 line=BORDER_GRAY, line_width_pt=0.75)
        add_text(slide, x + 0.4, y + 2, card_w - 0.8, 7.5, desc, 12,
                 color=TEXT, anchor=MSO_ANCHOR.TOP)

    # プレゼンストーリー例
    note_top = y + 10.4
    add_rect(slide, 1.5, note_top, 26.7, 1.8, fill=ORANGE_LIGHT,
             line=ORANGE, line_width_pt=0.75)
    add_text(slide, 1.8, note_top + 0.15, 26.2, 0.8,
             '🎯 プレゼンストーリー例（5分版）', 13, bold=True, color=NAVY)
    add_text(slide, 1.8, note_top + 0.9, 26.2, 1,
             '「現状はこうです（1）」→「特にここを直すべきです（2）」→「こう変えます（3）」の流れで、\n'
             '意思決定に必要な情報を5分で伝えられます。',
             11, color=TEXT, anchor=MSO_ANCHOR.TOP)


# ==========================================
# 最終ページ：おわりに（暫定・15Pガイド完成前の案内）
# ==========================================
def build_placeholder(prs):
    slide = add_blank_slide(prs)
    add_page_header(slide, 9, 15, 'ご案内')
    add_page_title(slide, 'Section 3 以降について',
                   '応用・トラブル対処・アップデートは順次追加します')

    add_text(slide, 1.5, 5, 26.7, 4,
             'このガイドは v1.0（Section 1-2 完成版）です。\n\n'
             '以下のセクションは、7月17日以降に順次追加されます。',
             13, color=TEXT)

    upcoming = [
        ('Section 3：応用と対話', ORANGE,
         '09. 生成後にAIと相談する方法\n'
         '10. 業種別の活用パターン\n'
         '11. クライアント提案への転用'),
        ('Section 4：トラブルシューティング', ORANGE,
         '12. URLアクセス失敗時の対処\n'
         '13. 出力が期待と異なる場合\n'
         '14. サポート窓口の使い方'),
        ('おわりに', ORANGE,
         '15. アップデート情報の受け取り方'),
    ]

    y = 10
    for section, color, items in upcoming:
        add_rect(slide, 1.5, y, 8, 2.5, fill=color)
        add_text(slide, 1.5, y, 8, 2.5, section, 13, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, 10, y, 18, 2.5, items, 12, color=TEXT,
                 anchor=MSO_ANCHOR.MIDDLE)
        y += 2.8

    # 更新通知
    add_rect(slide, 1.5, 19.2, 26.7, 1, fill=NAVY)
    add_text(slide, 1.5, 19.2, 26.7, 1,
             '📮 アップデート通知：note @Yasuaki_Irie / X @Yasuaki_Irie にてご案内予定',
             12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ==========================================
# メイン
# ==========================================
def main():
    prs = create_guide_presentation()
    build_cover(prs)
    build_toc(prs)
    build_p01(prs)
    build_p02(prs)
    build_p03(prs)
    build_p04(prs)
    build_p05(prs)
    build_p06(prs)
    build_p07(prs)
    build_p08(prs)
    build_placeholder(prs)
    out = '/tmp/UI診断ディレクター_使い方ガイド_v1.0_Section1-2.pptx'
    prs.save(out)
    print(f'Generated: {out}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
