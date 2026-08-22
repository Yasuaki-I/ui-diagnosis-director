# -*- coding: utf-8 -*-
"""v17 P1｜21組合せ動作テスト（3パターン×7テーマ）＋境界/フォールバック検証"""
import io, os, re, sys, traceback

BUILDER = os.environ.get('BUILDER', 'builder_v16_5.py')
V17 = os.environ.get('V17', 'v17/v17_p1_draw_patterns_20260823.py')
OUT = os.environ.get('OUT', 'v17/test_out')

g = {'__name__': 'builder_v16_5'}
exec(compile(open(BUILDER, encoding='utf-8').read(), BUILDER, 'exec'), g)
exec(compile(open(V17, encoding='utf-8').read(), V17, 'exec'), g)

THEMES = ['SolidGray', 'Blue', 'LightBlue', 'Green', 'Cyan', 'Red', 'Orange']

CATEGORY_DATA = {
    'title': '診断結果｜改善提案のカテゴリ分類',
    'categories': [
        {'label': 'ファーストビュー', 'score': 82, 'description': '訴求と導線が一致している'},
        {'label': '情報設計',       'score': 64, 'description': '見出し階層に飛びがある'},
        {'label': '導線設計',       'score': 38, 'description': 'CTAが画面外に落ちている'},
        {'label': '可読性',         'score': 71, 'description': '行間は適正・字間に余地'},
    ],
}
BREAKDOWN_DATA = {
    'title': '総合スコアの内訳',
    'whole': {'label': '総合スコア', 'value': 100},
    'components': [
        {'label': '視認性',   'value': 32, 'score': 78, 'note': '良好'},
        {'label': '導線',     'value': 24, 'score': 36, 'note': '要改善'},
        {'label': '情報設計', 'value': 18, 'score': 62, 'note': '中位'},
        {'label': '信頼性',   'value': 14, 'score': 70, 'note': '中位'},
        {'label': '表示速度', 'value': 12, 'score': 55, 'note': '中位'},
    ],
}
COMPARISON_DATA = {
    'title': '改善前後のUX比較',
    'comparison_axis': '改善前 / 改善後',
    'attribute_labels': ['CTA到達', '離脱率', '平均滞在'],
    'items': [
        {'label': 'Before', 'score': 38,
         'attributes': {'CTA到達': '2.1%', '離脱率': '68%', '平均滞在': '41秒'}},
        {'label': 'After', 'score': 76,
         'attributes': {'CTA到達': '5.4%', '離脱率': '43%', '平均滞在': '96秒'}},
    ],
}

results = []


def run(label, fn):
    try:
        fn()
        results.append((label, 'PASS', ''))
    except Exception as e:
        results.append((label, 'FAIL', '%s: %s' % (type(e).__name__, e)))
        traceback.print_exc()


# --- T1: 21組合せ（3パターン×7テーマ）---------------------------------
for theme in THEMES:
    prs = g['create_presentation']()
    pal = g['get_theme_palette'](theme)
    for pname, data, fn in (
        ('category', CATEGORY_DATA, g['draw_category']),
        ('breakdown', BREAKDOWN_DATA, g['draw_breakdown']),
        ('comparison', COMPARISON_DATA, g['draw_comparison']),
    ):
        def _do(fn=fn, prs=prs, pal=pal, data=data, pname=pname, theme=theme):
            sl = g['_blank_slide'](prs)
            g['_add_header'](sl, '%s / %s' % (pname, theme), pname.upper())
            g['_add_footer'](sl, 1, 3)
            rep = fn(sl, pal, data)
            assert rep['pattern'] in (pname, 'category'), rep
        run('T1｜%s × %s' % (pname, theme), _do)
    prs.save(os.path.join(OUT, 'v17_p1_%s.pptx' % theme))

# --- T2: category 要素数境界 3/4/5/6 と 2件・7件（丸め込み）------------
for n in (2, 3, 4, 5, 6, 7):
    def _do(n=n):
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        pal = g['get_theme_palette']('Blue')
        data = {'title': 'category n=%d' % n,
                'categories': [{'label': 'C%d' % (i + 1), 'score': 50 + i,
                                'description': 'desc'} for i in range(n)]}
        rep = g['draw_category'](sl, pal, data)
        exp = min(n, 6)
        assert rep['elements_drawn'] == exp, (n, rep)
    run('T2｜category 要素数 n=%d' % n, _do)

# --- T3: breakdown 境界 3〜7 / 範囲外2・8 / 数値不正 → フォールバック ---
for n in (2, 3, 5, 7, 8):
    def _do(n=n):
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        pal = g['get_theme_palette']('Green')
        data = {'title': 'breakdown n=%d' % n,
                'whole': {'label': '全体', 'value': None},
                'components': [{'label': 'B%d' % (i + 1), 'value': 10 + i,
                                'score': 60, 'note': ''} for i in range(n)]}
        rep = g['draw_breakdown'](sl, pal, data)
        if 3 <= n <= 7:
            assert rep['pattern'] == 'breakdown', rep
        else:
            assert rep['pattern'] == 'category' and rep['fallback_from'] == 'breakdown', rep
    run('T3｜breakdown 要素数 n=%d' % n, _do)


def _bad_value():
    prs = g['create_presentation']()
    sl = g['_blank_slide'](prs)
    pal = g['get_theme_palette']('Red')
    data = {'title': 'breakdown 数値不正',
            'components': [{'label': 'X', 'value': 'abc'},
                           {'label': 'Y', 'value': 10},
                           {'label': 'Z', 'value': 5}]}
    rep = g['draw_breakdown'](sl, pal, data)
    assert rep['fallback_from'] == 'breakdown', rep


run('T3｜breakdown 数値不正→category', _bad_value)


def _zero_total():
    prs = g['create_presentation']()
    sl = g['_blank_slide'](prs)
    pal = g['get_theme_palette']('Red')
    data = {'title': 'breakdown 合計0',
            'components': [{'label': 'X', 'value': 0}, {'label': 'Y', 'value': 0},
                           {'label': 'Z', 'value': 0}]}
    rep = g['draw_breakdown'](sl, pal, data)
    assert rep['fallback_from'] == 'breakdown', rep


run('T3｜breakdown 合計0→category', _zero_total)

# --- T4: comparison 2/3要素 と 範囲外1・4 -----------------------------
for n in (1, 2, 3, 4):
    def _do(n=n):
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        pal = g['get_theme_palette']('Orange')
        data = {'title': 'comparison n=%d' % n,
                'comparison_axis': '軸' if n == 2 else None,
                'attribute_labels': ['A', 'B'],
                'items': [{'label': 'I%d' % (i + 1), 'score': 40 + i * 10,
                           'attributes': {'A': '1', 'B': '2'}} for i in range(n)]}
        rep = g['draw_comparison'](sl, pal, data)
        if 2 <= n <= 3:
            assert rep['pattern'] == 'comparison', rep
        else:
            assert rep['fallback_from'] == 'comparison', rep
    run('T4｜comparison 要素数 n=%d' % n, _do)

# --- T5: 警告オーバーライド（score<40 → warning色）--------------------
def _warn():
    pal = g['get_theme_palette']('Blue')
    assert g['_tier_fill'](pal, 39, 'primary') == pal['warning']
    assert g['_tier_fill'](pal, 40, 'primary') == pal['primary']
    assert g['_tier_fill'](pal, None, 'primary') == pal['primary']


run('T5｜警告オーバーライド閾値40', _warn)

# --- T6: score=None 許容（警告判定スキップ）---------------------------
def _none_score():
    prs = g['create_presentation']()
    sl = g['_blank_slide'](prs)
    pal = g['get_theme_palette']('Cyan')
    data = {'title': 'score省略',
            'categories': [{'label': 'A', 'score': None, 'description': 'x'},
                           {'label': 'B', 'score': None, 'description': 'y'},
                           {'label': 'C', 'score': None, 'description': 'z'}]}
    rep = g['draw_category'](sl, pal, data)
    assert rep['elements_drawn'] == 3


run('T6｜score=None 許容', _none_score)

# --- T7: 決定論性（同一入力→同一XML）100回 ---------------------------
def _determinism():
    sigs = set()
    for _ in range(100):
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        pal = g['get_theme_palette']('Blue')
        g['draw_category'](sl, pal, CATEGORY_DATA)
        from lxml import etree
        sigs.add(etree.tostring(sl._element))
    assert len(sigs) == 1, len(sigs)


run('T7｜決定論性 category 100回→1種', _determinism)


def _determinism2():
    sigs = set()
    from lxml import etree
    for _ in range(50):
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        pal = g['get_theme_palette']('Green')
        g['draw_breakdown'](sl, pal, BREAKDOWN_DATA)
        sigs.add(etree.tostring(sl._element))
    assert len(sigs) == 1, len(sigs)


run('T7｜決定論性 breakdown 50回→1種', _determinism2)

# --- T8: 描画範囲（ヘッダ60〜フッター660px内）------------------------
def _bounds():
    from pptx.util import Emu
    prs = g['create_presentation']()
    sl = g['_blank_slide'](prs)
    pal = g['get_theme_palette']('Blue')
    g['draw_category'](sl, pal, CATEGORY_DATA)
    bad = []
    for sh in sl.shapes:
        l = sh.left / 9525.0
        t = sh.top / 9525.0
        r = l + sh.width / 9525.0
        b = t + sh.height / 9525.0
        if l < 0 or t < 60 or r > 1280 or b > 660:
            bad.append((sh.shape_type, round(l), round(t), round(r), round(b)))
    assert not bad, bad


run('T8｜描画範囲 60〜660px 内', _bounds)

# --- T9: フォント最小14pt / メイリオ強制（set_run 経由の担保）--------
def _font():
    from pptx.oxml.ns import qn
    prs = g['create_presentation']()
    sl = g['_blank_slide'](prs)
    pal = g['get_theme_palette']('Blue')
    g['draw_comparison'](sl, pal, COMPARISON_DATA)
    small = []
    nonmeiryo = []
    for sh in sl.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size is not None and r.font.size.pt < 14:
                    small.append(r.text)
                rPr = r._r.find(qn('a:rPr'))
                if rPr is not None:
                    for tag in ('a:latin', 'a:ea', 'a:cs'):
                        el = rPr.find(qn(tag))
                        if el is None or el.get('typeface') != 'メイリオ':
                            nonmeiryo.append((r.text, tag))
    assert not small, small
    assert not nonmeiryo, nonmeiryo[:5]


run('T9｜最小14pt・メイリオ3スクリプト強制', _font)

# --- T10: resolve_pattern（原本 DIAGNOSIS_TO_PATTERN 参照）-----------
def _resolve():
    assert g['resolve_pattern']('proposal_categorization') == 'category'
    assert g['resolve_pattern']('score_breakdown') == 'breakdown'
    assert g['resolve_pattern']('before_after') == 'comparison'
    assert g['resolve_pattern']('unknown_key') == 'category'


run('T10｜resolve_pattern マッピング', _resolve)

# --- T11: draw_pattern 未実装キー→category 退避 ----------------------
def _dispatch():
    prs = g['create_presentation']()
    sl = g['_blank_slide'](prs)
    pal = g['get_theme_palette']('Blue')
    rep = g['draw_pattern'](sl, pal and pal, {'title': 't', 'items': [
        {'label': 'a'}, {'label': 'b'}, {'label': 'c'}]}) if False else None
    rep = g['draw_pattern'](sl, 'pyramid', pal,
                            {'title': 't',
                             'items': [{'label': 'a', 'score': 50},
                                       {'label': 'b', 'score': 60},
                                       {'label': 'c', 'score': 70}]})
    assert rep['fallback_from'] == 'pyramid', rep


run('T11｜draw_pattern 未実装キー退避', _dispatch)

# --- T12: 既存C-1〜C-3の後方互換（同時使用で保存できる）--------------
def _compat():
    prs = g['create_presentation']()
    g['add_cover'](prs, title='v17 後方互換確認', date='2026年8月',
                   author='UI診断ディレクター', page_total=3)
    sl = g['_blank_slide'](prs)
    g['_add_header'](sl, 'v17 図解パターン', 'DIAGRAM')
    g['_add_footer'](sl, 2, 3)
    g['draw_category'](sl, g['get_theme_palette']('Blue'), CATEGORY_DATA)
    g['add_closing'](prs, message='Thank you.', page_num=3, total=3)
    prs.save(os.path.join(OUT, 'v17_p1_compat.pptx'))


run('T12｜既存レイアウトとの後方互換', _compat)

# --- 集計 -------------------------------------------------------------
npass = sum(1 for _, s, _ in results if s == 'PASS')
nfail = len(results) - npass
print('\n===== v17 P1 テスト結果 =====')
for label, st, msg in results:
    print('%-40s %s %s' % (label, st, msg))
print('-----')
print('合計 %d 項目｜PASS %d｜FAIL %d' % (len(results), npass, nfail))
sys.exit(1 if nfail else 0)
