# -*- coding: utf-8 -*-
"""v17 P2｜21組合せ動作テスト（3パターン×7テーマ）＋境界値・フォールバック検証"""
import os, sys, traceback

BUILDER = os.environ.get('BUILDER', 'builder_v17_10.py')
P2 = os.environ.get('P2', '')
OUT = os.environ.get('OUT', 'v17/p2_out')
os.makedirs(OUT, exist_ok=True)

g = {'__name__': 'builder'}
exec(compile(open(BUILDER, encoding='utf-8').read(), BUILDER, 'exec'), g)
if P2:
    exec(compile(open(P2, encoding='utf-8').read(), P2, 'exec'), g)

THEMES = ['SolidGray', 'Blue', 'LightBlue', 'Green', 'Cyan', 'Red', 'Orange']

PYRAMID = {
    'title': '改善施策の優先順位',
    'levels': [
        {'label': 'FVに申込CTAを追加', 'score': 82, 'description': '工数小・効果大'},
        {'label': '料金直下に導線を置く', 'score': 64, 'description': '工数小'},
        {'label': '信頼材料を増やす',   'score': 48, 'description': '工数中'},
        {'label': '表示速度の最適化',   'score': 36, 'description': '工数中'},
    ],
}
SEQUENCE = {
    'title': 'ユーザー行動フローと課題',
    'steps': [
        {'label': '流入',   'score': 78, 'description': 'FVを見る'},
        {'label': '理解',   'score': 66, 'description': '特徴を確認'},
        {'label': '比較',   'score': 58, 'description': '料金を確認'},
        {'label': '検討',   'score': 38, 'description': '声を見る'},
        {'label': '申込',   'score': 30, 'description': 'CTAを探す'},
    ],
}
FRAMEWORK = {
    'title': 'インパクト × コスト マトリクス',
    'axis_x_label': 'インパクト', 'axis_y_label': 'コスト',
    'axis_x_low': '小', 'axis_x_high': '大',
    'axis_y_low': '低', 'axis_y_high': '高',
    'cells': [
        {'row': 0, 'col': 0, 'label': '再設計',   'score': 45, 'items': ['LP全面改修']},
        {'row': 0, 'col': 1, 'label': '要検討',   'score': 62, 'items': ['動画追加', '事例充実']},
        {'row': 1, 'col': 0, 'label': '後回し',   'score': 35, 'items': ['配色微調整']},
        {'row': 1, 'col': 1, 'label': '即着手',   'score': 88, 'items': ['CTA追加', '導線整理']},
    ],
}

results = []


def run(label, fn):
    try:
        fn()
        results.append((label, 'PASS', ''))
    except Exception as e:
        results.append((label, 'FAIL', '%s: %s' % (type(e).__name__, e)))
        traceback.print_exc()


# --- T1: 21組合せ（3パターン×7テーマ）------------------------------
for theme in THEMES:
    prs = g['create_presentation']()
    pal = g['get_theme_palette'](theme)
    for pname, data in (('pyramid', PYRAMID), ('sequence', SEQUENCE),
                        ('framework', FRAMEWORK)):
        def _do(pname=pname, data=data, prs=prs, pal=pal):
            sl, rep = g['add_diagram_slide'](prs, pname, pal, data,
                                             page_num=1, total=3)
            assert rep['pattern'] == pname, rep
            assert rep['fallback_from'] is None, rep
        run('T1｜%s × %s' % (pname, theme), _do)
    prs.save(os.path.join(OUT, 'v17_p2_%s.pptx' % theme))

# --- T2: pyramid 境界 3/4/5 と範囲外 2/6 ----------------------------
for n in (2, 3, 4, 5, 6):
    def _do(n=n):
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        pal = g['get_theme_palette']('Blue')
        rep = g['draw_pyramid'](sl, pal, {'title': 'pyramid n=%d' % n,
            'levels': [{'label': 'L%d' % (i + 1), 'score': 50 + i,
                        'description': 'd'} for i in range(n)]})
        if 3 <= n <= 5:
            assert rep['pattern'] == 'pyramid', rep
        else:
            assert rep['fallback_from'] == 'pyramid', rep
    run('T2｜pyramid 段数 n=%d' % n, _do)

# --- T3: sequence 境界 3/6 と範囲外 2/7 -----------------------------
for n in (2, 3, 4, 6, 7):
    def _do(n=n):
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        pal = g['get_theme_palette']('Green')
        rep = g['draw_sequence'](sl, pal, {'title': 'sequence n=%d' % n,
            'steps': [{'label': 'S%d' % (i + 1), 'score': 60,
                       'description': 'd'} for i in range(n)]})
        if 3 <= n <= 6:
            assert rep['pattern'] == 'sequence', rep
        else:
            assert rep['fallback_from'] == 'sequence', rep
    run('T3｜sequence ステップ n=%d' % n, _do)

# --- T4: framework セル数 4〜9 と範囲外 3/10 ------------------------
EXPECT_GRID = {4: '2x2', 5: '3x2', 6: '3x2', 7: '3x3', 8: '3x3', 9: '3x3'}
for n in (3, 4, 5, 6, 7, 8, 9, 10):
    def _do(n=n):
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        pal = g['get_theme_palette']('Orange')
        rep = g['draw_framework'](sl, pal, {'title': 'framework n=%d' % n,
            'axis_x_label': 'X', 'axis_y_label': 'Y',
            'cells': [{'label': 'C%d' % (i + 1), 'score': 55} for i in range(n)]})
        if 4 <= n <= 9:
            assert rep['pattern'] == 'framework', rep
            assert EXPECT_GRID[n] in ' '.join(rep['notes']).replace('グリッド ', ''), rep
        else:
            assert rep['fallback_from'] == 'framework', rep
    run('T4｜framework セル n=%d' % n, _do)

# --- T5: framework 軸ラベル欠落 → category（例外を投げない）---------
for miss in ('axis_x_label', 'axis_y_label', 'both'):
    def _do(miss=miss):
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        pal = g['get_theme_palette']('Cyan')
        d = {'title': 'axes欠落', 'axis_x_label': 'X', 'axis_y_label': 'Y',
             'cells': [{'label': 'C%d' % i, 'score': 50} for i in range(4)]}
        if miss == 'both':
            d.pop('axis_x_label'); d.pop('axis_y_label')
        else:
            d.pop(miss)
        rep = g['draw_framework'](sl, pal, d)
        assert rep['fallback_from'] == 'framework', rep
        assert 'requires_axes' in ' '.join(rep['notes']), rep
    run('T5｜framework 軸欠落(%s)→category' % miss, _do)

# --- T6: 警告オーバーライド（score<40）------------------------------
def _warn():
    pal = g['get_theme_palette']('Blue')
    assert g['_tier_fill'](pal, 39, 'primary') == pal['warning']
    assert g['_tier_fill'](pal, 40, 'primary') == pal['primary']
    assert g['_tier_fill'](pal, None, 'primary') == pal['primary']
run('T6｜警告オーバーライド閾値40', _warn)

# --- T7: score=None 許容（3パターン）-------------------------------
def _none_score():
    prs = g['create_presentation']()
    pal = g['get_theme_palette']('Red')
    sl1 = g['_blank_slide'](prs)
    r1 = g['draw_pyramid'](sl1, pal, {'title': 't', 'levels': [
        {'label': 'a', 'score': None}, {'label': 'b', 'score': None},
        {'label': 'c', 'score': None}]})
    sl2 = g['_blank_slide'](prs)
    r2 = g['draw_sequence'](sl2, pal, {'title': 't', 'steps': [
        {'label': 'a'}, {'label': 'b'}, {'label': 'c'}]})
    sl3 = g['_blank_slide'](prs)
    r3 = g['draw_framework'](sl3, pal, {'title': 't', 'axis_x_label': 'X',
        'axis_y_label': 'Y', 'cells': [{'label': 'c%d' % i} for i in range(4)]})
    assert (r1['elements_drawn'], r2['elements_drawn'], r3['elements_drawn']) == (3, 3, 4)
run('T7｜score=None 許容（P2 3種）', _none_score)

# --- T8: 決定論性（同一入力→同一XML）-------------------------------
from lxml import etree
def _det(fnname, data, times, theme='Blue'):
    def _do():
        sigs = set()
        for _ in range(times):
            prs = g['create_presentation']()
            sl = g['_blank_slide'](prs)
            g[fnname](sl, g['get_theme_palette'](theme), data)
            sigs.add(etree.tostring(sl._element))
        assert len(sigs) == 1, len(sigs)
    return _do
run('T8｜決定論性 pyramid 100回→1種', _det('draw_pyramid', PYRAMID, 100))
run('T8｜決定論性 sequence 50回→1種', _det('draw_sequence', SEQUENCE, 50))
run('T8｜決定論性 framework 50回→1種', _det('draw_framework', FRAMEWORK, 50))

# --- T9: 描画範囲（ヘッダ60〜フッター660px内）----------------------
def _bounds(fnname, data):
    def _do():
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        g[fnname](sl, g['get_theme_palette']('Blue'), data)
        bad = []
        for sh in sl.shapes:
            l = sh.left / 9525.0; t = sh.top / 9525.0
            r = l + sh.width / 9525.0; b = t + sh.height / 9525.0
            if l < 0 or t < 60 or r > 1280.5 or b > 660.5:
                bad.append((round(l), round(t), round(r), round(b)))
        assert not bad, bad
    return _do
run('T9｜描画範囲 pyramid', _bounds('draw_pyramid', PYRAMID))
run('T9｜描画範囲 sequence', _bounds('draw_sequence', SEQUENCE))
run('T9｜描画範囲 framework', _bounds('draw_framework', FRAMEWORK))

# --- T10: 最小14pt・メイリオ3スクリプト強制 ------------------------
from pptx.oxml.ns import qn
def _font(fnname, data):
    def _do():
        prs = g['create_presentation']()
        sl = g['_blank_slide'](prs)
        g[fnname](sl, g['get_theme_palette']('Blue'), data)
        small, nonm = [], []
        for sh in sl.shapes:
            if not sh.has_text_frame:
                continue
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is not None and r.font.size.pt < 14:
                        small.append(r.text)
                    rPr = r._r.find(qn('a:rPr'))
                    if rPr is not None:
                        for tag in ('a:latin', 'a:ea', 'a:cs'):
                            el = rPr.find(qn(tag))
                            if el is None or el.get('typeface') != 'メイリオ':
                                nonm.append((r.text, tag))
        assert not small, small
        assert not nonm, nonm[:5]
    return _do
run('T10｜フォント pyramid', _font('draw_pyramid', PYRAMID))
run('T10｜フォント sequence', _font('draw_sequence', SEQUENCE))
run('T10｜フォント framework', _font('draw_framework', FRAMEWORK))

# --- T11: resolve_pattern（原本マッピング）--------------------------
def _resolve():
    assert g['resolve_pattern']('priority_ranking') == 'pyramid'
    assert g['resolve_pattern']('user_flow') == 'sequence'
    assert g['resolve_pattern']('impact_cost_matrix') == 'framework'
run('T11｜resolve_pattern P2マッピング', _resolve)

# --- T12: draw_pattern 6種対応 + P3未実装キーの退避 -----------------
def _dispatch():
    prs = g['create_presentation']()
    pal = g['get_theme_palette']('Blue')
    for k, d in (('pyramid', PYRAMID), ('sequence', SEQUENCE),
                 ('framework', FRAMEWORK)):
        sl = g['_blank_slide'](prs)
        assert g['draw_pattern'](sl, k, pal, d)['pattern'] == k
    for k in ('funnel', 'timeline', 'contrast', 'cycle', 'network', 'integration'):
        sl = g['_blank_slide'](prs)
        rep = g['draw_pattern'](sl, k, pal, {'title': 't',
            'levels': [{'label': 'a'}, {'label': 'b'}, {'label': 'c'}]})
        assert rep['fallback_from'] == k, (k, rep)
        assert 'P3' in ' '.join(rep['notes']), rep
run('T12｜draw_pattern 6種＋P3退避6件', _dispatch)

# --- T13: add_diagram_slide のヘッダ重複なし ------------------------
def _no_dup():
    prs = g['create_presentation']()
    pal = g['get_theme_palette']('Blue')
    sl, _ = g['add_diagram_slide'](prs, 'pyramid', pal, PYRAMID, 1, 3)
    texts = [sh.text_frame.text.strip() for sh in sl.shapes
             if sh.has_text_frame and sh.text_frame.text.strip()
             and sh.top / 9525 < 130]
    # ヘッダ帯は原本 ja（'ピラミッド'）、内見出しは data['title']
    assert 'ピラミッド' in texts, texts
    assert PYRAMID['title'] in texts, texts
    assert texts.count(PYRAMID['title']) == 1, texts
run('T13｜ヘッダ帯と内見出しの重複なし', _no_dup)

# --- T14: 既存C-1〜C-3との後方互換（同時使用）----------------------
def _compat():
    prs = g['create_presentation']()
    g['add_cover'](prs, title='v17.1.0 後方互換', date='2026年8月',
                   author='UI診断ディレクター', page_total=4)
    pal = g['get_theme_palette']('Blue')
    for i, (k, d) in enumerate([('pyramid', PYRAMID), ('sequence', SEQUENCE),
                                ('framework', FRAMEWORK)], start=2):
        g['add_diagram_slide'](prs, k, pal, d, page_num=i, total=4)
    prs.save(os.path.join(OUT, 'v17_p2_compat.pptx'))
run('T14｜既存レイアウトとの後方互換', _compat)

npass = sum(1 for _, s, _ in results if s == 'PASS')
nfail = len(results) - npass
print('\n===== v17 P2 テスト結果 =====')
for label, st, msg in results:
    print('%-40s %s %s' % (label, st, msg))
print('-----')
print('合計 %d 項目｜PASS %d｜FAIL %d' % (len(results), npass, nfail))
sys.exit(1 if nfail else 0)
