# -*- coding: utf-8 -*-
# =====================================================================
# ▼▼▼ v17 P3 追加ブロック ここから ▼▼▼
#     対象: funnel / cycle / contrast / timeline / network（5種）
#     配置: 03_pptx_builder.py の末尾（v17 P2 ブロックの直後）に追記
#     版数: 17.1.0 → 17.2.0（__version__ を更新すること）
#
#     設計根拠:
#       - 拡張定義集約表（v35_core_extended_pattern_definitions.md）P3-1〜P3-5 行
#       - 集約表と実装記録（8/13〜8/15）の照合結果: ⭐ 相違 0件
#       - 原本 DIAGRAM_PATTERNS は無改変。拡張層は DIAGRAM_PATTERN_SPEC に分離。
#
#     横断設計原則（v3.5コア）:
#       原則① 描画不能時も例外を投げず劣化描画（category フォールバック）で通す
#       原則② score < 40 は警告色でオーバーライドする
#       原則③ 同一入力から同一出力（決定論性｜乱数・力学配置を使わない）
#       原則④ 基本図形を優先し、環境依存図形を増やさない
#
#     ⚠️ cycle の図形選定について（8/25 実測にもとづく設計変更）
#       実装記録 8/14 判断④は MSO_SHAPE.BLOCK_ARC ＋ rotation を採ると
#       していたが、8/25 に実機同等環境で実測した結果:
#         (a) adjustments は「度」単位（既定値 [108.0, 0.0, 0.25]）であり、
#             実装記録が想定した 0.0〜1.0 比率ではない。
#             → 比率で指定すると極小弧になり視覚的に「線」になる。
#         (b) 正しく度で指定すれば分割弧は描けるが、⚠️ 弧の内部に
#             テキストを格納できないため、他10パターンが持つ
#             「ラベル＋スコア＋説明」の情報構造を表現できない。
#       よって本実装は角丸矩形の円周配置＋RIGHT_ARROW（基本図形のみ）で
#       円環構造を構成する。原本 shape='circle_arrow'（円環＋矢印）との
#       整合は矢印の実装によって維持される。原則④とも整合する。
#       （判断原理14「テストが通ることは実機で描画されることの代替に
#         ならない」を設計段階に先取り適用して得た結論）
# =====================================================================

# ---------------------------------------------------------------------
# v17-P3-0｜拡張仕様の登録（原本 DIAGRAM_PATTERNS は無改変）
# ---------------------------------------------------------------------
DIAGRAM_PATTERN_SPEC.update({
    'funnel': {
        'min_elements': 3, 'max_elements': 6, 'requires_axes': False,
        'direction': 'vertical', 'color_gradation': 'progressive_narrowing',
        # 上段90%→下段40%（原本 use「上から下へ絞り込むファネル型」）
        'width_ratio_top': 0.90, 'width_ratio_bottom': 0.40,
    },
    'cycle': {
        'min_elements': 3, 'max_elements': 6, 'requires_axes': False,
        'direction': 'clockwise', 'color_gradation': 'uniform_cyclic',
        # 12時起点・時計回り固定（集約表 P3-2 行）
        'start_angle_deg': -90.0,
    },
    'contrast': {
        'min_elements': 2, 'max_elements': 2, 'requires_axes': False,
        'direction': 'horizontal', 'color_gradation': 'polarized_contrast',
        'divider_w': 6,
    },
    'timeline': {
        'min_elements': 3, 'max_elements': 7, 'requires_axes': True,
        'direction': 'horizontal', 'color_gradation': 'progressive',
        'axis_h': 8, 'dot_d': 26,
    },
    'network': {
        'min_elements': 3, 'max_elements': 7, 'requires_axes': False,
        'direction': 'hierarchical_top_down',
        'color_gradation': 'depth_hierarchical',
        'node_w': 210, 'node_h': 66, 'edge_w': 3,
    },
})

# 段階減衰色（progressive_narrowing）｜funnel の段数別色キー列
V17_FUNNEL_TIERS = {
    3: ['primary', 'secondary', 'midtone'],
    4: ['primary', 'secondary', 'midtone', 'light'],
    5: ['primary', 'secondary', 'midtone', 'light', 'lightest'],
    6: ['primary', 'primary', 'secondary', 'midtone', 'light', 'lightest'],
}

# 均等循環色（uniform_cyclic）｜4色を循環させる
# ⚠️ 段階減衰は「終わりがある」誤読を生むため使わない（集約表 P3-2 行の根拠）
V17_CYCLE_TIERS = ['primary', 'secondary', 'midtone', 'light']

# 2極化色（polarized_contrast）｜中間色を排し落差を可視化
V17_CONTRAST_TIERS = ['lightest', 'primary']

# 深度別階層色（depth_hierarchical）｜network の depth→色キー
V17_NETWORK_TIERS = ['primary', 'secondary', 'midtone', 'light']


# ---------------------------------------------------------------------
# v17-P3-1｜funnel（絞り込み）｜vertical / progressive_narrowing / 3〜6段
# ---------------------------------------------------------------------
def draw_funnel(slide, palette, data):
    """段階的な絞り込みを上から下へ幅を狭めて表現する。

    Args:
        data : {'title': str,
                'stages': [{'label': str, 'score': int|None,
                            'description': str}, ...]}   # 3〜6・index0が最上段
    Returns:
        dict : 描画レポート

    仕様: 集約表 P3-1 行（min3／max6／vertical／progressive_narrowing／
          原本shape=trapezoid）
    ⚠️ 台形は MSO_SHAPE.TRAPEZOID を使わず、幅可変の矩形段で近似する。
       pyramid（P2-1）と同一方針（原則④）。pyramid は上が狭いが、
       funnel は上が広い（絞り込み方向が逆）。
    範囲外は draw_category へフォールバック（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['funnel']
    notes = []
    stages = _v17_normalize(data.get('stages', []))
    n = len(stages)

    if not (spec['min_elements'] <= n <= spec['max_elements']):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': s.get('label', ''), 'score': s.get('score'),
                            'description': str(s.get('description', ''))}
                           for s in stages],
        }, _fallback_from='funnel')
        rep['notes'].append(
            '段数 %d は funnel の範囲外（min %d／max %d）｜category へフォールバック'
            % (n, spec['min_elements'], spec['max_elements']))
        return rep

    _v17_title(slide, data.get('title', ''), palette)

    top = V17_AREA['body_top']
    avail_h = V17_AREA['body_bottom'] - top
    gap = 10
    band_h = int((avail_h - gap * (n - 1)) / n)
    band_h = min(band_h, 118)

    r_top = spec['width_ratio_top']
    r_bot = spec['width_ratio_bottom']
    tiers = V17_FUNNEL_TIERS.get(n, V17_FUNNEL_TIERS[6])
    cx = V17_AREA['left'] + V17_AREA['width'] // 2

    drawn = 0
    y = top
    for i, st in enumerate(stages):
        ratio = r_top + (r_bot - r_top) * (i / float(max(n - 1, 1)))
        w = int(V17_AREA['width'] * ratio)
        x = cx - w // 2
        score = st.get('score')
        fill_hex = _tier_fill(palette, score, tiers[min(i, len(tiers) - 1)])
        _v17_card(slide, x, y, w, band_h, fill_hex, palette)

        # 段番号（左端｜描画領域内に収める）
        num_x = max(V17_AREA['left'], x - 34)
        add_text(slide, num_x, y + band_h // 2 - 12, 30, str(i + 1), 14,
                 bold=True, color=hex_to_rgb(palette['midtone']), height_px=24)

        txt_color = _text_color_on(fill_hex)
        pad = 12
        add_text(slide, x + pad, y + 8, w - pad * 2, str(st.get('label', '')),
                 16, bold=True, color=txt_color, height_px=24,
                 align=_pp_center())
        line2 = []
        if score is not None:
            line2.append('%s%%' % score)
        desc = str(st.get('description', '')).strip()
        if desc:
            line2.append(desc)
        if line2 and band_h >= 56:
            add_text(slide, x + pad, y + 34, w - pad * 2, '　'.join(line2),
                     14, color=txt_color, height_px=band_h - 42,
                     align=_pp_center())
        drawn += 1
        y += band_h + gap

    notes.append('段数 %d｜幅比 %.2f→%.2f で絞り込みを表現' % (n, r_top, r_bot))
    return _v17_report('funnel', drawn, None, notes)


# ---------------------------------------------------------------------
# v17-P3-2｜cycle（循環）｜clockwise / uniform_cyclic / 3〜6段
# ---------------------------------------------------------------------
def _compute_cycle_positions(n, cx, cy, radius, box_w, box_h, start_deg):
    """円周上の n 個のボックス左上座標を決定論的に返す（原則③）。

    12時起点・時計回り。三角関数のみを使い乱数・力学配置は使わない。
    """
    import math
    step = 360.0 / n
    out = []
    for i in range(n):
        ang = start_deg + i * step
        rad = math.radians(ang)
        px = cx + radius * math.cos(rad) - box_w / 2.0
        py = cy + radius * math.sin(rad) - box_h / 2.0
        out.append((int(round(px)), int(round(py)), ang))
    return out


def _draw_cycle_arrow(slide, palette, cx, cy, radius, angle_deg):
    """円周の接線方向へ向く矢印を1本描く（基本図形 RIGHT_ARROW ＋ rotation）。

    ⚠️ BLOCK_ARC を使わない理由は本ブロック冒頭のコメントを参照。
    """
    import math
    aw, ah = 56, 20
    rad = math.radians(angle_deg)
    ax = cx + radius * math.cos(rad) - aw / 2.0
    ay = cy + radius * math.sin(rad) - ah / 2.0
    sh = add_shape(slide, _mso_right_arrow(), int(round(ax)), int(round(ay)),
                   aw, ah, fill=hex_to_rgb(palette['midtone']))
    try:
        sh.rotation = angle_deg + 90.0
    except Exception:
        pass          # rotation 非対応環境でも矢印自体は残す（原則①）
    return sh


def draw_cycle(slide, palette, data):
    """反復プロセスを円環で表現する（12時起点・時計回り固定）。

    Args:
        data : {'title': str,
                'cycle_name': str,                        # 円環中心の名称（任意）
                'phases': [{'label': str, 'score': int|None,
                            'description': str}, ...]}    # 3〜6
    Returns:
        dict : 描画レポート

    仕様: 集約表 P3-2 行（min3／max6／clockwise／uniform_cyclic／
          原本shape=circle_arrow）
    ⚠️ 色は段階減衰させず4色を循環させる。段階減衰は「終わりがある」
       誤読を生むため（集約表 P3-2 行の根拠）。
    範囲外は draw_category へフォールバック（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['cycle']
    notes = []
    phases = _v17_normalize(data.get('phases', []))
    n = len(phases)

    if not (spec['min_elements'] <= n <= spec['max_elements']):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': p.get('label', ''), 'score': p.get('score'),
                            'description': str(p.get('description', ''))}
                           for p in phases],
        }, _fallback_from='cycle')
        rep['notes'].append(
            '段数 %d は cycle の範囲外（min %d／max %d）｜category へフォールバック'
            % (n, spec['min_elements'], spec['max_elements']))
        return rep

    _v17_title(slide, data.get('title', ''), palette)

    top = V17_AREA['body_top']
    field_h = V17_AREA['body_bottom'] - top
    cx = V17_AREA['left'] + V17_AREA['width'] // 2
    cy = top + field_h // 2

    # v17-P3-fix1（8/25 自己目視）: 説明文が2行に折返してカード内が窮屈になる
    # 事象を検出。ボックス幅を広げ高さも確保する（充填率と可読性の両立）。
    box_w = 268 if n <= 4 else 212
    box_h = 100 if n <= 4 else 90
    radius = int(min(V17_AREA['width'] / 2 - box_w / 2 - 16,
                     field_h / 2 - box_h / 2 - 6))

    positions = _compute_cycle_positions(n, cx, cy, radius, box_w, box_h,
                                         spec['start_angle_deg'])

    # 矢印は各ボックスの中間角に置く（時計回りの進行を示す）
    step = 360.0 / n
    for i in range(n):
        mid = spec['start_angle_deg'] + (i + 0.5) * step
        _draw_cycle_arrow(slide, palette, cx, cy, int(radius * 0.62), mid)

    drawn = 0
    for i, (x, y, _ang) in enumerate(positions):
        ph = phases[i]
        score = ph.get('score')
        base_key = V17_CYCLE_TIERS[i % len(V17_CYCLE_TIERS)]
        fill_hex = _tier_fill(palette, score, base_key)
        _v17_card(slide, x, y, box_w, box_h, fill_hex, palette)

        txt_color = _text_color_on(fill_hex)
        pad = 10
        head = '%d. %s' % (i + 1, str(ph.get('label', '')))
        add_text(slide, x + pad, y + 7, box_w - pad * 2, head, 15, bold=True,
                 color=txt_color, height_px=22)
        sub = []
        if score is not None:
            sub.append('%s%%' % score)
        desc = str(ph.get('description', '')).strip()
        if desc:
            sub.append(desc)
        if sub:
            add_text(slide, x + pad, y + 31, box_w - pad * 2, '　'.join(sub),
                     14, color=txt_color, height_px=box_h - 38)
        drawn += 1

    # 判断⑤｜円環中心にサイクル名を配置
    cname = str(data.get('cycle_name', '')).strip()
    if cname:
        add_text(slide, cx - 110, cy - 16, 220, cname, 18, bold=True,
                 color=hex_to_rgb(palette['primary']), height_px=32,
                 align=_pp_center())

    notes.append('段数 %d｜12時起点・時計回り｜4色循環（uniform_cyclic）' % n)
    notes.append('円環は角丸矩形の円周配置＋RIGHT_ARROW で構成'
                 '（BLOCK_ARC 不使用｜8/25 実測にもとづく設計判断）')
    return _v17_report('cycle', drawn, None, notes)


# ---------------------------------------------------------------------
# v17-P3-3｜contrast（対比）｜horizontal / polarized_contrast / 2固定
# ---------------------------------------------------------------------
def draw_contrast(slide, palette, data):
    """対照的な2要素を左右分割で並列強調する。

    Args:
        data : {'title': str,
                'sides': [{'label': str, 'score': int|None,
                           'items': [str]}, ...]}     # ⚠️ 2固定
    Returns:
        dict : 描画レポート

    仕様: 集約表 P3-3 行（min2／max2 固定／horizontal／polarized_contrast／
          原本shape=split_screen）
    ⚠️ 中間色を排し2極化（薄→濃）する。目的は「推移」ではなく「落差」の
       可視化であるため（集約表 P3-3 行の根拠）。
    ⚠️ requires_axes=False のため軸ラベルは描かない（framework との差異）。
    2要素以外は draw_category へフォールバック（原則①）。
    """
    spec = DIAGRAM_PATTERN_SPEC['contrast']
    notes = []
    sides = _v17_normalize(data.get('sides', []))
    n = len(sides)

    if n != 2:
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': s.get('label', ''), 'score': s.get('score'),
                            'description': ''} for s in sides],
        }, _fallback_from='contrast')
        rep['notes'].append(
            '要素数 %d は contrast の範囲外（2固定）｜category へフォールバック' % n)
        return rep

    _v17_title(slide, data.get('title', ''), palette)

    top = V17_AREA['body_top']
    dv = spec['divider_w']
    col_w = (V17_AREA['width'] - dv - V17_AREA['gap'] * 2) // 2

    # v17-P3-fix2（8/25 自己目視）: items が少ないとカード下部が大きく空き、
    # 充填率が低下する事象を検出（P2 framework の V17_CARD_H_MAX と同趣旨）。
    # 内容量（ラベル+スコア+items 行数）から必要高を算出し上限で丸める。
    # 行高の実測値: 14pt×line_height1.6 ≒ 32px／折返し1行を許容して 44px 見込む
    max_items = max(len([t for t in (sd.get('items') or [])][:5]) for sd in sides)
    has_score = any(sd.get('score') is not None for sd in sides)
    need_h = (14 + 28) + (40 if has_score else 0) + max_items * 44 + 24
    # v17-P3-fix3（8/25 自己目視）: 内容追従だけだとカード下部ではなく
    # スライド下部に大きな空白が残り、他10パターンと重心が揃わない。
    # 描画領域の62%を下限として確保する（決定論的｜原則③）。
    avail_h = V17_AREA['body_bottom'] - top
    h = min(avail_h, max(need_h, int(avail_h * 0.62)))

    drawn = 0
    for i, sd in enumerate(sides):
        x = V17_AREA['left'] + i * (col_w + dv + V17_AREA['gap'] * 2)
        score = sd.get('score')
        fill_hex = _tier_fill(palette, score, V17_CONTRAST_TIERS[i])
        _v17_card(slide, x, top, col_w, h, fill_hex, palette)

        txt_color = _text_color_on(fill_hex)
        pad = 18
        add_text(slide, x + pad, top + 14, col_w - pad * 2,
                 str(sd.get('label', '')), 18, bold=True, color=txt_color,
                 height_px=28)
        cur = top + 48
        if score is not None:
            add_text(slide, x + pad, cur, col_w - pad * 2, '%s%%' % score, 24,
                     bold=True, color=txt_color, height_px=34)
            cur += 40
        items = [str(t) for t in (sd.get('items') or [])][:5]
        if items:
            add_paragraph_box(
                slide, x + pad, cur, col_w - pad * 2,
                [{'text': '・' + t, 'size': 14} for t in items],
                default_color=txt_color, default_size=14, line_height=1.6,
                height_px=max(top + h - cur - 14, len(items) * 44))
        drawn += 1

    # 中央の分割罫（split_screen の視覚的境界）
    dx = V17_AREA['left'] + col_w + V17_AREA['gap']
    _v17_rule(slide, dx, top, dv, h, palette['primary'])

    notes.append('2要素固定｜2極化色（%s → %s）で落差を可視化'
                 % (V17_CONTRAST_TIERS[0], V17_CONTRAST_TIERS[1]))
    return _v17_report('contrast', drawn, None, notes)


# ---------------------------------------------------------------------
# v17-P3-4｜timeline（時間軸）｜horizontal / progressive / 3〜7
#           ⚠️ requires_axes=True（統括厳守事項｜8/13）
# ---------------------------------------------------------------------
def draw_timeline(slide, palette, data):
    """期間別のマイルストーンを水平時間軸上に配置する。

    Args:
        data : {'title': str,
                'axis_label': str,                          # ⚠️ 必須
                'milestones': [{'label': str, 'axis': str,  # ⚠️ axis も必須
                                'score': int|None,
                                'description': str}, ...]}  # 3〜7
    Returns:
        dict : 描画レポート

    仕様: 集約表 P3-4 行（min3／max7／requires_axes=True／horizontal／
          progressive／原本shape=horizontal_bar）
    ⚠️ requires_axes=True は timeline の本質（統括厳守事項｜8/13）。
       時間軸ラベル（axis_label／各 milestone の axis）が欠落した場合は
       例外を投げず draw_category へフォールバックする（原則①）。
       framework と同一の姿勢。
    """
    spec = DIAGRAM_PATTERN_SPEC['timeline']
    notes = []
    ms = _v17_normalize(data.get('milestones', []))
    n = len(ms)

    def _fallback(reason):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': m.get('label', ''), 'score': m.get('score'),
                            'description': str(m.get('description', ''))}
                           for m in ms],
        }, _fallback_from='timeline')
        rep['notes'].append(reason)
        return rep

    if not (spec['min_elements'] <= n <= spec['max_elements']):
        return _fallback(
            'マイルストーン数 %d は timeline の範囲外（min %d／max %d）'
            '｜category へフォールバック'
            % (n, spec['min_elements'], spec['max_elements']))

    # ⚠️ requires_axes=True の検査（厳守事項）
    axis_label = str(data.get('axis_label') or '').strip()
    missing = [i + 1 for i, m in enumerate(ms)
               if not str(m.get('axis') or '').strip()]
    if not axis_label:
        return _fallback('timeline は requires_axes=True｜axis_label が未指定'
                         '｜category へフォールバック')
    if missing:
        return _fallback('timeline は requires_axes=True｜%s 番目の axis が未指定'
                         '｜category へフォールバック'
                         % '／'.join(str(i) for i in missing))

    _v17_title(slide, data.get('title', ''), palette)
    _v17_axis_label(slide, axis_label, palette)

    top = V17_AREA['body_top'] + 24
    field_h = V17_AREA['body_bottom'] - top
    axis_y = top + field_h // 2
    dot_d = spec['dot_d']

    # 時間軸本体（水平バー）
    _v17_rule(slide, V17_AREA['left'], axis_y - spec['axis_h'] // 2,
              V17_AREA['width'], spec['axis_h'], palette['midtone'])

    slot_w = V17_AREA['width'] // n
    card_w = min(slot_w - V17_AREA['gap'], 220)
    card_h = 104
    tiers = V17_SEQUENCE_TIERS.get(n, V17_SEQUENCE_TIERS[6])

    drawn = 0
    for i, m in enumerate(ms):
        slot_cx = V17_AREA['left'] + slot_w * i + slot_w // 2
        score = m.get('score')
        fill_hex = _tier_fill(palette, score, tiers[min(i, len(tiers) - 1)])

        # マイルストーンドット（軸上）
        add_shape(slide, _mso_oval(), slot_cx - dot_d // 2,
                  axis_y - dot_d // 2, dot_d, dot_d,
                  fill=hex_to_rgb(fill_hex))

        # 交互配置（上下）で重なりを避ける｜決定論的（原則③）
        above = (i % 2 == 0)
        cy = axis_y - dot_d // 2 - 14 - card_h if above else axis_y + dot_d // 2 + 32
        cx = slot_cx - card_w // 2
        cx = max(V17_AREA['left'], min(cx, V17_AREA['right'] - card_w))
        _v17_card(slide, cx, cy, card_w, card_h, fill_hex, palette)

        txt_color = _text_color_on(fill_hex)
        pad = 10
        add_text(slide, cx + pad, cy + 7, card_w - pad * 2,
                 str(m.get('label', '')), 15, bold=True, color=txt_color,
                 height_px=22)
        cur = cy + 31
        if score is not None:
            add_text(slide, cx + pad, cur, card_w - pad * 2, '%s%%' % score,
                     18, bold=True, color=txt_color, height_px=26)
            cur += 28
        desc = str(m.get('description', '')).strip()
        if desc and cur + 20 <= cy + card_h - 6:
            add_text(slide, cx + pad, cur, card_w - pad * 2, desc, 14,
                     color=txt_color, height_px=cy + card_h - cur - 6)

        # 時間軸ラベル（軸の反対側｜requires_axes=True の実体）
        lab_y = axis_y + dot_d // 2 + 6 if above else axis_y - dot_d // 2 - 26
        add_text(slide, slot_cx - slot_w // 2, lab_y, slot_w,
                 str(m.get('axis', '')), 14, bold=True,
                 color=hex_to_rgb(palette['secondary']), height_px=22,
                 align=_pp_center())
        drawn += 1

    notes.append('マイルストーン %d｜requires_axes=True 充足（軸ラベル＋各期間ラベル）' % n)
    notes.append('カードは軸に対し交互配置（決定論的｜重なり回避）')
    return _v17_report('timeline', drawn, None, notes)


# ---------------------------------------------------------------------
# v17-P3-5｜network（ネットワーク）｜階層型 / depth_hierarchical / 3〜7
# ---------------------------------------------------------------------
def _network_depths(nodes, edges):
    """各ノードの深度を決定論的に算出する（原則③）。

    ルート（入次数0）を深度0とし、エッジをたどって深度を確定する。
    循環がある場合も無限ループしないよう訪問済みを管理する。
    """
    ids = [str(nd.get('id') or nd.get('label') or i) for i, nd in enumerate(nodes)]
    idx = {k: i for i, k in enumerate(ids)}
    indeg = dict((k, 0) for k in ids)
    adj = dict((k, []) for k in ids)
    for e in edges:
        a = str(e.get('from', ''))
        b = str(e.get('to', ''))
        if a in idx and b in idx:
            adj[a].append(b)
            indeg[b] += 1
    roots = [k for k in ids if indeg[k] == 0] or [ids[0]]
    depth = dict((k, None) for k in ids)
    queue = [(r, 0) for r in roots]
    while queue:
        k, d = queue.pop(0)
        if depth[k] is not None and depth[k] <= d:
            continue
        depth[k] = d
        for nb in adj[k]:
            queue.append((nb, d + 1))
    for k in ids:
        if depth[k] is None:
            depth[k] = 0          # 孤立ノードは深度0に置く（劣化描画｜原則①）
    return ids, idx, depth


def _draw_network_edge(slide, palette, x1, y1, x2, y2, width_px):
    """ノード間のエッジを描く。

    ⚠️ add_connector（MSO_CONNECTOR）を第一選択とし、非対応環境では
       極細矩形の回転で代替する（原則①④｜cycle の教訓を適用）。
    """
    try:
        from pptx.enum.shapes import MSO_CONNECTOR
        from pptx.util import Pt
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            px(x1), px(y1), px(x2), px(y2))
        conn.line.color.rgb = hex_to_rgb(palette['midtone'])
        conn.line.width = Pt(max(width_px * 0.75, 1))
        return conn
    except Exception:
        import math
        length = int(round(math.hypot(x2 - x1, y2 - y1))) or 1
        ang = math.degrees(math.atan2(y2 - y1, x2 - x1))
        mx = (x1 + x2) // 2 - length // 2
        my = (y1 + y2) // 2 - width_px // 2
        sh = _v17_rule(slide, mx, my, length, max(width_px, 2),
                       palette['midtone'])
        try:
            sh.rotation = ang
        except Exception:
            pass
        return sh


def draw_network(slide, palette, data):
    """ノード間の関係性を階層レイアウト（上→下）で描画する。

    Args:
        data : {'title': str,
                'nodes': [{'id': str, 'label': str,
                           'score': int|None}, ...],        # 3〜7
                'edges': [{'from': str, 'to': str}, ...]}
    Returns:
        dict : 描画レポート

    仕様: 集約表 P3-5 行（min3／max7／hierarchical_top_down／
          depth_hierarchical／原本shape=node_edge）
    ⚠️ 再現性担保のため力学配置（force-directed）は採らず、
       決定論的な階層レイアウトを採用（集約表 P3-5 行の根拠｜原則③）。
    ⚠️ フォールバックは3起動条件（統括指示｜8/12議題4）:
       (1) ノード数が範囲外
       (2) エッジ参照が不整合（存在しないノードIDを参照）
       (3) エッジが1本もない（node_edge の本質を満たさない）
    """
    spec = DIAGRAM_PATTERN_SPEC['network']
    notes = []
    nodes = _v17_normalize(data.get('nodes', []))
    edges = [e for e in (data.get('edges') or []) if isinstance(e, dict)]
    n = len(nodes)

    def _fallback(reason):
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': nd.get('label', ''), 'score': nd.get('score'),
                            'description': ''} for nd in nodes],
        }, _fallback_from='network')
        rep['notes'].append(reason)
        return rep

    # 起動条件(1)｜ノード数
    if not (spec['min_elements'] <= n <= spec['max_elements']):
        return _fallback(
            'ノード数 %d は network の範囲外（min %d／max %d）'
            '｜category へフォールバック'
            % (n, spec['min_elements'], spec['max_elements']))

    ids, idx, depth = _network_depths(nodes, edges)

    # 起動条件(2)｜エッジ参照の整合
    bad = [(str(e.get('from', '')), str(e.get('to', '')))
           for e in edges
           if str(e.get('from', '')) not in idx or str(e.get('to', '')) not in idx]
    if bad:
        return _fallback('エッジ参照が不整合（%s）｜category へフォールバック'
                         % '／'.join('%s→%s' % b for b in bad[:3]))

    # 起動条件(3)｜エッジ0本
    if not edges:
        return _fallback('エッジが0本｜node_edge の本質を満たさない'
                         '｜category へフォールバック')

    _v17_title(slide, data.get('title', ''), palette)

    top = V17_AREA['body_top']
    field_h = V17_AREA['body_bottom'] - top
    nw, nh = spec['node_w'], spec['node_h']

    # 深度ごとにグルーピング（決定論的）
    groups = {}
    for k in ids:
        groups.setdefault(depth[k], []).append(k)
    max_depth = max(groups.keys())
    row_h = field_h / float(max_depth + 1)

    center = {}
    for d in sorted(groups.keys()):
        row = groups[d]
        cnt = len(row)
        span = V17_AREA['width'] / float(cnt)
        y = int(top + row_h * d + (row_h - nh) / 2.0)
        for j, k in enumerate(row):
            x = int(V17_AREA['left'] + span * j + (span - nw) / 2.0)
            x = max(V17_AREA['left'], min(x, V17_AREA['right'] - nw))
            center[k] = (x, y)

    # エッジを先に描く（ノードの下に来るように）
    for e in edges:
        a, b = str(e.get('from')), str(e.get('to'))
        ax, ay = center[a]
        bx, by = center[b]
        _draw_network_edge(slide, palette,
                           ax + nw // 2, ay + nh, bx + nw // 2, by,
                           spec['edge_w'])

    drawn = 0
    for k in ids:
        nd = nodes[idx[k]]
        x, y = center[k]
        score = nd.get('score')
        base_key = V17_NETWORK_TIERS[min(depth[k], len(V17_NETWORK_TIERS) - 1)]
        fill_hex = _tier_fill(palette, score, base_key)
        _v17_card(slide, x, y, nw, nh, fill_hex, palette)

        txt_color = _text_color_on(fill_hex)
        pad = 10
        add_text(slide, x + pad, y + 8, nw - pad * 2,
                 str(nd.get('label', '')), 15, bold=True, color=txt_color,
                 height_px=22)
        if score is not None:
            add_text(slide, x + pad, y + 32, nw - pad * 2, '%s%%' % score, 16,
                     bold=True, color=txt_color, height_px=24)
        drawn += 1

    notes.append('ノード %d／エッジ %d｜深度 0〜%d の階層レイアウト（決定論的）'
                 % (n, len(edges), max_depth))
    notes.append('エッジ交差の回避は v3.5 範囲外（統括承認｜8/12議題4）')
    return _v17_report('network', drawn, None, notes)


# ---------------------------------------------------------------------
# v17-P3-6｜ディスパッチの拡張（P1 3種＋P2 3種＋P3 5種＝11種）
# ---------------------------------------------------------------------
def draw_pattern(slide, pattern_key, palette, data):
    """パターンキーで描画関数を振り分ける（11種に対応）。

    ⚠️ v3.5コアの実装対象は11パターン。原本 DIAGNOSIS_TO_PATTERN に
       対応診断カテゴリを持たない 'integration' は構造的に対象外であり、
       category へ退避して notes に理由を記録する（例外は送出しない）。
    """
    table = {
        # P1
        'category':   draw_category,
        'breakdown':  draw_breakdown,
        'comparison': draw_comparison,
        # P2
        'pyramid':    draw_pyramid,
        'sequence':   draw_sequence,
        'framework':  draw_framework,
        # P3
        'funnel':     draw_funnel,
        'cycle':      draw_cycle,
        'contrast':   draw_contrast,
        'timeline':   draw_timeline,
        'network':    draw_network,
    }
    fn = table.get(pattern_key)
    if fn is None:
        src = (data.get('categories') or data.get('items')
               or data.get('components') or data.get('levels')
               or data.get('steps') or data.get('cells')
               or data.get('stages') or data.get('phases')
               or data.get('sides') or data.get('milestones')
               or data.get('nodes') or [])
        rep = draw_category(slide, palette, {
            'title': data.get('title', ''),
            'categories': [{'label': str(x.get('label', '')) if isinstance(x, dict) else str(x),
                            'score': x.get('score') if isinstance(x, dict) else None,
                            'description': ''} for x in src],
        }, _fallback_from=pattern_key)
        if pattern_key == 'integration':
            rep['notes'].append(
                'パターン "integration" は原本 DIAGNOSIS_TO_PATTERN に対応診断'
                'カテゴリがないため v3.5 の実装対象外（構造的に対象外）')
        else:
            rep['notes'].append('パターン "%s" は未定義キー｜category へ退避'
                                % pattern_key)
        return rep
    return fn(slide, palette, data)


def _mso_oval():
    from pptx.enum.shapes import MSO_SHAPE
    return MSO_SHAPE.OVAL

# =====================================================================
# ▲▲▲ v17 P3 追加ブロック ここまで ▲▲▲
# =====================================================================
