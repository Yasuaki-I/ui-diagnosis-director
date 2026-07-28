# -*- coding: utf-8 -*-
"""
紺＆クリーン スライド作成 ─ python-pptx ビルダー
==================================================

GPTs の Knowledge にアップロードして使用する PPTX 生成ライブラリ。
ChatGPT の Code Interpreter から exec() で読み込み、関数を呼び出して .pptx を組み立てます。

使い方:
    exec(open('/mnt/data/03_pptx_builder.py').read())
    prs = create_presentation()
    add_cover(prs, title='Q2 販売戦略', date='2026年7月', author='○○株式会社 営業企画部')
    add_agenda(prs, items=['市場環境', '現状課題', '施策', '計画'])
    # ... 他のレイアウト関数 ...
    prs.save('output.pptx')

必達13条項（Instructions参照）を厳格に遵守。

バージョン履歴:
    v15 (2026-07-12): リスト表示の「折返し行間」と「項目間余白」を分離制御。
        - _add_multi_run_box / add_paragraph_box に space_after_pt /
          space_before_pt / vertical_anchor パラメータを追加。
        - リスト表示6箇所（C-1強み/課題、C-3 slide1 Top3/強み、C-3 slide2
          コンパクト強み/課題）の line_height を 1.6-1.7 → 1.2 に、
          space_after_pt=8pt を項目間に付与。「1文の連続感」と「項目間の分離感」
          を独立制御し、折返し時の可読性を根本改善。
        - C-3 slide1 総評帯：SUMMARY_H を 62→68 に拡大、
          vertical_anchor=MIDDLE で天地余白を均等化。
    v14 (2026-07-09): 総評 1行運用確定、17字上限、SUMMARY_H=62 に最適化。
    v13 (2026-07-08): 総評本文18字、SUMMARY_H=76、ISSUE_H=190。
    v12 (2026-07-07): 強み項目 40→28字厳格化（語途中折返し防止）。
    v11 (2026-07-06): 総評⇄Top3レイアウト崩れ根治（文字数上限厳格化）。
    v10 (2026-07-05): C-1/C-2/C-3 の120字統一・word_wrap=True で文字切れ根絶。
"""
__version__ = '15.0.0'
__version_date__ = '2026-07-12'

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from copy import deepcopy

# =====================================================================
# 配色パレット（02_design_spec.md 準拠・変更禁止）
# =====================================================================
NAVY        = RGBColor(0x1C, 0x36, 0x6C)   # メインナビー
RED         = RGBColor(0xD0, 0x02, 0x1B)   # アクセント赤
LIGHT_GRAY  = RGBColor(0xF4, 0xF5, 0xF8)   # カード背景
BORDER_GRAY = RGBColor(0xD0, 0xD4, 0xDC)   # 罫線
TEXT        = RGBColor(0x40, 0x40, 0x40)   # 本文
SUB_TEXT    = RGBColor(0x60, 0x60, 0x60)   # 注釈
NAVY_LIGHT  = RGBColor(0x9D, 0xB0, 0xD6)   # ナビー帯内補助
NAVY_E6     = RGBColor(0xE6, 0xEA, 0xF3)   # ナビー帯内本文
GOLD        = RGBColor(0xFF, 0xD5, 0x4F)   # 紺帯内強調
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BORDER = RGBColor(0xCC, 0xCC, 0xCC)   # フッター境界線
PAGE_NUM    = RGBColor(0x26, 0x26, 0x26)   # ページ番号
STRIPE      = RGBColor(0xF8, 0xF9, 0xFB)   # 縞模様セル
# C-3 / UI診断ボード用 アクセントカラー
ORANGE      = RGBColor(0xF9, 0x73, 0x16)   # アクセントオレンジ（After/方向性）
ORANGE_LIGHT = RGBColor(0xFE, 0xF1, 0xE3)  # オレンジ薄帯（After背景）
RED_LIGHT   = RGBColor(0xFD, 0xEC, 0xEC)   # 赤薄帯（Before背景）
PRIO_RED    = RGBColor(0xD0, 0x02, 0x1B)   # 優先度高
PRIO_ORANGE = RGBColor(0xE8, 0x8B, 0x1F)   # 優先度中
PRIO_GRAY   = RGBColor(0x88, 0x88, 0x88)   # 優先度低

# =====================================================================
# ▼▼▼ Phase A 追加ブロック（v16.5・2026-07-28） ▼▼▼
# 目的：機能強化の下ごしらえ（辞書のみ・利用ロジック未実装）
# 既存のC-1〜C-3描画ロジックには一切影響しない（後方互換完全維持）
# 詳細設計：phase_a_design_20260727_rev2.md 参照
# =====================================================================

# ============================================================
# A-3：UI診断ディレクターの10項目評価軸の背景知識
# 出典：入江さんご自身の知見の集約（web-director.skill）
# 併用プロジェクト：マーケティングオーケストレーター
# 対応マッピング：10項目中8項目が web-director.skill 6大機能領域と「高」対応
# 詳細：phase_a_design_20260727_rev2.md A-3セクション参照
# ============================================================

# ============================================================
# A-1：デジタル庁 公式カラーパレット
# 出典：デジタル庁 ダッシュボードデザインの実践ガイドブックとデザインテンプレート
# https://www.digital.go.jp/resources/dashboard-guidebook
# ライセンス：PDL1.0（公共データ利用規約 第1.0版）
# GitHub：https://github.com/digital-go-jp/policy-dashboard-assets
# 取得日：2026-07-27（Claude-Chat先行調査）
# 用途：v3.5時点は内部保持のみ・利用ロジック未実装
#       将来（Phase B以降）で10項目スコア色分けへ応用予定
# ============================================================
DIGITAL_AGENCY_PALETTE = {
    "SolidGray": {
        "primary":  "#4D4D4D",
        "secondary":"#767676",
        "midtone":  "#999999",
        "light":    "#CCCCCC",
        "lightest": "#F2F2F2",
        "accent":   "#3460FB",  # 強調用（青）
        "warning":  "#FE3939",  # 警告用（赤）
        "bg":       "#F8F8FB",
    },
    "Blue": {
        "primary":  "#0017C1",
        "secondary":"#3460FB",
        "midtone":  "#7096F8",
        "light":    "#C5D7FB",
        "lightest": "#E8F1FE",
        "accent":   "#FE3939",
        "warning":  "#FFBBBB",
        "bg":       "#F8F8FB",
    },
    "LightBlue": {
        "primary":  "#0055AD",
        "secondary":"#008BF2",
        "midtone":  "#57B8FF",
        "light":    "#C0E4FF",
        "lightest": "#F0F9FF",
        "accent":   "#FE3939",
        "warning":  "#FFBBBB",
        "bg":       "#F8F8FB",
    },
    "Green": {
        "primary":  "#115A36",
        "secondary":"#259D63",
        "midtone":  "#51B883",
        "light":    "#9BD4B5",
        "lightest": "#E6F5EC",
        "accent":   "#666666",
        "warning":  "#CCCCCC",
        "bg":       "#F8F8FB",
    },
    "Cyan": {
        "primary":  "#006F83",
        "secondary":"#00A3BF",
        "midtone":  "#2BC8E4",
        "light":    "#99F2FF",
        "lightest": "#E9F7F9",
        "accent":   "#666666",
        "warning":  "#CCCCCC",
        "bg":       "#F8F8FB",
    },
    "Red": {
        "primary":  "#CE0000",
        "secondary":"#FE3939",
        "midtone":  "#FF7171",
        "light":    "#FFBBBB",
        "lightest": "#FDEEEE",
        "accent":   "#666666",
        "warning":  "#CCCCCC",
        "bg":       "#F8F8FB",
    },
    "Orange": {
        "primary":  "#AC3E00",
        "secondary":"#FB5B01",
        "midtone":  "#FF8D44",
        "light":    "#FFC199",
        "lightest": "#FFEEE2",
        "accent":   "#666666",
        "warning":  "#CCCCCC",
        "bg":       "#F8F8FB",
    },
}

# 全テーマ共通の閾値色（good/bad判定用）
DIGITAL_AGENCY_THRESHOLD = {
    "center":  "#E6E6E6",  # 中央値・中立表示
    # maximum / minimum は各テーマの4番目色に準ずる（テーマ依存）
}

# ============================================================
# A-2：パーツ図鑑・図解集 由来 図解パターン辞書
# 出典：うちた氏「パーツ図鑑_120種」「図解集_50種」「テンプレ大全_100枚」
# 商用利用許諾：2026-07-26 note返信にて確認済み
# ライセンス記録：UCHITA_LICENSE_RECORD_20260726.md 参照
# 用途：v3.5時点は辞書のみ・描画実装なし（Phase B以降で実装）
# ============================================================
DIAGRAM_PATTERNS = {
    'category':     {'ja': '分類',           'use': '要素を並列カテゴリで整理',      'shape': 'grid'},
    'pyramid':      {'ja': 'ピラミッド',     'use': '階層・優先順位を上下で表現',    'shape': 'triangle'},
    'comparison':   {'ja': '比較',           'use': '2〜3要素の対比',                'shape': 'side_by_side'},
    'sequence':     {'ja': '順序',           'use': 'ステップ・時系列を左→右',       'shape': 'arrow_chain'},
    'cycle':        {'ja': '循環',           'use': '反復プロセスを円環で表現',      'shape': 'circle_arrow'},
    'funnel':       {'ja': '絞り込み',       'use': '上から下へ絞り込むファネル型',  'shape': 'trapezoid'},
    'timeline':     {'ja': '時間軸',         'use': '期間別のマイルストーン',        'shape': 'horizontal_bar'},
    'breakdown':    {'ja': '分解',           'use': '全体を構成要素に分解',          'shape': 'tree'},
    'contrast':     {'ja': '対比',           'use': '対照的な2要素の並列強調',       'shape': 'split_screen'},
    'integration':  {'ja': '統合',           'use': '複数要素の統合結果',            'shape': 'merge'},
    'framework':    {'ja': 'フレームワーク', 'use': '4象限マトリクス等',             'shape': 'quadrant'},
    'network':      {'ja': 'ネットワーク',   'use': 'ノード間の関係性',              'shape': 'node_edge'},
}

# 診断結果 → 推奨図解パターンのマッピング
DIAGNOSIS_TO_PATTERN = {
    'proposal_categorization': 'category',
    'priority_ranking':        'pyramid',
    'before_after':            'comparison',
    'user_flow':               'sequence',
    'improvement_cycle':       'cycle',
    'conversion_funnel':       'funnel',
    'schedule':                'timeline',
    'score_breakdown':         'breakdown',
    'ux_contrast':             'contrast',
    'impact_cost_matrix':      'framework',
    'site_structure':          'network',
}

# =====================================================================
# ▲▲▲ Phase A 追加ブロック ここまで ▲▲▲
# =====================================================================

# キャンバス（1280×720 想定）
CANVAS_W_PX = 1280
CANVAS_H_PX = 720

# 共通フォント
FONT = 'メイリオ'

# =====================================================================
# 単位変換・ユーティリティ
# =====================================================================
def px(n):
    """1280×720 想定の px を EMU に変換 (1 px = 9525 EMU)"""
    return Emu(int(n * 9525))


def set_run(run, text, size_pt, bold=False, color=TEXT, italic=False):
    """run にテキスト・フォント設定を適用（メイリオ・14pt以上を強制）

    フォント設定は run の OOXML rPr 内の latin / ea / cs すべてに
    明示的にメイリオを書き込み、英数字も含めて完全に統一する。
    """
    if size_pt < 14:
        raise ValueError(
            f"font_size must be >= 14pt (got {size_pt}pt) ─ 条項2違反"
        )
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = False  # 条項13: 斜体強調禁止
    run.font.color.rgb = color
    _force_all_script_fonts(run, FONT)


def _force_all_script_fonts(run, font_name):
    """run の latin / eastAsia / complexScript すべてにフォントを強制設定。

    PowerPoint OOXML は run のフォントを latin（英数字）/ ea（日本語等）/
    cs（複雑表記）の3スクリプトで個別に持つ。run.font.name は latin の
    みを設定するため、日本語文字＋数字混在のテキストではフォントが
    分かれて見えることがある。このヘルパーで3スクリプトすべて明示的に
    同じフォントへ統一する。
    """
    rPr = run._r.get_or_add_rPr()
    for tag in ('a:latin', 'a:ea', 'a:cs'):
        # 既存タグを削除してから追加（重複防止）
        for el in rPr.findall(qn(tag)):
            rPr.remove(el)
        el = rPr.makeelement(qn(tag), {'typeface': font_name})
        rPr.append(el)


# 後方互換のエイリアス
_force_east_asian_font = _force_all_script_fonts


# =====================================================================
# 文字数バリデーション（design_system.md §1 準拠・絶対遵守）
# =====================================================================
# UI診断ディレクター用の文字数規定。
# スライドはコンテナサイズ固定のため、テキスト量を物理的に制御しないと
# 必ずレイアウトが崩壊する。規定文字数を超えた場合は ValueError で停止し、
# GPTs 側で「末尾…省略ではなく、規定文字内で要約し直す」運用とする。
# 詳細は gpts-package/design_system.md §1 を参照。

LIMITS = {
    # ─── C-1 スコアカード ─────────────────────
    'service_name':       30,
    'c1_comment':         40,   # 一言所見（10項目×1コメント）
    'c1_strength':        60,   # 強み（3項目）
    'c1_issue':           60,   # 最優先課題（3項目）
    'c1_conclusion':     120,   # 結論・総評
    # ─── C-2 改善提案リスト ──────────────────
    'c2_title':           30,   # 提案タイトル
    'c2_point':          150,   # POINT本文
    'c2_priority':         6,   # 優先度バッジ
    'c2_category':        12,   # カテゴリ
    # ─── C-3 ビジュアル診断ボード ────────────
    'c3_summary':         80,   # 総評（スライド1中央）
    'c3_top_issue':       40,   # 最優先課題タイトル
    'c3_direction':      120,   # 改善方向帯（2行折り返し許容）
    'c3_section_label':   20,   # LP構造マップ各セクション
    'c3_flow_step':       25,   # 行動フローステップ
    'c3_highlight_title': 30,   # Before/Afterタイトル
    'c3_before':          40,   # Before本文
    'c3_after':           40,   # After本文
}


def validate_length(text, limit_key, field_label=None, *, allow_none=True):
    """規定文字数を超えた場合 ValueError で停止する。

    Parameters
    ----------
    text : str or None
        検査対象テキスト。None / 空文字は許容（allow_none=False で禁止）。
    limit_key : str
        LIMITS 辞書のキー。例: 'c1_comment', 'c2_point', 'c3_summary'。
    field_label : str, optional
        エラーメッセージに表示する人間向けの項目名。
        未指定なら limit_key をそのまま使う。
    allow_none : bool
        True なら None / 空文字を素通り。False なら必須扱いで ValueError。

    Returns
    -------
    str
        検査済みテキスト（None なら空文字に正規化）。

    Raises
    ------
    KeyError
        limit_key が LIMITS に存在しない場合（開発者ミス）。
    ValueError
        文字数超過、または allow_none=False で空入力の場合。
    """
    if limit_key not in LIMITS:
        raise KeyError(
            f"validate_length: 未定義の limit_key '{limit_key}'。"
            f" LIMITS 辞書に追加してください。"
        )
    max_len = LIMITS[limit_key]
    label = field_label or limit_key

    if text is None or text == '':
        if allow_none:
            return ''
        raise ValueError(
            f"\n[文字数バリデーションエラー]\n"
            f"  項目: {label}\n"
            f"  問題: 必須項目が空欄です。\n"
            f"→ GPTs側でテキストを生成してください。"
        )

    text_str = str(text)
    n = len(text_str)
    if n > max_len:
        over = n - max_len
        preview = text_str[:30] + ('...' if n > 30 else '')
        raise ValueError(
            f"\n[文字数バリデーションエラー]\n"
            f"  項目: {label}\n"
            f"  上限: {max_len}文字（design_system.md §1 準拠）\n"
            f"  実際: {n}文字（{over}文字オーバー）\n"
            f"  内容: 「{preview}」\n"
            f"→ GPTs側で規定文字数内に要約し直してください。\n"
            f"  ※末尾「…」での省略は禁止です。要点を絞り込んでください。"
        )
    return text_str


def add_text(slide, left_px, top_px, width_px, text, size_pt, *,
             bold=False, color=TEXT, align=PP_ALIGN.LEFT,
             height_px=None, line_height=None, letter_spacing=None,
             vertical_anchor=None):
    """1行のテキストボックスを追加（内容追従、高さ自動・条項5準拠）

    vertical_anchor: MSO_ANCHOR.TOP / MIDDLE / BOTTOM（既定None=PPTXデフォルト）。
                     帯内の縦中央配置などで使用。
    """
    if height_px is None:
        height_px = max(int(size_pt * 1.6), 24)
    box = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.auto_size = MSO_AUTO_SIZE.NONE  # 高さ調整は word_wrap で
    if vertical_anchor is not None:
        tf.vertical_anchor = vertical_anchor
    para = tf.paragraphs[0]
    para.alignment = align
    if line_height:
        para.line_spacing = line_height
    run = para.add_run()
    set_run(run, text, size_pt, bold=bold, color=color)
    return box


def add_paragraph_box(slide, left_px, top_px, width_px, paragraphs, *,
                      height_px=None, default_size=16, default_color=TEXT,
                      line_height=1.6,
                      space_after_pt=None, space_before_pt=None):
    """
    複数行（段落 or 箇条書き）を1つのテキストボックスに集約（条項6準拠）

    paragraphs: list of dict or str
      dict: {'text': str, 'size': int, 'bold': bool, 'color': RGBColor,
             'align': PP_ALIGN, 'bullet': bool,
             'space_after_pt': float, 'space_before_pt': float}
      str:  既定スタイルで追加

    line_height : float
        段落**内部**の折返し行間（既定 1.6）。
    space_after_pt : float, optional
        [v15/2026-07-12] 全段落に適用する段落末尾余白（pt）。
        リスト表示で項目間の分離感を演出する用途。
    space_before_pt : float, optional
        全段落に適用する段落先頭余白（pt）。通常は未使用。
    """
    if height_px is None:
        height_px = 400  # 大きめに確保（内容で自動縮小される表示）
    box = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for i, p in enumerate(paragraphs):
        if isinstance(p, str):
            p = {'text': p}
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        para.line_spacing = p.get('line_height', line_height)
        # [v15] 段落間余白の分離制御
        sa = p.get('space_after_pt', space_after_pt)
        if sa is not None and sa != 0:
            para.space_after = Pt(sa)
        sb = p.get('space_before_pt', space_before_pt)
        if sb is not None and sb != 0:
            para.space_before = Pt(sb)
        run = para.add_run()
        set_run(
            run,
            p.get('text', ''),
            p.get('size', default_size),
            bold=p.get('bold', False),
            color=p.get('color', default_color),
        )
    return box


def add_shape(slide, shape_type, left_px, top_px, width_px, height_px, *,
              fill=None, line=None, line_width_pt=None, fill_alpha=None):
    """図形を追加（条項7: 図形には文字を入れない）

    Parameters
    ----------
    fill_alpha : int or None
        塗りつぶしの不透明度（0-100000、None=不透明）。
        例: 95000=95%不透明=5%透過、80000=80%不透明=20%透過。
        design_system.md §3.2 の透明度バリエーション実装に使用。
    """
    shape = slide.shapes.add_shape(
        shape_type, px(left_px), px(top_px), px(width_px), px(height_px)
    )
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        # 透明度指定がある場合は alpha 属性を OOXML に直接書き込む
        if fill_alpha is not None:
            _set_shape_fill_alpha(shape, fill_alpha)
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    # 図形内のテキストは空に
    if shape.has_text_frame:
        shape.text_frame.text = ''
    shape.shadow.inherit = False
    return shape


def _set_shape_fill_alpha(shape, alpha):
    """図形の塗りつぶし色に透明度（alpha）を設定する。

    python-pptx は標準で fill の alpha 属性をサポートしないため、
    OOXML を直接操作して <a:srgbClr><a:alpha val=".."/></a:srgbClr>
    を組み立てる。

    Parameters
    ----------
    alpha : int
        不透明度（0=完全透過、100000=完全不透明）。
        例: 95000 = 95%不透明 = 5%透過
            80000 = 80%不透明 = 20%透過
            12000 = 12%不透明 = 88%透過（薄帯背景）
    """
    spPr = shape.fill._xPr.find(qn('a:solidFill')) if False else None
    # 上記は型ヒント目的。実体は以下：
    solidFill = shape.fill._xPr.find(qn('a:solidFill'))
    if solidFill is None:
        return  # solid 塗りでない場合はスキップ
    srgbClr = solidFill.find(qn('a:srgbClr'))
    if srgbClr is None:
        return
    # 既存の alpha タグを削除してから追加（重複防止）
    for el in srgbClr.findall(qn('a:alpha')):
        srgbClr.remove(el)
    alpha_el = srgbClr.makeelement(qn('a:alpha'), {'val': str(int(alpha))})
    srgbClr.append(alpha_el)


def _add_bg_frame(slide, left_px, top_px, width_px, height_px, *,
                  fill=None, fill_alpha=None, line=None, line_width_pt=None,
                  radius_px=8):
    """カード背景フレーム（角丸 + 任意で透明度フィル + 任意で罫線）。

    design_system.md §4.3 カード規定準拠：
    - 角丸標準 8px（radius_px で変更可）
    - フラットデザイン（影なし）
    - 罫線 1pt / BORDER_GRAY を推奨

    ビジュアル強化（フェーズ1.5）で典型的に使うパターン：

        # NAVY 5%透過カード薄塗り
        _add_bg_frame(slide, 40, 90, 1200, 200,
                      fill=NAVY, fill_alpha=5000,
                      line=BORDER_GRAY, line_width_pt=1)

        # NAVY 18%透過カード強塗り
        _add_bg_frame(slide, 40, 90, 1200, 200,
                      fill=NAVY, fill_alpha=18000)

        # 不透明白カード（標準）
        _add_bg_frame(slide, 40, 90, 1200, 200,
                      fill=WHITE, line=BORDER_GRAY, line_width_pt=1)
    """
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      left_px, top_px, width_px, height_px,
                      fill=fill, fill_alpha=fill_alpha,
                      line=line, line_width_pt=line_width_pt)
    # 角丸半径を radius_px に合わせる（PowerPoint は短辺比で 0-0.5）
    short_side = min(width_px, height_px)
    if short_side > 0:
        ratio = max(0.0, min(0.5, radius_px / short_side))
        try:
            shape.adjustments[0] = ratio
        except (IndexError, AttributeError):
            pass
    return shape


# =====================================================================
# 共通要素：ヘッダ・フッター
# =====================================================================
def _add_header(slide, title, sub_label_en=''):
    """ヘッダ帯（高さ60px）を追加"""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W_PX, 60, fill=NAVY)
    add_text(slide, 40, 18, 800, title, 22, bold=True, color=WHITE)
    if sub_label_en:
        # サブラベルは右寄せ。長い英文（'Improvement Proposals' 等）が
        # 折り返さないよう、幅を400pxに広げる（左端= CANVAS_W_PX - 40 - 400）
        add_text(slide, CANVAS_W_PX - 40 - 400, 22, 400,
                 sub_label_en, 14,
                 color=NAVY_LIGHT, align=PP_ALIGN.RIGHT,
                 height_px=24)


def _add_footer(slide, page_num, total, author='紺＆クリーン スライド作成'):
    """フッター帯（ロゴ＋ページ番号のみ・条項8準拠）"""
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 660, CANVAS_W_PX, 1, fill=GRAY_BORDER)
    add_text(slide, 40, 682, 400, author, 14, bold=True, color=NAVY)
    add_text(slide, 1140, 685, 120, f'{page_num} / {total}', 14,
             color=PAGE_NUM, align=PP_ALIGN.RIGHT)


def _add_lead(slide, text, top_px=90):
    """リード文（ヘッダ直下の導入文）"""
    add_text(slide, 40, top_px, 1200, text, 18, color=TEXT,
             height_px=60, line_height=1.7)


def _add_conclusion_band(slide, top_px, body_text, label='CONCLUSION',
                          color_label=NAVY_LIGHT, color_body=WHITE,
                          height_px=92, bg=NAVY, body_size=18):
    """結論帯（本文末尾の主張帯）

    本文の下端切れを防ぐため：
      - 帯の高さを十分（既定92px）に確保
      - 本文 textbox は ラベル分(32px)を引いた残量を高さに割当
      - body_text に '\\n' が含まれる場合は複数行段落として配置
    """
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 40, top_px, 1200, height_px,
              fill=bg)
    # ラベル（上部）
    add_text(slide, 60, top_px + 10, 1160, label, 14, bold=True,
             color=color_label, letter_spacing=2, height_px=20)
    # 本文（複数行対応）
    body_top = top_px + 34
    body_h = max(int(body_size * 1.5), height_px - 38)
    # body_text に改行が含まれていれば複数段落、そうでなければ単一行
    lines = [ln for ln in body_text.split('\n') if ln]
    if not lines:
        lines = [body_text]
    paragraphs = [{'text': line, 'size': body_size, 'bold': True,
                    'color': color_body, 'line_height': 1.35}
                   for line in lines]
    add_paragraph_box(slide, 60, body_top, 1160, paragraphs,
                       height_px=body_h, line_height=1.35)


# =====================================================================
# Presentation 初期化
# =====================================================================
def create_presentation():
    """1280×720 px キャンバスの空 Presentation を生成"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _blank_slide(prs):
    """白紙レイアウトでスライドを追加"""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


# =====================================================================
# Layout 1: 表紙（add_cover）
# =====================================================================
def add_cover(prs, title, date, author, subtitle='', page_total=1):
    """
    表紙スライド。

    Args:
        title: メインタイトル
        date: 日付（"2026年7月15日" など）
        author: 著者表記（"○○株式会社 営業企画部"）
        subtitle: サブタイトル（省略可）
    """
    slide = _blank_slide(prs)

    # ナビー背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W_PX, 660, fill=NAVY)

    # プロジェクトラベル
    add_text(slide, 80, 160, 900, 'PRESENTATION', 16, bold=True,
             color=NAVY_LIGHT, letter_spacing=3)

    # メインタイトル（複数行対応）
    title_lines = title.split('\n')
    title_h = 70 * len(title_lines)
    add_paragraph_box(slide, 80, 200, 1120,
        [{'text': line, 'size': 40, 'bold': True, 'color': WHITE,
          'line_height': 1.35} for line in title_lines],
        height_px=title_h)

    # 赤ライン
    add_shape(slide, MSO_SHAPE.RECTANGLE, 80, 200 + title_h + 20, 480, 3,
              fill=RED)

    # サブタイトル
    if subtitle:
        add_text(slide, 80, 200 + title_h + 40, 1120, subtitle, 20,
                 color=NAVY_E6, line_height=1.6, height_px=60)

    # 日付＋著者（右下右揃え・条項9・幅を広く）
    add_text(slide, 480, 560, 720, f'{date}　{author}', 20, bold=True,
             color=WHITE, align=PP_ALIGN.RIGHT, height_px=70, line_height=1.4)

    # 白フッター背景（表紙はページ番号を省略）
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 661, CANVAS_W_PX, 59, fill=WHITE)
    add_text(slide, 40, 682, 400, '紺＆クリーン スライド作成', 14, bold=True,
             color=NAVY)

    return slide


# =====================================================================
# Layout 2: アジェンダ（add_agenda）
# =====================================================================
def add_agenda(prs, items, page_num=2, total=10,
               title='Agenda', lead='本資料の構成は以下のとおりです。'):
    """
    アジェンダスライド。

    Args:
        items: list of str or dict
            str: 章タイトルのみ
            dict: {'title': str, 'desc': str}
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Agenda')
    _add_lead(slide, lead)

    # 項目を縦に並べる（最大8項目まで）
    n = len(items)
    top_start = 170
    available_h = 460
    item_h = max(60, min(80, available_h // max(n, 1)))
    num_size = 32 if n >= 7 else 36
    title_size = 20 if n >= 7 else 22

    for i, item in enumerate(items):
        if isinstance(item, str):
            item = {'title': item, 'desc': ''}
        y = top_start + i * item_h

        # 番号（紺・大きく）
        add_text(slide, 40, y, 80, f'{i+1:02d}', num_size, bold=True,
                 color=NAVY, line_height=1.0, height_px=num_size + 4)

        # タイトル
        add_text(slide, 130, y + 4, 600, item['title'], title_size, bold=True,
                 color=NAVY, height_px=title_size + 8)

        # 説明（任意、タイトルの右側に配置して高さを節約）
        if item.get('desc'):
            add_text(slide, 740, y + 10, 500, item['desc'], 14,
                     color=SUB_TEXT, line_height=1.5, height_px=20)

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 3: 課題整理（add_issue_summary）
# =====================================================================
def add_issue_summary(prs, title, cards, page_num=3, total=10,
                      lead='現状の主な課題を整理します。',
                      conclusion=None):
    """
    課題整理スライド。3カード横並び。

    Args:
        cards: list of dict, 各 dict は {'no': str, 'heading': str, 'body': str}
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Issues')
    _add_lead(slide, lead)

    n = len(cards)
    card_w = (1200 - 15 * (n - 1)) // n  # カード間 15px
    card_top = 160
    card_h = 390 if conclusion else 460

    # ヘッダ帯（紺）の高さ
    HEADER_BAND_H = 88

    for i, card in enumerate(cards):
        left = 40 + i * (card_w + 15)
        # ---- カード背景（白＋薄ボーダー） ----
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  left, card_top, card_w, card_h,
                  fill=WHITE, line=BORDER_GRAY, line_width_pt=1)

        # ---- 上端ヘッダ帯（紺ベタ・大型ナンバー） ----
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  left + 1, card_top + 1,
                  card_w - 2, HEADER_BAND_H, fill=NAVY)
        # ナンバー（左端、大型）
        no_text = card.get('no', f'{i + 1:02d}')
        add_text(slide, left + 22, card_top + 14, 90, no_text,
                 40, bold=True, color=WHITE, line_height=1.0,
                 height_px=52)
        # 縦罫線（ナンバーと見出しの区切り）
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  left + 110, card_top + 22, 1, HEADER_BAND_H - 44,
                  fill=NAVY_LIGHT)
        # 見出し（ヘッダ帯内、白文字）
        add_text(slide, left + 124, card_top + 20,
                 card_w - 140, card.get('heading', ''),
                 18, bold=True, color=WHITE,
                 height_px=HEADER_BAND_H - 32, line_height=1.35)

        # ---- 本文領域（カード白地、縦中央寄せ） ----
        body_top = card_top + HEADER_BAND_H + 18
        body_bottom = card_top + card_h - 18
        body_h = body_bottom - body_top
        body = card.get('body', '')
        body_paragraphs = [
            {'text': '● ' + line.strip(), 'size': 15,
             'color': TEXT, 'line_height': 1.6}
            for line in body.split('\n') if line.strip()
        ]
        add_paragraph_box(slide, left + 22, body_top, card_w - 40,
                           body_paragraphs, height_px=body_h,
                           line_height=1.6)

    if conclusion:
        _add_conclusion_band(slide, 560, conclusion, label='POINT')

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 4: 優先度マトリクス（add_priority_matrix）
# =====================================================================
def add_priority_matrix(prs, title, items, page_num=4, total=10,
                        x_label_low='低', x_label_high='高',
                        y_label_low='低', y_label_high='高',
                        x_axis_name='重要度', y_axis_name='緊急度',
                        lead='4象限で施策を整理します。'):
    """
    優先度マトリクス。

    Args:
        items: list of dict, 各 {'label': str, 'x': float, 'y': float, 'priority': 'S'|'A'|'B'|'C'}
               x, y は 0.0〜1.0 で象限内の相対位置
               priority: S/A=赤, B/C=紺
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Priority Matrix')
    _add_lead(slide, lead)

    # マトリクス領域 (left=220, top=190, w=880, h=420)
    # リード文(y=90〜150)・Y軸名(mx_t-36)が重ならないよう mx_t を下にずらす
    mx_l, mx_t, mx_w, mx_h = 220, 190, 880, 420

    # 4象限の背景（右上が最優先）
    qw, qh = mx_w // 2, mx_h // 2
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l + qw, mx_t, qw, qh,
              fill=RGBColor(0xFF, 0xE8, 0xE8))  # 右上 = 最優先（薄赤）
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l, mx_t, qw, qh,
              fill=LIGHT_GRAY)  # 左上
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l, mx_t + qh, qw, qh,
              fill=STRIPE)  # 左下
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l + qw, mx_t + qh, qw, qh,
              fill=LIGHT_GRAY)  # 右下

    # 枠線
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l, mx_t, mx_w, mx_h,
              fill=None, line=BORDER_GRAY, line_width_pt=1)
    # 十字線
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l + qw, mx_t, 1, mx_h,
              fill=BORDER_GRAY)
    add_shape(slide, MSO_SHAPE.RECTANGLE, mx_l, mx_t + qh, mx_w, 1,
              fill=BORDER_GRAY)

    # 軸ラベル（高/低）
    add_text(slide, mx_l - 60, mx_t - 8, 60, y_label_high, 14, bold=True,
             color=NAVY, align=PP_ALIGN.RIGHT, height_px=22)
    add_text(slide, mx_l - 60, mx_t + mx_h - 20, 60, y_label_low, 14,
             bold=True, color=NAVY, align=PP_ALIGN.RIGHT, height_px=22)
    add_text(slide, mx_l - 8, mx_t + mx_h + 8, 60, x_label_low, 14,
             bold=True, color=NAVY, height_px=22)
    add_text(slide, mx_l + mx_w - 30, mx_t + mx_h + 8, 60, x_label_high, 14,
             bold=True, color=NAVY, align=PP_ALIGN.RIGHT, height_px=22)

    # 軸の名称（X軸名は下に、十分な幅でセンタリング）
    x_axis_text = f'{x_axis_name} →'
    add_text(slide, mx_l, mx_t + mx_h + 36, mx_w,
             x_axis_text, 14, color=SUB_TEXT,
             align=PP_ALIGN.CENTER, height_px=22)
    # Y軸名：マトリクス左外、横長領域に十分な幅で表示（折返し防止）
    # 14pt全角約19px、全角6文字「↑ インパクト」→ 114px必要 → 幅180pxで余裕
    add_text(slide, mx_l - 200, mx_t + mx_h // 2 - 12, 190,
             f'↑ {y_axis_name}', 14, color=SUB_TEXT,
             align=PP_ALIGN.RIGHT, height_px=22)

    # プロット
    # ラベル領域の最大幅を計算（マトリクス右端を超えないよう動的調整）
    for i, it in enumerate(items):
        cx = mx_l + int(it['x'] * mx_w)
        cy = mx_t + int((1 - it['y']) * mx_h)
        d = 32
        priority = it.get('priority', 'B')
        fill = RED if priority in ('S', 'A') else NAVY
        add_shape(slide, MSO_SHAPE.OVAL, cx - d // 2, cy - d // 2, d, d,
                  fill=fill)
        # 番号（円の中・白）
        add_text(slide, cx - d // 2, cy - 14, d, f'{i+1:02d}', 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, height_px=22)
        # ラベル位置の自動判定：
        #  - 円の右側に十分な余白があれば右配置
        #  - 右余白が少ない（マトリクス右端まで180px未満）なら左配置
        label_w = 200
        label_left_default = cx + d // 2 + 6
        # マトリクス右端 mx_l + mx_w を超えないか判定
        if label_left_default + label_w > mx_l + mx_w:
            # 円の左側にラベル配置（右寄せ）
            label_left = cx - d // 2 - 6 - label_w
            label_align = PP_ALIGN.RIGHT
        else:
            label_left = label_left_default
            label_align = PP_ALIGN.LEFT
        add_text(slide, label_left, cy - 12, label_w, it['label'], 14,
                 bold=True, color=TEXT, align=label_align, height_px=22)

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 5: OK・NG例の対比（add_ok_ng_pair）
# =====================================================================
def add_ok_ng_pair(prs, title, ng, ok, page_num=5, total=10,
                   lead='良い例と悪い例を対比します。',
                   conclusion=None):
    """
    OK/NG対比カード（スタイリッシュ版）。

    レイアウト：
      ┌─────────────────────┐
      │ ▌NG ✕  見出し       │ ← 上端帯（色帯）
      │ ─────────────────── │
      │ 本文テキスト         │ ← 中央領域
      │                     │
      │ ─────────────────── │
      │ ⚠ キャプション帯    │ ← 下端帯（淡色）
      └─────────────────────┘

    Args:
        ng: dict {'heading': str, 'body': str, 'caption': str}
        ok: dict {'heading': str, 'body': str, 'caption': str}
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'OK / NG')
    _add_lead(slide, lead)

    col_w = 580
    col_h = 380 if conclusion else 460
    card_top = 160

    # 視覚的に間延びを防ぐため、上端ヘッダ帯／中央本文／下端キャプション帯を
    # それぞれ独立した「帯」として配置する
    HEADER_BAND_H = 56   # 上端帯（バッジ+見出し）
    CAPTION_BAND_H = 44  # 下端帯（一言サマリ）

    for i, (col, badge_text, badge_color, label_color, accent_bg) in enumerate(
        [(ng, 'NG ✕', RGBColor(0xDC, 0x35, 0x45), RED,
          RGBColor(0xFD, 0xEC, 0xEC)),                  # NG淡赤帯
         (ok, 'OK ✓', RGBColor(0x10, 0x88, 0x4A),
          RGBColor(0x10, 0x88, 0x4A),
          RGBColor(0xE6, 0xF4, 0xEA))]                  # OK淡緑帯
    ):
        left = 40 + i * (col_w + 40)
        # ---- カード背景（薄いボーダー） ----
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top,
                  col_w, col_h, fill=WHITE,
                  line=BORDER_GRAY, line_width_pt=1)

        # ---- 上端ヘッダ帯（色付き） ----
        add_shape(slide, MSO_SHAPE.RECTANGLE, left + 1, card_top + 1,
                  col_w - 2, HEADER_BAND_H,
                  fill=accent_bg, line=accent_bg)
        # 色帯（左4px、強アクセント）
        add_shape(slide, MSO_SHAPE.RECTANGLE, left, card_top,
                  6, HEADER_BAND_H + 1, fill=badge_color)
        # バッジ
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  left + 20, card_top + 12, 96, 32,
                  fill=badge_color)
        add_text(slide, left + 20, card_top + 16, 96, badge_text, 16,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=24)
        # 見出し（バッジの右）
        add_text(slide, left + 130, card_top + 16, col_w - 150,
                 col.get('heading', ''), 18, bold=True,
                 color=label_color, height_px=28)

        # ---- 中央本文領域（縦中央寄せ） ----
        body_top = card_top + HEADER_BAND_H + 16
        body_bottom = card_top + col_h - CAPTION_BAND_H - 10
        body_h = body_bottom - body_top
        body = col.get('body', '')
        body_paragraphs = [{'text': line, 'size': 16, 'color': TEXT,
                            'line_height': 1.6}
                            for line in body.split('\n') if line.strip()]
        if body_paragraphs:
            add_paragraph_box(slide, left + 24, body_top, col_w - 48,
                               body_paragraphs, height_px=body_h,
                               line_height=1.6)

        # ---- 下端キャプション帯（淡色背景＋一言サマリ） ----
        if col.get('caption'):
            caption_top = card_top + col_h - CAPTION_BAND_H
            # 帯背景
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      left + 1, caption_top, col_w - 2, CAPTION_BAND_H - 1,
                      fill=accent_bg, line=accent_bg)
            # 左マーカー（▶）
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      left + 1, caption_top, 4, CAPTION_BAND_H - 1,
                      fill=badge_color)
            add_text(slide, left + 20, caption_top + 11,
                     col_w - 40, col['caption'], 14, bold=True,
                     color=label_color, height_px=22)

    if conclusion:
        _add_conclusion_band(slide, 558, conclusion, label='POINT')

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 6: 施策一覧表（add_action_table）
# =====================================================================
def add_action_table(prs, title, columns, rows, page_num=6, total=10,
                     lead='主要な施策を一覧化します。',
                     source=None, emphasize=None,
                     col_widths=None):
    """
    表組スライド（条項4準拠）。

    Args:
        columns: list of str（ヘッダ）
        rows: list of list（各行のセル値）
        emphasize: list of (row_idx, col_idx, original_text, emphasis_text)
                   セル内で original_text を emphasis_text に置換し、赤で大きく表示
        col_widths: list of int（各列の幅%。合計100。Noneなら均等）
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Action List')
    _add_lead(slide, lead)

    n_cols = len(columns)
    n_rows = len(rows) + 1  # ヘッダ含む

    table_left, table_top = 40, 160
    table_w = 1200
    row_h = 38

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, px(table_left), px(table_top),
        px(table_w), px(row_h * n_rows)
    )
    table = table_shape.table

    # 列幅
    if col_widths is None:
        col_widths = [100 // n_cols] * n_cols
    for i, w in enumerate(col_widths):
        table.columns[i].width = px(int(table_w * w / 100))

    # ヘッダ行
    for c, col_name in enumerate(columns):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ''
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        set_run(run, col_name, 16, bold=True, color=WHITE)
        cell.margin_left = Pt(8)
        cell.margin_right = Pt(8)
        cell.margin_top = Pt(6)
        cell.margin_bottom = Pt(6)

    # データ行
    for r, row in enumerate(rows, start=1):
        bg = WHITE if r % 2 == 1 else STRIPE
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            cell.text = ''
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c > 0 else PP_ALIGN.CENTER

            # 強調指定があれば run を分割
            emphasis_for_cell = None
            if emphasize:
                for em in emphasize:
                    if em[0] == r - 1 and em[1] == c:
                        emphasis_for_cell = em
                        break

            if emphasis_for_cell:
                _, _, original, target = emphasis_for_cell
                parts = str(val).split(target)
                for k, part in enumerate(parts):
                    if part:
                        run = para.add_run()
                        set_run(run, part, 14, color=TEXT)
                    if k < len(parts) - 1:
                        emp_run = para.add_run()
                        set_run(emp_run, target, 18, bold=True, color=RED)
            else:
                run = para.add_run()
                set_run(run, str(val), 14, color=TEXT)

            cell.margin_left = Pt(8)
            cell.margin_right = Pt(8)
            cell.margin_top = Pt(6)
            cell.margin_bottom = Pt(6)

    # 表終端の y を実測（条項11準拠）
    table_end_y = table_top + row_h * n_rows

    if source:
        # 出所注記（フッター帯の上、表終端 + 30px）
        add_text(slide, 40, table_end_y + 30, 1200, source, 14,
                 color=SUB_TEXT, line_height=1.5)

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 7: KPI カード（add_kpi_card）
# =====================================================================
def add_kpi_card(prs, title, kpis, page_num=7, total=10,
                 lead='主要KPIの進捗を確認します。',
                 conclusion=None):
    """
    KPIカード（3〜4個横並び）。

    Args:
        kpis: list of dict
              {'label': str, 'value': str, 'unit': str, 'desc': str, 'color': 'red'|'navy'}
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Key Metrics')
    _add_lead(slide, lead)

    n = len(kpis)
    gap = 20
    card_w = (1200 - gap * (n - 1)) // n
    card_h = 280
    top = 160

    for i, kpi in enumerate(kpis):
        left = 40 + i * (card_w + gap)
        is_red = kpi.get('color', 'red') == 'red'
        value_color = RED if is_red else NAVY

        # カード（白＋紺枠）
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  left, top, card_w, card_h,
                  fill=WHITE, line=NAVY, line_width_pt=2)
        # ラベル（上部）
        add_text(slide, left + 24, top + 26, card_w - 48,
                 kpi.get('label', ''), 14, bold=True, color=NAVY,
                 letter_spacing=2)
        # 大きな数値（条項13: サイズ拡大＋色変更）
        value_size = 48 if len(kpi.get('value', '')) <= 6 else 36
        add_text(slide, left + 24, top + 80, card_w - 48,
                 kpi.get('value', '') + kpi.get('unit', ''),
                 value_size, bold=True, color=value_color,
                 line_height=1.0, height_px=70)
        # 説明（複数行対応：改行 \n を段落として保持）
        desc = kpi.get('desc', '')
        desc_paragraphs = [
            {'text': line.strip(), 'size': 14, 'color': SUB_TEXT,
             'line_height': 1.55}
            for line in desc.split('\n') if line.strip()
        ] or [{'text': '', 'size': 14, 'color': SUB_TEXT,
               'line_height': 1.55}]
        add_paragraph_box(slide, left + 24, top + 170, card_w - 48,
                           desc_paragraphs, height_px=90,
                           line_height=1.55)

    if conclusion:
        _add_conclusion_band(slide, 470, conclusion, label='SUMMARY')

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 8: スケジュール／ガント風（add_schedule_gantt）
# =====================================================================
def add_schedule_gantt(prs, title, months, tasks, page_num=8, total=10,
                       lead='実行スケジュールは以下のとおりです。'):
    """
    ガント風スケジュール。

    Args:
        months: list of str（例: ['4月','5月','6月','7月','8月','9月']）
        tasks: list of dict {'name': str, 'start': int, 'end': int, 'milestone': bool}
               start/end は months のインデックス（0始まり、end は inclusive）
               milestone=True ならバーではなく◆マーク
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Schedule')
    _add_lead(slide, lead)

    chart_left = 280
    chart_top = 170
    chart_w = 960
    n_months = len(months)
    col_w = chart_w // n_months

    # 月ヘッダ
    for i, m in enumerate(months):
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  chart_left + i * col_w, chart_top, col_w, 36,
                  fill=NAVY)
        add_text(slide, chart_left + i * col_w, chart_top + 8, col_w, m,
                 14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # タスク行
    n_tasks = len(tasks)
    row_h = min(48, (560 - chart_top - 50) // max(n_tasks, 1))

    for t_idx, t in enumerate(tasks):
        y = chart_top + 40 + t_idx * row_h
        # タスク名
        add_text(slide, 40, y + 8, 230, t['name'], 16, bold=True,
                 color=TEXT, height_px=row_h - 4)
        # 列の縞模様
        for i in range(n_months):
            stripe_bg = STRIPE if t_idx % 2 == 1 else WHITE
            add_shape(slide, MSO_SHAPE.RECTANGLE,
                      chart_left + i * col_w, y, col_w, row_h - 4,
                      fill=stripe_bg)
        # バー or マイルストーン
        if t.get('milestone'):
            mx = chart_left + t['start'] * col_w + col_w // 2 - 12
            my = y + (row_h - 4) // 2 - 12
            add_shape(slide, MSO_SHAPE.DIAMOND, mx, my, 24, 24,
                      fill=RED)
        else:
            bar_left = chart_left + t['start'] * col_w + 4
            bar_right = chart_left + (t['end'] + 1) * col_w - 4
            bar_w = bar_right - bar_left
            bar_h = row_h - 18
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                      bar_left, y + 6, bar_w, bar_h,
                      fill=NAVY)

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 9: 運用フロー比較（add_flow_compare）
# =====================================================================
def add_flow_compare(prs, title, before_steps, after_steps,
                     page_num=9, total=10,
                     lead='現状フローと改善後フローを対比します。',
                     conclusion=None):
    """
    Before/After フロー比較。

    Args:
        before_steps: list of str（現状のステップ名）
        after_steps: list of str（改善後のステップ名）
    """
    slide = _blank_slide(prs)
    _add_header(slide, title, 'Flow Compare')
    _add_lead(slide, lead)

    def _draw_flow(top, label, steps, badge_color, card_fill, card_line):
        # バッジ
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 40, top + 20, 100, 36,
                  fill=badge_color)
        add_text(slide, 40, top + 26, 100, label, 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)

        # ステップカード（横並び）
        n = len(steps)
        avail_w = 1080
        card_w = min(180, (avail_w - 30 * (n - 1)) // n)
        gap = (avail_w - card_w * n) // max(n - 1, 1) if n > 1 else 0
        start_x = 160

        for i, step in enumerate(steps):
            x = start_x + i * (card_w + gap)
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, top, card_w, 76,
                      fill=card_fill, line=card_line, line_width_pt=2)
            add_text(slide, x + 10, top + 22, card_w - 20, step, 15,
                     bold=True, color=TEXT, align=PP_ALIGN.CENTER,
                     height_px=40, line_height=1.4)
            # 矢印
            if i < n - 1:
                arrow_x = x + card_w + gap // 2 - 10
                add_text(slide, arrow_x, top + 28, 20, '▸', 22,
                         bold=True, color=RED, align=PP_ALIGN.CENTER)

    _draw_flow(170, 'BEFORE', before_steps,
               RGBColor(0x88, 0x88, 0x88), LIGHT_GRAY, BORDER_GRAY)
    _draw_flow(320, 'AFTER', after_steps, NAVY, WHITE, NAVY)

    if conclusion:
        _add_conclusion_band(slide, 470, conclusion, label='IMPROVEMENT')

    _add_footer(slide, page_num, total)
    return slide


# =====================================================================
# Layout 10: クロージング（add_closing）
# =====================================================================
def add_closing(prs, message='Thank you.', next_step='',
                contact='', page_num=10, total=10):
    """
    クロージングスライド。

    Args:
        message: 大きく表示するメッセージ
        next_step: 次のアクション説明
        contact: 問い合わせ先
    """
    slide = _blank_slide(prs)

    # ナビー背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, CANVAS_W_PX, 660, fill=NAVY)

    # メッセージ
    add_text(slide, 80, 200, 1120, message, 56, bold=True, color=WHITE,
             line_height=1.2, height_px=80)

    # 赤ライン
    add_shape(slide, MSO_SHAPE.RECTANGLE, 80, 300, 480, 3, fill=RED)

    # ネクストステップ
    if next_step:
        add_text(slide, 80, 340, 1120, next_step, 22, color=NAVY_E6,
                 line_height=1.6, height_px=120)

    # 問い合わせ先（右下右揃え）
    if contact:
        add_text(slide, 740, 560, 460, contact, 16, color=WHITE,
                 align=PP_ALIGN.RIGHT, line_height=1.6, height_px=60)

    # 白フッター背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 661, CANVAS_W_PX, 59, fill=WHITE)
    add_text(slide, 40, 682, 400, '紺＆クリーン スライド作成', 14, bold=True,
             color=NAVY)
    add_text(slide, 1140, 685, 120, f'{page_num} / {total}', 14,
             color=PAGE_NUM, align=PP_ALIGN.RIGHT)

    return slide


# =====================================================================
# C-1 / Layout 11: UI診断スコアカード概要（add_scorecard_overview）
# =====================================================================
def add_scorecard_overview(prs, diagnosis, page_num=1, total=3,
                            author='UI/UX診断 by GPTs'):
    """
    UI診断スコアカード スライド1（概要）。

    総合スコア（44〜56pt）＋ランク＋サービス情報＋強み（箇条書き）＋
    最優先課題（箇条書き）を1枚に集約。

    Args:
        diagnosis: dict
            - service_name (str): 診断対象サービス名
            - url (str): URL
            - input_type (str): "URL入力" / "スクショ入力"
            - total_score (int): 50点満点
            - rank (str): S/A/B/C/D
            - rank_label (str): "優秀" など
            - strengths (list of str): 2〜4項目
            - priority_issues (list of str): 1〜3項目
    """
    slide = _blank_slide(prs)
    _add_header(slide, 'UI診断スコアカード', 'UI Diagnosis Scorecard')

    # ── 上段左：サービス情報（left=40, top=88, width=720）
    add_text(slide, 40, 90, 720, '診断対象', 14, bold=True,
             color=NAVY, letter_spacing=2)
    add_text(slide, 40, 114, 720, diagnosis.get('service_name', ''), 22,
             bold=True, color=TEXT, height_px=34)
    # URL とinput_type を1行で
    url = diagnosis.get('url', '')
    input_type = diagnosis.get('input_type', '')
    meta = f'{url}　／　{input_type}' if url else input_type
    add_text(slide, 40, 150, 720, meta, 14, color=SUB_TEXT, height_px=22)

    # ── 上段右：総合スコア＋ランクカード（left=800, top=88, width=440, height=130）
    score = diagnosis.get('total_score', 0)
    rank = diagnosis.get('rank', '-')
    rank_label = diagnosis.get('rank_label', '')
    # ランクに応じた色
    rank_color = NAVY if rank in ('S', 'A', 'B') else RED

    # スコアカード背景
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 800, 88, 440, 130,
              fill=LIGHT_GRAY)
    # 左帯
    add_shape(slide, MSO_SHAPE.RECTANGLE, 800, 88, 6, 130, fill=rank_color)

    # ラベル
    add_text(slide, 822, 100, 200, 'TOTAL SCORE', 14, bold=True,
             color=NAVY, letter_spacing=2)
    # 総合スコア（条項13：核心数値、48ptで強調）
    add_text(slide, 822, 124, 280,
             f'{score}', 48, bold=True, color=rank_color,
             line_height=1.0, height_px=64)
    # 「/50」を小さく
    add_text(slide, 822 + 100, 162, 80, '／ 50', 18, bold=True,
             color=SUB_TEXT, height_px=24)

    # ランクバッジ（右半分）
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1060, 110, 160, 86,
              fill=rank_color)
    add_text(slide, 1060, 118, 160, 'RANK', 14, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, letter_spacing=2)
    add_text(slide, 1060, 138, 160, rank, 40, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, line_height=1.0,
             height_px=52)
    if rank_label:
        add_text(slide, 800, 226, 440, rank_label, 16, color=SUB_TEXT,
                 align=PP_ALIGN.CENTER)

    # ── 下段左：強み（left=40, top=260, width=590）
    # 見出し帯（紺ベタ＋白文字）でスタイリッシュに
    SECTION_TOP = 260
    SECTION_W = 590
    LABEL_H = 36
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, SECTION_TOP, SECTION_W,
              LABEL_H, fill=NAVY)
    add_text(slide, 56, SECTION_TOP + 8, SECTION_W - 32,
             '◎ 強み （Strengths）', 16, bold=True, color=WHITE,
             height_px=24, letter_spacing=2)

    # 各項目をカード化して並べる（最大3項目想定）
    # カード幅 590px、本文width = 590 - 14(左padding) - 28(badge) - 8(gap) - 14(右padding) = 526px
    # 14pt全角約19px → 1行=27.6全角文字 → 30文字までの本文を1行表示できる
    strengths = diagnosis.get('strengths', [])[:3]
    item_top = SECTION_TOP + LABEL_H + 10
    item_h = 104
    item_gap = 10
    BADGE_D = 28  # バッジ直径（小さくして本文領域を確保）
    BADGE_X_OFFSET = 14
    BODY_LEFT_OFFSET = BADGE_X_OFFSET + BADGE_D + 10  # 52px
    BODY_RIGHT_PADDING = 14

    def _draw_item(left_x, y, idx, s, badge_color):
        # カード背景
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left_x, y,
                  SECTION_W, item_h, fill=LIGHT_GRAY)
        # 番号バッジ（サークル、縦中央）
        bx = left_x + BADGE_X_OFFSET
        by = y + (item_h - BADGE_D) // 2
        add_shape(slide, MSO_SHAPE.OVAL, bx, by, BADGE_D, BADGE_D,
                  fill=badge_color)
        add_text(slide, bx, by + 4, BADGE_D, f'{idx + 1:02d}', 14,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=22)
        # 本文（縦中央寄せ）
        body_left = left_x + BODY_LEFT_OFFSET
        body_w = SECTION_W - BODY_LEFT_OFFSET - BODY_RIGHT_PADDING
        add_paragraph_box(slide, body_left, y + 18, body_w,
            [{'text': s, 'size': 14, 'color': TEXT, 'line_height': 1.6}],
            height_px=item_h - 32, line_height=1.6)

    for idx, s in enumerate(strengths):
        y = item_top + idx * (item_h + item_gap)
        _draw_item(40, y, idx, s, NAVY)

    # ── 下段右：最優先課題（left=650, top=260, width=590）
    add_shape(slide, MSO_SHAPE.RECTANGLE, 650, SECTION_TOP, SECTION_W,
              LABEL_H, fill=RED)
    add_text(slide, 666, SECTION_TOP + 8, SECTION_W - 32,
             '⚠ 最優先課題 （Priority Issues）', 16, bold=True,
             color=WHITE, height_px=24, letter_spacing=2)

    issues = diagnosis.get('priority_issues', [])[:3]
    for idx, s in enumerate(issues):
        y = item_top + idx * (item_h + item_gap)
        _draw_item(650, y, idx, s, RED)

    _add_footer(slide, page_num, total, author=author)
    return slide


# =====================================================================
# C-1 / Layout 12: UI診断 10項目スコア表（add_scorecard_table）
# =====================================================================
def add_scorecard_table(prs, diagnosis, page_num=2, total=3,
                         author='UI/UX診断 by GPTs',
                         high_threshold=0.75, low_threshold=0.50):
    """
    UI診断スコアカード スライド2（10項目スコア一覧）。

    スコア列のみ色分け：
      - スコア / max >= high_threshold → 紺（高評価）
      - スコア / max <= low_threshold → 赤（低評価）
      - その他 → 通常テキスト色

    Args:
        diagnosis: scores キーに list of {category, score, max, comment}
        high_threshold: 高評価とみなす比率（既定 0.75）
        low_threshold: 低評価とみなす比率（既定 0.50）
    """
    slide = _blank_slide(prs)
    _add_header(slide, 'UI診断スコアカード', '10-Item Score Detail')
    _add_lead(slide, f'{diagnosis.get("service_name", "")} の10項目評価詳細です。'
                       'スコア列は高評価=紺、低評価=赤で色分けしています。')

    scores = diagnosis.get('scores', [])
    columns = ['No', '評価項目', 'スコア', 'コメント']
    col_widths_pct = [8, 28, 14, 50]

    # 凡例帯（高さ36px）を表の下に確実に置くため、表の縦サイズを管理
    table_left, table_top = 40, 160
    table_w = 1200
    n_rows = len(scores) + 1  # ヘッダ含む
    # フッター(660)・余白(12px)・凡例帯(36px)・余白(12px) を引いた領域を使用
    # 表下端 = 660 - 12 - 36 - 12 = 600
    table_max_bottom = 600
    table_avail_h = table_max_bottom - table_top
    # ヘッダ行はやや高め、データ行は均等
    header_h = 38
    data_avail = table_avail_h - header_h
    row_h = max(28, min(36, data_avail // max(1, len(scores))))

    table_shape = slide.shapes.add_table(
        n_rows, len(columns),
        px(table_left), px(table_top),
        px(table_w), px(header_h + row_h * len(scores))
    )
    table = table_shape.table
    # 行高を明示設定（PowerPoint側の自動拡大を抑制）
    table.rows[0].height = px(header_h)
    for r in range(1, n_rows):
        table.rows[r].height = px(row_h)

    # 列幅
    for i, w in enumerate(col_widths_pct):
        table.columns[i].width = px(int(table_w * w / 100))

    # ヘッダ行
    for c, col_name in enumerate(columns):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ''
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        set_run(run, col_name, 15, bold=True, color=WHITE)
        cell.margin_left = Pt(8)
        cell.margin_right = Pt(8)
        cell.margin_top = Pt(6)
        cell.margin_bottom = Pt(6)

    # データ行
    for r, item in enumerate(scores, start=1):
        bg = WHITE if r % 2 == 1 else STRIPE
        score_val = item.get('score', 0)
        max_val = item.get('max', 5)
        ratio = score_val / max_val if max_val else 0
        if ratio >= high_threshold:
            score_color = NAVY
            score_bold = True
        elif ratio <= low_threshold:
            score_color = RED
            score_bold = True
        else:
            score_color = TEXT
            score_bold = False

        # No
        _fill_cell(table.cell(r, 0), str(r), 14, color=TEXT,
                    bg=bg, align=PP_ALIGN.CENTER, margin_v=3, no_wrap=True)
        # 評価項目
        # [FB対応 2026-06-30] scores の `name` キーを優先（schema統一）
        _fill_cell(table.cell(r, 1),
                    item.get('name') or item.get('category', ''), 14,
                    bold=True, color=TEXT, bg=bg, align=PP_ALIGN.LEFT,
                    margin_v=3, no_wrap=True)
        # スコア（色分け＋大きめ）
        score_text = f'{score_val} / {max_val}'
        _fill_cell(table.cell(r, 2), score_text, 16, bold=score_bold,
                    color=score_color, bg=bg, align=PP_ALIGN.CENTER,
                    margin_v=3, no_wrap=True)
        # コメント（長文は省略して1行で確実に収める）
        comment_raw = item.get('comment', '')
        comment = comment_raw
        total_w = 0.0
        for idx, ch in enumerate(comment_raw):
            total_w += 0.5 if ord(ch) < 128 else 1.0
            if total_w > 28:
                comment = comment_raw[:idx].rstrip() + '…'
                break
        _fill_cell(table.cell(r, 3), comment, 14,
                    color=TEXT, bg=bg, align=PP_ALIGN.LEFT,
                    margin_v=3, no_wrap=True)

    # 凡例（表の下に色チップ付きで配置）
    table_end_y = table_top + header_h + row_h * len(scores)
    legend_y = table_end_y + 12
    # 凡例帯背景（淡グレー、高さ36px、14pt下限厳守）
    legend_h = 36
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
              40, legend_y, 1200, legend_h,
              fill=LIGHT_GRAY, line=BORDER_GRAY, line_width_pt=0.5)
    # 凡例ラベル「凡例」
    add_text(slide, 56, legend_y + 8, 60, '凡例', 14,
             bold=True, color=SUB_TEXT, height_px=22)
    # 紺チップ＋説明
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              120, legend_y + 12, 16, 12, fill=NAVY)
    add_text(slide, 142, legend_y + 8, 360,
             f'紺＝高評価（達成率 {int(high_threshold*100)}% 以上）',
             14, color=SUB_TEXT, height_px=22)
    # 赤チップ＋説明
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              520, legend_y + 12, 16, 12, fill=RED)
    add_text(slide, 542, legend_y + 8, 360,
             f'赤＝低評価（達成率 {int(low_threshold*100)}% 以下）',
             14, color=SUB_TEXT, height_px=22)
    # 通常チップ＋説明
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              920, legend_y + 12, 16, 12, fill=BORDER_GRAY)
    add_text(slide, 942, legend_y + 8, 260,
             'グレー＝中位', 14, color=SUB_TEXT, height_px=22)

    _add_footer(slide, page_num, total, author=author)
    return slide


def _truncate_full(s, max_w):
    """文字列を全角換算 max_w で切り、超えたら末尾を '…' に置換。

    全角=1.0、半角(ASCII)=0.5 でカウント。
    """
    total = 0.0
    for idx, ch in enumerate(s):
        total += 0.5 if ord(ch) < 128 else 1.0
        if total > max_w:
            return s[:idx].rstrip() + '…'
    return s


def _add_multi_run_box(slide, left_px, top_px, width_px, height_px,
                        paragraphs, *, line_height=None,
                        space_after_pt=None, space_before_pt=None,
                        vertical_anchor=None):
    """1テキストボックス内に複数段落＋段落内の複数runを配置するヘルパ。

    1コンテキスト1ブロック原則の核となる関数。
    タイトル+サブタイトル、ラベル+値、ヘッダ+本文など、論理的に
    1つの文脈に属するテキスト群を「1つのテキストボックス」に集約する。

    paragraphs: list of dict
      各dictは段落を表す。以下のキーをサポート：
        - text (str): 単一run。size, bold, color, line_height を併用可。
        - runs (list of dict): 段落内に複数run。各runは {text, size, bold, color}。
        - line_height (float): 段落の line_spacing（＝段落**内部**の折返し行間）
        - space_after_pt (float): 段落末尾に追加する余白（pt）。項目間の分離余白に使う。
        - space_before_pt (float): 段落先頭に追加する余白（pt）。
        - align (PP_ALIGN): 段落の整列

    Parameters
    ----------
    line_height : float, optional
        全段落に適用する既定の line_spacing（各段落dictで上書き可）。
        **段落内の折返し行間** を制御する。項目間の余白ではない。
    space_after_pt : float, optional
        全段落に適用する既定の段落末尾余白（pt）。**リスト項目間の分離感** に使う。
        [v15/2026-07-12] リスト表示（強み・課題Top3等）で
        line_height と項目間余白を分離制御するために新設。
    space_before_pt : float, optional
        全段落に適用する既定の段落先頭余白（pt）。通常は使わない。

    Notes
    -----
    - line_height と space_after_pt を分離する設計思想：
      折返し行間（line_height）は詰めて「1文章の連続感」を担保し、
      項目間余白（space_after_pt）で「別項目である分離感」を明示する。
      これによりリスト表示の可読性が根本改善する（v15対応）。
    """
    box = slide.shapes.add_textbox(
        px(left_px), px(top_px), px(width_px), px(height_px))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    # [v15/2026-07-12] 縦中央寄せ（総評帯など「1行運用でボックス高さより
    # テキストが小さい」場合に、上下余白の不均衡を解消するため）
    if vertical_anchor is not None:
        try:
            tf.vertical_anchor = vertical_anchor
        except Exception:
            pass

    for i, p in enumerate(paragraphs):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = p.get('align', PP_ALIGN.LEFT)
        lh = p.get('line_height', line_height)
        if lh:
            para.line_spacing = lh
        # [v15] 段落間余白の分離制御
        sa = p.get('space_after_pt', space_after_pt)
        if sa is not None and sa != 0:
            para.space_after = Pt(sa)
        sb = p.get('space_before_pt', space_before_pt)
        if sb is not None and sb != 0:
            para.space_before = Pt(sb)
        if 'runs' in p:
            for r_def in p['runs']:
                run = para.add_run()
                set_run(run, r_def.get('text', ''),
                        r_def.get('size', 14),
                        bold=r_def.get('bold', False),
                        color=r_def.get('color', TEXT))
        else:
            run = para.add_run()
            set_run(run, p.get('text', ''),
                    p.get('size', 14),
                    bold=p.get('bold', False),
                    color=p.get('color', TEXT))
    return box


def _fill_cell(cell, text, size_pt, *, bold=False, color=TEXT, bg=WHITE,
                align=PP_ALIGN.LEFT, margin_v=6, no_wrap=False):
    """テーブルセルを単一スタイルで埋めるヘルパー

    Args:
        margin_v: セル上下マージン(pt)。行高をタイトに保ちたいときは小さく。
        no_wrap: True で word_wrap を OFF（長文を1行で描画、はみ出しはクリップ）
    """
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.text = ''
    tf = cell.text_frame
    if no_wrap:
        tf.word_wrap = False
    # 行高の自動拡大を抑制（LibreOffice 描画対策）
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass
    # 縦中央寄せ（行高に余裕がない場合でも文字が枠に収まる）
    try:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    except Exception:
        pass
    para = tf.paragraphs[0]
    para.alignment = align
    # 段落 line_spacing を pt で明示固定（LibreOffice の行高拡大を抑制）
    para.line_spacing = Pt(size_pt + 2)
    para.space_before = Pt(0)
    para.space_after = Pt(0)
    run = para.add_run()
    set_run(run, str(text), size_pt, bold=bold, color=color)
    cell.margin_left = Pt(8)
    cell.margin_right = Pt(8)
    cell.margin_top = Pt(margin_v)
    cell.margin_bottom = Pt(margin_v)


# =====================================================================
# C-1 / Layout 13: UI診断 結論（add_scorecard_conclusion）
# =====================================================================
def add_scorecard_conclusion(prs, diagnosis, page_num=3, total=3,
                              author='UI/UX診断 by GPTs',
                              cta='詳細な改善提案は別途お送りします。'):
    """
    UI診断スコアカード スライド3（結論）。

    結論テキストを大きく表示＋CONCLUSION結論帯＋CTA。

    Args:
        diagnosis: conclusion キー（総評3〜5文）
        cta: 結論帯に表示する次のアクション文
    """
    slide = _blank_slide(prs)
    _add_header(slide, 'UI診断スコアカード', 'Conclusion')

    # ラベル
    add_text(slide, 40, 90, 1200, '総評', 14, bold=True,
             color=NAVY, letter_spacing=2, height_px=22)

    # サービス名 + ランク要約
    score = diagnosis.get('total_score', 0)
    rank = diagnosis.get('rank', '-')
    rank_label = diagnosis.get('rank_label', '')
    summary_line = (f'{diagnosis.get("service_name", "")} の総合評価は '
                     f'{score}/50点・ランク {rank}（{rank_label}）')
    add_text(slide, 40, 116, 1200, summary_line, 18, bold=True,
             color=TEXT, height_px=28)

    # 区切り赤ライン
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, 156, 200, 3, fill=RED)

    # 結論テキスト（18〜20pt）
    conclusion_text = diagnosis.get('conclusion', '')
    # 改行（句点で改行を入れて読みやすく）
    paragraphs = [p.strip() for p in conclusion_text.split('\n') if p.strip()]
    if not paragraphs:
        paragraphs = [conclusion_text]

    paragraph_specs = []
    for p in paragraphs:
        paragraph_specs.append({
            'text': p, 'size': 19, 'color': TEXT, 'line_height': 1.8,
        })

    add_paragraph_box(slide, 40, 180, 1200, paragraph_specs,
                      height_px=340, line_height=1.8)

    # 結論帯（CONCLUSION）
    _add_conclusion_band(slide, 540, cta, label='NEXT STEP')

    _add_footer(slide, page_num, total, author=author)
    return slide


# =====================================================================
# C-2 / Layout 14: 改善提案 一覧（add_proposal_overview）
# =====================================================================
def add_proposal_overview(prs, proposals_data, page_num=1, total=3,
                            author='UI/UX診断 by GPTs'):
    """
    改善提案 スライド1（一覧表）。

    add_action_table を流用したPPTX表組。
    優先度列は S/A=赤、B=紺 で色分け。

    Args:
        proposals_data: dict
            - service_name (str)
            - proposals (list of dict): {no, title, category, priority, ...}
    """
    slide = _blank_slide(prs)
    _add_header(slide, '改善提案リスト', 'Improvement Proposals')
    _add_lead(slide, f'{proposals_data.get("service_name", "")} に対する改善提案の一覧です。'
                       '詳細は次スライド以降に記載しています。')

    proposals = proposals_data.get('proposals', [])
    columns = ['No', '改善タイトル', '対象カテゴリ', '優先度']
    col_widths_pct = [8, 55, 22, 15]

    table_left, table_top = 40, 160
    table_w = 1200
    n_rows = len(proposals) + 1
    row_h = 44

    table_shape = slide.shapes.add_table(
        n_rows, len(columns),
        px(table_left), px(table_top),
        px(table_w), px(row_h * n_rows)
    )
    table = table_shape.table

    for i, w in enumerate(col_widths_pct):
        table.columns[i].width = px(int(table_w * w / 100))

    # ヘッダ
    for c, col_name in enumerate(columns):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text = ''
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        set_run(run, col_name, 16, bold=True, color=WHITE)
        cell.margin_left = Pt(8); cell.margin_right = Pt(8)
        cell.margin_top = Pt(8); cell.margin_bottom = Pt(8)

    # データ行
    for r, p in enumerate(proposals, start=1):
        bg = WHITE if r % 2 == 1 else STRIPE
        priority = p.get('priority', 'B')
        priority_color = RED if priority in ('S', 'A') else NAVY

        # No
        _fill_cell(table.cell(r, 0), str(p.get('no', r)), 16,
                    bold=True, color=NAVY, bg=bg, align=PP_ALIGN.CENTER)
        # タイトル
        _fill_cell(table.cell(r, 1), p.get('title', ''), 15,
                    bold=True, color=TEXT, bg=bg, align=PP_ALIGN.LEFT)
        # カテゴリ
        _fill_cell(table.cell(r, 2), p.get('category', ''), 14,
                    color=SUB_TEXT, bg=bg, align=PP_ALIGN.CENTER)
        # 優先度（バッジ風：背景色付き）
        cell = table.cell(r, 3)
        cell.fill.solid()
        cell.fill.fore_color.rgb = priority_color
        cell.text = ''
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        set_run(run, priority, 18, bold=True, color=WHITE)
        cell.margin_left = Pt(8); cell.margin_right = Pt(8)
        cell.margin_top = Pt(8); cell.margin_bottom = Pt(8)

    _add_footer(slide, page_num, total, author=author)
    return slide


# =====================================================================
# C-2 / Layout 15: 改善提案 詳細（add_proposal_detail）
# =====================================================================
def add_proposal_detail(prs, proposal, page_num=2, total=3,
                         author='UI/UX診断 by GPTs',
                         summary=None,
                         force_two_column=False):
    """
    改善提案 詳細スライド。

    **デフォルト：1スライド1提案（全幅1200pxカード）**
      - 「現状」「アクション」「期待効果」の各セクション本文を
        省略なしでフル表示
      - 原文の改行（\\n）はそのまま段落として尊重

    Args:
        proposal: dict（推奨）または list of dict（互換）
            dict 形式: {no, title, category, priority,
                        current, action, expected}
            list 形式: 旧API互換。1件なら1スライド表示、
                       2件 + force_two_column=True なら2列表示。
        summary: 結論帯に表示する総括コメント（最終スライドに付与）
        force_two_column: True のとき、list で2件渡された場合に
                          従来の2列レイアウトを強制（後方互換用）。
                          デフォルト False（1件1スライド推奨）。
    """
    # ---- 引数正規化：dict も list も受け付ける ----
    if isinstance(proposal, dict):
        proposals_slice = [proposal]
    elif isinstance(proposal, (list, tuple)):
        proposals_slice = list(proposal)
    else:
        raise TypeError('proposal は dict または list を渡してください')

    if len(proposals_slice) == 0:
        raise ValueError('proposal が空です')
    if len(proposals_slice) > 2:
        raise ValueError(
            'add_proposal_detail は1スライド最大2件まで。'
            '3件以上は分割して複数回呼び出してください。')

    # 2件渡されたとき：force_two_column=True なら2列、未指定なら警告
    n = len(proposals_slice)
    if n == 2 and not force_two_column:
        raise ValueError(
            '2件渡されましたが、1スライド1提案がデフォルトです。'
            '2列表示が必要なら force_two_column=True を指定してください。')

    slide = _blank_slide(prs)
    _add_header(slide, '改善提案リスト', 'Proposal Detail')

    # ---- リード文 ----
    nums = [str(p.get('no', i + 1)) for i, p in enumerate(proposals_slice)]
    if n == 1:
        lead_text = (f'提案 #{nums[0]} ── '
                     f'現状 → アクション → 期待効果 の3段で整理。')
    else:
        lead_text = (f'提案 #{nums[0]} / #{nums[1]} ── '
                     f'現状 → アクション → 期待効果 の3段で整理。')
    add_text(slide, 40, 90, 1200, lead_text, 16,
             color=SUB_TEXT, height_px=24)

    # ---- カード寸法 ----
    card_top = 130
    # 結論帯ありなら top=568, height=82 → 下端650
    # 結論帯なしならフッター（top=660）まで使用
    card_max_bottom = 552 if summary else 640
    card_h = card_max_bottom - card_top
    card_w = (1200 - 30) // 2 if n == 2 else 1200

    # フォントサイズ（2列時は少し小さく；14pt下限厳守）
    title_size = 17 if n == 2 else 22
    body_size = 14
    label_size = 14
    no_size = 26 if n == 2 else 36

    # カード内領域配分
    PAD_TOP = 14
    HEADER_H = 32 if n == 2 else 40
    # タイトル：1列なら最大3行折返しOK、2列なら1行に制限（…で省略）
    TITLE_H = 30 if n == 2 else 72
    CATEGORY_H = 22
    DIVIDER_GAP = 8
    SECTION_LABEL_H = 22
    SECTION_BODY_LINE_H = int(body_size * 1.6)  # 22px / 1行
    SECTION_GAP = 12  # セクション間の明示的余白

    # ---- タイトル省略ヘルパ（2列のみ使用） ----
    def _truncate_title(s, max_w):
        total = 0.0
        for idx, ch in enumerate(s):
            add = 0.5 if ord(ch) < 128 else 1.0
            if total + add > max_w:
                return s[:idx].rstrip() + '…'
            total += add
        return s

    for i, p in enumerate(proposals_slice):
        left = 40 + i * (card_w + 30) if n == 2 else 40
        priority = p.get('priority', 'B')
        priority_color = RED if priority in ('S', 'A') else NAVY

        # カード背景＋左帯
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top,
                  card_w, card_h, fill=LIGHT_GRAY)
        add_shape(slide, MSO_SHAPE.RECTANGLE, left, card_top, 6, card_h,
                  fill=priority_color)

        # 番号
        add_text(slide, left + 22, card_top + PAD_TOP, 200,
                 f'#{p.get("no", "")}', no_size, bold=True,
                 color=priority_color, line_height=1.0,
                 height_px=no_size + 4)
        # 優先度バッジ（右上）
        badge_w = 110 if n == 1 else 90
        badge_h = 34 if n == 1 else 30
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                  left + card_w - badge_w - 16, card_top + PAD_TOP + 2,
                  badge_w, badge_h, fill=priority_color)
        add_text(slide, left + card_w - badge_w - 16,
                 card_top + PAD_TOP + (10 if n == 1 else 8),
                 badge_w, f'優先度 {priority}',
                 16 if n == 1 else 14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, height_px=20)

        # タイトル
        title_y = card_top + PAD_TOP + HEADER_H + 4
        raw_title = p.get('title', '')
        display_title = _truncate_title(raw_title, max_w=14) if n == 2 else raw_title
        add_text(slide, left + 22, title_y, card_w - 44,
                 display_title, title_size, bold=True, color=NAVY,
                 height_px=TITLE_H, line_height=1.3)

        # カテゴリ
        cat_y = title_y + TITLE_H + 2
        add_text(slide, left + 22, cat_y, card_w - 44,
                 p.get('category', ''), 14,
                 color=SUB_TEXT, height_px=CATEGORY_H)

        # 区切り線
        divider_y = cat_y + CATEGORY_H + DIVIDER_GAP
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  left + 22, divider_y, card_w - 44, 1, fill=BORDER_GRAY)

        # 3段構成
        sections = [
            ('現状',     p.get('current', ''),  SUB_TEXT),
            ('アクション', p.get('action', ''),  NAVY),
            ('期待効果', p.get('expected', ''), RED),
        ]
        section_top_start = divider_y + DIVIDER_GAP + 4
        section_avail = (card_h - (section_top_start - card_top)
                         - 14 - SECTION_GAP * 2)
        section_h = section_avail // 3

        if n == 1:
            # ---- 1列モード（フル幅）：省略なし、原文改行を尊重 ----
            for s_idx, (label, body, label_color) in enumerate(sections):
                sy = section_top_start + s_idx * (section_h + SECTION_GAP)
                # ラベル
                add_text(slide, left + 22, sy, 240, label, label_size,
                         bold=True, color=label_color, letter_spacing=1,
                         height_px=SECTION_LABEL_H)
                # 本文（原文改行を段落として保持）
                body_top = sy + SECTION_LABEL_H + 4
                body_h = section_h - SECTION_LABEL_H - 6
                body_paragraphs = [
                    {'text': line.strip(),
                     'size': body_size,
                     'color': TEXT,
                     'line_height': 1.55}
                    for line in (body or '').split('\n') if line.strip()
                ] or [{'text': '', 'size': body_size,
                       'color': TEXT, 'line_height': 1.55}]
                add_paragraph_box(slide, left + 22, body_top,
                                  card_w - 44,
                                  body_paragraphs,
                                  height_px=body_h, line_height=1.55)
        else:
            # ---- 2列モード（互換）：文字数制限で省略 ----
            body_inner_w = card_w - 44
            approx_fullwidth_per_line = max(7, body_inner_w // 22)
            raw_lines = max(1, (section_h - SECTION_LABEL_H - 6)
                            // SECTION_BODY_LINE_H)
            safety_factor = 0.82
            max_fullwidth_chars = max(6, int(approx_fullwidth_per_line
                                              * raw_lines * safety_factor))

            def _truncate_to_fullwidth(s, max_w):
                total = 0.0
                for idx, ch in enumerate(s):
                    add = 0.5 if ord(ch) < 128 else 1.0
                    if total + add > max_w:
                        return s[:idx].rstrip() + '…'
                    total += add
                return s

            for s_idx, (label, body, label_color) in enumerate(sections):
                sy = section_top_start + s_idx * (section_h + SECTION_GAP)
                add_text(slide, left + 22, sy, 200, label, label_size,
                         bold=True, color=label_color, letter_spacing=1,
                         height_px=SECTION_LABEL_H)
                body_top = sy + SECTION_LABEL_H + 2
                body_h = section_h - SECTION_LABEL_H - 4
                flat_body = (body or '').replace('\n', ' ')
                flat_body = _truncate_to_fullwidth(flat_body,
                                                   max_fullwidth_chars)
                add_paragraph_box(slide, left + 22, body_top,
                                  card_w - 44,
                                  [{'text': flat_body, 'size': body_size,
                                    'color': TEXT, 'line_height': 1.55}],
                                  height_px=body_h, line_height=1.55)

    # 総括（最終スライドのみ）
    if summary:
        _add_conclusion_band(slide, 568, summary,
                             label='SUMMARY', height_px=82)

    _add_footer(slide, page_num, total, author=author)
    return slide


# =====================================================================
# C-1 / Layout 16: UI診断スコアカード（2枚構成）
# （add_scorecard_onepager）
#   1/2 サマリ：上段カード + 強み/最優先課題（2列） + 結論
#   2/2 詳細スコア表：10項目を横幅1200px全面活用、一言所見は全文表示
# =====================================================================
def _scorecard_title(slide, slide_no, subtitle):
    """両スライド共通のタイトル帯描画（左赤縦帯10px + タイトル + サブタイトル）

    [ビジュアル強化 2026-06-29]
    - 左赤縦帯：6px → 10px に拡幅して視認性向上
    - タイトル位置調整：60pxへシフト（縦帯余白）
    [FB対応 2026-06-30]
    - タイトル帯高さ：56 → 64px に拡大
    - サブタイトル top：+8px下げて可読性向上
    """
    TITLE_TOP = 16
    TITLE_H = 64
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, TITLE_TOP, 10, TITLE_H,
              fill=RED)
    add_text(slide, 60, TITLE_TOP + 4, 1140,
             f'【{slide_no}】UI診断スコアカード', 24, bold=True,
             color=NAVY, height_px=34)
    add_text(slide, 60, TITLE_TOP + 42, 1140, subtitle,
             14, color=SUB_TEXT, height_px=22)


def _scorecard_info_cards(slide, diagnosis):
    """両スライド共通：上段4分割カード（診断対象/入力種別/総合スコア/ランク）

    [ビジュアル強化 2026-06-29]
    - 総合スコア値：26pt → 32pt に拡大
    - ランク値：26pt → 32pt に拡大
    - 左色帯：4px → 6px に拡幅
    - カード背景：NAVY 5%透過の極薄背景で奥行き演出
    - ランクをバッジ風に強調（数字大型 + ラベル小型）
    [FB対応 2026-06-30]
    - INFO_TOP：82 → 90px（タイトル帯64px拡大に追従、+8pxシフト）
    """
    INFO_TOP = 90
    INFO_H = 80
    card_w = (1200 - 12 * 3) // 4  # 291px

    score_val = diagnosis.get('total_score', 0)
    rank = diagnosis.get('rank', '-')
    rank_label = diagnosis.get('rank_label', '')
    rank_color = NAVY if rank in ('S', 'A', 'B') else RED

    # [FB対応 2026-06-30] 左色帯を全4枚NAVYに統一
    # 装飾目的の色分けは廃止。スコア値・ランク値は内部の数字フォント色で
    # 個別強調（rank_color = NAVY or RED）するため左色帯の機能的役割は薄い。
    info_cards = [
        {'bar': NAVY, 'paragraphs': [
            {'text': '診断対象', 'size': 14, 'bold': True,
             'color': SUB_TEXT, 'line_height': 1.0},
            {'text': diagnosis.get('service_name', '') or '-',
             'size': 18, 'bold': True, 'color': TEXT, 'line_height': 1.2},
        ]},
        {'bar': NAVY, 'paragraphs': [
            {'text': '入力種別', 'size': 14, 'bold': True,
             'color': SUB_TEXT, 'line_height': 1.0},
            {'text': diagnosis.get('input_type', '') or '-',
             'size': 16, 'bold': False, 'color': TEXT, 'line_height': 1.2},
        ]},
        {'bar': NAVY, 'paragraphs': [
            {'text': '総合スコア', 'size': 14, 'bold': True,
             'color': SUB_TEXT, 'line_height': 1.0},
            {'runs': [
                {'text': str(score_val), 'size': 32, 'bold': True,
                 'color': rank_color},
                {'text': '  / 50', 'size': 14, 'bold': False,
                 'color': SUB_TEXT},
            ], 'line_height': 1.0},
        ]},
        {'bar': NAVY, 'paragraphs': [
            {'text': 'ランク', 'size': 14, 'bold': True,
             'color': SUB_TEXT, 'line_height': 1.0},
            {'runs': [
                {'text': rank, 'size': 32, 'bold': True, 'color': rank_color},
                {'text': '  ' + (rank_label or ''),
                 'size': 14, 'bold': False, 'color': SUB_TEXT},
            ], 'line_height': 1.0},
        ]},
    ]
    for i, card in enumerate(info_cards):
        x = 40 + i * (card_w + 12)
        # カード本体：白＋極薄NAVYの2層構造で奥行き演出
        _add_bg_frame(slide, x, INFO_TOP, card_w, INFO_H,
                      fill=WHITE,
                      line=BORDER_GRAY, line_width_pt=1)
        _add_bg_frame(slide, x, INFO_TOP, card_w, INFO_H,
                      fill=NAVY, fill_alpha=4000)
        # 左色帯：4px → 6px に拡幅
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, INFO_TOP, 6, INFO_H,
                  fill=card['bar'])
        _add_multi_run_box(slide, x + 18, INFO_TOP + 10,
                           card_w - 26, INFO_H - 14, card['paragraphs'])


def add_scorecard_onepager(prs, diagnosis, page_num=1, total=2,
                            author='UI/UX診断 by GPTs',
                            high_threshold=0.75, low_threshold=0.50,
                            slide_no='1'):
    """
    UI診断スコアカードを2枚に分割して出力。

    1/2 サマリ：
      - タイトル帯
      - 上段4分割カード（診断対象/入力種別/総合スコア/ランク）
      - 強み / 最優先課題（横並び2列、各3項目を line_height 1.7 で余裕表示）
      - 結論帯
    2/2 詳細スコア表：
      - タイトル帯
      - 上段4分割カード（同上：コンテキストを引継ぎ）
      - 10項目スコア表（横幅1200px全面、一言所見は省略なしで折返し表示）

    Returns: (slide_summary, slide_detail) tuple
    """
    # ==============================================
    # 文字数バリデーション（design_system.md §1.3 準拠）
    # ==============================================
    validate_length(diagnosis.get('service_name'), 'service_name',
                    'C-1 サービス名 (diagnosis.service_name)')
    for i, sc in enumerate(diagnosis.get('scores', [])[:10], start=1):
        if isinstance(sc, dict):
            validate_length(sc.get('comment'), 'c1_comment',
                            f'C-1 一言所見 #{i} (scores[{i-1}].comment / 項目: {sc.get("name","?")})')
    for i, st in enumerate(diagnosis.get('strengths', [])[:3], start=1):
        validate_length(st, 'c1_strength',
                        f'C-1 強み #{i} (strengths[{i-1}])')
    for i, iss in enumerate(diagnosis.get('priority_issues', [])[:3], start=1):
        validate_length(iss, 'c1_issue',
                        f'C-1 最優先課題 #{i} (priority_issues[{i-1}])')
    validate_length(diagnosis.get('conclusion'), 'c1_conclusion',
                    'C-1 結論 (diagnosis.conclusion)')

    # ==============================================
    # スライド 1/2 : サマリ
    # ==============================================
    slide = _blank_slide(prs)
    _scorecard_title(slide, slide_no,
                     'サマリ：診断対象・スコア・強み・最優先課題・結論')
    _scorecard_info_cards(slide, diagnosis)

    # ---- 強み と 最優先課題（横並び 2分割、余裕ある行高） ----
    # [ビジュアル強化 2026-06-29]
    # 情報カード下端=82+80=162 → 余白20px → LOWER_TOP=182
    # 結論帯top=608 まで → LOWER_H=426（ヘッダ36 + 本文390px）
    # 本文 line_height=1.6、16pt実描画~28px×折返し2行×3項目 = 168px、余裕あり
    # [FB対応 2026-06-30] タイトル帯64px化に追従、+8pxシフト
    # 情報カード下端=90+80=170 → 余白20px → LOWER_TOP=190
    # [FB対応 2026-07-02] 結論帯62px化に伴い top 614→594、下段下端を
    # 586 (=594-8余白) に揃えるため LOWER_H を 416→396 に縮小。
    FULL_W = 1200
    LOWER_TOP = 190
    LOWER_H = 396  # 606→586 に短縮（結論帯top=594 の直前に安全余白8px）
    GAP = 16
    HALF_W = (FULL_W - GAP) // 2  # 592px
    HEADER_BAR_H = 36  # 32→36：余白拡大
    LEFT_X = 40
    RIGHT_X = 40 + HALF_W + GAP

    # 強み（左半分） - NAVYヘッダ
    add_shape(slide, MSO_SHAPE.RECTANGLE, LEFT_X, LOWER_TOP, HALF_W,
              HEADER_BAR_H, fill=NAVY)
    add_text(slide, LEFT_X + 18, LOWER_TOP, HALF_W - 36,
             '◎ 強み（活かすべき点）', 16, bold=True, color=WHITE,
             height_px=HEADER_BAR_H, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    str_body_top = LOWER_TOP + HEADER_BAR_H
    str_body_h = LOWER_H - HEADER_BAR_H
    # 本文背景：LIGHT_GRAY → 白＋NAVY 3%極薄で上品さ演出
    add_shape(slide, MSO_SHAPE.RECTANGLE, LEFT_X, str_body_top,
              HALF_W, str_body_h, fill=WHITE,
              line=BORDER_GRAY, line_width_pt=0.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, LEFT_X, str_body_top,
              HALF_W, str_body_h, fill=NAVY, fill_alpha=3000)
    # [v15/2026-07-12] リスト内折返しの連続感と項目間の分離感を分離制御：
    # line_height 1.6→1.2（折返し行間タイト）＋ space_after_pt=8pt（項目間余白）
    strengths = diagnosis.get('strengths', [])[:3]
    str_paragraphs = []
    for s in strengths:
        str_paragraphs.append({
            'text': '● ' + s, 'size': 16, 'bold': False,
            'color': TEXT, 'line_height': 1.2, 'space_after_pt': 8,
        })
    _add_multi_run_box(slide, LEFT_X + 20, str_body_top + 18,
                       HALF_W - 40, str_body_h - 32,
                       str_paragraphs)

    # 最優先課題（右半分） - REDヘッダ
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, LOWER_TOP, HALF_W,
              HEADER_BAR_H, fill=RED)
    add_text(slide, RIGHT_X + 18, LOWER_TOP, HALF_W - 36,
             '⚠ 最優先で直すべき点', 16, bold=True, color=WHITE,
             height_px=HEADER_BAR_H, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    issue_body_top = LOWER_TOP + HEADER_BAR_H
    issue_body_h = LOWER_H - HEADER_BAR_H
    # 課題側はRED 3%極薄で警告色トーンを醸成
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, issue_body_top,
              HALF_W, issue_body_h, fill=WHITE,
              line=BORDER_GRAY, line_width_pt=0.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, issue_body_top,
              HALF_W, issue_body_h, fill=RED, fill_alpha=3000)
    # [v15/2026-07-12] リスト内折返しの連続感と項目間の分離感を分離制御
    issues = diagnosis.get('priority_issues', [])[:3]
    issue_paragraphs = []
    for i, s in enumerate(issues):
        issue_paragraphs.append({
            'text': f'{i + 1}. ' + s, 'size': 16, 'bold': False,
            'color': TEXT, 'line_height': 1.2, 'space_after_pt': 8,
        })
    _add_multi_run_box(slide, RIGHT_X + 20, issue_body_top + 18,
                       HALF_W - 40, issue_body_h - 32,
                       issue_paragraphs)

    # 結論帯（最下部）─ [ビジュアル強化] ORANGE系へ変更
    # ORANGE基調＋NAVYラベルで「行動喚起」の視覚的トーン演出
    # [FB対応 2026-07-02] 文字切れ根絶のため帯高さ 38→62px（2行対応）、
    #   最大120字＋word_wrap=True、top を 614→594 に上シフトしフッター境界
    #   （660）まで6pxの安全余白確保。ラベル箱は帯全高で垂直センター視認性向上。
    CONCL_TOP = 594
    CONCL_H = 62  # 帯本体高さ（2行対応・+24px）
    LABEL_W = 88  # ラベル箱幅
    GAP = 40  # ラベル⇔本文間隔（案B推奨値）
    conclusion_text = diagnosis.get('conclusion', '')
    if conclusion_text:
        flat = conclusion_text.replace('\n', ' ')
        # 本体：ORANGE薄塗背景＋ORANGE枠
        add_shape(slide, MSO_SHAPE.RECTANGLE, 40, CONCL_TOP, 1200,
                  CONCL_H, fill=WHITE,
                  line=ORANGE, line_width_pt=1)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 40, CONCL_TOP, 1200,
                  CONCL_H, fill=ORANGE, fill_alpha=10000)
        # ラベル部：ORANGEベタ、帯全高で垂直センター（視認性優先）
        add_shape(slide, MSO_SHAPE.RECTANGLE,
                  40, CONCL_TOP, LABEL_W, CONCL_H,
                  fill=ORANGE)
        add_text(slide, 40, CONCL_TOP, LABEL_W, '結論', 14,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=CONCL_H, letter_spacing=2,
                 vertical_anchor=MSO_ANCHOR.MIDDLE)
        # 120字上限で切り詰め（設計方針：GPTs側で120字以内に生成させる）
        flat_disp = _truncate_full(flat, 120)
        # 本文 left = 40 + LABEL_W + GAP = 168、幅 = 40 + 1200 - 168 = 1072
        body_left = 40 + LABEL_W + GAP
        body_w = 40 + 1200 - body_left
        body_box = add_text(slide, body_left, CONCL_TOP, body_w, flat_disp, 14,
                            bold=True, color=NAVY, height_px=CONCL_H,
                            line_height=1.35,
                            vertical_anchor=MSO_ANCHOR.MIDDLE)
        # word_wrap=True で2行折り返しを許容（文字切れ根絶）
        body_box.text_frame.word_wrap = True

    _add_footer(slide, page_num, total, author=author)
    slide_summary = slide

    # ==============================================
    # スライド 2/2 : 詳細スコア表
    # ==============================================
    slide = _blank_slide(prs)
    _scorecard_title(slide, slide_no,
                     '10項目評価：各項目のスコアと一言所見（全文表示）')
    _scorecard_info_cards(slide, diagnosis)

    # 表セクション（横幅1200px全面、一言所見は省略なしで折返し）
    # 情報カード下端=160 → 余白20px → TABLE_TOP=180
    # フッター帯top=660 まで → 表領域=460px → ヘッダ28 + 10行×43px = 458px
    FULL_W = 1200
    TABLE_TOP = 180
    header_h = 28
    n_items = len(diagnosis.get('scores', [])[:10])
    # 残り432px ÷ 10行 = 43px/行（一言所見の2行折返しに十分）
    row_h = 30  # 設定値30px、LibreOffice実描画で約40-43px
    # No:項目:スコア:一言所見 = 5:17:7:71 → 60:204:84:852px
    # 「ファーストビュー」など全角8文字を1行で収めるため項目幅を拡大
    col_widths_pct_local = [5, 17, 7, 71]
    columns_local = ['No', '項目', 'スコア', '一言所見']
    n_rows = n_items + 1

    table_shape = slide.shapes.add_table(
        n_rows, len(columns_local),
        px(40), px(TABLE_TOP),
        px(FULL_W), px(header_h + row_h * n_items))
    table = table_shape.table
    for i, wpct in enumerate(col_widths_pct_local):
        table.columns[i].width = px(int(FULL_W * wpct / 100))
    table.rows[0].height = px(header_h)
    for r in range(1, n_rows):
        table.rows[r].height = px(row_h)

    # ヘッダ行
    header_aligns = [PP_ALIGN.CENTER, PP_ALIGN.LEFT,
                     PP_ALIGN.CENTER, PP_ALIGN.LEFT]
    for c, (col_name, a) in enumerate(zip(columns_local, header_aligns)):
        _fill_cell(table.cell(0, c), col_name, 14, bold=True,
                   color=WHITE, bg=NAVY, align=a,
                   margin_v=2, no_wrap=True)

    # データ行
    scores = diagnosis.get('scores', [])[:10]
    for r, item in enumerate(scores, start=1):
        bg = WHITE if r % 2 == 1 else STRIPE
        sv = item.get('score', 0)
        mv = item.get('max', 5)
        ratio = sv / mv if mv else 0
        if ratio >= high_threshold:
            score_color = NAVY; score_bold = True
        elif ratio <= low_threshold:
            score_color = RED; score_bold = True
        else:
            score_color = TEXT; score_bold = False

        _fill_cell(table.cell(r, 0), str(r), 14, color=TEXT,
                   bg=bg, align=PP_ALIGN.CENTER,
                   margin_v=2, no_wrap=True)
        # [FB対応 2026-06-30] scores の `name` キーを優先（schema統一）
        cat_raw = item.get('name') or item.get('category', '')
        cat_disp = _truncate_full(cat_raw, 10)
        _fill_cell(table.cell(r, 1), cat_disp, 14, bold=True,
                   color=TEXT, bg=bg, align=PP_ALIGN.LEFT,
                   margin_v=2, no_wrap=True)
        score_text = f'{sv}/{mv}'
        _fill_cell(table.cell(r, 2), score_text, 14, bold=score_bold,
                   color=score_color, bg=bg, align=PP_ALIGN.CENTER,
                   margin_v=2, no_wrap=True)
        # 一言所見：全文表示（折返し許可）
        comment_raw = item.get('comment', '')
        _fill_cell(table.cell(r, 3), comment_raw, 14, color=TEXT,
                   bg=bg, align=PP_ALIGN.LEFT,
                   margin_v=2, no_wrap=False)

    _add_footer(slide, page_num + 1, total, author=author)
    slide_detail = slide

    return (slide_summary, slide_detail)


# =====================================================================
# C-2 / Layout 17: 改善提案リスト（2枚構成）
# （add_proposal_onepager）
#   1/2 前半提案：5件中の前半3件をカード形式で表示
#   2/2 後半提案：残り2件のカード + POINT帯（最優先の一手まとめ）
# =====================================================================
def _proposal_legend(slide, title_top):
    """両スライド共通：右上の優先度凡例（高/中/低バッジ）"""
    LEGEND_TOP = title_top + 14
    add_text(slide, 880, LEGEND_TOP, 80, '優先度', 14, bold=True,
             color=SUB_TEXT, align=PP_ALIGN.RIGHT, height_px=22)
    PRIO_RED = RGBColor(0xD0, 0x02, 0x1B)
    PRIO_ORANGE = RGBColor(0xE8, 0x8B, 0x1F)
    PRIO_GRAY = RGBColor(0x88, 0x88, 0x88)
    for i, (lbl, c) in enumerate([('高', PRIO_RED),
                                   ('中', PRIO_ORANGE),
                                   ('低', PRIO_GRAY)]):
        bx = 970 + i * 80
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, bx, LEGEND_TOP,
                  70, 28, fill=c)
        add_text(slide, bx, LEGEND_TOP + 3, 70, lbl, 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER, height_px=22)


def _proposal_title(slide, slide_no, subtitle):
    """両スライド共通：タイトル帯

    [ビジュアル強化 2026-06-29]
    - 左赤縦帯：6px → 10px に拡幅
    [FB対応 2026-06-30]
    - タイトル帯高さ：56 → 64px に拡大（サブタイトル収納問題解消）
    - タイトル段落の line_height 拡大で視覚的間隔確保
    - サブタイトル space_before 相当の余白追加
    """
    TITLE_TOP = 16
    TITLE_H = 64
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, TITLE_TOP, 10, TITLE_H,
              fill=RED)
    # タイトル（24pt）とサブタイトル（14pt）を別々の add_text で配置
    # → 1ボックスにまとめる方式だと line_height でしか間隔調整できないため
    add_text(slide, 60, TITLE_TOP + 4, 820,
             f'【{slide_no}】改善提案リスト', 24, bold=True,
             color=NAVY, height_px=34)
    add_text(slide, 60, TITLE_TOP + 42, 820, subtitle,
             14, color=SUB_TEXT, height_px=22)
    _proposal_legend(slide, TITLE_TOP)


def _proposal_card(slide, p, y, card_h, prio_color, prio_label):
    """1提案カードを描画（共通ヘルパ）。

    [ビジュアル強化 2026-06-29]
    - カード背景：白＋NAVY 3%極薄の2層構造で奥行き演出
    - 左色帯：6px → 8px に拡幅
    - 番号サークル：30→34px、文字16pt化
    - 「改善後」帯：淡赤 → ORANGE薄塗（行動喚起トーン）
    """
    # カード背景：白＋NAVY 3%極薄
    _add_bg_frame(slide, 40, y, 1200, card_h,
                  fill=WHITE, line=BORDER_GRAY, line_width_pt=1)
    _add_bg_frame(slide, 40, y, 1200, card_h,
                  fill=NAVY, fill_alpha=3000)
    # 左色帯：6 → 8px に拡幅
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, y, 8, card_h,
              fill=prio_color)

    # ヘッダ行（番号サークル + タイトル + 優先度バッジ + 工数）
    head_top = y + 10
    # 番号サークル：30→34px、文字16pt化
    add_shape(slide, MSO_SHAPE.OVAL, 58, head_top, 34, 34,
              fill=prio_color)
    add_text(slide, 58, head_top + 4, 34,
             str(p.get('no', '')), 16, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER, height_px=26,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    raw_title = p.get('title', '')
    t_disp = _truncate_full(raw_title, 36)
    add_text(slide, 102, head_top + 2, 770, t_disp, 18,
             bold=True, color=NAVY, height_px=32)
    # 優先度バッジ：中央揃え強化（vertical_anchor=MIDDLE）
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 980, head_top + 1,
              70, 30, fill=prio_color)
    add_text(slide, 980, head_top + 1, 70, prio_label, 14,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             height_px=30, vertical_anchor=MSO_ANCHOR.MIDDLE)
    effort = p.get('effort', '中')
    add_text(slide, 1060, head_top + 5, 170,
             f'工数：{effort}', 14, color=SUB_TEXT, height_px=22)

    # 区切り線
    add_shape(slide, MSO_SHAPE.RECTANGLE, 56, head_top + 36,
              1170, 1, fill=BORDER_GRAY)

    # 改善箇所 / 課題（2行表示・各1ブロック）
    BODY_TOP = head_top + 44
    target_area = p.get('target_area', '')
    issue = p.get('issue', '')
    _add_multi_run_box(slide, 56, BODY_TOP, 568, 48, [
        {'runs': [
            {'text': '改善箇所　', 'size': 14, 'bold': True, 'color': NAVY},
            {'text': target_area, 'size': 14, 'bold': False, 'color': TEXT},
        ], 'line_height': 1.5},
    ])
    _add_multi_run_box(slide, 640, BODY_TOP, 584, 48, [
        {'runs': [
            {'text': '課題　', 'size': 14, 'bold': True, 'color': NAVY},
            {'text': issue, 'size': 14, 'bold': False, 'color': TEXT},
        ], 'line_height': 1.5},
    ])

    # 現状 / 改善後（2行表示・色違い背景帯）
    # [ビジュアル強化 2026-06-29]
    # - 現状：淡グレー（変更なし、現状を中立的に表現）
    # - 改善後：淡赤 → ORANGE 12%薄塗（行動喚起トーン）、ラベルORANGE化
    bottom_top = BODY_TOP + 52
    # 現状（左、淡グレー帯）
    add_shape(slide, MSO_SHAPE.RECTANGLE, 56, bottom_top, 568, 48,
              fill=LIGHT_GRAY, line=BORDER_GRAY, line_width_pt=0.5)
    _add_multi_run_box(slide, 64, bottom_top + 4, 552, 40, [
        {'runs': [
            {'text': '現状　', 'size': 14, 'bold': True, 'color': SUB_TEXT},
            {'text': p.get('before', ''), 'size': 14,
             'bold': False, 'color': TEXT},
        ], 'line_height': 1.5},
    ])
    # 改善後（右、ORANGE 12%薄塗 帯）
    # [FB対応 2026-07-02] 色ルール明確化：
    #   ORANGE は「行動喚起帯（結論・POINT・改善方向）」専用とし、
    #   ここの識別枠は NAVY_LIGHT に統一。塗りの ORANGE 12% は
    #   Before との対比を出す視覚要素として維持（枠線のみ変更）。
    add_shape(slide, MSO_SHAPE.RECTANGLE, 640, bottom_top, 584, 48,
              fill=WHITE, line=NAVY_LIGHT, line_width_pt=0.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 640, bottom_top, 584, 48,
              fill=ORANGE, fill_alpha=12000)
    _add_multi_run_box(slide, 648, bottom_top + 4, 568, 40, [
        {'runs': [
            {'text': '改善後　', 'size': 14, 'bold': True, 'color': ORANGE},
            {'text': p.get('after', ''), 'size': 14,
             'bold': False, 'color': TEXT},
        ], 'line_height': 1.5},
    ])


def _proposal_point_band(slide, summary, point_band_top=594, point_band_h=62):
    """POINT帯描画（スライド最下部）

    [ビジュアル強化 2026-06-29]
    - 本体：NAVYベタ → ORANGE薄塗（行動喚起トーン）
    - ラベル：RED → ORANGEベタ
    - 本文：WHITE → NAVY太字（コントラスト確保＋ブランド統一）
    [FB対応 2026-06-30]
    - ラベル⇔本文間隔 12→24px拡大
    [FB対応 2026-07-01 案B]
    - ラベル箱を上下2pxずつ縮めて立体感を演出
    - ラベル⇔本文間隔 24→40pxに拡大
    [FB対応 2026-07-02]
    - 帯高さ 32→62px（2行対応）、top 620→594 に上シフト
    - 文字上限 56→120字＋word_wrap=True（文字切れ根絶）
    - ラベル箱は帯全高で垂直センター視認性向上
    """
    if not summary:
        return
    flat = summary.replace('\n', ' ')
    LABEL_W = 90
    GAP = 40
    # 本体：ORANGE薄塗背景＋ORANGE枠
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, point_band_top,
              1200, point_band_h, fill=WHITE,
              line=ORANGE, line_width_pt=1)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, point_band_top,
              1200, point_band_h, fill=ORANGE, fill_alpha=10000)
    # ラベル部：ORANGEベタ、帯全高で垂直センター（視認性優先）
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              40, point_band_top, LABEL_W, point_band_h,
              fill=ORANGE)
    add_text(slide, 40, point_band_top, LABEL_W, 'POINT', 14,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             height_px=point_band_h, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    # 120字上限で切り詰め（GPTs側で120字以内に生成させる）
    flat_disp = _truncate_full(flat, 120)
    # [FB対応 2026-07-16] 句点孤立防止：末尾の「。」直前に word joiner
    # (U+2060) を挿入し、直前文字と結束させる。これにより「。」だけが
    # 2行目に取り残される日本語組版として不自然な折返しを回避する。
    # 副作用：word joiner は zero-width の非表示文字のため見た目に影響なし。
    if flat_disp.endswith('。'):
        flat_disp = flat_disp[:-1] + '\u2060。'
    # 文中の句点も同様に処理（「〜こと。〜」のように文中で句点＋続きがある場合）
    # は稀なので、末尾のみ対応で実運用の大半をカバー
    # 本文 left = 40 + LABEL_W + GAP = 170、幅 = 1200 + 40 - 170 = 1070
    body_left = 40 + LABEL_W + GAP
    body_w = 40 + 1200 - body_left
    body_box = add_text(slide, body_left, point_band_top, body_w, flat_disp, 14,
                        bold=True, color=NAVY, height_px=point_band_h,
                        line_height=1.35,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)
    # word_wrap=True で2行折り返しを許容（文字切れ根絶）
    body_box.text_frame.word_wrap = True


def add_proposal_onepager(prs, proposals_data, page_num=1, total=None,
                           author='UI/UX診断 by GPTs',
                           slide_no='2'):
    """
    改善提案リストを件数に応じて自動的に1or2枚で出力。

    - 4件以上：2枚構成（1/2 前半3件 / 2/2 後半残り + POINT帯）
    - 3件以下：1枚構成（全件 + POINT帯を1枚に集約）

    Args:
        proposals_data: dict
            service_name (str), proposals (list of max 5 dicts):
              no, title, priority ('高'/'中'/'低' or 'S'/'A'/'B'/'C'),
              effort ('小'/'中'/'大'), target_area, issue, before, after,
              target_score_item (optional)
            summary (str): POINT帯本文
        total: 全体総ページ数（None指定時は提案件数から自動算出）
    Returns:
        slide1 (3件以下) or (slide1, slide2) tuple (4件以上)
    """
    PRIO_RED = RGBColor(0xD0, 0x02, 0x1B)
    PRIO_ORANGE = RGBColor(0xE8, 0x8B, 0x1F)
    PRIO_GRAY = RGBColor(0x88, 0x88, 0x88)

    def _prio_color(p):
        if p in ('S', 'A', '高'):
            return PRIO_RED
        if p in ('B', '中'):
            return PRIO_ORANGE
        return PRIO_GRAY

    prio_label_map = {'S': '高', 'A': '高', 'B': '中', 'C': '低',
                      '高': '高', '中': '中', '低': '低'}

    # ==============================================
    # 文字数バリデーション（design_system.md §1.4 準拠）
    # ==============================================
    validate_length(proposals_data.get('service_name'), 'service_name',
                    'C-2 サービス名 (proposals_data.service_name)')
    for i, p in enumerate(proposals_data.get('proposals', [])[:5], start=1):
        if isinstance(p, dict):
            validate_length(p.get('title'), 'c2_title',
                            f'C-2 提案#{i} タイトル (proposals[{i-1}].title)')
            # POINT本文：issue / before / after / target_area などを連結したもの。
            # 各フィールド単体ではなく point があれば優先、なければ issue を上限内に。
            point_text = p.get('point') or p.get('issue') or ''
            validate_length(point_text, 'c2_point',
                            f'C-2 提案#{i} POINT (proposals[{i-1}].point|issue)')
            if p.get('priority') is not None:
                validate_length(str(p.get('priority')), 'c2_priority',
                                f'C-2 提案#{i} 優先度 (proposals[{i-1}].priority)')
            if p.get('category') is not None:
                validate_length(p.get('category'), 'c2_category',
                                f'C-2 提案#{i} カテゴリ (proposals[{i-1}].category)')

    proposals = proposals_data.get('proposals', [])[:5]
    n = len(proposals)
    summary = proposals_data.get('summary', '')

    # 全体ページ数を自動決定（1スライド or 2スライド）
    is_split = n >= 4
    auto_total = 2 if is_split else 1
    if total is None:
        total = auto_total

    # [FB対応 2026-06-30] タイトル帯64px化に追従、92→100に+8pxシフト
    card_top = 100
    card_gap = 10

    if not is_split:
        # ===============================================
        # 1スライド構成（3件以下：全件 + POINT帯を1枚に集約）
        # ===============================================
        slide = _blank_slide(prs)
        _proposal_title(slide, slide_no,
                        f'優先度の高い改善アクションを整理（全{n}件）')

        # POINT帯は最下部、カード領域はその上
        # [FB対応 2026-07-02] POINT帯を 62px に拡大し top を 620→594 に上シフト。
        # AREA_BOTTOM も 594-8=586 に縮小（カード高さは件数に応じ自動均等割）。
        POINT_BAND_TOP = 594
        AREA_BOTTOM = POINT_BAND_TOP - 8  # 586
        avail = AREA_BOTTOM - card_top
        if n > 0:
            card_h = (avail - card_gap * (n - 1)) // n
        else:
            card_h = 200

        for i, p in enumerate(proposals):
            y = card_top + i * (card_h + card_gap)
            priority = p.get('priority', 'B')
            _proposal_card(slide, p, y, card_h,
                           _prio_color(priority),
                           prio_label_map.get(priority, '中'))

        _proposal_point_band(slide, summary, point_band_top=POINT_BAND_TOP)
        _add_footer(slide, page_num, total, author=author)
        return slide

    # ===============================================
    # 2スライド構成（4件以上：前半3件 + 後半残り+POINT）
    # ===============================================
    first_n = 3
    proposals_a = proposals[:first_n]
    proposals_b = proposals[first_n:]

    # --- スライド 1/2 : 前半提案 ---
    slide = _blank_slide(prs)
    _proposal_title(slide, slide_no,
                    f'優先度の高い改善アクションを整理（1/2 前半 {first_n}件）')

    AREA_BOTTOM_A = 658  # フッター帯top=660 直前
    avail_a = AREA_BOTTOM_A - card_top
    card_h_a = (avail_a - card_gap * (first_n - 1)) // first_n

    for i, p in enumerate(proposals_a):
        y = card_top + i * (card_h_a + card_gap)
        priority = p.get('priority', 'B')
        _proposal_card(slide, p, y, card_h_a,
                       _prio_color(priority),
                       prio_label_map.get(priority, '中'))

    _add_footer(slide, page_num, total, author=author)
    slide1 = slide

    # --- スライド 2/2 : 後半提案 + POINT帯 ---
    slide = _blank_slide(prs)
    second_n = len(proposals_b)
    _proposal_title(slide, slide_no,
                    f'優先度の高い改善アクションを整理（2/2 後半 {second_n}件・POINT）')

    # [FB対応 2026-07-02] POINT帯を 62px に拡大し top を 620→594 に上シフト。
    # [FB対応 2026-07-16] 2枚目のカード高さを1枚目と統一し上詰め配置。
    #   従来：件数で領域を均等割 → 2件時カードが肥大化し下部余白が不自然に大
    #   改善：card_h_b = card_h_a（1枚目と同じ）で上詰め、下部余白は自然な
    #         「資料的余白」として機能させる。POINT帯位置は594px据え置き。
    POINT_BAND_TOP = 594
    AREA_BOTTOM_B = POINT_BAND_TOP - 8  # 586（現状は参考値、上詰めのため未使用）
    card_h_b = card_h_a  # 1枚目と統一（上詰め・下余白自然化）

    for i, p in enumerate(proposals_b):
        y = card_top + i * (card_h_b + card_gap)
        priority = p.get('priority', 'B')
        _proposal_card(slide, p, y, card_h_b,
                       _prio_color(priority),
                       prio_label_map.get(priority, '中'))

    _proposal_point_band(slide, summary, point_band_top=POINT_BAND_TOP)
    _add_footer(slide, page_num + 1, total, author=author)
    slide2 = slide

    return (slide1, slide2)


# =====================================================================
# C-3 / Layout 18: ビジュアル診断ボード（3枚構成）
# （add_visual_board）— フェーズ1：レーダーは表形式代替
#   1/3 LP構造マップ + 総評・最重要課題 + 6ステップ行動フロー
#   2/3 スコア視覚化（10項目表＋強み/課題コンパクトカード）
#   3/3 Before/After Top3 ハイライト（赤帯Before／オレンジ帯After）
# =====================================================================
def _visual_title(slide, slide_no, sub_title):
    """C-3共通：タイトル帯（左赤縦帯 + タイトル + サブタイトル）

    [ビジュアル強化 2026-06-29]
    - 左赤縦帯：6px → 10px に拡幅
    [FB対応 2026-06-30]
    - タイトル帯高さ：56 → 64px に拡大
    - サブタイトル top：+8px下げて可読性向上
    """
    TITLE_TOP = 16
    TITLE_H = 64
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, TITLE_TOP, 10, TITLE_H, fill=RED)
    add_text(slide, 60, TITLE_TOP + 4, 1140,
             f'【{slide_no}】ビジュアル診断ボード', 24, bold=True,
             color=NAVY, height_px=34)
    add_text(slide, 60, TITLE_TOP + 42, 1140, sub_title,
             14, color=SUB_TEXT, height_px=22)


def _visual_meta_header(slide, visual_data, top=90):
    """C-3共通：目的/対象/診断日/スコアの4分割メタカード

    [ビジュアル強化 2026-06-29]
    - 左色帯：4px → 6px に拡幅
    - カード背景：白＋NAVY 4%極薄の2層構造で奥行き演出
    [FB対応 2026-06-30]
    - top デフォルト 82 → 90px（タイトル帯64px化に追従）
    """
    INFO_H = 64
    card_w = (1200 - 12 * 3) // 4  # 291px
    score_val = visual_data.get('total_score', 0)
    rank = visual_data.get('rank', '-')
    rank_label = visual_data.get('rank_label', '')

    # [FB対応 2026-07-04] 全4カード NAVY 統一（装飾目的の色分けは廃止）
    # [FB対応 2026-07-16 v15.4] 総合スコアカードのみ訴求強化：
    #   スコア値を14pt→26pt に、ランク値を14pt→32pt に大型化。
    #   販売資産として最も訴求すべきポイントを視覚的に主役化する。
    rank_color = NAVY if rank in ('S', 'A', 'B') else RED
    meta_cards = [
        {'bar': NAVY, 'label': '目的',
         'value': visual_data.get('purpose', '-'), 'size': 14},
        {'bar': NAVY, 'label': '対象',
         'value': visual_data.get('target', '-'), 'size': 14},
        {'bar': NAVY, 'label': '診断日',
         'value': visual_data.get('diagnosis_date', '-'), 'size': 14},
        # 総合スコアは runs で「スコア + / 50 + ランク」を大型化・色分け
        {'bar': NAVY, 'label': '総合スコア',
         'is_score_card': True},
    ]

    for i, card in enumerate(meta_cards):
        x = 40 + i * (card_w + 12)
        # カード本体：白＋極薄NAVYの2層構造
        _add_bg_frame(slide, x, top, card_w, INFO_H,
                      fill=WHITE, line=BORDER_GRAY, line_width_pt=1)
        _add_bg_frame(slide, x, top, card_w, INFO_H,
                      fill=NAVY, fill_alpha=4000)
        # 左色帯：4 → 6px に拡幅
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, top, 6, INFO_H,
                  fill=card['bar'])
        # 総合スコアカードのみ、C-1側と同じ大型化仕様を適用
        if card.get('is_score_card'):
            _add_multi_run_box(slide, x + 18, top + 8, card_w - 26,
                               INFO_H - 14, [
                {'text': card['label'], 'size': 14, 'bold': True,
                 'color': SUB_TEXT, 'line_height': 1.0},
                {'runs': [
                    {'text': str(score_val), 'size': 26, 'bold': True,
                     'color': rank_color},
                    {'text': '  / 50    ', 'size': 14, 'bold': False,
                     'color': SUB_TEXT},
                    {'text': rank, 'size': 26, 'bold': True,
                     'color': rank_color},
                ], 'line_height': 1.0},
            ])
        else:
            _add_multi_run_box(slide, x + 18, top + 8, card_w - 26,
                               INFO_H - 14, [
                {'text': card['label'], 'size': 14, 'bold': True,
                 'color': SUB_TEXT, 'line_height': 1.0},
                {'text': card['value'],
                 'size': card.get('size', 16),
                 'bold': card.get('value_bold', False),
                 'color': card.get('value_color', TEXT),
                 'line_height': 1.1},
            ])


def _visual_direction_footer(slide, direction, top=594):
    """C-3共通：フッター上の改善方向性帯

    [ビジュアル強化 2026-06-29]
    - 帯高さ：28 → 36px に拡大（視認性向上）
    - ORANGE_LIGHT → ORANGE 12%透過に変更
    - 本文 14pt → 16pt 太字 NAVY色（メリハリ強化）
    [FB対応 2026-06-30]
    - ラベル⇔本文間隔 20→28px拡大
    [FB対応 2026-07-01 案B]
    - ラベル箱を上下2pxずつ縮めて立体感を演出
    - ラベル⇔本文間隔 28→40pxに拡大
    [FB対応 2026-07-02]
    - 帯高さ 36→62px（2行対応）、top を 594 へ統一
    - 文字上限 44→120字＋word_wrap=True（文字切れ根絶）
    - 本文14pt太字に変更（2行構成時の可読性優先）
    - ラベル箱は帯全高で垂直センター視認性向上
    - 後方互換：top=620 指定時も 594 として扱う（旧呼出しの吸収）
    """
    if not direction:
        return
    # 後方互換：旧デフォルト620/612を渡された場合も新top 594 に強制上書き
    BAND_TOP = 594 if top in (620, 612) else top
    BAND_H = 62
    LABEL_W = 110
    GAP = 40
    # 本体：ORANGE薄塗＋枠
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, BAND_TOP, 1200, BAND_H,
              fill=WHITE, line=ORANGE, line_width_pt=1)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, BAND_TOP, 1200, BAND_H,
              fill=ORANGE, fill_alpha=12000)
    # ラベル部：ORANGEベタ、帯全高で垂直センター（視認性優先）
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              40, BAND_TOP, LABEL_W, BAND_H,
              fill=ORANGE)
    add_text(slide, 40, BAND_TOP, LABEL_W, '改善方向', 14,
             bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             height_px=BAND_H, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    flat = direction.replace('\n', ' ')
    # 120字上限で切り詰め（GPTs側で120字以内に生成させる）
    flat = _truncate_full(flat, 120)
    # 本文 left = 40 + LABEL_W + GAP = 190、幅 = 40 + 1200 - 190 = 1050
    body_left = 40 + LABEL_W + GAP
    body_w = 40 + 1200 - body_left
    body_box = add_text(slide, body_left, BAND_TOP, body_w, flat, 14,
                        bold=True, color=NAVY, height_px=BAND_H,
                        line_height=1.35,
                        vertical_anchor=MSO_ANCHOR.MIDDLE)
    # word_wrap=True で2行折り返しを許容（文字切れ根絶）
    body_box.text_frame.word_wrap = True


def _normalize_flow_status(s):
    """flow_steps.status の表記揺れを正規化"""
    if s in ('✓', 'OK', True, 'success', '成功'):
        return '✓'
    if s in ('✕', '×', 'NG', False, 'fail', '離脱'):
        return '✕'
    return '？'


def add_visual_board(prs, visual_data, page_num=1, total=3, slide_no='3',
                     author='UI/UX診断 by GPTs'):
    """
    ビジュアル診断ボードを3枚構成で出力。

    Args:
        visual_data: dict （詳細は visual_data_schema.md 参照）
        page_num: 開始ページ番号（個別出力時=1、統合版では既存ページ数+1）
        total: 全体ページ数（個別=3、統合版では C-1+C-2+3 を渡す）
        slide_no: 章番号（既定 '3'）
    Returns:
        (slide1, slide2, slide3) tuple
    """
    # ==============================================
    # 文字数バリデーション（design_system.md §1.5 準拠）
    # ==============================================
    validate_length(visual_data.get('service_name'), 'service_name',
                    'C-3 サービス名 (visual_data.service_name)')
    validate_length(visual_data.get('summary'), 'c3_summary',
                    'C-3 総評 (visual_data.summary)')
    validate_length(visual_data.get('top_issue'), 'c3_top_issue',
                    'C-3 最優先課題 (visual_data.top_issue)')
    validate_length(visual_data.get('direction'), 'c3_direction',
                    'C-3 改善方向 (visual_data.direction)')
    for i, sec in enumerate(visual_data.get('sections', []), start=1):
        if isinstance(sec, dict):
            validate_length(sec.get('label'), 'c3_section_label',
                            f'C-3 LP構造マップ#{i} (sections[{i-1}].label)')
    for i, step in enumerate(visual_data.get('flow_steps', [])[:6], start=1):
        if isinstance(step, dict):
            # ステップ全体（label + note）は1ステップ枠（25文字以内）
            combined = (step.get('label', '') + ' ' + step.get('note', '')).strip()
            validate_length(combined, 'c3_flow_step',
                            f'C-3 行動フロー#{i} (flow_steps[{i-1}].label+note)')
    for i, hl in enumerate(visual_data.get('highlights', [])[:3], start=1):
        if isinstance(hl, dict):
            validate_length(hl.get('title'), 'c3_highlight_title',
                            f'C-3 ハイライト#{i} タイトル (highlights[{i-1}].title)')
            validate_length(hl.get('before'), 'c3_before',
                            f'C-3 ハイライト#{i} Before (highlights[{i-1}].before)')
            validate_length(hl.get('after'), 'c3_after',
                            f'C-3 ハイライト#{i} After (highlights[{i-1}].after)')

    # ==============================================================
    # スライド 1/3 : LP構造マップ + 総評・最重要課題 + 行動フロー
    # ==============================================================
    slide = _blank_slide(prs)
    _visual_title(slide, slide_no,
                  '構造マップ・総評・最重要課題・ユーザー行動フロー')
    _visual_meta_header(slide, visual_data, top=90)

    # ---- 左：LP構造マップ ----
    # [FB対応 2026-07-03] 番号サークルの色（NAVY=良好 / RED=要改善）が
    # 凡例なしで意味不明だったため、ヘッダ帯の直下にコンパクトな凡例チップを追加。
    # ヘッダタイトルは幅を切って右側に凡例スペースを確保。
    LEFT_X = 40
    LEFT_W = 380
    MAP_TOP = 158
    MAP_HEADER_H = 28
    add_shape(slide, MSO_SHAPE.RECTANGLE, LEFT_X, MAP_TOP, LEFT_W,
              MAP_HEADER_H, fill=NAVY)
    # タイトル短縮（読み進み順 → 順）で凡例と共存
    add_text(slide, LEFT_X + 16, MAP_TOP, 220,
             '◤ LP構造マップ', 14, bold=True, color=WHITE,
             height_px=MAP_HEADER_H, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    # 凡例チップ（ヘッダ右端に配置）：●良好 / ●要改善
    LEGEND_Y = MAP_TOP + 6
    LEGEND_H = 16
    # 良好チップ（NAVY_LIGHT の丸 + "良好"）
    add_shape(slide, MSO_SHAPE.OVAL, LEFT_X + LEFT_W - 148, LEGEND_Y + 3,
              10, 10, fill=NAVY_LIGHT)
    add_text(slide, LEFT_X + LEFT_W - 132, LEGEND_Y, 40, '良好',
             14, bold=False, color=WHITE, height_px=LEGEND_H,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    # 要改善チップ（RED の丸 + "要改善"）
    add_shape(slide, MSO_SHAPE.OVAL, LEFT_X + LEFT_W - 78, LEGEND_Y + 3,
              10, 10, fill=RED)
    add_text(slide, LEFT_X + LEFT_W - 62, LEGEND_Y, 60, '要改善',
             14, bold=False, color=WHITE, height_px=LEGEND_H,
             vertical_anchor=MSO_ANCHOR.MIDDLE)

    sections = visual_data.get('sections', [])[:9]
    # [FB対応 2026-07-20 v15.6] 空値フォールバック実装
    # sectionsが空の場合、代替テキストを1件だけ描画
    if not sections:
        sections = [{'no': 1, 'label': 'データ未取得',
                     'desc': 'LP構造データが渡されませんでした',
                     'has_issue': False}]
    n_sec = len(sections)
    MAP_BODY_TOP = MAP_TOP + MAP_HEADER_H
    MAP_BODY_H = 276 - MAP_HEADER_H  # =248 (158+28+248=434, 行動フロー top=444 直前)
    add_shape(slide, MSO_SHAPE.RECTANGLE, LEFT_X, MAP_BODY_TOP, LEFT_W,
              MAP_BODY_H, fill=LIGHT_GRAY,
              line=BORDER_GRAY, line_width_pt=0.5)

    if n_sec > 0:
        # 行高は最小26pxを確保し、name｜desc を1行 runs 構造で表示
        row_h = max(26, MAP_BODY_H // n_sec)
        for i, sec in enumerate(sections):
            y = MAP_BODY_TOP + 4 + i * row_h
            # [FB対応 2026-06-30] has_issue 判定を status からも派生
            # status='!' or '✕' → 課題あり、'✓' → 課題なし
            status_val = sec.get('status', '')
            has_issue = (
                sec.get('has_issue', False)
                or status_val in ('!', '✕', '×', 'NG', 'fail', '離脱')
            )
            # 番号サークル（課題ありは赤、なしは紺）
            circle_color = RED if has_issue else NAVY
            add_shape(slide, MSO_SHAPE.OVAL, LEFT_X + 8, y, 22, 22,
                      fill=circle_color)
            add_text(slide, LEFT_X + 8, y, 22,
                     str(sec.get('no', i + 1)), 14, bold=True,
                     color=WHITE, align=PP_ALIGN.CENTER,
                     height_px=22, vertical_anchor=MSO_ANCHOR.MIDDLE)
            # 名前｜説明（1ブロック内に runs で混在、1行で表示）
            # [FB対応 2026-06-30] スキーマ不整合バグ修正
            # visual_data_schema.md と test データは `label` を使うため、
            # まず `label` を優先し、後方互換で `name` も読む。
            # [FB対応 2026-07-03] status 記号（✓/!/○）は凡例なしで意味不明のため、
            # 状態はサークル色（NAVY=良好/RED=要改善）で伝える設計に変更。
            # desc フィールドがあればそれを表示、なければ空文字（記号は使わない）。
            # [FB対応 2026-07-08 v12] LP構造マップで name=「FV/メインビジュアル」(10字)
            # が 8字上限で「FV/メインビジュ…」に末尾途切れを発生。
            # 実描画基準（14pt × 幅約300px, LEFT_W - 48）で 8→ 12字へ緩和。
            name = _truncate_full(
                sec.get('label') or sec.get('name', ''), 12)
            desc_raw = sec.get('desc', '')  # status は表示から除外
            desc = _truncate_full(desc_raw, 16)
            _add_multi_run_box(slide, LEFT_X + 38, y, LEFT_W - 48, row_h - 2, [
                {'runs': [
                    {'text': name, 'size': 14, 'bold': True, 'color': TEXT},
                    {'text': f'  {desc}' if desc else '',
                     'size': 14, 'bold': False, 'color': SUB_TEXT},
                ], 'line_height': 1.1},
            ])

    # ---- 中央上：総評 ----
    # [FB対応 2026-07-05 v10] v9で SUMMARY_H=66 に縮めた結果、
    # 内部領域 66-36=30px に対し summary_text=40字は 14pt×1.35≈
    # 26px×2行=52px 必要で 22px はみ出し→ ISSUE_TOP(234) に衝突。
    # 根本対策：
    #   1) summary上限を 40→ 24字 (1行収容確定)
    #   2) SUMMARY_H を 66→ 76（+10px）で内部領域 40px を確保
    #      → 1行(26px)+上下余白の安全値
    #   3) ISSUE_H を 200→ 190（-10px）で縦合計を維持
    #      SUMMARY_END = 158+76 = 234（旧ISSUE_TOPと同値）
    #      ISSUE_TOP = 234+10 = 244、ISSUE_END = 244+190 = 434（同一）
    #      内部領域 190-28-14 = 148px → 22字×3件×26px = 78px 安全
    # [FB対応 2026-07-09 v14・1行運用確定] 大原則：総評は"1行運用"を絶対条件とする。
    # 2行折返しが起きたら被り事故が発生するため、コード側で「1行に確実収容できる
    # 最大文字数」を上限に設定し、Instructions側でAI生成コピーがその範囲内で
    # 書くよう指示する。折返しを許容するのではなく発生させない設計に統一。
    #
    # 実測結果（14pt Meiryo × width 352px, 2026-07-09）：
    #   全角 17字（換算17.0）→ 1行内収容 ✅ 限界値
    #   全角 18字（換算18.0）→ 2行折返し ❌
    #   半角混 18字（換算16.0）→ 1行内収容 ✅ 安全
    #   半角混 19字（換算17.5）→ 2行折返し ❌
    # ⇒ 全角換算 17 が明確な境界値。_truncate_full(17) が正解。
    #
    # 座標最適化（v13 SUMMARY_H=84 は2行想定で過剰。1行分に最適化）：
    #   SUMMARY_H 84→ 62（ヘッダ28px + 本文帯34px = 1行安全収容）
    #   ISSUE_TOPギャップ +14→ +12（元設計値）
    #   ISSUE_H 178→ 200（削減分22を吸収、Top3内部余白拡大＝可読性向上）
    #   下端 158+62+12+200 = 432（右カラム下端434とほぼ同値・±2px許容）
    # [v15/2026-07-12] 総評本文帯の天地余白拡大（62→68）：
    # v14 では本文帯 34px、テキスト vertical_anchor=MIDDLE でも
    # 「詰まっている感」が残っていた。本文帯を 40px に拡大して天地に
    # 各6px の余裕を確保する。下流連動：ISSUE_TOP=238、ISSUE_END=438
    # （行動フロー帯 top=444 との安全余白 6px）。
    # [FB対応 2026-07-16 v15.5・入江さん提案] 3カラム再配分：
    # 中央（総評・最重要課題）を主役化するため、右（強み）から80px移譲。
    # 強み：3点固定・28字上限で折返し許容の設計 → 336px幅でも表示可
    # 中央：内部描画幅 380→460、20字（380px）に52pxの安全マージン確保
    CENTER_X = 432
    CENTER_W = 460  # 380 → 460（+80px拡張）
    SUMMARY_TOP = 158
    SUMMARY_H = 68  # v14=62 → v15=68（本文帯 34→40px, 天地余白 +6px）
    add_shape(slide, MSO_SHAPE.RECTANGLE, CENTER_X, SUMMARY_TOP, CENTER_W,
              28, fill=ORANGE)
    add_text(slide, CENTER_X + 16, SUMMARY_TOP, CENTER_W - 32,
             '◤ 総評', 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, CENTER_X, SUMMARY_TOP + 28, CENTER_W,
              SUMMARY_H - 28, fill=ORANGE_LIGHT,
              line=BORDER_GRAY, line_width_pt=0.5)
    # [v14] 全角換算17が実測境界値（全角17字ギリギリ1行、18字で折返し発生）。
    # Instructions側でも「総評は17字以内・1行運用」をルール化。
    # [v15/2026-07-12] 本文帯（34px）内でテキスト縦位置が上寄り（上6px/下2px）
    # だったため、テキストボックスを本文帯フル高（34px）で作成し
    # vertical_anchor=MIDDLE で縦中央寄せ。上下余白の視覚的不均衡を解消。
    summary_text = _truncate_full(visual_data.get('summary', ''), 17)
    _add_multi_run_box(slide, CENTER_X + 14, SUMMARY_TOP + 28,
                       CENTER_W - 28, SUMMARY_H - 28, [
        {'text': summary_text, 'size': 14, 'bold': False,
         'color': TEXT, 'line_height': 1.35},
    ], vertical_anchor=MSO_ANCHOR.MIDDLE)

    # ---- 中央下：最重要課題（Top3 or Top4） ----
    # [FB対応 2026-07-03] Top3レイアウト崩れ修正：line_height 1.35 + 22字切詰
    # [FB対応 2026-07-04 v9・案C] ISSUE_H 172→200（内部158px確保）
    # [FB対応 2026-07-16] 22字境界LibreOffice折返し回避で20字に厳格化
    # [FB対応 2026-07-16 v15.3] Top3/Top4動的対応：
    #   - Instructions v3.2 で C=4項目化に伴い、GPTs側が独自書換して
    #     ヘッダーが黒文字化する副作用が実測で確認された
    #   - コード側で件数判定→ヘッダー文字列自動生成でGPTs書換余地を排除
    #   - 4項目時は ISSUE_H を 200→240 に拡張し尾切れを根本解消
    # top_issues (複数) と top_issue (単数) の両方を扱える後方互換実装
    top_issues_all = (
        visual_data.get('top_issues')
        or ([visual_data.get('top_issue')] if visual_data.get('top_issue') else [])
        or visual_data.get('priority_issues', [])
    )
    # ランクに応じて表示件数を決定（B=3, C=4, D=5 は将来対応）
    # 現状は Top3/Top4 の2パターンをサポート
    n_issues = min(len(top_issues_all), 4)  # 最大4項目まで
    top_issues = top_issues_all[:n_issues]
    # ヘッダー文字列をコード側で動的生成（GPTs書換によるカラー崩壊防止）
    issue_header_text = f'⚠ 最重要課題（Top{n_issues}）'
    # [FB対応 2026-07-16 v15.5] 中央カラム幅拡張（380→460）により20字が
    # 1行完結する境界（内部432px確保）。ISSUE_H は元の200pxに復帰。
    # 3項目：26px×3 + space8pt×2 = 約94px（内部158pxに余裕64px）
    # 4項目：26px×4 + space2pt×3 = 約112px（内部158pxに余裕46px）
    ISSUE_H = 200
    space_after = 8 if n_issues <= 3 else 2
    ISSUE_TOP = SUMMARY_TOP + SUMMARY_H + 12
    add_shape(slide, MSO_SHAPE.RECTANGLE, CENTER_X, ISSUE_TOP, CENTER_W,
              28, fill=RED)
    add_text(slide, CENTER_X + 16, ISSUE_TOP, CENTER_W - 32,
             issue_header_text, 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, CENTER_X, ISSUE_TOP + 28, CENTER_W,
              ISSUE_H - 28, fill=RED_LIGHT,
              line=BORDER_GRAY, line_width_pt=0.5)
    # [FB対応 2026-07-16 v15.5] 中央カラム460px化で内部432px確保。
    # 20字（380px）が1行完結する境界。安全マージン+52pxで折返し防止。
    issue_paragraphs = [
        {'text': f'{i + 1}. ' + _truncate_full(s, 20),
         'size': 14, 'bold': False, 'color': TEXT,
         'line_height': 1.2, 'space_after_pt': space_after}
        for i, s in enumerate(top_issues)
    ]
    _add_multi_run_box(slide, CENTER_X + 14, ISSUE_TOP + 34,
                       CENTER_W - 28, ISSUE_H - 42, issue_paragraphs)

    # ---- 右：強み3点 ----
    # [FB対応 2026-07-16 v15.5・入江さん提案] 右カラム縮小・80px中央へ移譲。
    # 強み3点固定・28字上限の設計上、内部304pxで折返し許容表示は成立。
    RIGHT_X = 904  # 824 → 904（80px右シフト）
    RIGHT_W = 336  # 416 → 336（-80px縮小）
    STR_TOP = 158
    STR_H = 276  # 元の値に復帰（ISSUE_H拡張が不要になったため）
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, STR_TOP, RIGHT_W,
              28, fill=NAVY)
    add_text(slide, RIGHT_X + 16, STR_TOP, RIGHT_W - 32,
             '◎ 強み（活かすべき点）', 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, STR_TOP + 28, RIGHT_W,
              STR_H - 28, fill=LIGHT_GRAY,
              line=BORDER_GRAY, line_width_pt=0.5)
    # [FB対応 2026-07-08 v12] 強み項目の「ファーストビュー設/計」のような
    # 語の途中での不自然な折返しを避けるため、実描画基準で 40→ 28字に厳格化。
    # 14pt Meiryo × width 384px（実領域 352px）の実測描画で 1行に収まる上限。
    # [v15/2026-07-12] 折返し発生時の「1文の連続感」と「項目間の分離感」を分離制御：
    #   - line_height 1.7→1.2（折返し行間を詰める＝連続感）
    #   - space_after_pt=8pt を項目間に追加（項目間の分離感）
    # これにより「明確な CTA ボタンとファーストビュー設計」のような28字文が
    # 折返しても、2行目が「同一文の続き」と視認できる。
    strengths = visual_data.get('strengths', [])[:3]
    str_paragraphs = [
        {'text': '● ' + _truncate_full(s, 28),
         'size': 14, 'bold': False, 'color': TEXT,
         'line_height': 1.2, 'space_after_pt': 8}
        for s in strengths
    ]
    _add_multi_run_box(slide, RIGHT_X + 16, STR_TOP + 36,
                       RIGHT_W - 32, STR_H - 44, str_paragraphs)

    # ---- 下：行動フロー6ステップ ----
    # [ビジュアル強化 2026-06-29] ステップ円36→44px、ステータス16→18pt
    # [FB対応 2026-07-02] direction帯62px化に伴い FLOW_H 156→146
    # [FB対応 2026-07-16 v15.5] 中央カラム拡張で情報表示問題を解決したため
    # FLOW_TOP/FLOW_H は元の値（444/146）に復帰。レイアウト整合性を保持。
    FLOW_TOP = 444
    FLOW_H = 146
    FLOW_W = 1200
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, FLOW_TOP, FLOW_W, 28, fill=NAVY)
    add_text(slide, 56, FLOW_TOP, 500,
             '◤ ユーザー行動フローと課題', 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    # [FB対応 2026-07-02] フロー総括をヘッダ右端に移設（従来はステップnoteと
    # 重なっていたため）。40字上限は据置き。GOLDでナビー帯内の視認性確保。
    flow_summary_early = visual_data.get('flow_summary', '')
    if flow_summary_early:
        add_text(slide, 560, FLOW_TOP, FLOW_W - 560 - 16,
                 '→ ' + _truncate_full(flow_summary_early, 40), 14, bold=True,
                 color=GOLD, align=PP_ALIGN.RIGHT,
                 height_px=28, vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, FLOW_TOP + 28, FLOW_W,
              FLOW_H - 28, fill=WHITE,
              line=BORDER_GRAY, line_width_pt=0.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, FLOW_TOP + 28, FLOW_W,
              FLOW_H - 28, fill=NAVY, fill_alpha=3000)

    flow_steps = visual_data.get('flow_steps', [])[:6]
    n_flow = len(flow_steps)
    if n_flow > 0:
        step_gap = 8
        step_w = (FLOW_W - 20 - step_gap * (n_flow - 1) - 60) // n_flow
        arrow_w = 60 // max(1, n_flow - 1) if n_flow > 1 else 0
        step_top = FLOW_TOP + 36
        STEP_D = 44  # 36→44に拡大
        for i, step in enumerate(flow_steps):
            x = 50 + i * (step_w + step_gap + arrow_w)
            status = _normalize_flow_status(step.get('status', '✓'))
            step_color = NAVY if status == '✓' else RED
            # ステップ円
            add_shape(slide, MSO_SHAPE.OVAL, x + step_w // 2 - STEP_D // 2,
                      step_top, STEP_D, STEP_D, fill=step_color)
            add_text(slide, x + step_w // 2 - STEP_D // 2, step_top,
                     STEP_D, status, 18, bold=True, color=WHITE,
                     align=PP_ALIGN.CENTER, height_px=STEP_D,
                     vertical_anchor=MSO_ANCHOR.MIDDLE)
            # ラベル + note（1ブロック2段落）
            # ステップ円44pxに合わせてラベル開始位置を+8px調整
            label = _truncate_full(step.get('label', ''), 6)
            note = _truncate_full(step.get('note', ''), 14)
            _add_multi_run_box(slide, x, step_top + STEP_D + 4, step_w, 56, [
                {'text': label, 'size': 14, 'bold': True,
                 'color': TEXT, 'align': PP_ALIGN.CENTER,
                 'line_height': 1.1},
                {'text': note, 'size': 14, 'bold': False,
                 'color': SUB_TEXT, 'align': PP_ALIGN.CENTER,
                 'line_height': 1.2},
            ])
            # 矢印（最終ステップ以外）
            if i < n_flow - 1:
                arrow_x = x + step_w + step_gap
                add_text(slide, arrow_x, step_top + 8, arrow_w, '→',
                         20, bold=True, color=ORANGE,
                         align=PP_ALIGN.CENTER, height_px=STEP_D,
                         vertical_anchor=MSO_ANCHOR.MIDDLE)

    # フロー総括はヘッダ右端に移設済み（上のフロー帯ヘッダ内 GOLD 表示）。
    # 下段位置はステップ note と重なっていたため廃止。

    # 改善方向フッター帯
    _visual_direction_footer(slide, visual_data.get('direction', ''))
    _add_footer(slide, page_num, total, author=author)
    slide1 = slide

    # ==============================================================
    # スライド 2/3 : スコア視覚化（10項目表 + 強み/課題コンパクト）
    # ==============================================================
    slide = _blank_slide(prs)
    _visual_title(slide, slide_no,
                  'スコア視覚化：10項目評価と強み・最優先課題の俯瞰')
    _visual_meta_header(slide, visual_data, top=90)

    # ---- 左：10項目スコア表（フェーズ1=表形式、フェーズ2でレーダー差し替え） ----
    SCORE_TOP = 158
    SCORE_W = 580
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, SCORE_TOP, SCORE_W, 28, fill=NAVY)
    add_text(slide, 56, SCORE_TOP, SCORE_W - 32,
             '◤ 10項目スコア（満点5・ギャップ可視化）',
             14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)

    scores = visual_data.get('scores', [])[:10]
    # [FB対応 2026-07-20 v15.6] 空値フォールバック実装
    # scoresが空の場合、代替表示を提供（バーは描かず、警告文を1行）
    n_scores = len(scores)
    SCORE_BODY_TOP = SCORE_TOP + 28
    # [FB対応 2026-07-03] 改善方向帯top=594 と18px重複していたため
    # SCORE_BODY_H を 426→400 に短縮。end=158+28+400=586（帯top594との
    # 間に8px安全余白確保）。
    SCORE_BODY_H = 400
    add_shape(slide, MSO_SHAPE.RECTANGLE, 40, SCORE_BODY_TOP, SCORE_W,
              SCORE_BODY_H, fill=WHITE,
              line=BORDER_GRAY, line_width_pt=0.5)

    if n_scores > 0:
        row_h = min(42, SCORE_BODY_H // n_scores)
        for i, item in enumerate(scores):
            y = SCORE_BODY_TOP + 6 + i * row_h
            sv = item.get('score', '－')
            mv = item.get('max', 5)
            is_na = (sv == '－' or sv == '－/5')
            ratio = (sv / mv) if (not is_na and mv) else 0
            # スコア色分け：2点以下=赤、4点以上=紺、それ以外=グレー
            if is_na:
                bar_color = BORDER_GRAY
                score_color = SUB_TEXT
                score_bold = False
            elif sv <= 2:
                bar_color = RED
                score_color = RED
                score_bold = True
            elif sv >= 4:
                bar_color = NAVY
                score_color = NAVY
                score_bold = True
            else:
                bar_color = ORANGE
                score_color = TEXT
                score_bold = False
            # 項目名
            # [FB対応 2026-06-30] スキーマ不整合バグ修正
            # scores 配列は C-1 と同じ `name` キーを使う（schema 統一）
            # `category` は後方互換のため最終フォールバック
            cat_disp = _truncate_full(
                item.get('name') or item.get('category', ''), 10)
            add_text(slide, 56, y + 8, 180, cat_disp, 14, bold=True,
                     color=TEXT, height_px=22)
            # ★バー視覚化（max=5、塗りつぶし数=sv）
            bar_x = 240
            bar_w = 28
            for star_i in range(5):
                star_color = bar_color if (not is_na and star_i < sv) else BORDER_GRAY
                add_shape(slide, MSO_SHAPE.RECTANGLE,
                          bar_x + star_i * (bar_w + 4), y + 12,
                          bar_w, 16, fill=star_color)
            # スコア数値
            score_text = '－/5' if is_na else f'{sv}/5'
            add_text(slide, bar_x + 5 * (bar_w + 4) + 8, y + 8, 50, score_text,
                     14, bold=score_bold, color=score_color, height_px=22)
    else:
        # [FB対応 2026-07-20 v15.6] scoresが空の時のフォールバック描画
        add_text(slide, 56, SCORE_BODY_TOP + 20, SCORE_W - 32,
                 '⚠ 10項目スコアデータが渡されていません', 14, bold=True,
                 color=RED, height_px=22)
        add_text(slide, 56, SCORE_BODY_TOP + 50, SCORE_W - 32,
                 'visual_dataに`scores`キー（C-1のscoresと同一10件）を必ず含めてください',
                 12, color=SUB_TEXT, height_px=40)

    # ---- 右上：強み3点（コンパクト） ----
    RIGHT_X = 640
    RIGHT_W = 600
    STR_TOP = 158
    STR_H = 218
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, STR_TOP, RIGHT_W,
              28, fill=NAVY)
    add_text(slide, RIGHT_X + 16, STR_TOP, RIGHT_W - 32,
             '◎ 強み（活かすべき点）', 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, STR_TOP + 28, RIGHT_W,
              STR_H - 28, fill=LIGHT_GRAY,
              line=BORDER_GRAY, line_width_pt=0.5)
    # [v15/2026-07-12] リスト内折返しの連続感と項目間の分離感を分離制御
    strengths = visual_data.get('strengths', [])[:3]
    # [FB対応 2026-07-20 v15.6] 空値フォールバック実装
    if not strengths:
        strengths = ['強みデータが渡されませんでした（visual_dataにstrengthsキー必須）']
    str_paragraphs = [
        {'text': '● ' + _truncate_full(s, 40),
         'size': 14, 'bold': False, 'color': TEXT,
         'line_height': 1.2, 'space_after_pt': 8}
        for s in strengths
    ]
    _add_multi_run_box(slide, RIGHT_X + 16, STR_TOP + 36,
                       RIGHT_W - 32, STR_H - 44, str_paragraphs)

    # ---- 右下：最優先課題3点（コンパクト） ----
    ISSUE_TOP = STR_TOP + STR_H + 12
    # [FB対応 2026-07-02] direction帯62px化(top=594)に伴い ISSUE_H を 220→198
    # に縮小し下端を 586 に (安全余白8px)。
    #   ISSUE_TOP = STR_TOP(158) + STR_H(218) + 12 = 388
    #   ISSUE_TOP + ISSUE_H = 388 + 198 = 586 < 594（direction帯 top）
    ISSUE_H = 198
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, ISSUE_TOP, RIGHT_W,
              28, fill=RED)
    add_text(slide, RIGHT_X + 16, ISSUE_TOP, RIGHT_W - 32,
             '⚠ 最優先で直すべき点', 14, bold=True, color=WHITE,
             height_px=28, letter_spacing=2,
             vertical_anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, RIGHT_X, ISSUE_TOP + 28, RIGHT_W,
              ISSUE_H - 28, fill=RED_LIGHT,
              line=BORDER_GRAY, line_width_pt=0.5)
    # [v15/2026-07-12] リスト内折返しの連続感と項目間の分離感を分離制御
    priority_issues = visual_data.get('priority_issues', [])[:3]
    # [FB対応 2026-07-20 v15.6] 空値フォールバック実装
    if not priority_issues:
        priority_issues = ['最優先課題データが渡されませんでした（visual_dataにpriority_issuesキー必須）']
    issue_paragraphs = [
        {'text': f'{i + 1}. ' + _truncate_full(s, 50),
         'size': 14, 'bold': False, 'color': TEXT,
         'line_height': 1.2, 'space_after_pt': 8}
        for i, s in enumerate(priority_issues)
    ]
    _add_multi_run_box(slide, RIGHT_X + 16, ISSUE_TOP + 36,
                       RIGHT_W - 32, ISSUE_H - 44, issue_paragraphs)

    # 改善方向フッター帯
    _visual_direction_footer(slide, visual_data.get('direction', ''))
    _add_footer(slide, page_num + 1, total, author=author)
    slide2 = slide

    # ==============================================================
    # スライド 3/3 : Before/After Top3 ハイライト
    # ==============================================================
    slide = _blank_slide(prs)
    _visual_title(slide, slide_no,
                  'Before/After ハイライト：優先度の高い改善Top3')

    highlights = visual_data.get('highlights', [])[:3]
    n_h = len(highlights)

    HL_TOP = 92
    HL_GAP = 12
    # [FB対応 2026-07-02] direction帯62px化(top=594)に伴い HL_AREA_H を
    # 520→494 に縮小し下端を 586 に (安全余白8px)。92 + 494 = 586。
    HL_AREA_H = 494
    hl_h = (HL_AREA_H - HL_GAP * max(0, n_h - 1)) // max(1, n_h)

    prio_color_map = {'高': PRIO_RED, '中': PRIO_ORANGE, '低': PRIO_GRAY,
                      'S': PRIO_RED, 'A': PRIO_RED,
                      'B': PRIO_ORANGE, 'C': PRIO_GRAY}

    for i, h in enumerate(highlights):
        y = HL_TOP + i * (hl_h + HL_GAP)
        priority = h.get('priority', '中')
        prio_col = prio_color_map.get(priority, PRIO_ORANGE)
        effort = h.get('effort', '中')

        # カード背景 + 左色帯
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 40, y, 1200,
                  hl_h, fill=WHITE,
                  line=BORDER_GRAY, line_width_pt=1)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 40, y, 6, hl_h, fill=prio_col)

        # ヘッダ：番号 + タイトル + 優先度バッジ + 工数
        head_top = y + 10
        add_shape(slide, MSO_SHAPE.OVAL, 60, head_top, 32, 32, fill=prio_col)
        add_text(slide, 60, head_top + 4, 32, str(h.get('no', i + 1)),
                 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=24, vertical_anchor=MSO_ANCHOR.MIDDLE)
        title = _truncate_full(h.get('title', ''), 30)
        # [FB対応 2026-07-16 v15.3] 優先度x=820→780 に伴い、タイトル幅を
        # 700→660 に縮小（右端 762、優先度左端 780 との安全余白 18px）
        add_text(slide, 102, head_top + 4, 660, title, 18,
                 bold=True, color=NAVY, height_px=28)
        # [FB対応 2026-07-16 v15.3] 優先度・工数・箇所のx座標配分を再設計。
        # 従来：優先度820/工数900/箇所1040-1220（箇所幅180pxで日本語折返し多発）
        # 改善：優先度780/工数860/箇所960-1220（箇所幅260pxに拡張）
        # 優先度バッジ
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 780, head_top + 2,
                  70, 28, fill=prio_col)
        add_text(slide, 780, head_top + 5, 70, priority, 14,
                 bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=22)
        # 工数（幅130→90に縮小、コンパクト化）
        add_text(slide, 860, head_top + 5, 90, f'工数：{effort}', 14,
                 color=SUB_TEXT, height_px=22)
        # 改善箇所（幅180→260に拡張、日本語テキストの折返し防止）
        target_area = _truncate_full(h.get('target_area', ''), 20)
        add_text(slide, 960, head_top + 5, 260,
                 f'箇所：{target_area}', 14, color=SUB_TEXT, height_px=22)

        # 区切り線
        add_shape(slide, MSO_SHAPE.RECTANGLE, 60, head_top + 38, 1170, 1,
                  fill=BORDER_GRAY)

        # Before/After 2カラム
        body_top = head_top + 50
        body_h = hl_h - 60
        col_w = 568

        # Before（赤帯、letter_spacing削除でBEFORE全文表示）
        add_shape(slide, MSO_SHAPE.RECTANGLE, 60, body_top, 90, 28, fill=RED)
        add_text(slide, 60, body_top, 90, 'BEFORE', 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=28,
                 vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 150, body_top,
                  col_w - 90, 28, fill=RED_LIGHT)
        before_text = _truncate_full(h.get('before', ''), 38)
        add_text(slide, 158, body_top, col_w - 106, before_text, 14,
                 color=TEXT, height_px=28,
                 vertical_anchor=MSO_ANCHOR.MIDDLE)

        # After（オレンジ帯、letter_spacing削除）
        add_shape(slide, MSO_SHAPE.RECTANGLE, 644, body_top, 90, 28,
                  fill=ORANGE)
        add_text(slide, 644, body_top, 90, 'AFTER', 14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER,
                 height_px=28,
                 vertical_anchor=MSO_ANCHOR.MIDDLE)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 734, body_top,
                  col_w - 90, 28, fill=ORANGE_LIGHT)
        after_text = _truncate_full(h.get('after', ''), 38)
        add_text(slide, 742, body_top, col_w - 106, after_text, 14,
                 color=TEXT, height_px=28,
                 vertical_anchor=MSO_ANCHOR.MIDDLE)

    # 改善方向フッター帯
    _visual_direction_footer(slide, visual_data.get('direction', ''))
    _add_footer(slide, page_num + 2, total, author=author)
    slide3 = slide

    return (slide1, slide2, slide3)


# =====================================================================
# build_full_report：C-1 + C-2 + C-3 を1ファイルに統合
# =====================================================================
def build_full_report(diagnosis, proposals, visual_data,
                      author='UI/UX診断 by GPTs'):
    """
    UI診断統合レポートを生成。C-1（2枚）+ C-2（1-2枚）+ C-3（3枚）を
    1つの Presentation に格納して返す。

    Args:
        diagnosis: dict（add_scorecard_onepager と同じスキーマ）
        proposals: dict（add_proposal_onepager と同じスキーマ）
        visual_data: dict（visual_data_schema.md 参照）
        author: フッター著者名

    Returns:
        prs: Presentation オブジェクト（.save() してダウンロード可能）
    """
    prs = create_presentation()

    # 提案件数からC-2のスライド数を推定
    n_proposals = len(proposals.get('proposals', []))
    c2_slides = 2 if n_proposals >= 4 else 1

    # 全体ページ数 = C-1(2) + C-2(1or2) + C-3(3)
    total = 2 + c2_slides + 3

    # C-1: スコアカード (2枚)
    add_scorecard_onepager(prs, diagnosis,
                            page_num=1, total=total,
                            author=author, slide_no='1')

    # C-2: 改善提案リスト (1 or 2枚)
    add_proposal_onepager(prs, proposals,
                          page_num=3, total=total,
                          author=author, slide_no='2')

    # C-3: ビジュアル診断ボード (3枚)
    add_visual_board(prs, visual_data,
                     page_num=3 + c2_slides, total=total,
                     slide_no='3', author=author)

    return prs


# =====================================================================
# 動作確認用（直接実行時）
# =====================================================================
if __name__ == '__main__':
    prs = create_presentation()

    add_cover(prs,
              title='Q2 販売戦略レビュー',
              date='2026年7月15日',
              author='○○株式会社 営業企画部',
              subtitle='上期実績の振り返りと下期に向けた重点施策')

    add_agenda(prs, items=[
        {'title': '市場環境の整理', 'desc': 'マクロ・競合・顧客動向'},
        {'title': '現状の課題', 'desc': '上期で見えた3つの論点'},
        {'title': '優先施策', 'desc': '4象限で整理した優先順位'},
        {'title': '実行スケジュール', 'desc': '下期6か月のロードマップ'},
    ], page_num=2, total=4)

    add_issue_summary(prs,
        title='現状の3つの課題',
        cards=[
            {'no': '01', 'heading': 'リード獲得の伸び悩み',
             'body': '主要チャネルで前年比横ばい。\n新規流入経路の追加検討が必要。'},
            {'no': '02', 'heading': '商談化率の低下',
             'body': 'リードから商談への転換が\n前年比10pt低下。\nスコアリング再設計が必要。'},
            {'no': '03', 'heading': '受注後の解約率増',
             'body': '導入後3か月の解約が増加。\nオンボーディング強化が急務。'},
        ],
        page_num=3, total=4,
        lead='上期実績の分析から、下期に向けた3つの課題を整理しました。',
        conclusion='3課題は独立ではなく「リード→商談→受注→定着」の流れで連動している。')

    add_closing(prs,
                message='Thank you.',
                next_step='次回ミーティング：7月22日（火）\n各施策の責任者から進捗報告をお願いします。',
                contact='問い合わせ：営業企画部 山田',
                page_num=4, total=4)

    prs.save('test_output.pptx')
    print('Generated: test_output.pptx')

    # =====================================================================
    # C-1 / C-2 動作確認サンプル
    # =====================================================================
    # ── C-1: UI診断スコアカード（3枚） ──
    diagnosis_sample = {
        'service_name': 'サンプルECサイト',
        'url': 'https://example.com',
        'input_type': 'URL入力',
        'total_score': 32,
        'rank': 'B',
        'rank_label': '標準的（改善余地あり）',
        'scores': [
            {'category': 'ファーストビュー訴求',  'score': 4.0, 'max': 5, 'comment': 'メインビジュアルは明快。サブコピーで価値訴求を補強したい。'},
            {'category': 'ナビゲーション設計',    'score': 3.5, 'max': 5, 'comment': 'グローバルナビは整理されているが、カテゴリ深度がやや深い。'},
            {'category': '視覚的階層',            'score': 4.5, 'max': 5, 'comment': '見出し階層が明確で読みやすい。'},
            {'category': '検索・絞り込み',        'score': 2.0, 'max': 5, 'comment': '絞り込み条件の保存・組み合わせができない。'},
            {'category': '商品詳細の情報量',      'score': 3.0, 'max': 5, 'comment': 'サイズ感や使用シーンの情報が不足している。'},
            {'category': 'カート/購入導線',       'score': 4.0, 'max': 5, 'comment': 'カート遷移はスムーズ。ゲスト購入も可能。'},
            {'category': 'フォーム使いやすさ',    'score': 2.5, 'max': 5, 'comment': '入力エラー表示が遅延し、修正箇所が分かりにくい。'},
            {'category': 'モバイル最適化',        'score': 3.5, 'max': 5, 'comment': 'タップ領域は十分。横スクロールが一部発生。'},
            {'category': '表示速度',              'score': 3.0, 'max': 5, 'comment': '初回読み込みが3.2秒。画像最適化が必要。'},
            {'category': 'アクセシビリティ',      'score': 2.0, 'max': 5, 'comment': 'コントラスト比不足の箇所が散見される。alt属性も欠落多い。'},
        ],
        'strengths': [
            'ファーストビューのメインビジュアルが明快で価値が伝わる',
            '見出し階層が整理されており情報の優先順位が分かりやすい',
            'カート→購入の導線がスムーズで離脱が起きにくい',
        ],
        'priority_issues': [
            '検索・絞り込みの組み合わせ条件が保存できず再操作の手間が大きい',
            'フォーム入力時のエラー表示が遅延し、入力者が修正箇所を特定しにくい',
            'アクセシビリティ（コントラスト・alt属性）の改善が必要',
        ],
        'conclusion': (
            '総合32/50点・ランクBで、基本的なUI設計は整っているが改善余地が大きい状態です。\n'
            '特に「検索・絞り込み」「フォーム」「アクセシビリティ」の3領域は緊急度が高く、'
            'いずれもCVR・離脱率に直接影響する論点です。\n'
            '一方、ファーストビュー訴求や視覚的階層は強みであり、これらを軸に他領域の改善を進めることで'
            '効率的なUX向上が見込めます。'
        ),
    }

    prs2 = create_presentation()
    add_scorecard_overview(prs2, diagnosis_sample, page_num=1, total=3)
    add_scorecard_table(prs2, diagnosis_sample, page_num=2, total=3)
    add_scorecard_conclusion(prs2, diagnosis_sample, page_num=3, total=3,
        cta='続けて C-2 で改善提案リストをご確認ください。')
    prs2.save('test_scorecard.pptx')
    print('Generated: test_scorecard.pptx')

    # ── C-2: 改善提案リスト（3枚） ──
    proposals_sample = {
        'service_name': 'サンプルECサイト',
        'proposals': [
            {'no': 1, 'title': '絞り込み条件の保存機能と複数条件の組み合わせ対応',
             'category': '検索・絞り込み', 'priority': 'S',
             'current': '条件選択後に他ページへ移動すると条件がリセットされ、再度入力が必要。\n複数カテゴリの同時選択にも対応していない。',
             'action': 'URLパラメータ／LocalStorageで条件を保持し、戻る操作でも復元する。\nチェックボックスUIに変更し、複数条件のAND選択を可能にする。',
             'expected': '再操作率を約40%削減、検索ページ滞在時間を-25%。\n結果として商品詳細到達率の向上が見込める。'},
            {'no': 2, 'title': 'フォーム入力エラーのリアルタイム表示',
             'category': 'フォーム使いやすさ', 'priority': 'S',
             'current': '送信ボタン押下後にまとめてエラーが表示され、修正箇所が分かりにくい。\nエラー文言も「入力が誤っています」と汎用的で具体性が低い。',
             'action': '各フィールドのblur時にバリデーションを実行し、エラー文言を「具体的に何がおかしいか」「どう直すか」の2点を含めて表示。',
             'expected': 'フォーム離脱率を-15pt、購入完了率を+5pt見込み。'},
            {'no': 3, 'title': 'アクセシビリティ最低基準（WCAG AA）への準拠',
             'category': 'アクセシビリティ', 'priority': 'A',
             'current': '主要テキストのコントラスト比が4.5未満の箇所が多数。\n商品画像のalt属性が空欄またはファイル名のままで、SR利用者が情報を取得できない。',
             'action': 'デザインシステムのカラートークンを見直し、AA基準を満たすよう調整。\nalt属性は「商品名＋色／柄」を必須項目化し、CMS入力時にバリデーションする。',
             'expected': '潜在顧客層（視覚障碍ユーザ・高齢者）への到達拡大。\nSEO面でも画像検索流入の改善が見込める。'},
            {'no': 4, 'title': '商品詳細ページの情報量強化',
             'category': '商品詳細', 'priority': 'A',
             'current': 'サイズ表記が「M」「L」のみで実寸が不明。使用シーン写真も少なく、購入後イメージが湧きにくい。',
             'action': 'サイズ実寸表（cm）と着用モデル身長を必須化。シーン別の使用写真を最低3枚追加する。',
             'expected': '商品詳細→カート遷移率を+8pt。返品率の低下も期待される。'},
        ],
        'summary': '優先度SとAの計4件を3か月以内に着手することで、ランクBからAへの引き上げが現実的。',
    }

    prs3 = create_presentation()
    # 一覧（1枚）＋ 詳細スライド4枚（1スライド1提案）＝ 計5枚
    proposals_list = proposals_sample['proposals']
    total_pages = 1 + len(proposals_list)  # = 5
    add_proposal_overview(prs3, proposals_sample,
                          page_num=1, total=total_pages)
    for idx, proposal in enumerate(proposals_list):
        is_last = (idx == len(proposals_list) - 1)
        add_proposal_detail(
            prs3, proposal,
            page_num=2 + idx,
            total=total_pages,
            summary=proposals_sample['summary'] if is_last else None,
        )
    prs3.save('test_proposals.pptx')
    print('Generated: test_proposals.pptx')

    # =====================================================================
    # C-1 onepager / C-2 onepager 動作確認サンプル
    # （添付参照デザイン: SpeakUpEnglish）
    # =====================================================================
    onepager_diagnosis = {
        'service_name': 'スピークアップ英会話',
        'input_type': 'スクリーンショット画像',
        'total_score': 30,
        'rank': 'B',
        'rank_label': '標準的（改善余地あり）',
        'scores': [
            {'category': 'ファーストビュー',   'score': 3, 'max': 5,
             'comment': '清潔感はあるが、体験価値と成果が一瞬で伝わりにくい'},
            {'category': 'キャッチコピー',     'score': 3, 'max': 5,
             'comment': '開始喚起は明快だが、選ばれる理由が弱い'},
            {'category': 'CTA設計',           'score': 1, 'max': 5,
             'comment': 'FV・中盤に主要CTAがなく行動導線が弱い'},
            {'category': '信頼性・権威性',     'score': 3, 'max': 5,
             'comment': '実績数値はあるが、根拠や証拠の厚みが不足'},
            {'category': 'フォーム設計',       'score': 1, 'max': 5,
             'comment': '問い合わせ導線はあるが入力前の不安解消が少ない'},
            {'category': 'レスポンシブ',       'score': 3, 'max': 5,
             'comment': 'PC表示は整っているが、SPでCTA固定が必要'},
            {'category': '読みやすさ',         'score': 5, 'max': 5,
             'comment': '余白・行間・色数が整理され読みやすい'},
            {'category': '情報設計',           'score': 3, 'max': 5,
             'comment': '特徴と料金は見やすいが、比較・成果・不安解消が不足'},
            {'category': 'ブランド一貫性',     'score': 5, 'max': 5,
             'comment': '青基調で教育サービスらしい安心感がある'},
            {'category': '表示速度・技術',     'score': 3, 'max': 5,
             'comment': '画像中心のため最適化余地がある'},
        ],
        'strengths': [
            '余白が広く、読みやすいレイアウトで安心感がある',
            '料金・特徴・講師情報が整理され、サービス概要を把握しやすい',
            '青基調の配色で教育サービスらしい信頼感が出ている',
        ],
        'priority_issues': [
            'ファーストビューに主要CTAがなく、次の行動が分かりにくい',
            '受講後の成果・具体的なベネフィットが弱く、比較検討時の決め手に欠ける',
            '受講生の声が短く、信頼材料としての説得力が不足している',
        ],
        'conclusion': '興味を持ったユーザーも、無料体験への問い合わせ前に離脱する可能性が高い。FV内CTAと成果訴求・信頼材料の3点を最優先で改善する。',
    }

    prs4 = create_presentation()
    add_scorecard_onepager(prs4, onepager_diagnosis,
                            page_num=1, total=1, slide_no='1')
    prs4.save('test_onepager_scorecard.pptx')
    print('Generated: test_onepager_scorecard.pptx')

    onepager_proposals = {
        'service_name': 'スピークアップ英会話',
        'proposals': [
            {'no': 1, 'title': 'ファーストビューに無料体験CTAを追加',
             'priority': '高', 'effort': '小',
             'target_area': 'ファーストビュー',
             'issue': '興味を持った直後に行動できない',
             'before': '見出しと説明文はあるが、主要ボタンが見当たらない',
             'after': '無料体験・料金確認の2ボタンをFV内に配置',
             'target_score_item': 'CTA設計'},
            {'no': 2, 'title': 'キャッチコピーを成果訴求型に変更',
             'priority': '高', 'effort': '小',
             'target_area': 'メインコピー',
             'issue': '学習後のメリットが弱く差別化しにくい',
             'before': '「英語、始めませんか？」で汎用的な印象',
             'after': '忙しい人でも続く、成果が見えるオンライン英会話として訴求',
             'target_score_item': 'キャッチコピー'},
            {'no': 3, 'title': '実績数値に根拠と補足を追加',
             'priority': '中', 'effort': '中',
             'target_area': '実績帯',
             'issue': '数値はあるが信頼の裏付けが不足',
             'before': '開校年・採用率・段階数・時間帯のみ表示',
             'after': '累計受講者数、継続率、満足度など検討材料を追加',
             'target_score_item': '信頼性・権威性'},
            {'no': 4, 'title': '受講生の声を具体化',
             'priority': '中', 'effort': '小',
             'target_area': '受講生の声',
             'issue': '短文のみで利用シーンや成果が伝わりにくい',
             'before': '「続けられています」など一言コメント中心',
             'after': '課題・受講理由・成果を含む3行レビューへ拡張',
             'target_score_item': '信頼性・権威性'},
            {'no': 5, 'title': '料金セクション下に比較とCTAを追加',
             'priority': '高', 'effort': '中',
             'target_area': '料金プラン',
             'issue': 'プラン選択後の次アクションが弱い',
             'before': '料金カードの提示で終わり、相談導線が離れている',
             'after': 'おすすめプラン表示と無料相談CTAを直下に配置',
             'target_score_item': 'CTA設計'},
        ],
        'summary': '最優先の一手：ファーストビューに無料体験CTAを追加し、直下に成果・安心材料を3点で提示する。',
    }

    prs5 = create_presentation()
    add_proposal_onepager(prs5, onepager_proposals,
                           page_num=1, total=1, slide_no='2')
    prs5.save('test_onepager_proposals.pptx')
    print('Generated: test_onepager_proposals.pptx')
